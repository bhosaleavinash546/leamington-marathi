#!/usr/bin/env python3
"""
CostVision — Implementation Blueprint presentation (secure deployment + CAPEE
integration) for senior management. Light professional theme, logo top-left on
every slide, native shapes/diagrams, speaker notes per slide.

Content is grounded in docs/CostVision-Secure-Deployment-CAPEE-Integration.md,
which itself is grounded in a line-by-line audit of the codebase.

Regenerate:  python3 build_blueprint_pptx.py
Output:      CostVision-Implementation-Blueprint.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
DARK    = RGBColor(0x0F, 0x17, 0x2A)
BODY    = RGBColor(0x33, 0x41, 0x55)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
BG      = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xF1, 0xF5, 0xF9)
PANEL2  = RGBColor(0xEF, 0xF6, 0xFF)
GREENBG = RGBColor(0xEC, 0xFD, 0xF5)
AMBERBG = RGBColor(0xFF, 0xFB, 0xEB)
GREEN   = RGBColor(0x05, 0x96, 0x69)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
RED     = RGBColor(0xDC, 0x26, 0x26)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
CYAN    = RGBColor(0x08, 0x91, 0xB2)
LINE    = RGBColor(0xE2, 0xE8, 0xF0)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def box(slide, x, y, w, h, fill=None, line=None, round_=False, radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for r in para:
            t, size, color, bold = r[0], r[1], r[2], r[3]
            italic = r[4] if len(r) > 4 else False
            run = p.add_run(); run.text = t
            f = run.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold
            f.italic = italic; f.name = 'Calibri'
    return tb

def logo(slide, x=Inches(0.35), y=Inches(0.22), scale=1.0):
    s = scale
    badge = box(slide, x, y, Inches(0.42 * s), Inches(0.42 * s), fill=INDIGO, round_=True, radius=0.28)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'cv'
    r.font.size = Pt(17 * s); r.font.bold = True; r.font.color.rgb = BG; r.font.name = 'Calibri'
    text(slide, x + Inches(0.52 * s), y - Inches(0.03 * s), Inches(2.6), Inches(0.32),
         [[('CostVision', 18 * s, BLUE, True)]])
    text(slide, x + Inches(0.52 * s), y + Inches(0.24 * s), Inches(2.8), Inches(0.22),
         [[('AI  COST  INTELLIGENCE', 7.5 * s, MUTED, False)]])

def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt

def header(title, kicker=None):
    slide = prs.slides.add_slide(BLANK)
    box(slide, 0, 0, W, H, fill=BG)
    logo(slide)
    box(slide, 0, Inches(0.78), W, Pt(2.2), fill=INDIGO)
    if kicker:
        text(slide, Inches(0.45), Inches(0.95), Inches(11.5), Inches(0.3),
             [[(kicker.upper(), 11, BLUE, True)]])
        ty = Inches(1.22)
    else:
        ty = Inches(1.02)
    text(slide, Inches(0.45), ty, Inches(12.4), Inches(0.6), [[(title, 27, DARK, True)]])
    return slide

def flow_box(slide, x, y, w, h, title, sub, color, fill=None):
    b = box(slide, x, y, w, h, fill=(fill or PANEL), round_=True, radius=0.12)
    box(slide, x, y, Inches(0.07), h, fill=color)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = color; r.font.name = 'Calibri'
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(9.5); r2.font.color.rgb = BODY; r2.font.name = 'Calibri'
    return b

def down_arrow(slide, x, y, h=Inches(0.32)):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.32), h)
    a.fill.solid(); a.fill.fore_color.rgb = LINE
    a.line.fill.background(); a.shadow.inherit = False
    return a


# ═══════════════ 1 — TITLE ═══════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=BG)
box(s, 0, 0, W, Inches(0.16), fill=INDIGO)
logo(s, x=Inches(0.5), y=Inches(0.45), scale=1.25)
text(s, Inches(0.9), Inches(2.15), Inches(11.8), Inches(1.0),
     [[('CostVision Implementation Blueprint', 36, DARK, True)]])
text(s, Inches(0.9), Inches(3.1), Inches(11.4), Inches(0.9),
     [[('Secure enterprise deployment inside our network — and integration with CAPEE.', 19, BODY, False)],
      [('All CAD models, drawings and images stay inside the company. Verified in the code.', 19, BODY, False)]])
for i, (t, c) in enumerate([('100% CAD stays internal', GREEN), ('AI controls BUILT & tested', BLUE), ('~3–5 weeks remaining', VIOLET), ('6-phase rollout', CYAN)]):
    x = Inches(0.9 + i * 2.95)
    chip = box(s, x, Inches(4.55), Inches(2.7), Inches(0.52), fill=PANEL, round_=True, radius=0.5)
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t; r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = c; r.font.name = 'Calibri'
text(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
     [[('Management briefing  ·  July 2026  ·  Grounded in a line-by-line audit of the platform code', 12, MUTED, False)]])
box(s, 0, H - Inches(0.16), W, Inches(0.16), fill=INDIGO)
notes(s, "Welcome. This session is the implementation blueprint for CostVision: how we deploy it securely so that "
         "no CAD data ever leaves our network, and how we integrate it with our existing CAPEE costing tool. "
         "Everything in this deck comes from an actual audit of the platform's code — not vendor promises. "
         "One important update since the written plan: the two security controls at the heart of this blueprint — "
         "private AI routing and the air-gapped switch — have now been BUILT into the platform and tested live. "
         "By the end I'll ask for one decision: approval to start Phase 1 and 2.")

# ═══════════════ 2 — THE DECISION ═══════════════
s = header('The decision we are asking for today', 'Purpose of this meeting')
box(s, Inches(0.45), Inches(2.1), Inches(12.45), Inches(1.5), fill=PANEL2, round_=True, radius=0.08)
text(s, Inches(0.8), Inches(2.35), Inches(11.8), Inches(1.1),
     [[('Approve Phases 1–2: ', 17, BLUE, True),
       ('the IT-Security assessment and the architecture design for deploying CostVision inside our network '
        'and connecting it to CAPEE.  (~5–7 weeks, existing teams, no licence spend.)', 17, BODY, False)]],
     line_spacing=1.2)
pts = [
    ('Why now', 'The tool is built, tested (1,005 automated tests) and proven on live runs — the value is waiting on deployment, not development.', BLUE),
    ('Why it is safe', 'CAD processing already runs fully inside the server. The private AI routing and the air-gapped switch are now BUILT and tested live.', GREEN),
    ('Why CAPEE wins', 'CAPEE keeps the workflow; CostVision adds the AI, physics and organisational memory behind it. No tool replacement.', VIOLET),
]
for i, (t, d, c) in enumerate(pts):
    x = Inches(0.45 + i * 4.25)
    box(s, x, Inches(4.0), Inches(4.0), Inches(2.4), fill=PANEL, round_=True, radius=0.07)
    box(s, x, Inches(4.0), Inches(4.0), Inches(0.09), fill=c)
    text(s, x + Inches(0.25), Inches(4.25), Inches(3.55), Inches(0.4), [[(t, 15.5, c, True)]])
    text(s, x + Inches(0.25), Inches(4.75), Inches(3.55), Inches(1.55), [[(d, 12, BODY, False)]], line_spacing=1.18)
notes(s, "One clear ask: approve the first two phases — the security assessment and the architecture design. "
         "That's five to seven weeks with existing teams and no new licence spend. Why now? Because the tool itself is "
         "finished and proven; what remains is deployment. Why is it safe? Because CAD processing already happens "
         "entirely inside the server — we verified this in the code — and the two security controls this plan called "
         "for, private AI routing and the air-gapped switch, are now built into the platform and tested live. "
         "And why does CAPEE win? Because we are not replacing CAPEE — we're giving it an AI engine and a memory.")

# ═══════════════ 3 — WHAT COSTVISION IS ═══════════════
s = header('What CostVision is — a quick recap', 'Background')
caps = [
    ('18 cost engines', 'Machining, casting, injection, PCB, software and more — physics-based, every figure traceable', BLUE),
    ('CAD-to-Cost', 'Reads STEP/IGES models with a real geometry engine and auto-fills the costing', CYAN),
    ('Photo-to-BOM', 'Costs a PCB from photographs — detects components and builds the bill of materials', VIOLET),
    ('Agentic AI (upgraded 2026)', 'Learns from every quote; conformal confidence with a guarantee, a causal negotiation coach, outcome-weighted findings and what-if scenarios — all glass-box', GREEN),
]
for i, (t, d, c) in enumerate(caps):
    x = Inches(0.45 + (i % 2) * 6.35); y = Inches(2.05 + (i // 2) * 1.62)
    box(s, x, y, Inches(6.1), Inches(1.42), fill=PANEL, round_=True, radius=0.09)
    box(s, x, y, Inches(0.09), Inches(1.42), fill=c)
    text(s, x + Inches(0.28), y + Inches(0.14), Inches(5.6), Inches(0.4), [[(t, 14.5, c, True)]])
    text(s, x + Inches(0.28), y + Inches(0.58), Inches(5.6), Inches(0.75), [[(d, 11.5, BODY, False)]], line_spacing=1.12)
box(s, Inches(0.45), Inches(5.55), Inches(12.45), Inches(1.3), fill=GREENBG, round_=True, radius=0.08)
text(s, Inches(0.8), Inches(5.75), Inches(11.8), Inches(0.95),
     [[('Proven, not promised:  ', 13.5, GREEN, True),
       ('estimating error cut from 10.9% to 0.3% after learning from 3 real quotes  ·  £512k/yr of pricing issues '
        'found autonomously in the live demo  ·  1,005 automated tests protect it all.', 13.5, BODY, False)]],
     line_spacing=1.2)
notes(s, "Thirty seconds of background for anyone new. CostVision costs parts bottom-up with physics across 18 "
         "manufacturing processes. It reads CAD files directly, it can cost a circuit board from photographs, and — "
         "the newest layer — it learns: from every analysis and every real supplier quote, getting measurably more "
         "accurate and even raising savings findings on its own. The green bar shows verified results from live runs, "
         "including half a million pounds a year of findings surfaced autonomously in our demonstration.")

# ═══════════════ 3B — LATEST AGENTIC INTELLIGENCE (2026) ═══════════════
s = header('Latest intelligence — six upgrades that widen the moat', 'Background · new in 2026')
rows = [
    ('Self-audit & guardrails', 'A deterministic layer re-checks every estimate for known errors and applies bounded corrections — proven on five real automotive CAD parts (a fuel tank fell from £217 to £25). Geometry stays ground truth; the AI never overrules it.', INDIGO),
    ('Learns from your quotes', 'Bulk-import logged actuals; the engine builds per-segment calibration and watches for drift, so the number tightens as your own history grows — and status rides on every estimate.', GREEN),
    ('Negotiation coach', 'Names the cost driver and writes the counter: "a quote of £95 implies aluminium +14% above spot — hold at £86.34."', BLUE),
    ('Outcome-weighted findings', 'The agent learns which findings convert and ranks by money it can really recover — a £100k gap that closes beats a £200k gap that never does.', VIOLET),
    ('What-if engine', 'Drag a commodity ±20% and the price recomputes live; run it across the whole portfolio ("if steel +10%, 7 parts go underwater").', CYAN),
    ('Glass-box by design', 'Every learned or derived number stays auditable — no black-box weight ever touches the price. The competitor edge is the opposite.', AMBER),
]
for i, (t, d, c) in enumerate(rows):
    y = Inches(2.0 + i * 0.82)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.74), fill=PANEL if i % 2 == 0 else BG, round_=True, radius=0.1)
    box(s, Inches(0.45), y, Inches(0.09), Inches(0.74), fill=c)
    text(s, Inches(0.78), y + Inches(0.10), Inches(3.5), Inches(0.55), [[(t, 12.5, c, True)]])
    text(s, Inches(4.3), y + Inches(0.09), Inches(8.4), Inches(0.6), [[(d, 11, BODY, False)]], line_spacing=1.08)
notes(s, "One slide on what's newest, because it directly strengthens the security-and-defensibility case this deck "
         "makes. Six upgrades shipped in 2026, all built and tested. The one I'd start with is the self-audit: a "
         "deterministic layer that re-checks every single estimate against a list of mistakes we've actually seen, and "
         "corrects them within bounds — and I can prove it, because we ran it on five real automotive CAD parts and it "
         "caught every one, a fuel tank alone falling from two hundred and seventeen pounds to twenty-five. The tool now "
         "also learns from your own logged quotes, building per-segment calibration and watching for drift, so it gets "
         "tighter on your parts over time. Alongside those, the machine-sizing that used to be hard-coded is now "
         "universal across commodities. The negotiation coach names the commodity driving a cost and writes the buyer's "
         "counter-argument; the autonomous agent prioritises the savings we can really recover; and the what-if engine "
         "answers 'if a commodity moves, what happens' as a defensible conditional. Underneath all of it, everything "
         "stays glass-box — every learned number is auditable, nothing is a black box — and the whole thing, CAD engine "
         "included, ships in a container whose build is verified in our CI pipeline. That is exactly the combination — "
         "checked, learned-from-data, and auditable — that makes these numbers safe to deploy and defend.")

# ═══════════════ 4 — TECH STACK IN PLAIN WORDS ═══════════════
s = header('What the tool is made of — in plain words', 'Technology, simply explained')
stack = [
    ('Frontend', 'What users see', 'The web application in the browser — forms, dashboards, reports. Nothing is installed on laptops; it is served from our own server.', BLUE),
    ('Backend', 'The engine room', 'The server application that does the work: runs the 18 cost engines, the AI logic and all calculations. Runs on a standard company virtual machine.', INDIGO),
    ('Database', 'The memory', 'Where rate libraries, saved costings and the AI knowledge base live. Standard corporate database (PostgreSQL), encrypted, backed up by IT.', VIOLET),
    ('CAD engine', 'The 3D model reader', 'Specialist geometry software (OCCT — the same core used by CAD vendors) that measures the 3D model: volume, weight, walls, holes. Runs INSIDE the backend.', CYAN),
    ('AI layer', 'The language brain', 'The AI model used for vision and language tasks. NOW BUILT: one switchboard controls every AI call — route it to a private endpoint, or turn it fully OFF (air-gapped).', GREEN),
]
for i, (t, tag, d, c) in enumerate(stack):
    y = Inches(2.05 + i * 0.95)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.84), fill=PANEL, round_=True, radius=0.14)
    box(s, Inches(0.45), y, Inches(0.09), Inches(0.84), fill=c)
    text(s, Inches(0.75), y + Inches(0.10), Inches(2.0), Inches(0.4), [[(t, 15, c, True)]])
    text(s, Inches(0.75), y + Inches(0.47), Inches(2.0), Inches(0.3), [[(tag, 10, MUTED, False, True)]])
    text(s, Inches(2.95), y + Inches(0.12), Inches(9.7), Inches(0.65), [[(d, 11.5, BODY, False)]], line_spacing=1.1)
notes(s, "Before the architecture, the five building blocks in plain words. The FRONTEND is simply what users see in "
         "their browser — nothing installs on laptops. The BACKEND is the engine room on our server, where all "
         "calculations happen. The DATABASE is the memory — a standard corporate database our IT already knows how to "
         "run and back up. The GEOMETRY ENGINE is the specialist CAD reader — and note, it runs inside our backend, "
         "not in any cloud. The AI LAYER is the only piece that talks to an external service — and the control for "
         "that is now BUILT into the platform: a single switchboard that every AI call must pass through. Point it at "
         "a private endpoint with one setting, or flip the air-gapped switch and it makes no external calls at all. "
         "We tested all three modes live and the server announces its mode at start-up for audit.")

# ═══════════════ 5 — THE REQUIREMENT & KEY FINDING ═══════════════
s = header('The security requirement — and the key finding', 'Security')
box(s, Inches(0.45), Inches(2.05), Inches(12.45), Inches(1.15), fill=AMBERBG, round_=True, radius=0.08)
text(s, Inches(0.8), Inches(2.25), Inches(11.8), Inches(0.8),
     [[('IT-Security requirement:  ', 15, AMBER, True),
       ('no CAD files, engineering drawings or images may leave the company network. Ever.', 15, DARK, True)]],
     line_spacing=1.15)
box(s, Inches(0.45), Inches(3.5), Inches(12.45), Inches(2.0), fill=GREENBG, round_=True, radius=0.08)
text(s, Inches(0.8), Inches(3.72), Inches(11.8), Inches(1.6),
     [[('The key finding from the code audit:  ', 15, GREEN, True),
       ('CAD files already never leave the server.', 15, DARK, True)],
      [('The geometry engine runs inside our backend and processes CAD models in memory — files are not even '
        'written to disk, let alone sent anywhere. The only outbound flows are the AI layer (photos and derived '
        'summaries) and two optional public feeds. The controls that close this gap are now BUILT: private AI '
        'routing and a provable air-gapped switch, tested live in all three modes.', 13, BODY, False)]],
     space_after=8, line_spacing=1.2)
text(s, Inches(0.45), Inches(5.85), Inches(12.4), Inches(0.9),
     [[('Why you can trust this: ', 12.5, DARK, True),
       ('this is not a vendor claim — we audited every outbound network call in the platform\'s source code, '
        'line by line, and listed each one. The full inventory is in the written plan.', 12.5, BODY, False)]],
     line_spacing=1.2)
notes(s, "The requirement is absolute: CAD data must never leave our network. Here is the finding that makes this "
         "project straightforward: it already doesn't. We audited every outbound connection in the source code. "
         "CAD models are processed in memory, inside our server, by the local geometry engine — they are not even "
         "saved to disk. The only traffic that leaves today is the AI layer — photographs and derived summaries — "
         "plus two optional public feeds. And here is the update since the written plan: the controls that close "
         "this gap are already built. Every AI call now goes through one switchboard that can be pointed at a "
         "private endpoint or switched off entirely, and we tested all three modes live. So the job is not a "
         "redesign — the software side is done; what remains is infrastructure.")

# ═══════════════ 6 — DATA-FLOW MAP ═══════════════
s = header('Where data flows today — verified in the code', 'Data-flow map')
box(s, Inches(0.45), Inches(2.0), Inches(7.5), Inches(4.7), fill=GREENBG, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(2.2), Inches(6.9), Inches(0.4),
     [[('STAYS INSIDE — already ✓', 14, GREEN, True)]])
inside = [
    ('CAD / geometry processing', 'OCCT engine runs in the backend; files processed in memory, never stored'),
    ('All 18 cost engines', 'Physics + maths, fully local'),
    ('Knowledge base & learning loop', 'Similarity, calibration, autonomous findings — local database'),
    ('BOM file parsing, exports, reports', 'Local parsers, local PDF/Excel generation'),
]
for i, (a, b) in enumerate(inside):
    y = Inches(2.7 + i * 0.95)
    box(s, Inches(0.75), y, Inches(0.09), Inches(0.8), fill=GREEN)
    text(s, Inches(1.0), y, Inches(6.7), Inches(0.9),
         [[(a, 12.5, DARK, True)], [(b, 10.5, BODY, False)]], space_after=2, line_spacing=1.05)
box(s, Inches(8.3), Inches(2.0), Inches(4.6), Inches(4.7), fill=AMBERBG, round_=True, radius=0.05)
text(s, Inches(8.6), Inches(2.2), Inches(4.0), Inches(0.4),
     [[('LEAVES TODAY — controlled ✓', 14, AMBER, True)]])
outside = [
    ('AI layer calls', 'Photos + CAD-derived summaries. Private routing BUILT — one setting points every call at our endpoint', AMBER),
    ('Live part pricing (optional, OFF by default)', 'Part-number text only — silenced by the air-gap switch', MUTED),
    ('News ticker feeds', 'Public news, no company data — silenced by the air-gap switch', MUTED),
]
for i, (a, b, c) in enumerate(outside):
    y = Inches(2.7 + i * 1.28)
    box(s, Inches(8.6), y, Inches(0.09), Inches(1.1), fill=c)
    text(s, Inches(8.85), y, Inches(3.9), Inches(1.25),
         [[(a, 12, DARK, True)], [(b, 10.5, BODY, False)]], space_after=2, line_spacing=1.05)
notes(s, "The whole security story on one slide. Green, left: what already stays inside — CAD processing, all cost "
         "engines, the learning loop, reports. That is the overwhelming majority of the platform. Amber, right: the "
         "three flows that go out — and all three are now under our control in the shipped code. The AI layer is the "
         "one that matters — photos and derived summaries — and its private routing is built: one setting points "
         "every call at our own endpoint. The other two are optional conveniences, and the new air-gapped switch "
         "silences them along with everything else, which we verified live.")

# ═══════════════ 7 — DEPLOYMENT OPTIONS ═══════════════
s = header('Deployment options — our recommendation', 'Decision')
cols = [
    ('OPTION A', 'Fully air-gapped', 'Everything on our VMs. AI features off via the built-in AIR-GAPPED switch (already implemented and tested). Deterministic engines fully working.',
     'Zero external connections — provable today, switch is built', 'Reduced AI capability', PANEL, MUTED),
    ('OPTION B  ★ RECOMMENDED', 'Private AI, on-prem core', 'Platform on our VMs. AI calls go to Claude running in OUR OWN cloud tenancy over a private link — no public internet, no data retention.',
     'Full capability + data control — routing already BUILT (one setting)', 'Requires cloud tenancy sign-off', PANEL2, BLUE),
    ('OPTION C', 'Public AI API', 'Platform on-prem but AI calls to the public API.',
     'Simplest', 'Fails our requirement — photos would traverse a public API', PANEL, RED),
]
for i, (tag, t, d, pro, con, fill, c) in enumerate(cols):
    x = Inches(0.45 + i * 4.25)
    box(s, x, Inches(2.05), Inches(4.0), Inches(4.5), fill=fill, round_=True, radius=0.06)
    box(s, x, Inches(2.05), Inches(4.0), Inches(0.09), fill=c)
    text(s, x + Inches(0.25), Inches(2.25), Inches(3.55), Inches(0.35), [[(tag, 12, c, True)]])
    text(s, x + Inches(0.25), Inches(2.62), Inches(3.55), Inches(0.4), [[(t, 16, DARK, True)]])
    text(s, x + Inches(0.25), Inches(3.15), Inches(3.55), Inches(1.5), [[(d, 11.5, BODY, False)]], line_spacing=1.15)
    text(s, x + Inches(0.25), Inches(4.85), Inches(3.55), Inches(0.75),
         [[('+ ', 12, GREEN, True), (pro, 11, BODY, False)]], line_spacing=1.1)
    text(s, x + Inches(0.25), Inches(5.7), Inches(3.55), Inches(0.75),
         [[('– ', 12, RED, True), (con, 11, BODY, False)]], line_spacing=1.1)
text(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.5),
     [[('Recommendation: Option B — the software for BOTH A and B is already built; what remains is infrastructure sign-off.', 13, DARK, True)]])
notes(s, "Three ways to deploy. Option A: fully air-gapped — everything on our machines, zero external connections, "
         "and the switch that enforces it is already implemented and tested; but the vision AI features are reduced. "
         "Option C is the public API — I include it only to reject it, because photos would traverse a public "
         "endpoint. Option B is our recommendation: the platform runs on-premise, and the AI model runs in OUR OWN "
         "cloud tenancy, reached over a private link with a no-data-retention configuration. The routing for this is "
         "also already built — a single setting points every AI call at our endpoint; we demonstrated it live. "
         "So the software work for both A and B is done; the remaining effort is infrastructure and sign-off. "
         "And note — the two options combine: the most sensitive programmes can run air-gapped while the rest use B.")

# ═══════════════ 8 — TARGET ARCHITECTURE ═══════════════
s = header('Target architecture — everything inside our walls', 'Architecture')
box(s, Inches(0.45), Inches(1.95), Inches(9.2), Inches(4.95), fill=PANEL, round_=True, radius=0.04)
text(s, Inches(0.7), Inches(2.05), Inches(8.5), Inches(0.3), [[('COMPANY INTERNAL NETWORK', 11, MUTED, True)]])
flow_box(s, Inches(0.85), Inches(2.45), Inches(4.0), Inches(0.75), 'Engineers (browser)  +  CAPEE', 'single sign-on (Azure AD)', BLUE, fill=BG)
flow_box(s, Inches(5.15), Inches(2.45), Inches(4.2), Inches(0.75), 'Corporate API Gateway', 'authentication · rate limits · audit logs', INDIGO, fill=BG)
down_arrow(s, Inches(2.7), Inches(3.28))
down_arrow(s, Inches(7.1), Inches(3.28))
flow_box(s, Inches(0.85), Inches(3.68), Inches(8.5), Inches(1.05), 'CostVision server (frontend + backend on our VM)',
         'geometry engine (CAD, in-memory) · 18 cost engines · learning loop · report generation', VIOLET, fill=BG)
down_arrow(s, Inches(2.7), Inches(4.82))
down_arrow(s, Inches(7.1), Inches(4.82))
flow_box(s, Inches(0.85), Inches(5.22), Inches(4.0), Inches(0.85), 'Corporate database (PostgreSQL)', 'rates · knowledge base · encrypted at rest', CYAN, fill=BG)
flow_box(s, Inches(5.15), Inches(5.22), Inches(4.2), Inches(0.85), 'Internal AI gateway', 'ONLY allowed exit · inspected · logged', AMBER, fill=BG)
box(s, Inches(10.0), Inches(4.9), Inches(2.9), Inches(2.0), fill=PANEL2, round_=True, radius=0.08)
text(s, Inches(10.2), Inches(5.05), Inches(2.5), Inches(1.8),
     [[('Our cloud tenancy', 12.5, BLUE, True)],
      [('Claude on AWS Bedrock / Google Vertex', 10.5, BODY, False)],
      [('private link · no public internet · zero data retention', 10, MUTED, False)]],
     space_after=4, line_spacing=1.1)
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.42), Inches(5.5), Inches(0.55), Inches(0.3))
arr.fill.solid(); arr.fill.fore_color.rgb = AMBER; arr.line.fill.background(); arr.shadow.inherit = False
text(s, Inches(10.0), Inches(2.45), Inches(2.9), Inches(2.2),
     [[('Zero-trust rules', 12.5, DARK, True)],
      [('• Every request authenticated (Azure AD)', 10.5, BODY, False)],
      [('• Deny-all egress except the AI gateway', 10.5, BODY, False)],
      [('• Role-based access (engineer / lead / admin / CAPEE service)', 10.5, BODY, False)],
      [('• CAD pulled from PLM vault, never stored', 10.5, BODY, False)]],
     space_after=4, line_spacing=1.1)
notes(s, "The target architecture. Everything in the grey box is inside our network: engineers and CAPEE come through "
         "the corporate gateway with single sign-on; the CostVision server — frontend and backend — runs on our "
         "virtual machine, with the CAD geometry engine inside it; the database is our standard encrypted PostgreSQL. "
         "The only permitted exit is the internal AI gateway at the bottom — one door, inspected and logged — which "
         "connects over a private link to the AI model running in our own cloud tenancy with zero data retention. "
         "On the right, the zero-trust rules: authenticate everything, deny all other egress, role-based access, and "
         "CAD comes from the PLM vault and is never stored. If presenting live, a simple Fade animation on each layer "
         "of this diagram works well.")

# ═══════════════ 9 — CAD PROTECTION ═══════════════
s = header('How CAD data is protected', 'CAD data protection')
rows = [
    ('CAD model opened & measured', 'Inside our backend (OCCT engine)', 'Already local — verified'),
    ('Feature detection (holes, threads, walls)', 'Inside our backend', 'Already local — verified'),
    ('BOM extraction from files', 'Inside our backend (local parsers)', 'Already local — verified'),
    ('Cost calculation & learning', 'Inside our backend + our database', 'Already local — verified'),
    ('Photo analysis & AI narrative', 'Private AI endpoint (Option B)', 'Routing BUILT — set one value'),
    ('CAD file storage', 'Nowhere — processed in memory only', 'Files never written to disk'),
]
for i, (a, b, c) in enumerate(rows):
    y = Inches(2.1 + i * 0.73)
    bgc = PANEL if i % 2 == 0 else BG
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.62), fill=bgc, round_=True, radius=0.14)
    text(s, Inches(0.8), y + Inches(0.12), Inches(4.9), Inches(0.4), [[(a, 12.5, DARK, True)]])
    text(s, Inches(5.8), y + Inches(0.12), Inches(3.8), Inches(0.4), [[(b, 12, BODY, False)]])
    good = 'change' not in c and 'never' not in c
    col = GREEN if 'verified' in c or 'never' in c else AMBER
    text(s, Inches(9.7), y + Inches(0.12), Inches(3.0), Inches(0.4), [[(('✓  ' if col == GREEN else '→  ') + c, 11.5, col, True)]])
text(s, Inches(0.45), Inches(6.6), Inches(12.4), Inches(0.6),
     [[('Plus: logs never contain CAD payloads; metadata goes to the corporate SIEM; and the "air-gapped" switch is '
        'now BUILT & tested — security can PROVE zero egress in a witnessed firewall test.', 12, MUTED, False, True)]], line_spacing=1.15)
notes(s, "Function by function: where does CAD data actually go? Opening and measuring the model, detecting features, "
         "extracting BOMs, calculating cost — all of it already happens inside our backend, verified in code. "
         "The single amber row is photo analysis and AI narrative — and the routing for it is now built: one "
         "configuration value points it at the private endpoint. And note the last row — CostVision never stores CAD "
         "files at all; they are processed in memory and released. Finally, the air-gapped switch is no longer a "
         "promise — it is built and tested, so security can prove zero egress in a witnessed firewall test rather "
         "than take our word for it.")

# ═══════════════ 10 — CAPEE INTEGRATION ═══════════════
s = header('CAPEE + CostVision — better together', 'Integration')
box(s, Inches(0.45), Inches(2.0), Inches(5.6), Inches(2.1), fill=PANEL, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(2.2), Inches(5.0), Inches(1.8),
     [[('CAPEE  (stays the front door)', 15, DARK, True)],
      [('• Costing workflow, approvals, reporting', 12, BODY, False)],
      [('• System of record — unchanged for users', 12, BODY, False)],
      [('• Calls CostVision services behind the scenes', 12, BODY, False)]],
     space_after=5, line_spacing=1.12)
box(s, Inches(7.3), Inches(2.0), Inches(5.6), Inches(2.1), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(7.6), Inches(2.2), Inches(5.0), Inches(1.8),
     [[('CostVision  (the engine behind it)', 15, BLUE, True)],
      [('• 18 physics cost engines + CAD reading', 12, BODY, False)],
      [('• AI memory: similar parts, self-calibration', 12, BODY, False)],
      [('• Autonomous findings for sourcing', 12, BODY, False)]],
     space_after=5, line_spacing=1.12)
a1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.15), Inches(2.45), Inches(1.05), Inches(0.4))
a1.fill.solid(); a1.fill.fore_color.rgb = BLUE; a1.line.fill.background(); a1.shadow.inherit = False
a2 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(6.15), Inches(3.15), Inches(1.05), Inches(0.4))
a2.fill.solid(); a2.fill.fore_color.rgb = GREEN; a2.line.fill.background(); a2.shadow.inherit = False
text(s, Inches(5.95), Inches(2.1), Inches(1.5), Inches(0.3), [[('requests', 9.5, BLUE, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(5.95), Inches(3.6), Inches(1.5), Inches(0.3), [[('AI answers', 9.5, GREEN, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.45), Inches(4.5), Inches(12.45), Inches(2.3), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(4.7), Inches(11.9), Inches(0.35), [[('What flows over the connection (internal APIs that already exist)', 13.5, DARK, True)]])
flows = [
    ('Cost a part →', 'full should-cost with breakdown'),
    ('CAD file →', 'geometry + auto-filled inputs'),
    ('New part →', 'similar past parts + suggestions'),
    ('PO price →', 'feeds the learning loop automatically'),
    ('Dashboard ←', 'autonomous findings (£/yr)'),
    ('One shared →', 'rate library & knowledge database'),
]
for i, (a, b) in enumerate(flows):
    x = Inches(0.75 + (i % 3) * 4.1); y = Inches(5.2 + (i // 3) * 0.75)
    text(s, x, y, Inches(3.9), Inches(0.7),
         [[(a + ' ', 12, BLUE, True), (b, 11.5, BODY, False)]], line_spacing=1.1)
notes(s, "The integration philosophy: CAPEE stays the front door — the workflow, approvals and reporting our teams "
         "already know. CostVision becomes the engine behind it, called over internal APIs that already exist. "
         "The six flows that matter: CAPEE sends a part, gets a full should-cost back. It sends a CAD file, gets "
         "auto-filled inputs. It asks about a new part, gets similar history and suggestions. And the best one: "
         "every PO price CAPEE already captures feeds the learning loop automatically — CAPEE's data makes "
         "CostVision smarter without anyone lifting a finger. One shared rate library and knowledge database means "
         "one version of the truth.")

# ═══════════════ 11 — WHAT EACH SIDE NEEDS ═══════════════
s = header('What each side needs to change', 'Scope of work')
box(s, Inches(0.45), Inches(2.05), Inches(6.0), Inches(4.55), fill=PANEL2, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(2.25), Inches(5.4), Inches(0.4), [[('CostVision — 6 items (2 DONE)', 15, BLUE, True)]])
cv = [
    ('Private AI routing', 'DONE — one setting routes every AI call to our endpoint', True),
    ('"Air-gapped" switch', 'DONE — provably disables all external calls; tested live', True),
    ('Azure AD sign-on', 'replace local logins with company SSO', False),
    ('PostgreSQL', 'move from file database to the corporate DB estate', False),
    ('Cost API wrapper', 'one clean endpoint per commodity for CAPEE', False),
    ('Audit hardening', 'admin audit table + logs to SIEM', False),
]
for i, (a, b, done) in enumerate(cv):
    y = Inches(2.75 + i * 0.63)
    mark = '✓ ' if done else f'{i+1}. '
    tc = GREEN if done else DARK
    dc = GREEN if done else BODY
    text(s, Inches(0.75), y, Inches(5.5), Inches(0.6),
         [[(mark + a + ' — ', 12, tc, True), (b, 11, dc, False)]], line_spacing=1.05)
box(s, Inches(6.85), Inches(2.05), Inches(6.05), Inches(4.55), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(7.15), Inches(2.25), Inches(5.4), Inches(0.4), [[('CAPEE — small, additive changes', 15, VIOLET, True)]])
cap = [
    ('Backend', 'HTTP client to call CostVision + token handling; hook PO prices into the learning API', 'Small'),
    ('User interface', 'an "AI insights" panel on the costing screen; a CAD-upload button; a findings widget', 'Medium'),
    ('Data', 'adopt the shared rate library; one-off import of historical quotes to seed the memory', 'Small–Medium'),
]
for i, (a, b, sz) in enumerate(cap):
    y = Inches(2.8 + i * 1.25)
    text(s, Inches(7.15), y, Inches(5.5), Inches(1.2),
         [[(a + '  ', 13, DARK, True), ('· ' + sz, 11, VIOLET, True)],
          [(b, 11.5, BODY, False)]], space_after=3, line_spacing=1.12)
text(s, Inches(0.45), Inches(6.8), Inches(12.4), Inches(0.45),
     [[('Items 1–2 are built & tested. Remaining critical path: ~3–5 engineering weeks. No architectural surgery on either side.', 13.5, DARK, True)]])
notes(s, "The scope of work, honestly sized. CostVision needed six bounded changes — and the first two, private AI "
         "routing and the air-gap switch, are already done: built, covered by automated tests, and demonstrated live "
         "in all three deployment modes. What remains on the CostVision side is company single sign-on, the corporate "
         "database, a clean API wrapper for CAPEE, and audit hardening. CAPEE's changes are additive, not invasive: "
         "a client in the backend, an AI-insights panel in the UI, and adopting the shared rate library — plus a "
         "one-off import of historical quotes so the memory starts smart. Remaining critical path: roughly three to "
         "five engineering weeks. Nothing here is architectural surgery.")

# ═══════════════ 12 — SECURITY & COMPLIANCE ═══════════════
s = header('Security & compliance — how we tick the boxes', 'Compliance')
comp = [
    ('Encryption', 'TLS everywhere in transit; encrypted database at rest; CAD never stored at all', GREEN),
    ('Access control', 'Azure AD single sign-on; roles mapped to AD groups; CAPEE uses a least-privilege service account', BLUE),
    ('Data-loss prevention', 'One AI exit door, inspected and logged; firewall denies everything else; witnessed air-gap test (switch built)', AMBER),
    ('Audit trails', 'Every rate change, knowledge write and admin action attributable to a user; logs to corporate SIEM', VIOLET),
]
for i, (t, d, c) in enumerate(comp):
    x = Inches(0.45 + (i % 2) * 6.35); y = Inches(2.05 + (i // 2) * 1.55)
    box(s, x, y, Inches(6.1), Inches(1.35), fill=PANEL, round_=True, radius=0.09)
    box(s, x, y, Inches(0.09), Inches(1.35), fill=c)
    text(s, x + Inches(0.28), y + Inches(0.13), Inches(5.6), Inches(0.4), [[(t, 14, c, True)]])
    text(s, x + Inches(0.28), y + Inches(0.55), Inches(5.6), Inches(0.75), [[(d, 11.5, BODY, False)]], line_spacing=1.12)
box(s, Inches(0.45), Inches(5.35), Inches(12.45), Inches(1.35), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(5.55), Inches(11.9), Inches(1.05),
     [[('Standards mapping:  ', 13, DARK, True),
       ('ISO 27001 — covered by the controls above, add to ISMS scope  ·  SOC 2 — inherited from the cloud tenancy '
        '(Bedrock/Vertex are certified)  ·  GDPR — minimal personal data, zero-retention AI, EU region  ·  '
        'ISO/SAE 21434 — engineering IT tool, out of vehicle scope; supply-chain review (SBOM) in CI.', 12, BODY, False)]],
     line_spacing=1.25)
notes(s, "For the compliance-minded: encryption everywhere, and the strongest control of all for CAD — it is never "
         "stored. Access is company single sign-on with role mapping. Data-loss prevention comes down to one "
         "inspected exit door and a firewall that denies everything else — and security can witness the air-gap test "
         "themselves. Full audit trails to our SIEM. On standards: ISO 27001 is covered by these controls; SOC 2 is "
         "inherited from the certified cloud services; GDPR exposure is minimal with zero-retention AI in the EU "
         "region; and for automotive cyber, this is an engineering IT tool outside the vehicle scope, with "
         "supply-chain checks built into the build pipeline.")

# ═══════════════ 13 — ROLLOUT TIMELINE ═══════════════
s = header('Rollout — six phases, value from week eight', 'Plan')
phases = [
    ('1', 'Security assessment', '2–3 wks', BLUE),
    ('2', 'Architecture design', '3–4 wks', INDIGO),
    ('3', 'Pilot (dummy CAD)', '2 wks', CYAN),
    ('4', 'CAPEE pilot', '4–6 wks', VIOLET),
    ('5', 'Enterprise rollout', '4 wks', GREEN),
    ('6', 'Monitor & govern', 'ongoing', MUTED),
]
cw = Inches(2.24)
for i, (n, t, d, c) in enumerate(phases):
    shp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.4) + int(cw * 0.86) * i, Inches(2.3), cw, Inches(1.1))
    shp.adjustments[0] = 0.28
    shp.fill.solid(); shp.fill.fore_color.rgb = c
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = f'{n} · {t}'
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = BG; r.font.name = 'Calibri'
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = d
    r2.font.size = Pt(9.5); r2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xFF); r2.font.name = 'Calibri'
gates = [
    ('Phase 3 exit gate', 'Security-witnessed test: firewall blocks all egress, full test suite passes, zero unexpected traffic.'),
    ('Phase 4 exit gate', 'A cost engineer completes a real costing from inside CAPEE with AI insights; accuracy dashboard live.'),
    ('First value', 'Week ~8: pilot users get similar-part suggestions; sourcing sees the first autonomous findings.'),
]
for i, (a, b) in enumerate(gates):
    y = Inches(3.9 + i * 0.95)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.82), fill=PANEL, round_=True, radius=0.12)
    box(s, Inches(0.45), y, Inches(0.09), Inches(0.82), fill=BLUE)
    text(s, Inches(0.75), y + Inches(0.1), Inches(2.6), Inches(0.4), [[(a, 12.5, BLUE, True)]])
    text(s, Inches(3.5), y + Inches(0.1), Inches(9.2), Inches(0.65), [[(b, 11.5, BODY, False)]], line_spacing=1.1)
text(s, Inches(0.45), Inches(6.82), Inches(12.4), Inches(0.5),
     [[('After go-live we track: ', 11.5, DARK, True),
       ('accuracy vs PO prices · costings via CAPEE · £ findings actioned · adoption.   Approval today → pilot value by ~September 2026.', 11.5, BODY, False)]])
notes(s, "Six phases. The first two — assessment and design — are what we're asking approval for today. Phase three "
         "is a sealed pilot with dummy CAD data, and its exit gate is a security-witnessed test: firewall blocking "
         "everything, full test suite passing, zero unexpected traffic. Phase four connects CAPEE, with a real "
         "costing done end-to-end by a cost engineer as the gate. Rollout and then steady-state governance follow. "
         "First tangible value lands around week eight, when pilot users start getting similar-part suggestions "
         "and sourcing sees the first autonomous findings.")

# ═══════════════ 14 — RISKS ═══════════════
s = header('Risks — and how we manage them', 'Honest view')
risks = [
    ('Cloud tenancy not approved', 'Medium', 'Fall back to Option A (air-gapped) — deterministic engines keep full value today; add a self-hosted vision model later if needed.', AMBER),
    ('Learning depends on logged quotes', 'Medium', 'The CAPEE PO-price hook automates it — quotes flow in without anyone changing habits. Plus a one-off historical import to start smart.', BLUE),
    ('Single-team knowledge of the platform', 'Low–Med', '1,005 automated tests, written architecture docs, and a named CAPEE-side maintainer trained during Phase 4.', VIOLET),
    ('Adoption ("another tool")', 'Low', 'Users stay in CAPEE — CostVision works behind the scenes. Nothing new to learn except better answers appearing.', GREEN),
]
for i, (t, sev, m, c) in enumerate(risks):
    y = Inches(2.05 + i * 1.18)
    box(s, Inches(0.45), y, Inches(12.45), Inches(1.02), fill=PANEL, round_=True, radius=0.09)
    box(s, Inches(0.45), y, Inches(0.09), Inches(1.02), fill=c)
    text(s, Inches(0.75), y + Inches(0.12), Inches(4.4), Inches(0.45), [[(t, 13, DARK, True)]])
    text(s, Inches(0.75), y + Inches(0.55), Inches(2.0), Inches(0.35), [[('Likelihood: ' + sev, 10.5, c, True)]])
    text(s, Inches(5.3), y + Inches(0.12), Inches(7.4), Inches(0.85),
         [[('Mitigation:  ', 11.5, DARK, True), (m, 11.5, BODY, False)]], line_spacing=1.12)
notes(s, "The honest risk view. If the cloud tenancy isn't approved, we don't stall — we fall back to the fully "
         "air-gapped option, which still delivers the deterministic engines and the learning loop. The learning "
         "depends on real quotes being logged — but the CAPEE hook automates that, and a historical import means we "
         "don't start from zero. Platform knowledge is protected by the test suite, documentation and a trained "
         "CAPEE-side maintainer. And adoption risk is low precisely because users stay in CAPEE — for them, the "
         "change is simply that better answers start appearing.")

# ═══════════════ 15 — VERDICT & ASK ═══════════════
s = header('Feasibility verdict — and the ask', 'Decision')
verdicts = [
    ('Secure deployment, CAD fully internal', 'FEASIBLE — CAD never leaves; AI routing + air-gap switch already BUILT', GREEN),
    ('CAPEE + CostVision integration', 'FEASIBLE — API-first design; CAPEE stays the front door', GREEN),
    ('Changes required', 'Head start delivered: 2 of 6 items built — ~3–5 engineering weeks remain', BLUE),
    ('Long-term governance', 'Rate-library board · monthly findings review · quarterly access & egress audit', VIOLET),
]
for i, (a, b, c) in enumerate(verdicts):
    y = Inches(2.05 + i * 0.82)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.7), fill=(PANEL if i % 2 == 0 else BG), round_=True, radius=0.13)
    text(s, Inches(0.8), y + Inches(0.15), Inches(4.6), Inches(0.4), [[(a, 13, DARK, True)]])
    text(s, Inches(5.5), y + Inches(0.15), Inches(7.2), Inches(0.4), [[(b, 12.5, c, True)]])
box(s, Inches(0.45), Inches(5.6), Inches(12.45), Inches(1.25), fill=PANEL2, round_=True, radius=0.07)
text(s, Inches(0.8), Inches(5.8), Inches(11.8), Inches(0.9),
     [[('The ask today:  ', 16, BLUE, True),
       ('approve Phase 1 (IT-Security assessment) and Phase 2 (architecture design) — 5–7 weeks, existing teams. '
        'The full written plan and security checklist are ready for the security review.', 14, BODY, False)]],
     line_spacing=1.2)
box(s, 0, H - Inches(0.16), W, Inches(0.16), fill=INDIGO)
notes(s, "To summarise: secure deployment with CAD fully internal is feasible — in fact, the platform was built that "
         "way, and the two controls that complete it — private AI routing and the air-gapped switch — are already "
         "built, tested and demonstrated live. CAPEE integration is feasible and high-value, with CAPEE remaining "
         "the front door. The remaining changes are bounded — roughly three to five weeks, not months — and "
         "governance is defined. The ask today is simple: approve Phases 1 and 2, the security assessment and the "
         "architecture design. The full written plan, including the security checklist, is ready to hand to the "
         "IT-Security team. Thank you — questions welcome.")

# ═══════════════ 16 — BACKUP: MARKET LANDSCAPE ═══════════════
s = header('Has anyone done this? Yes — the market is real', 'Backup · Market landscape')
market = [
    ('aPriori', 'Market leader (US)', 'CAD-to-cost with physics-based models and regional cost databases. Used by Fortune-500 OEMs and Tier-1 suppliers — companies our size already pay for this category.', BLUE, 'apriori.com'),
    ('Siemens PCM', 'PLM giant (Teamcenter)', 'Bottom-up should-costing with process models and supplier collaboration, embedded in the Siemens PLM suite.', INDIGO, 'siemens.com/teamcenter'),
    ('Tset', 'AI challenger (Austria, 2018)', 'Automotive-focused costing from 3D models and BOMs, 50+ calculation modules plus CO2. Being acquired by A2MAC1 (June 2026) to build "AI-enabled costing intelligence".', VIOLET, 'tset.com · globenewswire.com (18 Jun 2026)'),
    ('Boothroyd DFMA', 'The classic (Dewhurst)', 'Design-for-manufacture and concurrent costing — the methodology textbooks are built on, used for decades.', CYAN, 'dfma.com'),
    ('New AI entrants', 'Start-ups', 'Razorlabs Cost Advisor, Emithran and others: AI quotes from CAD files, calibrated on live machining data.', MUTED, 'emithran.com'),
]
for i, (t, tag, d, c, src) in enumerate(market):
    y = Inches(1.98 + i * 0.87)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.78), fill=PANEL, round_=True, radius=0.14)
    box(s, Inches(0.45), y, Inches(0.09), Inches(0.78), fill=c)
    text(s, Inches(0.75), y + Inches(0.08), Inches(2.55), Inches(0.4), [[(t, 13, c, True)]])
    text(s, Inches(0.75), y + Inches(0.44), Inches(2.55), Inches(0.3), [[(tag, 9, MUTED, False, True)]])
    text(s, Inches(3.45), y + Inches(0.08), Inches(9.2), Inches(0.45), [[(d, 10.5, BODY, False)]], line_spacing=1.05)
    text(s, Inches(3.45), y + Inches(0.55), Inches(9.2), Inches(0.22), [[('Source: ' + src, 8.5, MUTED, False, True)]])
box(s, Inches(0.45), Inches(6.45), Inches(12.45), Inches(0.75), fill=PANEL2, round_=True, radius=0.1)
text(s, Inches(0.75), Inches(6.58), Inches(11.9), Inches(0.55),
     [[('Takeaway: ', 12, BLUE, True),
       ('the should-cost category is proven — large OEMs already pay for it. The question is not IF it works, '
        'but why OURS fits us better: security, learning on our own data, and CAPEE integration (next slide).', 12, BODY, False)]],
     line_spacing=1.12)
notes(s, "Backup slide, for the question 'has anyone done this before — are we the only ones?' The answer is "
         "reassuring in both directions. No, we are not inventing an unproven category: should-cost software is an "
         "established market. aPriori is the leader — CAD-to-cost with physics models, sold to exactly our kind of "
         "company. Siemens ships it inside the PLM suite. Tset is the automotive AI challenger — and notably A2MAC1 "
         "agreed to acquire them in June 2026 specifically to build AI-enabled costing intelligence, which tells you "
         "the whole market believes in the direction we've already built. Boothroyd Dewhurst is the classic "
         "methodology, and a wave of AI start-ups is doing CAD-to-quote. So the category is validated by the market. "
         "The real question is why our own tool rather than buying one — and that's the next slide.")

# ═══════════════ 17 — BACKUP: WHERE COSTVISION STANDS APART ═══════════════
s = header('Why ours — what no vendor offers today', 'Backup · Differentiation')
diffs = [
    ('PCB photo → costed BOM', 'Photograph a circuit board, get a costed bill of materials with live component pricing. Design tools generate BOMs; no mainstream costing suite costs a board from photos.', GREEN),
    ('Self-learning on OUR data', 'Calibrates itself on our real PO prices (error 10.9% → 0.3% after 3 quotes in testing) and raises findings autonomously. Incumbents are still marketing toward this.', GREEN),
    ('Runs inside our walls', 'Commercial tools are cloud SaaS — our CAD would live in their cloud. CostVision is on-prem with private AI routing and a provable air-gapped switch, both already built.', GREEN),
    ('One platform, no module licences', '18 commodity engines + automotive software costing + carbon + RFQ generation in one codebase. Vendors sell comparable breadth as separately licensed modules.', GREEN),
]
for i, (t, d, c) in enumerate(diffs):
    x = Inches(0.45 + (i % 2) * 6.35); y = Inches(2.0 + (i // 2) * 1.62)
    box(s, x, y, Inches(6.1), Inches(1.45), fill=GREENBG, round_=True, radius=0.09)
    box(s, x, y, Inches(0.09), Inches(1.45), fill=c)
    text(s, x + Inches(0.28), y + Inches(0.12), Inches(5.6), Inches(0.35), [[(t, 13.5, GREEN, True)]])
    text(s, x + Inches(0.28), y + Inches(0.5), Inches(5.6), Inches(0.9), [[(d, 10.5, BODY, False)]], line_spacing=1.1)
box(s, Inches(0.45), Inches(5.4), Inches(12.45), Inches(1.0), fill=AMBERBG, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(5.53), Inches(11.9), Inches(0.8),
     [[('Honest caveat: ', 12, AMBER, True),
       ('vendors like aPriori have decades of curated global cost data behind their numbers. Our answer is the '
        'learning loop: every quote CAPEE logs turns OUR history into a moat no vendor can buy — and our rates stay confidential.', 11.5, BODY, False)]],
     line_spacing=1.15)
text(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.35),
     [[('Positioning: not "better than aPriori" — an internal, secure, self-learning costing engine behind CAPEE that gets smarter on our own data.', 12, DARK, True)]])
text(s, Inches(0.45), Inches(6.95), Inches(12.4), Inches(0.25),
     [[('Sources: apriori.com · siemens.com/teamcenter · tset.com · globenewswire.com (A2MAC1-Tset, 18 Jun 2026) · dfma.com · emithran.com · flux.ai · circuitmind.io', 8.5, MUTED, False, True)]])
notes(s, "The second half of the backup answer: given the market exists, why build our own? Four things no vendor "
         "offers today. One — photo-to-cost for PCBs: photograph a supplier's or competitor's board and get a costed "
         "bill of materials; design tools generate BOMs during design, but no costing suite does this from photos. "
         "Two — the self-learning loop: CostVision calibrates on our own purchase-order prices, cutting error from "
         "about eleven percent to under one percent in testing after just three quotes, and raises findings on its "
         "own; the incumbents are only now marketing toward AI-enabled costing. Three — deployment: the commercial "
         "tools are cloud SaaS, meaning our CAD and our rates would sit in a vendor's cloud; CostVision runs inside "
         "our network with the private-AI and air-gap controls already built. Four — breadth without per-module "
         "licences. And the honest caveat management should hear: aPriori's strength is decades of curated cost "
         "data. Our counter is that every quote CAPEE logs makes CostVision smarter on OUR parts, OUR suppliers, "
         "OUR regions — data no vendor has. The positioning is not 'better than aPriori'; it is an internal, "
         "secure, self-learning engine behind CAPEE.")

# ═══════════════ 18 — BACKUP: COMPANY RATE DATA UPLOAD (PROOF) ═══════════════
s = header('Your rates, not vendor rates — already built', 'Backup · Company cost database')
pic = s.shapes.add_picture('docs/rate-library-upload-proof.png', Inches(0.45), Inches(2.0), width=Inches(6.9))
pic.line.color.rgb = LINE; pic.line.width = Pt(1)
text(s, Inches(0.45), Inches(6.72), Inches(6.9), Inches(0.5),
     [[('Live screenshot (July 2026): Rate Library screen after a company workbook upload — badge shows '
        '"Company rates active".', 9.5, MUTED, False, True)]], line_spacing=1.1)
text(s, Inches(7.6), Inches(2.0), Inches(5.3), Inches(2.4),
     [[('How admins load company data', 14, BLUE, True)],
      [('1.  Download the Excel template — six sheets: Materials · Machines · Labour · Energy · FX · Overhead.', 11.5, BODY, False)],
      [('2.  Fill in our rates and upload — every row is validated; the file becomes the active library instantly.', 11.5, BODY, False)],
      [('3.  Fine-tune any single cell in the tables — each change is logged with the user\'s name and a timestamp.', 11.5, BODY, False)]],
     space_after=7, line_spacing=1.15)
box(s, Inches(7.6), Inches(4.35), Inches(5.3), Inches(1.15), fill=GREENBG, round_=True, radius=0.1)
text(s, Inches(7.85), Inches(4.5), Inches(4.85), Inches(0.9),
     [[('Proven live today:  ', 11.5, GREEN, True),
       ('a full workbook uploaded and accepted — 328 materials, 155 machines, 42 labour, 11 energy, 9 FX and '
        '23 overhead rows — then activated.', 11, BODY, False)]], line_spacing=1.15)
text(s, Inches(7.6), Inches(5.7), Inches(5.3), Inches(1.4),
     [[('Also built in:', 12, DARK, True)],
      [('• 20 country rate sets, 8 labour categories each', 11, BODY, False)],
      [('• PCB country cost table — admin-editable', 11, BODY, False)],
      [('• Separate rate library for software costing', 11, BODY, False)],
      [('• One click back to built-in defaults, full audit trail', 11, BODY, False)]],
     space_after=3, line_spacing=1.12)
notes(s, "This backup slide answers 'whose numbers are these?' The answer: ours, whenever we want them to be. "
         "What you see is a live screenshot of the tool, not a mock-up. An administrator downloads an Excel "
         "template with six sheets — materials, machines, labour, energy, exchange rates and overheads — fills in "
         "our company rates, and uploads it. The file is validated row by row and becomes the active rate library "
         "instantly; the green badge confirms company rates are in force. After that, any single cell can be "
         "fine-tuned in the on-screen tables, and every change is recorded with the user's name and a timestamp — "
         "a full audit trail. We proved this live: a complete workbook with over five hundred rows across the six "
         "sheets uploaded, validated and activated. On top of that, the tool ships twenty country rate sets with "
         "eight labour categories each, the PCB country cost table is editable the same way, and the software-"
         "costing engine has its own separate library. And if anything goes wrong, one click returns to the "
         "built-in defaults. So management should hear this clearly: the tool calculates on OUR labour, material, "
         "machine and energy rates, by country — vendor tools can't offer that level of transparency and control.")

# ═══════════════ 19 — BACKUP: BUSINESS CASE ═══════════════
s = header('What it costs vs what it returns', 'Backup · Business case')
box(s, Inches(0.45), Inches(2.0), Inches(6.0), Inches(3.5), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(2.18), Inches(5.4), Inches(0.4), [[('What it costs', 15, DARK, True)]])
costs = [
    ('Engineering', '~3–5 weeks remaining on CostVision + small additive CAPEE changes — existing teams'),
    ('Infrastructure', 'one VM + corporate PostgreSQL — standard IT estate our teams already run'),
    ('AI usage', 'pay-per-use in our own cloud tenancy — no per-seat fees, off in air-gap mode'),
    ('Licence spend', '£0 — built in-house, nothing to buy or renew'),
]
for i, (a, b) in enumerate(costs):
    y = Inches(2.66 + i * 0.72)
    text(s, Inches(0.75), y, Inches(5.5), Inches(0.7),
         [[(a + ' — ', 12, DARK, True), (b, 11, BODY, False)]], line_spacing=1.1)
box(s, Inches(6.85), Inches(2.0), Inches(6.05), Inches(3.5), fill=GREENBG, round_=True, radius=0.05)
text(s, Inches(7.15), Inches(2.18), Inches(5.4), Inches(0.4), [[('What it returns', 15, GREEN, True)]])
gains = [
    ('Savings found', '£512k/yr of pricing issues surfaced autonomously in the live demo (indicative)'),
    ('Licence avoided', 'commercial should-cost suites run six figures per year, recurring'),
    ('Speed', 'first defensible should-cost in minutes, not days — every quote challengeable'),
    ('Leverage', 'bottom-up numbers on OUR rates that suppliers cannot wave away'),
]
for i, (a, b) in enumerate(gains):
    y = Inches(2.66 + i * 0.72)
    text(s, Inches(7.15), y, Inches(5.5), Inches(0.7),
         [[(a + ' — ', 12, GREEN, True), (b, 11, BODY, False)]], line_spacing=1.1)
ip = box(s, Inches(0.45), Inches(5.68), Inches(12.45), Inches(0.62), fill=GREENBG, round_=True, radius=0.12)
box(s, Inches(0.45), Inches(5.68), Inches(0.09), Inches(0.62), fill=GREEN)
text(s, Inches(0.8), Inches(5.83), Inches(11.9), Inches(0.4),
     [[('OUR IP:  ', 13.5, GREEN, True),
       ('built in-house — the code, the rate library and the knowledge base are entirely our intellectual property. '
        'No vendor owns any part of it.', 13.5, DARK, True)]])
box(s, Inches(0.45), Inches(6.44), Inches(12.45), Inches(0.85), fill=AMBERBG, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(6.55), Inches(11.9), Inches(0.65),
     [[('Honest numbers: ', 12, AMBER, True),
       ('the £512k/yr figure comes from demonstration data — we treat it as indicative and validate it against real '
        'programmes during the Phase 3–4 pilot before claiming it in any business case.', 11.5, BODY, False)]],
     line_spacing=1.15)
notes(s, "The business case in one view. Costs, left: three to five engineering weeks remaining plus small additive "
         "changes in CAPEE, all with existing teams; infrastructure is one virtual machine and our standard "
         "corporate database; AI usage is pay-per-use in our own tenancy with no per-seat fees; and licence spend "
         "is zero — this was built in-house, so the code, the rate library and the knowledge base are our "
         "intellectual property. Returns, right: the autonomous agent surfaced half a million pounds a year of "
         "pricing findings in the demonstration; a commercial should-cost suite would cost six figures every year, "
         "recurring; quotes get a defensible counter-number in minutes; and negotiations start from bottom-up "
         "figures built on our own rates. And the amber bar is deliberate honesty: the savings figure is from demo "
         "data — we validate it in the pilot before it goes into any business case. That candour is what makes the "
         "rest of this slide believable.")

# ═══════════════ 20 — BACKUP: EVIDENCE PACK ═══════════════
s = header('Does it actually work? The measured results', 'Backup · Evidence')
cd = CategoryChartData()
cd.categories = ['Machining error %', 'Casting error %', 'Uncertainty band ±%']
cd.add_series('Before learning', (10.9, 8.7, 20.4))
cd.add_series('After 3 real quotes', (0.3, 0.6, 2.8))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.45), Inches(2.0), Inches(7.3), Inches(4.4), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(10)
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(9.5)
plot.data_labels.font.bold = True
ch.category_axis.tick_labels.font.size = Pt(9.5)
ch.value_axis.tick_labels.font.size = Pt(9)
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = MUTED
ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = GREEN
text(s, Inches(0.45), Inches(6.5), Inches(7.3), Inches(0.6),
     [[('Live test runs, July 2026 — after the engine learned from just 3 real supplier quotes per commodity.', 10, MUTED, False, True)]],
     line_spacing=1.1)
stats = [
    ('99%', 'part-match accuracy — the AI memory finds the right similar past part', BLUE),
    ('1,005', 'automated tests protect every engine and feature on each change', VIOLET),
    ('3 quotes', 'is all the calibration needs before accuracy lands under 1%', GREEN),
    ('£512k/yr', 'findings raised by the autonomous agent in the live demo', CYAN),
]
for i, (n, d, c) in enumerate(stats):
    y = Inches(2.0 + i * 1.23)
    box(s, Inches(8.1), y, Inches(4.8), Inches(1.08), fill=PANEL, round_=True, radius=0.1)
    box(s, Inches(8.1), y, Inches(0.09), Inches(1.08), fill=c)
    text(s, Inches(8.4), y + Inches(0.12), Inches(1.7), Inches(0.5), [[(n, 19, c, True)]])
    text(s, Inches(10.05), y + Inches(0.14), Inches(2.75), Inches(0.85), [[(d, 10, BODY, False)]], line_spacing=1.1)
notes(s, "The evidence slide, for when someone asks 'does it actually work?' The chart shows measured accuracy from "
         "live test runs. Grey bars are the engine before learning: machining estimates were off by about eleven "
         "percent, casting by about nine, and the honest uncertainty band was around plus-or-minus twenty percent. "
         "Green bars are after the engine learned from just three real supplier quotes per commodity: errors drop "
         "below one percent and the band tightens to under three. That is the self-learning loop working — and it "
         "is why CAPEE's PO prices matter so much, because they are exactly the fuel this loop runs on. On the "
         "right: the part-matching memory finds the right historical part ninety-nine percent of the time; "
         "seven hundred and ninety-three automated tests protect the platform on every change; and the autonomous "
         "agent found half a million pounds a year of pricing issues in the demo without anyone asking it to.")

# ═══════════════ 21 — BACKUP: LIKELY QUESTIONS ═══════════════
s = header('Questions you may be asking', 'Backup · Straight answers')
qa = [
    ('Why not just buy aPriori?', 'A SaaS suite means our CAD and rates live in a vendor cloud, it cannot learn from our '
     'quote history, and it costs six figures every year. CostVision is on-prem, self-learning, and our IP.'),
    ('Is our data training a public AI?', 'No. AI calls run zero-retention in our own cloud tenancy (Option B) — and in '
     'air-gapped mode there are provably no external calls at all. Nothing is ever used to train public models.'),
    ('What if the key developer leaves?', '1,005 automated tests define how everything must behave, the architecture is '
     'documented in writing, and Phase 4 trains a named CAPEE-side maintainer.'),
    ('What does it cost to run?', 'One VM, a standard corporate database, and pay-per-use AI in our tenancy. '
     'No licences, no per-seat fees.'),
    ('How do we know the numbers are right?', 'Physics-based cost build-ups on OUR uploaded rates, calibrated against OUR '
     'real PO prices — and the accuracy dashboard tracks the error openly, part by part.'),
]
for i, (q, a) in enumerate(qa):
    y = Inches(2.0 + i * 0.98)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.86), fill=(PANEL if i % 2 == 0 else BG), round_=True, radius=0.1)
    text(s, Inches(0.8), y + Inches(0.1), Inches(3.9), Inches(0.7), [[(q, 12.5, BLUE, True)]], line_spacing=1.08)
    text(s, Inches(4.9), y + Inches(0.1), Inches(7.8), Inches(0.7), [[(a, 11, BODY, False)]], line_spacing=1.1)
notes(s, "Backup for the question-and-answer session — the five questions most likely to come up, with straight "
         "answers. Why not buy aPriori? Because a vendor SaaS puts our CAD and our rates in their cloud, can't learn "
         "from our quote history, and bills six figures a year forever — while CostVision is on-premise, "
         "self-learning and our own intellectual property. Is our data training a public AI? No — zero retention in "
         "our own tenancy, and the air-gapped switch makes 'no external calls' provable rather than promised. "
         "Key-person risk? Seven hundred ninety-three automated tests, written architecture docs, and a trained "
         "second maintainer from Phase 4. Running cost? A virtual machine, a standard database, pay-per-use AI — "
         "no licences. And how do we trust the numbers? Physics build-ups on our own rates, calibrated against our "
         "own purchase orders, with the error tracked openly on the accuracy dashboard.")

# ═══════════════ 22 — BACKUP: PHOTO-TO-BOM WORKFLOW ═══════════════
s = header('Photo to costed BOM — how it works, step by step', 'Backup · Feature deep-dive')
pcb_steps = [
    ('1 · Upload photos', 'Up to 5 photos of the board (+ an optional BOM or pick-and-place file). Photos are auto-downsized in the browser; nothing is stored.', BLUE),
    ('2 · Board classification', 'AI identifies what the board is — domain (automotive, industrial, consumer) and safety level (ASIL) — so the right cost assumptions apply.', INDIGO),
    ('3 · OCR text pass', 'Every printed marking is read: chip part numbers, connector types, silkscreen references — hard evidence, not guesswork.', VIOLET),
    ('4 · Full vision BOM extraction', 'The AI locates and identifies every visible component — type, package, value, manufacturer part number (~20 seconds).', CYAN),
    ('5 · Deterministic clean-up & costing', 'Code (not AI) normalises types, parses values, scores per-line confidence, grounds prices against the distributor catalogue, then adds board fabrication + SMT assembly + test on country rates.', GREEN),
    ('6 · Same photos, same answer', 'The result is cached by a digital fingerprint (SHA-256) — re-running identical photos returns the identical BOM. Export PDF/Excel, save to library, feed the learning loop.', AMBER),
]
y = Inches(1.95)
for i, (t, d, c) in enumerate(pcb_steps):
    flow_box(s, Inches(0.45), y, Inches(12.45), Inches(0.68), t, d, c, fill=PANEL)
    y += Inches(0.68)
    if i < len(pcb_steps) - 1:
        down_arrow(s, Inches(6.5), y + Inches(0.015), h=Inches(0.15))
        y += Inches(0.18)
notes(s, "The photo-to-BOM workflow in six steps. An engineer photographs a circuit board — up to five photos, and "
         "optionally a BOM or pick-and-place file if one exists. Step two: the AI first classifies what kind of "
         "board it is, including automotive safety level, so the right cost assumptions apply. Step three is an OCR "
         "pass — it literally reads the part numbers printed on the chips, which anchors the analysis in hard "
         "evidence. Step four is the full vision extraction: every visible component located, typed and identified, "
         "in about twenty seconds. Step five matters for trust: ordinary deterministic code — not AI — cleans the "
         "list, scores confidence line by line, grounds prices against a distributor catalogue, and then builds the "
         "full cost: bare board fabrication, component purchase, SMT assembly time and test, on our country rates. "
         "And step six: the result is fingerprinted and cached, so the same photos always give the same answer — "
         "an auditable property no ad-hoc AI chat can offer.")

# ═══════════════ 23 — BACKUP: PHOTO-TO-BOM TECH / ACCURACY / TIME ═══════════════
s = header('Photo to costed BOM — technology, accuracy, speed', 'Backup · Feature deep-dive')
cols3 = [
    ('Technology', BLUE, [
        'Claude vision model — via our controlled endpoint (Option B) ',
        '5-stage pipeline: classify → OCR → extract → cost → cache',
        'Deterministic TypeScript post-processor for types, values, MPNs',
        'SHA-256 result cache — repeatable, auditable output',
        'Optional live pricing (Nexar/RS/Farnell) — part numbers only, opt-in',
    ]),
    ('Why it is accurate', GREEN, [
        'OCR evidence + vision cross-check — read, not guessed',
        'Code normalises every line; implausible part numbers rejected',
        'Per-line confidence score — low lines flagged for engineer review',
        'Prices grounded against the distributor catalogue, shown as ranges',
        'Learning loop calibrates against real quotes over time',
    ]),
    ('Speed', VIOLET, [
        'Main AI pass: ~20 seconds',
        'Full costed BOM: typically under a minute',
        'Repeat of the same board: instant (cache)',
        'Manual alternative: hours to days per board',
        'Works from photos alone — no CAD, no drawings needed',
    ]),
]
for i, (t, c, items) in enumerate(cols3):
    x = Inches(0.45 + i * 4.25)
    box(s, x, Inches(2.0), Inches(4.0), Inches(4.15), fill=(PANEL2 if i == 0 else PANEL if i == 2 else GREENBG), round_=True, radius=0.06)
    box(s, x, Inches(2.0), Inches(4.0), Inches(0.09), fill=c)
    text(s, x + Inches(0.25), Inches(2.2), Inches(3.55), Inches(0.4), [[(t, 14.5, c, True)]])
    text(s, x + Inches(0.25), Inches(2.68), Inches(3.55), Inches(3.4),
         [[('• ' + it, 10.5, BODY, False)] for it in items], space_after=7, line_spacing=1.12)
box(s, Inches(0.45), Inches(6.42), Inches(12.45), Inches(0.85), fill=AMBERBG, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(6.53), Inches(11.9), Inches(0.65),
     [[('Honest limits: ', 12, AMBER, True),
       ('a vision BOM is a strong first pass, not gospel — hidden or underside parts need extra photos, and '
        'low-confidence lines are deliberately flagged for a human check rather than silently guessed.', 11.5, BODY, False)]],
     line_spacing=1.15)
notes(s, "The same feature from three angles. Technology: the vision model runs through our controlled endpoint, "
         "inside a five-stage pipeline, and everything after the AI step is plain deterministic code — including a "
         "fingerprint cache that makes results repeatable and auditable. Live distributor pricing is opt-in and "
         "sends only part-number text, never images. Accuracy: the OCR pass means part numbers are read off the "
         "chips, not guessed; the post-processor rejects implausible part numbers; every line carries a confidence "
         "score and low ones are flagged for review; and prices are grounded against a catalogue and shown as "
         "ranges, not false precision. Speed: the heavy AI pass is about twenty seconds, a full costed BOM lands in "
         "under a minute, and repeating the same board is instant. The honest limit to state up front: it can only "
         "see what is photographed — underside components need an underside photo — and where it is unsure it says "
         "so and asks for a human check. That is by design.")

# ═══════════════ 24 — BACKUP: CAD-TO-COST WORKFLOW ═══════════════
s = header('CAD file to cost — how it works, step by step', 'Backup · Feature deep-dive')
cad_steps = [
    ('1 · Upload a CAD model', 'STEP, IGES or STL from any CAD system — as the standalone CAD-to-Cost flow or inside any of the 13 commodity calculators.', BLUE),
    ('2 · Opened in memory, inside our server', 'The OCCT geometry kernel (the same engine class commercial CAD tools are built on) reads the model in memory. The file is never written to disk and never leaves the server.', INDIGO),
    ('3 · The model is measured', 'Exact volume, surface area, bounding box, mean wall thickness, hole count, free-form faces — measured from the geometry, not typed in from a drawing.', VIOLET),
    ('4 · Inputs auto-filled', 'Mass from volume × material density, stock/billet size, and suggested process parameters for the chosen commodity — the engineer reviews and adjusts.', CYAN),
    ('5 · Physics engines cost it', 'Material, cycle or machining time, labour, tooling amortisation, energy and overheads — computed bottom-up on OUR uploaded rates for the chosen country.', GREEN),
    ('6 · Defensible result', 'Full cost breakdown + DFM warnings (thin walls, tonnage limits) + an honest uncertainty band — and the learning loop tightens it as real quotes arrive.', AMBER),
]
y = Inches(1.95)
for i, (t, d, c) in enumerate(cad_steps):
    flow_box(s, Inches(0.45), y, Inches(12.45), Inches(0.68), t, d, c, fill=PANEL)
    y += Inches(0.68)
    if i < len(cad_steps) - 1:
        down_arrow(s, Inches(6.5), y + Inches(0.015), h=Inches(0.15))
        y += Inches(0.18)
notes(s, "The CAD-to-cost workflow. An engineer uploads a STEP, IGES or STL model — either in the dedicated "
         "CAD-to-Cost flow or directly inside any of thirteen commodity calculators. Step two is the security "
         "headline: the geometry engine opens the model in memory, inside our own server — the file is never "
         "written to disk and never leaves the network. Step three, the engine measures the part: exact volume, "
         "surface area, wall thickness, holes — numbers measured from the actual geometry rather than read off a "
         "drawing by eye. Step four, those measurements auto-fill the costing inputs — mass, stock size, suggested "
         "process parameters — which the engineer reviews rather than types. Step five, the physics engines build "
         "the cost bottom-up on our own rates. And step six, the output is a defensible breakdown with "
         "manufacturability warnings and an honest uncertainty band that tightens as the learning loop absorbs "
         "real quotes.")

# ═══════════════ 25 — BACKUP: CAD-TO-COST TECH / ACCURACY / TIME ═══════════════
s = header('CAD file to cost — technology, accuracy, speed', 'Backup · Feature deep-dive')
cols4 = [
    ('Technology', BLUE, [
        'OCCT (Open CASCADE) geometry kernel — runs inside our backend',
        'Pure-TypeScript STL fast path — no external process at all',
        '18 deterministic physics cost engines do the actual costing',
        'AI translates the measured geometry into process inputs — it never receives the CAD file itself',
        'Rate library + 20-country regional rates drive every figure',
    ]),
    ('Why it is accurate', GREEN, [
        'Measured, not estimated — exact volume → exact material mass',
        'Physics build-ups traceable line by line (no black box)',
        'DFM sanity checks: wall thickness, clamp tonnage, press force',
        'Honest uncertainty band shown with every result',
        'Calibration on 3 real quotes took machining error 10.9% → 0.3%',
    ]),
    ('Speed', VIOLET, [
        'Geometry read and measured in seconds',
        'Costing computes instantly once inputs are filled',
        'Model-to-first-price: a few minutes, mostly review time',
        'Manual alternative: hours of take-off per part',
        'Same flow in 13 commodities — one skill to learn',
    ]),
]
for i, (t, c, items) in enumerate(cols4):
    x = Inches(0.45 + i * 4.25)
    box(s, x, Inches(2.0), Inches(4.0), Inches(4.15), fill=(PANEL2 if i == 0 else PANEL if i == 2 else GREENBG), round_=True, radius=0.06)
    box(s, x, Inches(2.0), Inches(4.0), Inches(0.09), fill=c)
    text(s, x + Inches(0.25), Inches(2.2), Inches(3.55), Inches(0.4), [[(t, 14.5, c, True)]])
    text(s, x + Inches(0.25), Inches(2.68), Inches(3.55), Inches(3.4),
         [[('• ' + it, 10.5, BODY, False)] for it in items], space_after=7, line_spacing=1.12)
box(s, Inches(0.45), Inches(6.42), Inches(12.45), Inches(0.85), fill=GREENBG, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(6.53), Inches(11.9), Inches(0.65),
     [[('Security, restated: ', 12, GREEN, True),
       ('the CAD file is processed in memory inside our backend and is never stored or transmitted. Only a short '
        'numeric geometry summary reaches the AI layer — and in air-gapped mode, nothing leaves at all.', 11.5, BODY, True)]],
     line_spacing=1.15)
notes(s, "CAD-to-cost from the same three angles. Technology: the geometry work is done by OCCT — Open CASCADE, "
         "the same engineering kernel class that commercial CAD packages build on — running inside our backend, "
         "with a pure-TypeScript fast path for STL files. The costing itself is done by our eighteen deterministic "
         "physics engines. The AI's only job is translating the measured geometry into sensible process inputs — "
         "and it never receives the CAD file, only a short numeric summary. Accuracy: the inputs are measured, not "
         "estimated — exact volume gives exact mass; every cost line is traceable; manufacturability checks catch "
         "physically impossible set-ups; the uncertainty band is shown honestly; and calibration against just "
         "three real quotes took machining error from eleven percent to a third of a percent. Speed: seconds to "
         "measure, instant to cost, minutes to a defensible first price — against hours of manual take-off. And "
         "the green bar restates the point that matters most in this room: the CAD never leaves our walls, full "
         "stop.")

OUT = 'CostVision-Implementation-Blueprint.pptx'
prs.save(OUT)
print(f'Wrote {OUT} with {len(prs.slides._sldIdLst)} slides')
