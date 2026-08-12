#!/usr/bin/env python3
"""
Deck geometry and legibility guard.

A fast structural check, run before every send.

It exists because these faults survive a casual read of the builder source: the
numbers look reasonable in code and the shape lands off the slide. It is NOT a
substitute for looking at the deck — `render_decks.sh` converts to PDF, and
rendering caught a panel overlap on Workflow 44 that this file could not.

It found, on decks that had been presented:

  * Executive slide 10 — a 4-column grid of 3.88" cards on a 13.33" slide, so
    column 4 began at 12.39" and ended at 16.27". Two whole feature cards were
    off the screen: "Tooling & NRE Amortisation" and "Learning Curve", the two
    the slide's own speaker note tells the presenter to point at.
  * Executive slide 4 — two chips clipped 0.13" below the bottom edge.
  * Agentic slide 15 — "Expected realizable" 0.27" past the right edge.
  * 6.0–6.1 pt body text in two decks, which no projector makes legible.

    python3 check_decks.py            # report
    python3 check_decks.py --strict   # exit 1 on any fault (CI / pre-send)
"""
import sys
import re
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.dml import MSO_FILL

DECKS = [
    'CostVision-Workflow-Explained.pptx',
    'CostVision-Agentic-AI-Management-Presentation.pptx',
    'CostVision-Executive-Presentation.pptx',
    'CostVision-Implementation-Blueprint.pptx',
]

# Nothing a director reads across a meeting room should be smaller than this.
MIN_PT = 7.5
# Shapes may sit this far outside the slide before it counts — bleed strips and
# full-width rules are drawn deliberately flush, so a hairline is not a fault.
TOLERANCE_IN = 0.02

EMOJI = re.compile('[\U0001F000-\U0001FAFF☀-➿]')


def audit(path):
    prs = Presentation(path)
    w_in, h_in = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    offslide, small, no_notes, emoji_fonts = [], [], [], set()
    overlaps = []

    for idx, slide in enumerate(prs.slides, 1):
        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            no_notes.append(idx)
        for sh in slide.shapes:
            try:
                l, t = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except (TypeError, AttributeError):
                continue          # placeholders without explicit geometry
            r, b = l + w, t + h
            if (r > w_in + TOLERANCE_IN or b > h_in + TOLERANCE_IN
                    or l < -TOLERANCE_IN or t < -TOLERANCE_IN):
                label = sh.text_frame.text[:40].replace('\n', ' ') if sh.has_text_frame else ''
                over = max(r - w_in, b - h_in, -l, -t)
                offslide.append((idx, round(over, 2), label))
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if run.font.size and run.font.size.pt < MIN_PT:
                        small.append((idx, round(run.font.size.pt, 1), run.text[:34]))
                    # An emoji run with no explicit emoji font falls back to
                    # whatever the viewing machine offers — which is how a glyph
                    # becomes an empty box on somebody else's laptop.
                    if EMOJI.search(run.text) and (run.font.name or '') != 'Segoe UI Emoji':
                        emoji_fonts.add((idx, run.font.name or '(inherited)'))

        # ── panel-on-panel overlap ────────────────────────────────────────
        # Everything can sit inside the slide and still be wrong: on Workflow 44
        # a caption bar was laid 0.06" over the two panels above it and clipped a
        # bullet. Only rendering caught that, so the check now looks for filled
        # panels that PARTIALLY overlap. Containment is deliberate design (a card
        # on a background), so only partial intersection is reported.
        panels = []
        for sh in slide.shapes:
            try:
                l, t = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except (TypeError, AttributeError):
                continue
            if w < 0.6 or h < 0.25:
                continue                      # chips, rules and icons
            if w > w_in * 0.97 and h > h_in * 0.9:
                continue                      # full-bleed backgrounds
            try:
                if sh.fill.type is None or sh.fill.type == MSO_FILL.BACKGROUND:
                    continue
            except (AttributeError, ValueError, TypeError):
                continue
            panels.append((l, t, l + w, t + h))
        for a in range(len(panels)):
            for b in range(a + 1, len(panels)):
                ax1, ay1, ax2, ay2 = panels[a]
                bx1, by1, bx2, by2 = panels[b]
                ox = min(ax2, bx2) - max(ax1, bx1)
                oy = min(ay2, by2) - max(ay1, by1)
                if ox <= 0.02 or oy <= 0.02:
                    continue
                contained = ((ax1 >= bx1 - .01 and ax2 <= bx2 + .01
                              and ay1 >= by1 - .01 and ay2 <= by2 + .01)
                             or (bx1 >= ax1 - .01 and bx2 <= ax2 + .01
                                 and by1 >= ay1 - .01 and by2 <= ay2 + .01))
                if contained:
                    continue
                overlaps.append((idx, round(ox, 2), round(oy, 2)))

    return {
        'slides': len(prs.slides), 'offslide': offslide, 'small': small,
        'no_notes': no_notes, 'emoji_fonts': sorted(emoji_fonts),
        'overlaps': overlaps,
    }


def main():
    strict = '--strict' in sys.argv
    faults = 0
    for path in DECKS:
        r = audit(path)
        n_small = len(r['small'])
        n_off = len(r['offslide'])
        n_emoji = len(r['emoji_fonts'])
        n_lap = len(r['overlaps'])
        # Overlap is ADVISORY, not a hard fault. Decks overlap shapes on purpose
        # — accent bars on cards, badges on panels — so a geometric test cannot
        # tell intent from accident, and a noisy gate is one people learn to
        # ignore. Rendering to PDF is the real check for this class now that
        # libreoffice-impress is installed; see render_decks.sh.
        faults += n_off + n_small + len(r['no_notes'])
        status = 'OK' if not (n_off or n_small or r['no_notes']) else 'FAULTS'
        print(f"\n{path}  ({r['slides']} slides)  {status}")
        print(f"   off-slide shapes      : {n_off}")
        for s, over, label in r['offslide'][:6]:
            print(f"        slide {s:3d}  +{over}\"  \"{label}\"")
        print(f"   runs below {MIN_PT} pt      : {n_small}")
        for s, pt, txt in r['small'][:6]:
            print(f"        slide {s:3d}  {pt} pt  \"{txt}\"")
        print(f"   panel overlaps (info) : {n_lap}")
        for sl, ox, oy in r['overlaps'][:5]:
            print(f"        slide {sl:3d}  {ox}\" x {oy}\" intersection")
        print(f"   slides without notes  : {r['no_notes'] or 'none'}")
        print(f"   emoji runs w/o Segoe  : {n_emoji}"
              + (f"  e.g. {r['emoji_fonts'][:3]}" if n_emoji else ''))

    print(f"\n{'=' * 62}\nTotal hard faults (off-slide + tiny text + missing notes): {faults}")
    if strict and faults:
        print('FAILED — do not present until these are fixed.')
        return 1
    return 0


if __name__ == '__main__':
    sys.exit(main())
