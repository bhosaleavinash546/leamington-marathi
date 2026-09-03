#!/usr/bin/env python3
"""
CostVision and CAPEE — the business case, for the Cost Engineering Director.

Two options, in the order they would be done: Option 1 fills CAPEE's cost input
automatically and is proved on real JLR parts; Option 2 costs a whole basket in
one unattended run, driven by an agent. Option 1 comes first because it is the
cheaper thing to try and it produces the accuracy figure Option 2 rests on.

Written for a non-IT audience: plain language on the slides and in the speaker
notes, technical terms only where there is no honest substitute.

RULES THIS FILE FOLLOWS.

1. No durations. Not weeks, not months, not quarters. Nothing here has been
   scoped by JLR, so a date would be an invention presented as a plan. The order
   of work is real and is shown; how long each step takes is for JLR to size.

2. No claimed savings. A pound figure needs JLR rate data and the accuracy
   number Option 1 produces, and we have neither. Slide 16 states what changes
   in KIND and says explicitly that no saving is claimed.

3. Only what is in the software goes on a slide. Every count, value and timing
   was read off the codebase or measured in September 2026 and re-checked at
   build time. Slide 21 lists how each one was checked.

4. The accuracy harness figure is deliberately absent. The harness exists and
   grades honestly, but its sample file is a template of EXAMPLE rows, so any
   MAPE it prints today is meaningless. Slide 21 says so in the amber box.

Cross-references to slide numbers are load-bearing: if a slide moves, re-check
the references on slides 3, 6, 7, 9 and 13.

Regenerate:  python3 build_capee_options_pptx.py
Output:      CostVision-CAPEE-Implementation-Options.pptx
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pptx_fixup import finalise

# ── Palette: identical to build_blueprint_pptx.py so the pack reads as one ────
INDIGO  = RGBColor(0x1D, 0x6F, 0xB8)
BLUE    = RGBColor(0x1D, 0x6F, 0xB8)
DARK    = RGBColor(0x16, 0x32, 0x5C)
BODY    = RGBColor(0x3A, 0x43, 0x56)
MUTED   = RGBColor(0x6B, 0x72, 0x80)
BG      = RGBColor(0xF4, 0xF7, 0xFB)
PANEL   = RGBColor(0xFF, 0xFF, 0xFF)
PANEL2  = RGBColor(0xE8, 0xF1, 0xFA)
GREENBG = RGBColor(0xEA, 0xF6, 0xEF)
AMBERBG = RGBColor(0xFC, 0xF3, 0xE3)
REDBG   = RGBColor(0xFB, 0xEC, 0xEA)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
AMBER   = RGBColor(0xB7, 0x79, 0x1F)
RED     = RGBColor(0xB0, 0x3A, 0x2E)
VIOLET  = RGBColor(0x6B, 0x3F, 0xA0)
LINE    = RGBColor(0xDC, 0xE3, 0xEE)
NAVY    = RGBColor(0x16, 0x32, 0x5C)
ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)
HERO_SUB = RGBColor(0xCA, 0xDC, 0xFC)
HERO_DIM = RGBColor(0x8F, 0xA3, 0xCC)

TITLE_FONT = 'Cambria'

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def box(slide, x, y, w, h, fill=None, line=None, round_=False, radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _h(v, minimum=Inches(0.12)):
    """Clamp a derived height to a positive value.

    Several helpers below size an inner text box as `h - <chrome>`. When a
    caller passes a short card the result goes zero or negative, which is
    invalid OOXML: PowerPoint refuses the file and offers to repair it, while
    LibreOffice silently tolerates it. A text box does not clip its contents, so
    clamping changes nothing visually.
    """
    return v if v > minimum else minimum


_PICTO = re.compile('([\U0001F300-\U0001FAFF☀-⛿✀-✒✙-➿️]+)')
EMOJI_FONT = 'Segoe UI Emoji'


def _emit_runs(p, t, size, color, bold, italic, base_font='Calibri'):
    for part in _PICTO.split(t):
        if not part:
            continue
        run = p.add_run(); run.text = part
        f = run.font
        f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
        if _PICTO.fullmatch(part): f.name = EMOJI_FONT
        elif base_font: f.name = base_font


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for r in para:
            t, size, color, bold = r[0], r[1], r[2], r[3]
            italic = r[4] if len(r) > 4 else False
            _emit_runs(p, t, size, color, bold, italic, base_font=font)
    return tb


def logo(slide, x=Inches(0.35), y=Inches(0.22), scale=1.0, on_dark=False):
    s = scale
    badge = box(slide, x, y, Inches(0.42 * s), Inches(0.42 * s),
                fill=ON_DARK if on_dark else INDIGO, round_=True, radius=0.28)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'cv'
    r.font.size = Pt(17 * s); r.font.bold = True; r.font.name = 'Calibri'
    r.font.color.rgb = NAVY if on_dark else ON_DARK
    text(slide, x + Inches(0.52 * s), y - Inches(0.03 * s), Inches(2.6), Inches(0.32),
         [[('CostVision', 18 * s, ON_DARK if on_dark else BLUE, True)]])
    text(slide, x + Inches(0.52 * s), y + Inches(0.24 * s), Inches(2.8), Inches(0.22),
         [[('AI  COST  INTELLIGENCE', 7.5 * s, HERO_DIM if on_dark else MUTED, False)]])


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def header(title, kicker=None):
    slide = prs.slides.add_slide(BLANK)
    box(slide, 0, 0, W, H, fill=BG)
    logo(slide)
    box(slide, 0, Inches(0.78), W, Pt(2.2), fill=NAVY)
    if kicker:
        text(slide, Inches(0.45), Inches(0.95), Inches(11.5), Inches(0.3),
             [[(kicker.upper(), 11, BLUE, True)]])
        ty = Inches(1.22)
    else:
        ty = Inches(1.02)
    text(slide, Inches(0.45), ty, Inches(12.4), Inches(0.6), [[(title, 27, DARK, True)]],
         font=TITLE_FONT)
    return slide


def card(slide, x, y, w, h, accent, title, lines, title_size=13, body_size=10.5,
         fill=PANEL, title_color=None):
    """White card with a coloured left rule — the pack's workhorse block."""
    box(slide, x, y, w, h, fill=fill, line=LINE, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    text(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.3),
         [[(title, title_size, title_color or DARK, True)]])
    if lines:
        runs = [[(ln[0], body_size, ln[1] if len(ln) > 1 else BODY,
                  ln[2] if len(ln) > 2 else False)] for ln in lines]
        text(slide, x + Inches(0.22), y + Inches(0.48), w - Inches(0.4), _h(h - Inches(0.6)),
             runs, space_after=3, line_spacing=1.12)


def table(slide, x, y, w, cols, rows, col_w, head_fill=NAVY, row_h=Inches(0.34),
          size=9.5, head_size=9.5):
    """Simple native table drawn from shapes — keeps full colour control."""
    cx = x
    for i, c in enumerate(cols):
        box(slide, cx, y, col_w[i], Inches(0.36), fill=head_fill)
        text(slide, cx + Inches(0.09), y + Inches(0.09), col_w[i] - Inches(0.16), Inches(0.24),
             [[(c, head_size, ON_DARK, True)]])
        cx += col_w[i]
    ry = y + Inches(0.36)
    for r_i, row in enumerate(rows):
        cx = x
        band = PANEL if r_i % 2 == 0 else BG
        for i, cell in enumerate(row):
            val, colr, bold = (cell if isinstance(cell, tuple) else (cell, BODY, False))
            box(slide, cx, ry, col_w[i], row_h, fill=band, line=LINE)
            text(slide, cx + Inches(0.09), ry + Inches(0.07), col_w[i] - Inches(0.16),
                 row_h - Inches(0.1), [[(val, size, colr, bold)]])
            cx += col_w[i]
        ry += row_h
    return ry


def callout(slide, x, y, w, h, fill, accent, title, body, tsize=12, bsize=10.5):
    box(slide, x, y, w, h, fill=fill, line=None, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    text(slide, x + Inches(0.24), y + Inches(0.13), w - Inches(0.45), Inches(0.28),
         [[(title, tsize, accent, True)]])
    text(slide, x + Inches(0.24), y + Inches(0.44), w - Inches(0.45), _h(h - Inches(0.55)),
         [[(body, bsize, BODY, False)]], space_after=3, line_spacing=1.14)


def footer(slide, txt):
    text(slide, Inches(0.45), Inches(7.02), Inches(12.4), Inches(0.24),
         [[(txt, 8, MUTED, False)]])


# ── Icons ────────────────────────────────────────────────────────────────────
# 12 of the 16 icons in assets/workflow-deck/icons are WHITE artwork on a
# transparent background — measured, average RGB 255,255,255. Dropped straight
# onto a light card they are invisible. So a white icon always sits inside a
# filled circle. check / times / warn / arrow are coloured and work bare.
import os
ICON_DIR = 'assets/workflow-deck/icons'
_COLOURED_ICONS = {'check', 'times', 'warn', 'arrow'}


def icon_badge(slide, name, cx, cy, d=Inches(0.62), fill=INDIGO, pad=0.26):
    """Coloured circle with a white icon centred inside it."""
    path = os.path.join(ICON_DIR, f'{name}.png')
    if name not in _COLOURED_ICONS:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
        c.fill.solid(); c.fill.fore_color.rgb = fill
        c.line.fill.background(); c.shadow.inherit = False
    if os.path.exists(path):
        inset = int(d * pad)
        slide.shapes.add_picture(path, cx + inset, cy + inset, d - 2 * inset, d - 2 * inset)


def step_circle(slide, n, cx, cy, d=Inches(0.5), fill=INDIGO, size=17):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.margin_left = tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(size); r.font.bold = True; r.font.name = 'Calibri'
    r.font.color.rgb = ON_DARK


def chevron(slide, x, y, w, h, label, sub, fill, text_col=ON_DARK):
    """One block of a left-to-right process flow."""
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = text_col; r.font.name = 'Calibri'
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(8.5); r2.font.color.rgb = text_col; r2.font.name = 'Calibri'
    return shp


def flow_step(slide, x, y, w, h, icon, title, sub, accent):
    """Icon-topped step card used in the end-to-end summary slides."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, w, Inches(0.06), fill=accent)
    icon_badge(slide, icon, x + (w - Inches(0.62)) / 2, y + Inches(0.18), fill=accent)
    text(slide, x + Inches(0.08), y + Inches(0.92), w - Inches(0.16), Inches(0.34),
         [[(title, 10.5, DARK, True)]], align=PP_ALIGN.CENTER)
    text(slide, x + Inches(0.08), y + Inches(1.28), w - Inches(0.16), _h(h - Inches(1.34)),
         [[(sub, 8.8, BODY, False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)


def arrow_between(slide, x, y, w=Inches(0.3)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.22))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xB8, 0xC4, 0xD4)
    a.line.fill.background(); a.shadow.inherit = False


def lane(slide, x, y, w, h, label, colour, items, label_w=Inches(1.5)):
    """Swim-lane: who does what."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE)
    box(slide, x, y, label_w, h, fill=colour)
    tf_box = box(slide, x, y, label_w, h, fill=None)
    tf = tf_box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = ON_DARK; r.font.name = 'Calibri'
    cw = (w - label_w) / len(items)
    for i, it in enumerate(items):
        cx = x + label_w + cw * i
        if i:
            box(slide, cx, y + Inches(0.06), Pt(0.75), h - Inches(0.12), fill=LINE)
        text(slide, cx + Inches(0.12), y + Inches(0.13), cw - Inches(0.24), h - Inches(0.26),
             [[(it, 9.2, BODY, False)]], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)


# ─────────────────────────────────────────────────────────── 1 · TITLE ──────
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=NAVY)
box(s, 0, 0, Inches(0.09), H, fill=INDIGO)
logo(s, x=Inches(0.55), y=Inches(0.5), scale=1.15, on_dark=True)
text(s, Inches(0.6), Inches(1.95), Inches(11.6), Inches(1.3),
     [[('CostVision and CAPEE', 40, ON_DARK, True)]], font=TITLE_FONT)
text(s, Inches(0.6), Inches(2.95), Inches(11.4), Inches(0.6),
     [[('The business case for automating cost input, and then costing in bulk',
        16, HERO_SUB, False)]])
box(s, Inches(0.6), Inches(3.8), Inches(3.3), Pt(2.5), fill=INDIGO)
for i, (n, t_, sub) in enumerate([
        ('1', 'Feed CAPEE automatically, and prove it',
         'CAPEE keeps doing the costing. CostVision reads the numbers off the 3D model and the '
         'drawing and fills them in. We prove it on real JLR parts first.'),
        ('2', 'Cost in bulk, run by an agent',
         'A whole basket of parts costed in one run, unattended, with a person asked only where '
         'the geometry genuinely cannot decide.')]):
    y = Inches(4.15) + Inches(1.0) * i
    step_circle(s, n, Inches(0.62), y, d=Inches(0.55), fill=INDIGO)
    text(s, Inches(1.4), y + Inches(0.02), Inches(10.8), Inches(0.3), [[(t_, 15, ON_DARK, True)]])
    text(s, Inches(1.4), y + Inches(0.35), Inches(10.8), Inches(0.5), [[(sub, 11, HERO_SUB, False)]],
         line_spacing=1.15)
text(s, Inches(0.6), Inches(6.55), Inches(11.6), Inches(0.4),
     [[('JLR Cost Engineering  ·  September 2026', 10, HERO_DIM, False)]])
notes(s, "Thank you for the time. I want to take you through two things we could do with a tool we "
         "have built called CostVision, and ask for a decision on the first one. The first is to "
         "let it fill in the numbers CAPEE needs, automatically, off the 3D model and the drawing, "
         "and prove that works on our own parts. The second, once that is proven, is to cost a "
         "whole basket of parts in one unattended run. I will be straight about what the tool does "
         "today and what it does not. There are no dates in this pack. Nobody in IT or the business "
         "has sized any of this yet, so any timeline I put up would be made up. What I can show you "
         "is the order the work has to happen in.")

# ────────────────────────────────────────────── 2 · WHAT WE ARE ASKING ──────
s = header('What we are asking for', 'The decision')
callout(s, Inches(0.45), Inches(1.75), Inches(12.43), Inches(1.15), GREENBG, GREEN,
        'Approve a proof of concept: Option 1, inside CAPEE, on our own parts',
        'Let CostVision fill CAPEE\'s cost input automatically for one part type, on 30 to 50 parts '
        'we have already bought. At the end we will know how often the measurement is right, how '
        'often an engineer has to correct it, and what it saves. Option 2 follows once that is known.')
rows = [
    (('What we want to do', DARK, True), 'Fill CAPEE\'s cost input from the 3D model and the drawing, for one part type'),
    (('What we need from JLR', DARK, True), '30 to 50 parts with the price we paid, our own rate data, three answers from IT, and a team for the trial'),
    (('What it costs to try', DARK, True), 'Our own engineering time. No licence, no new tool for the business to learn, CAPEE unchanged'),
    (('What we will know at the end', DARK, True), 'A measured accuracy figure against real JLR prices, and a measured time saving per part'),
    (('What we are NOT asking for', DARK, True), 'A decision on Option 2 today. It needs the accuracy answer first, and some work that has not been done'),
]
table(s, Inches(0.45), Inches(3.05), Inches(12.43), ['', 'Detail'], rows,
      [Inches(3.5), Inches(8.93)], row_h=Inches(0.52), size=10.8)
callout(s, Inches(0.45), Inches(6.15), Inches(12.43), Inches(0.95), AMBERBG, AMBER,
        'The one thing to know before you decide',
        'CostVision has never been checked against a price JLR has actually paid. Not one part. '
        'That is exactly what the proof of concept is for, and until it is done I cannot tell you '
        'how accurate the tool is.')
notes(s, "This is the ask, so I will put it first. I want approval to run a proof of concept inside "
         "CAPEE, on one part type, using parts we have already bought. What we need from JLR is on "
         "the second row and I will come back to it at the end. What it costs is our own time. "
         "There is no licence and the business does not have to learn a new tool, because CAPEE "
         "does not change. The amber box is the honest bit and I would rather say it now than have "
         "it come out in questions. We have never compared this tool against a price we actually "
         "paid. That is the whole point of the proof of concept.")

# ──────────────────────────────────────── 3 · WHAT HAPPENS TODAY ────────────
s = header('What costing a part looks like today', 'The problem')
card(s, Inches(0.45), Inches(1.78), Inches(6.1), Inches(3.7), RED,
     'TODAY, BY HAND',
     [('The engineer opens the part in CAD, reads the drawing,',),
      ('and types the numbers into CAPEE.',),
      ('Boxes to fill, depending on the part type:', DARK, True),
      ('machined part 12 · casting 30 · moulding 39 · pressing 69',),
      ('',),
      ('Part weight, starting material weight, wall thickness,',),
      ('hole and pocket counts, machine times, surface area,',),
      ('tolerances, finish, heat treatment.',),
      ('',),
      ('Two engineers costing the same part will not type the', DARK, True),
      ('same numbers. So one part gets two answers.', DARK, True)], fill=REDBG)
card(s, Inches(6.78), Inches(1.78), Inches(6.1), Inches(3.7), GREEN,
     'WITH COSTVISION',
     [('The engineer hands over the 3D model and the drawing.',),
      ('Measured off the model:', DARK, True),
      ('volume, size, surface area, wall thickness, holes,',),
      ('pockets, bosses, gear teeth, bend count.',),
      ('',),
      ('Read off the drawing:', DARK, True),
      ('tolerances, surface finish, heat treatment, coating,',),
      ('masked features.',),
      ('',),
      ('165 costing rules fill the form. 149 write straight', DARK, True),
      ('into a named box. Same part, same numbers, always.', GREEN, True)], fill=GREENBG)
callout(s, Inches(0.45), Inches(5.55), Inches(12.43), Inches(1.25), PANEL2, INDIGO,
        'Measuring and reading are two different jobs',
        'The 3D model is measured by software, in the way a CMM measures a part. That is not AI. '
        'The AI only reads the words written on the drawing, and everything it reads is checked '
        'against the measured model before it can move a cost. Slide 8 covers those checks. On the '
        'deterministic setting the AI is switched off completely and the part still costs.')
notes(s, "This is the problem we are trying to fix. Today an engineer opens the model, reads the "
         "drawing, and types the numbers in. How many numbers depends on the part. A machined part "
         "is twelve boxes. A pressing is sixty-nine. The time matters, but the second point matters "
         "more. Two of our engineers costing the same part will not type the same numbers, so we "
         "end up with two answers for one part and no way to say which is right. On the right is "
         "what the tool does instead. A hundred and sixty-five costing rules turn the "
         "measurements into the values the form needs, and a hundred and forty-nine of those "
         "write straight into a named box. One part uses its own type's rules, between nine "
         "and eighteen of them. The blue box is the "
         "distinction I want people to hold on to. Measuring the model is ordinary software doing "
         "geometry. The AI only reads the writing on the drawing, and we check everything it reads.")

# ─────────────────────────────────── 4 · WHERE COSTVISION IS TODAY ──────────
s = header('Where CostVision is today', 'Current status')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('Read off the software this month. This is what the tool does today.',
        11.5, BODY, False)]])
card(s, Inches(0.45), Inches(2.1), Inches(6.1), Inches(3.6), GREEN,
     'WORKING NOW',
     [('19 part types cost through the engine, plus an assembly',),
      ('roll-up and a separate software model',),
      ('The costing gives the same answer every run',),
      ('Measures 3D models with no internet connection',),
      ('Reads drawings for tolerances, finish, heat treatment',),
      ('(needs the AI service; the default setting runs without',),
      ('it and still costs the part)',),
      ('Refuses geometry it cannot honestly measure',),
      ('Asks before it scales an inch model',),
      ('Same safety checks on every CAD route, printed in the',),
      ('report',),
      ('328 materials, 178 machines, 42 labour rates, 20 regions',),
      ('2,079 automatic tests across 156 files, all passing',)], fill=GREENBG)
card(s, Inches(6.78), Inches(2.1), Inches(6.1), Inches(3.6), AMBER,
     'NOT THERE YET',
     [('Never compared against a price JLR has paid',),
      ('No history kept when someone changes a rate',),
      ('Finished costings are not saved, so an old one cannot',),
      ('be reproduced',),
      ('No JLR sign-on; it has its own logins today',),
      ('Gear costs in the screens but the automated route',),
      ('cannot reach it',),
      ('The automated route cannot yet be told a region or an',),
      ('annual volume',),
      ('The geometry safety checks run on the CAD route, not',),
      ('yet on the automated one',),
      ('Still under active development',)], fill=AMBERBG)
callout(s, Inches(0.45), Inches(5.9), Inches(12.43), Inches(1.0), PANEL2, INDIGO,
        'What changed since the last time I showed you this',
        'Six rounds of work in the last few weeks. The tool now refuses bad geometry instead of '
        'costing it, asks about units, runs the same checks on every CAD route and prints them in '
        'the report, measures a file once instead of four times, and widens its confidence range '
        'when a number was typed rather than measured.')
notes(s, "Before the options, this is honestly where the tool is. Left side is working today and I "
         "checked every line of it against the software this week. Nineteen part types cost through "
         "the engine. The maths repeats. It measures models with no internet connection, which is "
         "the first thing IT security will ask about. Right side is what is not done. The top one "
         "is the important one and I will keep coming back to it. Below that, it does not keep a "
         "history when someone changes a rate and it does not save finished costings, so today you "
         "could not reproduce a costing from three months ago. Those two matter a lot for Option "
         "two and I will come back to them. The blue box is what has changed since I last showed "
         "you this. The headline is that it now refuses work it cannot do honestly.")

# ──────────────────────────────────────────────── 5 · THE PLAN ──────────────
s = header('The plan: prove it small, then scale it', 'Sequence')
steps = [
    ('clip',   'Option 1', 'Fill CAPEE automatically', INDIGO),
    ('check',  'Prove it',  'On 30 to 50 real JLR parts', GREEN),
    ('eye',    'Decide',    'With an accuracy figure in hand', AMBER),
    ('cog',    'Option 2',  'Cost the whole basket in bulk', VIOLET),
]
sw, gap = Inches(2.9), Inches(0.3)
x = Inches(0.6)
for i, (ic, t_, sub, c_) in enumerate(steps):
    flow_step(s, x, Inches(1.82), sw, Inches(1.85), ic, t_, sub, c_)
    if i < len(steps) - 1:
        arrow_between(s, x + sw + Inches(0.03), Inches(2.62), gap - Inches(0.06))
    x += sw + gap
card(s, Inches(0.45), Inches(3.86), Inches(6.1), Inches(2.3), INDIGO,
     'OPTION 1 · feed CAPEE automatically',
     [('CAPEE stays as it is and still does the costing.',),
      ('Instead of typing between twelve and seventy boxes,',),
      ('the engineer hands over the 3D model and the drawing',),
      ('and the boxes fill in.',),
      ('',),
      ('Moves five pieces of software, about 14,300 lines,', DARK, True),
      ('with the safety checks. Slide 7 lists them.', DARK, True)])
card(s, Inches(6.78), Inches(3.86), Inches(6.1), Inches(2.3), VIOLET,
     'OPTION 2 · cost in bulk, run by an agent',
     [('A list of parts goes in. Every one is measured, costed',),
      ('and reported without anyone sitting there.',),
      ('An engineer is asked only where the geometry',),
      ('genuinely cannot decide.',),
      ('',),
      ('Same engine and the same rates. The new work is the', DARK, True),
      ('run, its record, and wiring the checks onto it.', DARK, True)])
callout(s, Inches(0.45), Inches(6.2), Inches(12.43), Inches(0.82), GREENBG, GREEN,
        'Why this order',
        'Option 1 is the cheaper thing to try and it produces the accuracy number that Option 2 '
        'needs. Costing five hundred parts is only worth doing once we know the tool gets one part '
        'right.')
notes(s, "This is the shape of the whole thing. Do Option one first, prove it on our own parts, "
         "look at the number, then decide about Option two. Option one is the smaller change. CAPEE "
         "does not move. The engineer stops typing and starts checking. Option two is a different "
         "capability altogether. Instead of costing the parts we have time for, we cost the whole "
         "basket. The reason for this order is on the green strip. Option one is cheap to try and "
         "it gives us the accuracy figure. There is no sense costing five hundred parts until we "
         "know the tool gets one part right.")

# ────────────────────────────────── 6 · OPTION 1 · ONE PART ─────────────────
s = header('Option 1: what happens to one part', 'Option 1')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('What an engineer does each time they cost a part', 11.5, MUTED, True)]])
steps = [
    ('upload', 'Hands over files',    '3D model and drawing, from the CAPEE screen',  INDIGO),
    ('ruler',  'Software measures',   'Size, weight, surface area, holes, pockets',   INDIGO),
    ('clip',   'AI reads the drawing', 'Tolerances, finish, coating, heat treatment', VIOLET),
    ('shield', 'Checks run',          'What was read is compared with the model',     GREEN),
    ('person', 'Engineer confirms',   'Only the things a 3D model cannot show',       AMBER),
    ('calc',   'CAPEE costs the part', 'CAPEE gets the numbers and calculates',       GREEN),
]
sw, gap = Inches(1.87), Inches(0.19)
x = Inches(0.45)
for i, (ic, t_, sub, c_) in enumerate(steps):
    flow_step(s, x, Inches(2.0), sw, Inches(2.05), ic, t_, sub, c_)
    if i < len(steps) - 1:
        arrow_between(s, x + sw + Inches(0.02), Inches(2.9), gap - Inches(0.04))
    x += sw + gap
text(s, Inches(0.45), Inches(4.25), Inches(12.4), Inches(0.3),
     [[('What order the build has to happen in', 11.5, MUTED, True)]])
ph = ['Agree the connection', 'Stand up measuring', 'Hand over settings',
      'Add safety checks', 'Connect into CAPEE', 'Trial on one part type']
cw = Inches(2.05)
x = Inches(0.45)
for i, lbl in enumerate(ph):
    chevron(s, x, Inches(4.58), cw, Inches(0.6), lbl, '', INDIGO if i < 5 else GREEN)
    x += cw - Inches(0.04)
lane(s, Inches(0.45), Inches(5.38), Inches(12.43), Inches(0.6), 'WE DO', INDIGO,
     ['Package the measuring software', 'Write out every setting', 'Move the safety checks across',
      'Help with the CAPEE connection'])
lane(s, Inches(0.45), Inches(6.02), Inches(12.43), Inches(0.6), 'JLR DOES', VIOLET,
     ['Tell us how CAPEE is built', 'Provide a server', 'Change CAPEE to take the numbers',
      'Free up a team for the trial'])
callout(s, Inches(0.45), Inches(6.72), Inches(12.43), Inches(0.68), AMBERBG, AMBER,
        'One step we cannot size yet',
        'Connecting into CAPEE depends on how CAPEE is built, and we do not know that yet. Slide 10 '
        'sets out the three ways of doing it.')
notes(s, "Option one on one slide. Top row is what happens each time somebody costs a part. The "
         "engineer hands over the model and drawing from inside CAPEE. Software measures the model, "
         "the AI reads the drawing, the checks compare the two. A few things no model can show, "
         "like heat treatment, the engineer confirms. Then CAPEE gets the numbers and costs the "
         "part exactly as it does now. Underneath is the build order and who does what. I have "
         "deliberately not put weeks against these. Nobody has sized this yet, and the connection "
         "step in particular waits on an answer we do not have.")

# ────────────────────────────────── 7 · OPTION 1 · WHAT WE HAND OVER ────────
s = header('Option 1: what we hand over', 'Option 1 · handover')
text(s, Inches(0.45), Inches(1.7), Inches(12.4), Inches(0.35),
     [[('Sizes counted straight out of the software this month. Both languages are in common use '
        'and JLR IT will have people who know them.', 11.5, BODY, False)]])
rows = [
    ('Measuring engine', ('Python', DARK, True), '2,462 lines',
     'Opens the 3D model and measures it. Makes no internet connection.'),
    ('Reading the drawing', ('TypeScript', DARK, True), '830 lines',
     'Sends the drawing to the AI and sorts out what comes back.'),
    ('Safety checks', ('TypeScript', DARK, True), '536 lines',
     'The four checks on the next slide, in three files.'),
    ('Questions per part type', ('TypeScript', DARK, True), '8,509 lines',
     '165 costing rules that fill the form and decide what the engineer gets asked.'),
    ('Design-for-cost findings', ('TypeScript', DARK, True), '1,992 lines',
     'Flags features that are expensive to make and says why.'),
    ('Settings list', ('Plain list', DARK, True), 'one file',
     'Every fixed number pulled together in one place. Slide 9.'),
]
table(s, Inches(0.45), Inches(2.3), Inches(12.43),
      ['Part', 'Language', 'Size', 'What it does'], rows,
      [Inches(3.3), Inches(1.7), Inches(1.5), Inches(5.93)], row_h=Inches(0.5), size=10.5)
callout(s, Inches(0.45), Inches(5.82), Inches(6.1), Inches(1.15), PANEL2, INDIGO,
        'About the two languages',
        'Python is the normal choice for engineering and measurement work. TypeScript is widely '
        'used for business software. Neither is unusual and neither ties JLR to one supplier.')
callout(s, Inches(6.78), Inches(5.82), Inches(6.1), Inches(1.15), AMBERBG, AMBER,
        'One thing for IT to confirm',
        'The measuring engine needs a normal Linux server, not the cut-down kind used for small '
        'services. It decides where this can go, so it is worth settling early.')
notes(s, "You asked last time what would actually get handed over. Six things, and the sizes come "
         "straight out of the software. The measuring engine is Python, which is the normal choice "
         "for measurement work. The rest is TypeScript, common in business software. Neither is odd "
         "and neither ties us to one supplier. The fourth row is the biggest piece and it is worth "
         "a word: that is the hundred and sixty-five costing rules that fill the form and decide "
         "what to ask the engineer for. The amber box is the one practical thing worth raising with IT "
         "now, because it decides where this can live.")

# ─────────────────────────────────── 8 · OPTION 1 · SAFETY CHECKS ───────────
s = header('Option 1: the checks that have to come with it', 'Option 1 · required')
text(s, Inches(0.45), Inches(1.7), Inches(12.4), Inches(0.35),
     [[('The AI does not set a price. It does hand over numbers that feed into one, so a bad '
        'reading turns into a bad cost. These four checks are what stop that, and they have to '
        'come across with everything else.', 11.5, BODY, False)]])
checks = [
    ('cube', 'The measured model beats the AI', INDIGO,
     'If the AI reads 2 kg off the drawing and the model measures 1.4 kg, we use the measured '
     'figure. A disagreement the other way is flagged to the engineer, and a gap of more than '
     'half stops the costing until someone accepts it.'),
    ('eye', 'Everything read gets cross-checked', INDIGO,
     'Numbers taken off the drawing get compared against the measured model. Where they disagree, '
     'the engineer is told and has to accept it by name before the cost stands.'),
    ('press', 'Machining time has a ceiling', INDIGO,
     'On a casting or forging, machining is capped at what finishing that part could realistically '
     'take, so a bad estimate cannot run the cost up.'),
    ('person', 'Unknowns get asked, not guessed', VIOLET,
     'Heat treatment, tolerance class and material grade cannot be read off a solid model, so the '
     'engineer is asked. The AI suggestion is labelled as a suggestion and never comes pre-ticked.'),
]
y = Inches(2.45)
for ic, t_, c_, b_ in checks:
    box(s, Inches(0.45), y, Inches(12.43), Inches(0.92), fill=PANEL, line=LINE, round_=True)
    box(s, Inches(0.45), y, Inches(0.075), Inches(0.92), fill=c_)
    icon_badge(s, ic, Inches(0.72), y + Inches(0.16), d=Inches(0.6), fill=c_)
    text(s, Inches(1.55), y + Inches(0.14), Inches(11.1), Inches(0.28), [[(t_, 12.5, DARK, True)]])
    text(s, Inches(1.55), y + Inches(0.45), Inches(11.1), Inches(0.42),
         [[(b_, 10.2, BODY, False)]], line_spacing=1.12)
    y += Inches(1.02)
callout(s, Inches(0.45), Inches(6.6), Inches(12.43), Inches(0.82), REDBG, RED,
        'Two things we found in our own review, and fixed',
        'A route through the tool had no checks on it and returned a cost of "not a number" while '
        'reporting success. And a part exported as surfaces rather than a solid was costed at a '
        'plausible but wrong weight. Both now stop the costing instead of producing a figure.')
notes(s, "On a job like this the temptation is to take the measuring and the AI and leave what "
         "looks like plumbing until later. These four checks are what make the rest safe. The first "
         "one is the one to remember. If the AI and the measured model disagree, the model wins. "
         "The last one matters for a different reason. Where the tool cannot know something it asks "
         "rather than guessing, and it never pre-ticks its own answer. The red box is from our own "
         "review. We found a route with no checks that gave back a nonsense figure while saying it "
         "had worked, and we found that a part exported as surfaces got costed at a wrong weight. "
         "Both of those now stop rather than produce a number.")

# ────────────────────────────────── 9 · OPTION 1 · FIXED NUMBERS ────────────
s = header('Option 1: the fixed numbers inside the software', 'Option 1 · settings')
text(s, Inches(0.45), Inches(1.7), Inches(12.4), Inches(0.35),
     [[('You asked what the hardcoded values actually are. These are engineering judgements written '
        'into the code. Six real ones, read off the software this month:', 11.5, BODY, False)]])
rows = [
    ('Smallest face counted as machined', ('400 mm²', DARK, True),
     'Below this we treat it as an edge, not a face to machine'),
    ('Smallest pocket counted', ('80 mm²', DARK, True),
     'Stops small recesses getting costed as pockets'),
    ('Cost of a rejected plated part', ('4.5x a good one', DARK, True),
     'You strip it and plate it again, so a 3% reject rate is nearer 13% on cost'),
    ('Zinc price used for plating', ('USD 3.67 / kg', DARK, True),
     'Market reference, August 2026. Moves monthly and is meant to get updated'),
    ('Shape allowance, pressing', ('1.15', DARK, True),
     'A real pressing has more surface than a flat plate: edges, flanges, bends'),
    ('Shape allowance, casting and forging', ('1.25 / 1.10', DARK, True),
     'Castings carry ribs and bosses; forgings are simpler on the outside'),
]
table(s, Inches(0.45), Inches(2.35), Inches(12.43),
      ['What it controls', 'Value', 'Why'], rows,
      [Inches(4.3), Inches(2.3), Inches(5.83)], row_h=Inches(0.5), size=10.5)
callout(s, Inches(0.45), Inches(5.87), Inches(6.1), Inches(1.15), PANEL2, INDIGO,
        'How they get handed over',
        'Pulled out of the code into one list you can read and change without a developer. That '
        'list is a deliverable in its own right and it is a step in the Option 1 sequence.')
callout(s, Inches(6.78), Inches(5.87), Inches(6.1), Inches(1.15), AMBERBG, AMBER,
        'Where these came from',
        'They are our engineering estimates. They are not measurements from a JLR plant. Treat them '
        'as a starting point and expect to replace them as real data comes in.')
notes(s, "Hardcoded values is a phrase that means nothing on its own, so here are six real ones. "
         "Look at the third row. A rejected plated part costs about four and a half times a good "
         "one, because you strip it and start again. That is why a three percent reject rate is "
         "closer to thirteen percent on cost. That is an engineering judgement and exactly the sort "
         "of thing you would challenge in a review, so it should not be buried in code. We pull "
         "every one of them into a single list you can read and change. The amber box is the "
         "caveat. These are our estimates, not measurements from one of our plants.")

# ───────────────────────────────── 10 · OPTION 1 · CONNECTING ───────────────
s = header('Option 1: three ways to connect it to CAPEE', 'Option 1 · decision for IT')
text(s, Inches(0.45), Inches(1.7), Inches(12.4), Inches(0.35),
     [[('We do not know yet how CAPEE is built. All three routes work. The answer decides which one '
        'we use and how much effort it takes.', 11.5, BODY, False)]])
rows = [
    (('A · Put it inside CAPEE', DARK, True), 'CAPEE is built in JavaScript',
     'Our code goes straight in. Only the measuring engine sits separately.',
     ('Least effort', GREEN, True)),
    (('B · Run it alongside', DARK, True), 'CAPEE is Java, .NET or anything else',
     'It runs as its own service on a JLR server and CAPEE asks it for numbers over a secure link.',
     ('Most effort', AMBER, True)),
    (('C · Pass a file', DARK, True), 'CAPEE is a desktop or spreadsheet tool',
     'It writes a file that CAPEE imports. Simplest to build, but somebody has to move the file.',
     ('Least effort', GREEN, True)),
]
table(s, Inches(0.45), Inches(2.3), Inches(12.43),
      ['Route', 'Use when', 'What it means day to day', 'Relative effort'], rows,
      [Inches(2.9), Inches(2.9), Inches(5.13), Inches(1.5)], row_h=Inches(0.85), size=10.3)
callout(s, Inches(0.45), Inches(5.4), Inches(6.1), Inches(1.0), PANEL2, INDIGO,
        'If we had to pick blind',
        'Route B. It works whatever CAPEE turns out to be, and we can update CostVision without '
        'touching CAPEE every time.')
callout(s, Inches(6.78), Inches(5.4), Inches(6.1), Inches(1.0), AMBERBG, AMBER,
        'Three questions for IT',
        'What is CAPEE written in. Can it call another service on the network. Can it run on a '
        'normal Linux server.')
footer(s, 'Effort is shown against the other two routes. None of the three has been sized by JLR.')
notes(s, "I cannot give you one answer because nobody has told us how CAPEE is built. Rather than "
         "hold everything up, here are all three routes. All of them work. What changes is how much "
         "effort each takes, and I have shown that against each other rather than in weeks, because "
         "none of this has been sized. If you made me pick blind I would take route B, running it "
         "alongside, because it works whatever CAPEE is and we can improve our side without "
         "disturbing it. The three questions in amber are what I need from IT.")

# ──────────────────────────────── 11 · OPTION 1 · THE PROOF ─────────────────
s = header('Option 1: what the proof of concept proves', 'Option 1 · the trial')
text(s, Inches(0.45), Inches(1.68), Inches(12.4), Inches(0.32),
     [[('One part type, one team, 30 to 50 parts we have already bought. Four numbers come out of '
        'it, and none of them exist today.', 11.5, BODY, False)]])
meas = [
    ('1', 'How close is it?', GREEN,
     'Cost every part with the tool, compare against the price we actually paid, and report the '
     'average error. The test is whether the real error lands inside the range the tool itself '
     'predicted for each part. A tool that is wrong and says so is usable. One that is wrong and '
     'confident is not.'),
    ('2', 'How often does the engineer have to step in?', INDIGO,
     'Count the values the engineer changed after the tool filled them, and the parts where the '
     'tool stopped and asked. On our own eight parts, with the material family given, six costed '
     'straight through and two stopped. Without it all eight stop, because a solid model cannot '
     'tell steel from aluminium. We need that figure on JLR parts, not ours.'),
    ('3', 'How much time does it save?', VIOLET,
     'Time the same parts both ways: typed by hand into CAPEE, and filled by the tool with the '
     'engineer checking. The difference is the saving, measured rather than claimed.'),
    ('4', 'What does it refuse, and was it right to?', AMBER,
     'Every part the tool turned down, with the reason. If it refuses parts it should have costed, '
     'that is a fault we need to see early.'),
]
y = Inches(2.3)
for n, t_, c_, b_ in meas:
    box(s, Inches(0.45), y, Inches(12.43), Inches(1.02), fill=PANEL, line=LINE, round_=True)
    box(s, Inches(0.45), y, Inches(0.075), Inches(1.02), fill=c_)
    step_circle(s, n, Inches(0.72), y + Inches(0.18), d=Inches(0.48), fill=c_)
    text(s, Inches(1.45), y + Inches(0.12), Inches(11.2), Inches(0.28), [[(t_, 12.2, DARK, True)]])
    text(s, Inches(1.45), y + Inches(0.43), Inches(11.2), Inches(0.48),
         [[(b_, 10.0, BODY, False)]], line_spacing=1.12)
    y += Inches(1.08)
callout(s, Inches(0.45), Inches(6.6), Inches(12.43), Inches(0.84), GREENBG, GREEN,
        'Why it is worth doing even if the answer is disappointing',
        'A measured accuracy figure is worth having either way. Today nobody can state one, so '
        'every conversation about this tool stalls in the same place.')
notes(s, "This is what we would actually get out of the trial, and I want to be concrete because "
         "these four numbers do not exist today. First, how close is it. We cost the parts, compare "
         "against what we paid, and report the error. The test I would use is whether the real "
         "error lands inside the range the tool predicted for itself. A tool that is wrong and says "
         "so is usable. One that is wrong and confident is not. Second, how often the engineer has "
         "to step in. On our own eight test parts, once the material family was given, six went "
         "straight through and two stopped and asked. Given nothing at all, all eight stop and "
         "ask, because a solid model cannot tell steel from aluminium. In a bulk run that answer "
         "comes off the part master rather than from a person. I need that figure on JLR parts. Third, the time saving, measured by doing the "
         "same parts both ways. Fourth, what it refused and whether it was right to. And the green "
         "strip is the reason to do it even if the answer disappoints us. Right now nobody can "
         "state an accuracy figure at all, and every conversation about this tool stops there.")

# ─────────────────────────────── 12 · OPTION 1 · ORDER OF WORK ──────────────
s = header('Option 1: order of work', 'Option 1 · sequence')
plan1 = [
    ('1', 'Agree the connection', INDIGO,
     'Find out how CAPEE is built, pick one of the three routes, agree exactly which numbers CAPEE expects.'),
    ('2', 'Stand up the measuring software', INDIGO,
     'Install it on a JLR server. Measure 20 real JLR parts and check the results against their drawings.'),
    ('3', 'Hand over the settings list', VIOLET,
     'Pull every fixed number into one list. JLR review it and sign it off before anything gets costed with it.'),
    ('4', 'Move the safety checks across', GREEN,
     'The four checks and the questions the engineer confirms. In the scope from the start.'),
    ('5', 'Connect into CAPEE', INDIGO,
     'Wire it in. Route the drawing reading through a JLR-approved AI service. Cost one real part end to end.'),
    ('6', 'Run the trial', GREEN,
     'One team, one part type, 30 to 50 parts with prices we paid. Produce the four numbers from the last slide.'),
]
y = Inches(1.9)
for n, t_, c_, b_ in plan1:
    box(s, Inches(0.45), y, Inches(12.43), Inches(0.7), fill=PANEL, line=LINE, round_=True)
    box(s, Inches(0.45), y, Inches(0.075), Inches(0.7), fill=c_)
    step_circle(s, n, Inches(0.72), y + Inches(0.11), d=Inches(0.46), fill=c_)
    text(s, Inches(1.42), y + Inches(0.08), Inches(11.2), Inches(0.27), [[(t_, 12.2, DARK, True)]])
    text(s, Inches(1.42), y + Inches(0.38), Inches(11.2), Inches(0.28), [[(b_, 9.6, BODY, False)]])
    y += Inches(0.78)
callout(s, Inches(0.45), Inches(6.62), Inches(6.1), Inches(0.84), PANEL2, INDIGO,
        'Steps 1 and 2 can start together',
        'Standing up the measuring software does not depend on knowing how CAPEE is built. Only '
        'step 5 does.')
callout(s, Inches(6.78), Inches(6.62), Inches(6.1), Inches(0.84), AMBERBG, AMBER,
        'Start with one part type',
        'Machined parts or pressings first. Most of our volume and the shapes we understand best. '
        'Widen it once it holds up.')
notes(s, "Six steps in the order they have to happen. Two worth pulling out. Step three is the "
         "settings list, signed off before anything gets costed with it. Step four is the safety "
         "checks, and I have given them a step of their own so they do not get folded into the "
         "connection work and squeezed. The blue box is useful for planning. Steps one and two can "
         "run at the same time, because standing up the measuring software does not depend on "
         "knowing how CAPEE is built. Only step five does.")

# ─────────────────────────────────────── 13 · OPTION 2 · WHAT IT IS ─────────
s = header('Option 2: costing in bulk, run by an agent', 'Option 2')
card(s, Inches(0.45), Inches(1.72), Inches(6.1), Inches(2.75), RED,
     'WHAT WE DO TODAY',
     [('We cost the parts we have time to cost.',),
      ('',),
      ('A basket of several hundred parts gets sampled, or',),
      ('estimated by analogy against something similar, or',),
      ('taken from what the supplier last quoted.',),
      ('',),
      ('Nobody costs all of it from the geometry, because', DARK, True),
      ('there are not enough engineer-hours to do it.', DARK, True)], fill=REDBG)
card(s, Inches(6.78), Inches(1.72), Inches(6.1), Inches(2.75), GREEN,
     'WHAT OPTION 2 DOES',
     [('A list of parts goes in. Every one is measured,',),
      ('costed and reported in one run, unattended.',),
      ('',),
      ('An engineer is asked only where the geometry',),
      ('genuinely cannot decide, and only once per question.',),
      ('',),
      ('Same engine and the same rate book. The geometry', DARK, True),
      ('checks have to be wired onto this route.', DARK, True)], fill=GREENBG)
callout(s, Inches(0.45), Inches(4.55), Inches(12.43), Inches(1.1), PANEL2, INDIGO,
        'What the agent actually does, and what it does not',
        'The agent marshals the run. It picks the part type, answers what it can from the drawing, '
        'and escalates what it cannot. Every number still comes from the same deterministic engine '
        'at the same rates. The agent never sets a price, but it does choose the inputs, so the '
        'geometry checks that run on the CAD route today have to be wired onto this route before a '
        'run is left unattended.')
callout(s, Inches(0.45), Inches(5.8), Inches(12.43), Inches(1.15), AMBERBG, AMBER,
        'Where this stands today',
        'The engine is already callable on its own, and there is already an agent that can cost a '
        'part when asked. What does not exist is the run itself: a list going in, a record coming '
        'out, and a way to reproduce it later. Slide 15 sets out exactly what has to be built.')
notes(s, "Option two is a different thing from Option one, and I want to be clear about that. "
         "Option one saves the engineer typing. Option two gives us something we do not have at "
         "all. Today we cost the parts we have time to cost. A basket of several hundred gets "
         "sampled or estimated by analogy, because there are not enough engineer-hours. Option two "
         "costs the whole basket in one run. The blue box is important. The agent marshals the run "
         "but it never sets a price. Every number still comes out of the same engine at the same "
         "rates. It does choose the inputs, though, so the geometry checks that run on the CAD "
         "route today have to be wired onto this one, and that is on the list on slide fifteen. "
         "The amber box is the honest position. The engine already "
         "works standalone and an agent already exists that can cost a part. What does not exist is "
         "the run itself.")

# ────────────────────────────────── 14 · OPTION 2 · A BULK RUN ──────────────
s = header('Option 2: what a bulk run does', 'Option 2')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('What happens to a basket of parts, start to finish', 11.5, MUTED, True)]])
steps2 = [
    ('upload', 'A list goes in',      'Part numbers, CAD files, volumes, region',       VIOLET),
    ('ruler',  'Each one measured',   'Same engine as a single part, run in parallel',  VIOLET),
    ('cog',    'Rules fill the form', 'The part type\u2019s own rules run',              VIOLET),
    ('person', 'Only gaps escalate',  'One question, answered once, applied to all',    AMBER),
    ('calc',   'Engine costs them',   'Eight buckets, every figure traceable',          GREEN),
    ('clip',   'One report out',      'Per part and per basket, with what was assumed', GREEN),
]
sw, gap = Inches(1.87), Inches(0.19)
x = Inches(0.45)
for i, (ic, t_, sub, c_) in enumerate(steps2):
    flow_step(s, x, Inches(2.0), sw, Inches(2.05), ic, t_, sub, c_)
    if i < len(steps2) - 1:
        arrow_between(s, x + sw + Inches(0.02), Inches(2.9), gap - Inches(0.04))
    x += sw + gap
text(s, Inches(0.45), Inches(4.25), Inches(12.4), Inches(0.3),
     [[('How long a run takes. The 40-part row was measured this month; the rest carries that rate '
        'forward', 11.5, MUTED, True)]])
rows = [
    ('Measure one part, warm', ('0.9 to 3.1 seconds', DARK, True), 'Five real production parts, 0.6 MB to 3.0 MB. Bigger file, longer'),
    ('Run the rules and cost it', ('5 milliseconds', DARK, True), 'Rules plus the eight-bucket engine. Not the slow part'),
    ('40 parts', ('31 seconds', DARK, True), 'Measured. Two workers, nothing reused from cache'),
    ('100 parts', ('about 1.5 minutes', DARK, True), 'That measured rate, on two workers'),
    ('500 parts', ('about 6.5 minutes', DARK, True), 'Two workers. Roughly half that on four'),
    ('1,000 parts', ('about 13 minutes', DARK, True), 'Two workers. Roughly half that on four'),
]
table(s, Inches(0.45), Inches(4.52), Inches(12.43),
      ['What', 'Measured', 'Basis'], rows,
      [Inches(4.3), Inches(2.6), Inches(5.53)], row_h=Inches(0.30), size=10.2)
callout(s, Inches(0.45), Inches(6.72), Inches(12.43), Inches(0.66), PANEL2, INDIGO,
        'The machine time is not the constraint',
        'Getting the part list, the CAD files and the rate data together is the work. The computer '
        'finishes in minutes.')
notes(s, "Here is what a run actually does and what it costs in machine time. The forty-part row "
         "is measured, this month, on our own production parts. A part takes between one and three "
         "seconds to measure depending on how big the file is. Running the rules and costing it "
         "takes five thousandths of a second, so the measuring is the whole of it. Forty parts took "
         "thirty-one seconds on two workers with nothing reused, and the rows below that are that "
         "same rate carried forward, so five hundred parts is about six and a half minutes. The point of the blue strip at the bottom is that the "
         "computer is not the constraint here. Pulling together the part list, the CAD files and "
         "our own rate data is the work. The machine finishes while you are getting a coffee.")

# ──────────────────────────── 15 · OPTION 2 · WHAT MUST BE BUILT ────────────
s = header('Option 2: what has to be built first', 'Option 2 · prerequisites')
text(s, Inches(0.45), Inches(1.68), Inches(12.4), Inches(0.32),
     [[('These are the honest gaps. None of them is large on its own, and none of them is optional.',
        11.5, BODY, False)]])
rows = [
    (('A history of rate changes', DARK, True), 'Blocks audit',
     'Today the rate book is one row that gets overwritten. Change a rate and the old one is gone. '
     'A bulk run costed at rates nobody can recover is not defensible.',
     ('Must have', RED, True)),
    (('Saved costings', DARK, True), 'Blocks audit',
     'Finished costings are not stored anywhere. A five-hundred-part run would produce numbers with '
     'nowhere to live and no way to reproduce them.',
     ('Must have', RED, True)),
    (('Region and volume on the automated route', DARK, True), 'Blocks correctness',
     'The automated route always uses the UK rate book and never gets told the annual volume. Both '
     'already work in the screens; they are simply not passed through.',
     ('Must have', RED, True)),
    (('Gear on the automated route', DARK, True), 'Blocks coverage',
     'Gear costs in the screens and has its own rules, but the automated route cannot reach it. '
     'A basket with gears in it would silently skip them.',
     ('Must have', AMBER, True)),
    (('Geometry safety checks on the automated route', DARK, True), 'Blocks trust',
     'The checks that compare what was read against what was measured run on the CAD route today. '
     'The automated route validates the fields only. They must be wired on before a run is left '
     'unattended.',
     ('Must have', RED, True)),
    (('The run itself', DARK, True), 'The new work',
     'A list in, a job queue, a report out, and a way to answer one question once and apply it to '
     'every part it affects.',
     ('New build', INDIGO, True)),
]
table(s, Inches(0.45), Inches(2.15), Inches(12.43),
      ['What is missing', 'Why it matters', 'Detail', 'Type'], rows,
      [Inches(3.0), Inches(1.5), Inches(6.63), Inches(1.3)], row_h=Inches(0.7), size=9.8)
callout(s, Inches(0.45), Inches(6.74), Inches(12.43), Inches(0.66), GREENBG, GREEN,
        'The good news underneath all of that',
        'The costing engine itself needs no change. We ran it on its own, outside the application, '
        'and it produced the full eight-bucket answer with nothing adjusted.')
notes(s, "I would rather show you this list than have it come out later. Six things. The first two "
         "are the ones that would stop me signing off a bulk run for anything audited. The rate "
         "book today is a single row that gets overwritten, so change a rate and the old one is "
         "gone. And finished costings are not saved at all. Together that means a five-hundred-part "
         "run produces numbers you cannot reproduce next quarter. The third is smaller than it "
         "sounds: region and volume already work in the screens, they are just not passed through "
         "to the automated route. Fourth, gear does not reach the automated route at all, so a "
         "basket with gears in it would quietly skip them. Fifth is the geometry safety checks. They "
         "run on the CAD route today and the automated route only validates the fields, so they "
         "have to be wired on before anyone leaves a run unattended. Sixth is the run itself, which "
         "is the actual new build. The green strip is the good news. The engine needs no change. We "
         "checked that by running it on its own, outside the application.")

# ───────────────────────────── 16 · OPTION 2 · WHAT IT CHANGES ──────────────
s = header('Option 2: what it changes for the business', 'Option 2 · the case')
rows = [
    ('How many parts get a proper cost', ('The ones we have hours for', RED, False), ('All of them', GREEN, True)),
    ('How a basket gets priced', ('Sample and scale up', RED, False), ('Every part costed from its own geometry', GREEN, True)),
    ('Turnaround on a new basket', ('Days of engineer time', RED, False), ('Minutes of machine time, plus the data gathering', GREEN, True)),
    ('Consistency across a basket', ('Depends who costed it', RED, False), ('Same rules and rates for every part', GREEN, True)),
    ('What we can say to a supplier', ('This is our estimate', AMBER, False), ('Here is the build-up, per part, with the rates', GREEN, True)),
    ('Where an engineer spends the time', ('Typing numbers in', RED, False), ('Answering the questions only a person can', GREEN, True)),
    ('What it needs from a person', ('One engineer per part', RED, False), ('One answer per question, applied to every part', GREEN, True)),
]
table(s, Inches(0.45), Inches(1.9), Inches(12.43),
      ['', 'Today', 'With bulk costing'], rows,
      [Inches(4.13), Inches(3.7), Inches(4.6)], row_h=Inches(0.5), size=10.8)
callout(s, Inches(0.45), Inches(5.85), Inches(12.43), Inches(1.0), AMBERBG, AMBER,
        'The honest limit on all of this',
        'Every row assumes the tool is accurate enough to trust unattended, and today nobody can '
        'say whether it is. That is precisely what Option 1 measures, which is why it comes first.')
footer(s, 'No cost saving is claimed on this slide. The saving depends on our own rate data and on '
          'the accuracy figure Option 1 produces.')
notes(s, "This is the business case for Option two, and I have deliberately not put a pound figure "
         "on it, because an honest one needs two things we do not have yet: our own rate data and "
         "the accuracy number from Option one. What I can show is what changes in kind. Today we "
         "cost the parts we have hours for and scale up from a sample. With bulk costing every part "
         "in the basket gets costed from its own geometry, with the same rules and the same rates. "
         "The row I would draw your eye to is the last one. Today it is one engineer per part. With "
         "this it is one answer per question, applied across every part that question affects. And "
         "the amber box is the limit on all of it. Every row assumes the tool is accurate enough to "
         "leave running, and that is exactly what Option one measures.")

# ─────────────────────────── 17 · OPTION 2 · ORDER OF WORK ──────────────────
s = header('Option 2: order of work', 'Option 2 · sequence')
plan2 = [
    ('1', 'Load JLR rates', VIOLET,
     'Materials, machine rates, labour, energy, country factors. This sets the pace and it is ours to do.'),
    ('2', 'Keep a history of rate changes', RED,
     'Every change stamped and kept, so a costing can name the rate book it used. Prerequisite for anything audited.'),
    ('3', 'Save finished costings', RED,
     'A stored record per part: inputs, rates, answers, result. Without this a bulk run cannot be reproduced.'),
    ('4', 'Close the automated-route gaps', AMBER,
     'Pass region and annual volume through, and register gear. Small work, and a basket is wrong without it.'),
    ('5', 'Build the run', INDIGO,
     'A list in, a job queue, one question answered once, a report out per part and per basket.'),
    ('6', 'Trial on a real basket', GREEN,
     'One programme, one commodity family, run beside the way we cost today and compare.'),
]
y = Inches(1.88)
for n, t_, c_, b_ in plan2:
    box(s, Inches(0.45), y, Inches(12.43), Inches(0.7), fill=PANEL, line=LINE, round_=True)
    box(s, Inches(0.45), y, Inches(0.075), Inches(0.7), fill=c_)
    step_circle(s, n, Inches(0.72), y + Inches(0.11), d=Inches(0.46), fill=c_)
    text(s, Inches(1.42), y + Inches(0.08), Inches(11.2), Inches(0.27), [[(t_, 12.2, DARK, True)]])
    text(s, Inches(1.42), y + Inches(0.38), Inches(11.2), Inches(0.28), [[(b_, 9.6, BODY, False)]])
    y += Inches(0.78)
callout(s, Inches(0.45), Inches(6.6), Inches(6.1), Inches(0.84), PANEL2, INDIGO,
        'Steps 2 and 3 are the real prerequisites',
        'They are not part of the run, and they are what makes a run mean anything afterwards.')
callout(s, Inches(6.78), Inches(6.6), Inches(6.1), Inches(0.84), GREENBG, GREEN,
        'Step 1 can start now',
        'Gathering our own rate data helps both options and does not wait on any decision about '
        'Option 2.')
notes(s, "Six steps for Option two. Step one is our own rate data and it is worth starting whatever "
         "we decide, because both options need it. Steps two and three are the prerequisites I "
         "showed you two slides ago: a history of rate changes, and saved costings. They are not "
         "part of the run itself, they are what makes a run mean anything three months later. Step "
         "four is small but it is not optional, because a basket costed at UK rates when it should "
         "have been Poland is just wrong. Step five is the actual new build. Step six is a trial on "
         "a real basket, run beside how we cost today so we can compare.")

# ────────────────────────────────────── 18 · THE TWO TOGETHER ───────────────
s = header('The two options together', 'Comparison')
rows = [
    ('What it gives us', ('The engineer stops typing', DARK, True), ('We can cost a whole basket', DARK, True)),
    ('Which system does the costing', ('CAPEE, unchanged', GREEN, True), ('CostVision engine, beside CAPEE', AMBER, False)),
    ('Does anyone learn a new tool', ('No, same CAPEE screens', GREEN, True), ('A reviewer does, for the results', AMBER, False)),
    ('How much software work', ('Moderate', AMBER, False), ('More: the run, plus two prerequisites', AMBER, False)),
    ('How much JLR data work', ('Parts with prices we paid', AMBER, False), ('Our full rate book', RED, False)),
    ('Does it tell us how accurate we are', ('Yes, that is the point of it', GREEN, True), ('It relies on that answer', AMBER, False)),
    ('Can we stop part way', ('Yes, nothing in CAPEE has changed', GREEN, True), ('Yes, it runs alongside', GREEN, True)),
    ('Is it waiting on an unknown', ('Yes, how CAPEE is built', AMBER, False), ('Yes, the accuracy figure', AMBER, False)),
    ('Is it audit-ready today', ('Yes, CAPEE keeps the record', GREEN, True), ('No, until rate history and saved costings exist', RED, True)),
]
table(s, Inches(0.45), Inches(1.9), Inches(12.43),
      ['', 'OPTION 1 · feed CAPEE', 'OPTION 2 · bulk costing'], rows,
      [Inches(4.13), Inches(4.15), Inches(4.15)], row_h=Inches(0.41), size=10.4)
callout(s, Inches(0.45), Inches(6.1), Inches(12.43), Inches(0.82), PANEL2, INDIGO,
        'Read the last two rows together',
        'Option 1 produces the accuracy figure that Option 2 leans on, and Option 1 is audit-ready '
        'today because CAPEE keeps the record. Option 2 is not, until the rate history and the '
        'saved costings exist. That is the whole argument for this order.')
footer(s, 'No durations shown. Nothing in either option has been sized by JLR.')
notes(s, "The top half is fairly even and each option wins some rows. The two I would look at are "
         "the last two. Option one produces the accuracy figure that Option two leans on. And "
         "Option one is audit-ready today, because CAPEE keeps the record exactly as it does now, "
         "whereas Option two is not audit-ready until we have built the rate history and the saved "
         "costings. Those two rows are the whole argument for doing them in this order.")

# ─────────────────────────────────────── 19 · RISKS ─────────────────────────
s = header('Risks and open items', 'Risks')
rows = [
    (('We do not know how accurate it is', DARK, True), 'Both',
     'Never compared against a price JLR has paid. The Option 1 trial is what settles it.',
     ('High', RED, True)),
    (('Old costings cannot be reproduced', DARK, True), 'Option 2',
     'No rate history and no saved costings. Has to be built before bulk is used for anything audited.',
     ('High', RED, True)),
    (('We do not know how CAPEE is built', DARK, True), 'Option 1',
     'Holds up the choice of connection route. One answer from IT sorts it.',
     ('High', AMBER, True)),
    (('The automated route has gaps', DARK, True), 'Option 2',
     'No region, no annual volume, no gear. Small work, but a basket is quietly wrong without it.',
     ('Medium', AMBER, True)),
    (('The built-in numbers are our estimates', DARK, True), 'Both',
     'Starting points, not JLR plant measurements. They get replaced as real data comes in.',
     ('Medium', AMBER, True)),
    (('Drawing text goes to an AI service', DARK, True), 'Both',
     'Short extracts only, and it can go through a JLR service. The CAD file itself stays on the server.',
     ('Medium', AMBER, True)),
    (('Server type', DARK, True), 'Option 1',
     'The measuring engine needs a normal Linux server. Easy to confirm now, awkward if we miss it.',
     ('Low', AMBER, True)),
    (('The tool is still changing', DARK, True), 'Both',
     'Six rounds of work in the last few weeks. For any trial we pin a version and stay on it.',
     ('Low', AMBER, True)),
]
table(s, Inches(0.45), Inches(1.9), Inches(12.43),
      ['Item', 'Applies to', 'Detail and what sorts it', 'Rating'], rows,
      [Inches(3.4), Inches(1.15), Inches(6.58), Inches(1.3)], row_h=Inches(0.52), size=9.6)
callout(s, Inches(0.45), Inches(6.6), Inches(12.43), Inches(0.84), PANEL2, INDIGO,
        'One decision deals with the top item',
        'Approving the Option 1 trial turns "we do not know how accurate it is" into a '
        'number. Everything else on this list is smaller and has a known fix.')
notes(s, "Eight items. The top one is the one that would stop me putting this in front of a "
         "supplier today, and approving the trial is what fixes it. The second is the one that "
         "would stop me signing off a bulk run for anything audited, and it is a build, not a "
         "mystery. The third is waiting on IT. The fourth is small but I have put it up because a "
         "basket costed at the wrong region is wrong in a way nobody would notice. And the last "
         "row is worth saying plainly: the tool is still being worked on, so for any trial we pin a "
         "version and stay on it. Note also that the CAD file itself never leaves the server.")

# ─────────────────────────────── 20 · WHAT WE NEED TO START ─────────────────
s = header('What we need to get going', 'Next steps')
asks = [
    ('1', 'clip', '30 to 50 parts with the price we actually paid', RED,
     'Part, description, annual volume, where it was made, and the price paid. Existing purchase '
     'records will do. This is what turns "we think it is about right" into a number, and both '
     'options rest on it.'),
    ('2', 'coins', 'Our own rate data', VIOLET,
     'Material prices, machine rates per hour, labour by grade, electricity and gas. Worth starting '
     'now whatever is decided, because both options need it and it is ours to gather.'),
    ('3', 'cog', 'How CAPEE is built', INDIGO,
     'What it is written in, whether it can call another service on the network, and what kind of '
     'server it runs on. Three questions to IT. Unblocks Option 1.'),
    ('4', 'shield', 'A decision on the drawing-reading step', GREEN,
     'Either switch off the AI reading and keep the measuring, or run it through a JLR-approved '
     'service. The data-flow document is written and covers both.'),
    ('5', 'person', 'A team and one part type for the trial', AMBER,
     'One cost engineer and one commodity family for the duration. Machined parts or pressings '
     'would be my choice.'),
]
y = Inches(1.82)
for n, ic, t_, c_, b_ in asks:
    box(s, Inches(0.45), y, Inches(12.43), Inches(0.9), fill=PANEL, line=LINE, round_=True)
    box(s, Inches(0.45), y, Inches(0.075), Inches(0.9), fill=c_)
    step_circle(s, n, Inches(0.72), y + Inches(0.1), d=Inches(0.46), fill=c_)
    icon_badge(s, ic, Inches(1.38), y + Inches(0.1), d=Inches(0.46), fill=c_)
    text(s, Inches(2.06), y + Inches(0.09), Inches(10.5), Inches(0.28), [[(t_, 12.4, DARK, True)]])
    text(s, Inches(2.06), y + Inches(0.4), Inches(10.4), Inches(0.44),
         [[(b_, 10.0, BODY, False)]], line_spacing=1.12)
    y += Inches(0.98)
callout(s, Inches(0.45), Inches(6.8), Inches(12.43), Inches(0.66), GREENBG, GREEN,
        'The first two do not wait on any decision made today',
        'Purchase records and our own rates are useful whichever way this goes.')
notes(s, "Five things. The first is the one I will keep coming back to: thirty to fifty parts where "
         "we know what we paid. Purchase records are fine, I do not need anything new created. "
         "Without it I cannot tell anyone how accurate this is, and that is the first question we "
         "will get asked. The second is our own rates, and it is worth starting whatever you decide "
         "today because both options need it. The third is three questions to IT. The fourth is a "
         "security decision and either answer works for us. The fifth is a person and a part type "
         "for the trial. The green strip matters: the first two do not wait on anything decided in "
         "this room.")

# ────────────────────────────── 21 · WHERE THE NUMBERS COME FROM ────────────
s = header('Where the numbers come from', 'Appendix')
text(s, Inches(0.45), Inches(1.68), Inches(12.4), Inches(0.32),
     [[('Every count and timing in this pack was read off the software or measured this month, and '
        'checked again when these slides were built.', 11.5, BODY, False)]])
rows = [
    ('19 part types cost through the engine', 'Traced each one from its screen through to the costing'),
    ('165 costing rules, 149 write into a named box', 'Counted in the rule packs for the twelve part types that have them; 9 to 18 apply per part'),
    ('12 to 69 boxes typed by hand today', 'Counted the input fields on the machining, casting, moulding and pressing forms'),
    ('The costing gives the same answer every run', 'Ran one part five times, identical to the last decimal'),
    ('Measuring makes no internet connection', 'Checked every line for outbound calls; there are none'),
    ('CAD files stay on the server', 'Measuring happens locally; only drawing text goes to an AI service'),
    ('0.9 to 3.1 seconds to measure a part', 'Five real production parts, 0.6 MB to 3.0 MB, this month'),
    ('40 parts in 31 seconds on two workers', 'Measured wall-clock, nothing reused from cache; larger rows carry that rate forward'),
    ('Six of eight sample parts costed straight through', 'With the material family given. Given nothing, all eight ask for it first'),
    ('Line counts on slide 7', 'Counted directly in the source files'),
    ('Settings values on slide 9', 'Read off the software; each is a single defined value'),
    ('328 materials, 178 machines, 42 labour rates, 20 regions', 'Counted in the rate library; rates effective June and July 2026'),
    ('2,079 automatic tests across 156 files', 'Full test run at build time, all passing'),
    ('The engine runs on its own, outside the application', 'Called it standalone; produced the full eight-bucket answer unchanged'),
]
table(s, Inches(0.45), Inches(2.15), Inches(12.43), ['Statement', 'How it was checked'], rows,
      [Inches(5.9), Inches(6.53)], row_h=Inches(0.29), size=9.5)
callout(s, Inches(0.45), Inches(6.78), Inches(12.43), Inches(0.66), AMBERBG, AMBER,
        'One number we deliberately do not show',
        'The tool has an accuracy report, but no real purchase prices are loaded into it yet. That is ask number one.')
notes(s, "Keep this one for anyone who wants to know how we know. Every claim was checked against "
         "the software or measured this month rather than assumed. Two worth pointing at. The "
         "costing gives the same answer every time, which sounds obvious but is not true of every "
         "tool. And the fifth row is the answer to the first question IT security will ask: we "
         "checked every line of the measuring software for internet connections and there are none. "
         "The amber box at the bottom is the one number I have deliberately left off every slide. "
         "The tool has an accuracy report, but nothing real has been loaded into it, so what it "
         "prints today means nothing. That is why ask number one is what it is.")

# ─────────────────────────────────────────────────────────────────────────────
OUT = 'CostVision-CAPEE-Implementation-Options.pptx'
prs.save(OUT)
finalise(OUT)


def assert_powerpoint_can_open(path):
    """Fail the build on the OOXML faults that make PowerPoint offer to repair.

    LibreOffice opens files PowerPoint rejects, so converting to PDF proves
    nothing about whether the deck will open on a colleague's laptop. A shape
    with a zero or negative extent is the fault this build actually hit: a short
    callout made an inner text box `h - chrome` wide, which went negative, and
    PowerPoint refused the file while LibreOffice rendered it happily.
    """
    import zipfile, re
    import xml.etree.ElementTree as ET
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    z = zipfile.ZipFile(path)
    faults = []
    if z.testzip() is not None:
        faults.append('zip archive is corrupt')
    for n in sorted(x for x in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', x)):
        root = ET.fromstring(z.read(n))
        for ext in root.iter('{%s}ext' % A):
            cx, cy = int(ext.get('cx', '1')), int(ext.get('cy', '1'))
            if cx <= 0 or cy <= 0:
                faults.append(f'{n}: shape extent cx={cx} cy={cy}')
        ids = [int(e.get('id')) for e in root.iter('{%s}cNvPr' % NS) if e.get('id')]
        for d in sorted({i for i in ids if ids.count(i) > 1}):
            faults.append(f'{n}: duplicate shape id {d}')
    if faults:
        raise SystemExit('DECK IS INVALID - PowerPoint would ask to repair it:\n  '
                         + '\n  '.join(faults))
    return len([x for x in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', x)])


n_slides = assert_powerpoint_can_open(OUT)
print(f'{OUT}  -  {n_slides} slides, validated')
