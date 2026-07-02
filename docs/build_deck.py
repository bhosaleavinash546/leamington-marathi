#!/usr/bin/env python3
"""CostVision — management roadmap deck (light theme, plain language)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── Palette (light, professional) ─────────────────────────────────────────────
INK      = RGBColor(0x1F, 0x2A, 0x44)   # deep slate text
BLUE     = RGBColor(0x2A, 0x6A, 0xF2)   # primary
BLUE_DK  = RGBColor(0x1D, 0x4E, 0xD8)
TEAL     = RGBColor(0x0E, 0x9E, 0xA6)
GREEN    = RGBColor(0x1F, 0xA9, 0x6B)
AMBER    = RGBColor(0xE0, 0x8A, 0x1E)
RED      = RGBColor(0xD1, 0x4B, 0x4B)
PURPLE   = RGBColor(0x7C, 0x4D, 0xD6)
GREY     = RGBColor(0x63, 0x6B, 0x7A)
LIGHT    = RGBColor(0xF3, 0xF6, 0xFC)   # slide bg tint
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BD  = RGBColor(0xD9, 0xE2, 0xF1)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def box(s, x, y, w, h, fill=CARD, line=CARD_BD, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs, str): runs = [(runs, 18, INK, False)]
    first = True
    for item in runs:
        t, sz, col, bold = (item + (False,))[:4] if len(item) < 4 else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
        r.font.name = "Segoe UI"
        first = False
    return tb

def bullets(s, x, y, w, h, items, size=15, col=INK, gap=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple): txt, c, b = (it + (col, False))[:3] if len(it)>=2 else (it[0], col, False)
        else: txt, c, b = it, col, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  " + txt
        r.font.size = Pt(size); r.font.color.rgb = c; r.font.bold = b; r.font.name = "Segoe UI"
    return tb

def header(s, kicker, title, num):
    bar = box(s, 0, 0, 13.333, 0.10, fill=BLUE, line=None)
    text(s, 0.55, 0.30, 11.5, 0.4, [(kicker, 12, BLUE, True)])
    text(s, 0.55, 0.58, 11.8, 0.9, [(title, 27, INK, True)])
    text(s, 12.5, 0.30, 0.7, 0.4, [(f"{num:02d}", 12, GREY, True)], align=PP_ALIGN.RIGHT)

def chip(s, x, y, w, label, color, txtcol=WHITE, h=0.42, size=11):
    c = box(s, x, y, w, h, fill=color, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y+0.03, w, h-0.04, [(label, size, txtcol, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

def arrow(s, x, y, w=0.5, h=0.4, color=BLUE):
    a = box(s, x, y, w, h, fill=color, line=None, shape=MSO_SHAPE.RIGHT_ARROW)
    return a

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=RGBColor(0xFF,0xFF,0xFF), line=None)
box(s, 0, 0, 13.333, 2.6, fill=RGBColor(0xEA,0xF1,0xFE), line=None, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 2.55, 13.333, 0.08, fill=BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.8, 0.7, 11.7, 0.5, [("COSTVISION  ·  SHOULD-COST INTELLIGENCE PLATFORM", 14, BLUE, True)])
text(s, 0.8, 1.15, 11.7, 1.2, [("Bringing CostVision In-House", 40, INK, True)])
text(s, 0.8, 2.75, 11.7, 0.6, [("A plain-English roadmap to deploy the tool securely inside our organisation", 19, GREY, False)])
# three highlight cards
cards = [("100% On-Premise", "Runs inside our firewall — our data stays with us", GREEN),
         ("No External AI", "AI can be switched off or kept fully in-house", BLUE),
         ("Secure by Design", "Accounts, encryption, and audit built in", PURPLE)]
for i,(t,d,c) in enumerate(cards):
    x = 0.8 + i*4.05
    box(s, x, 3.9, 3.75, 2.4, fill=CARD, line=CARD_BD, line_w=1.2)
    box(s, x, 3.9, 3.75, 0.14, fill=c, line=None)
    text(s, x+0.25, 4.2, 3.3, 0.6, [(t, 17, INK, True)])
    text(s, x+0.25, 4.85, 3.3, 1.3, [(d, 13, GREY, False)])
text(s, 0.8, 6.7, 11.7, 0.4, [("Prepared for management review  ·  Implementation & data-security roadmap", 12, GREY, False)])
notes(s, "Opening slide. Set the frame: this is about taking a tool that already works and deploying it safely inside our own organisation, so our data never leaves the building. Three promises up front — it runs on our own servers, it does not depend on any outside AI service, and security is built in from day one. Keep it reassuring: we are not buying a risky external product; we are bringing something we control in-house.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is CostVision (context)
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "THE CONTEXT", "What CostVision does for us", 2)
text(s, 0.55, 1.6, 7.2, 0.9, [("It answers one question, fast and consistently:", 17, INK, True),
                              ("“What should this part or software cost to make?”", 17, BLUE, True)])
bullets(s, 0.55, 2.7, 7.2, 3.6, [
    ("Prices manufactured parts — machining, casting, sheet metal, moulding, PCB and more", INK, False),
    ("Prices automotive software programmes (the 6-step guided estimator)", INK, False),
    ("Turns a photo of a circuit board into a full, priced parts list", INK, False),
    ("Gives buyers and engineers one consistent, defensible number", INK, False),
    ("Replaces slow spreadsheets and gut-feel estimates", INK, False),
], size=15, gap=10)
# right visual: simple value funnel
bx = 8.1
box(s, bx, 1.7, 4.6, 4.7, fill=CARD, line=CARD_BD, line_w=1.2)
text(s, bx, 1.9, 4.6, 0.4, [("From inputs to a trusted price", 14, INK, True)], align=PP_ALIGN.CENTER)
steps = [("Part / software details", BLUE), ("CostVision engine", TEAL), ("Cost breakdown + benchmarks", GREEN), ("Decision-ready number", PURPLE)]
yy = 2.5
for i,(t,c) in enumerate(steps):
    chip(s, bx+0.5, yy, 3.6, t, c, h=0.55, size=13)
    if i < len(steps)-1:
        a = box(s, bx+2.15, yy+0.56, 0.3, 0.34, fill=GREY, line=None, shape=MSO_SHAPE.DOWN_ARROW)
    yy += 0.95
notes(s, "Remind the audience what the tool actually delivers, in business terms — not features, but outcomes. The key message: it gives us one consistent, defensible cost number, quickly, instead of every engineer or buyer estimating differently in their own spreadsheet. The little funnel on the right shows the simple idea: details go in, a trusted price comes out. This sets up WHY it is worth deploying properly.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Where does our data live today?
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "TODAY'S SITUATION", "Where our data lives right now", 3)
# left: on our server (good)
box(s, 0.55, 1.7, 6.0, 4.9, fill=RGBColor(0xEC,0xF7,0xF0), line=RGBColor(0xBF,0xE3,0xCE), line_w=1.2)
text(s, 0.8, 1.9, 5.5, 0.5, [("✅  Stays on our own server", 16, GREEN, True)])
bullets(s, 0.8, 2.55, 5.5, 3.9, [
    "User accounts (passwords are encrypted / hashed)",
    "Saved projects & cost scenarios — kept per user",
    "Rate library, supplier quotes, parts lists",
    "Uploaded board photos are held in memory only — never saved to disk",
], size=13.5, gap=9)
# right: goes outside (attention)
box(s, 6.75, 1.7, 6.0, 4.9, fill=RGBColor(0xFD,0xF0,0xEE), line=RGBColor(0xF2,0xCF,0xCA), line_w=1.2)
text(s, 7.0, 1.9, 5.5, 0.5, [("⚠️  Currently leaves our network", 16, RED, True)])
bullets(s, 7.0, 2.55, 5.5, 3.9, [
    ("Today the AI features send data to an outside AI service (Anthropic)", RED, True),
    "That includes board photos and cost details sent for AI analysis",
    "Optional: live parts-pricing and a news feed also reach the internet",
    "This is the one thing we must change before roll-out",
], size=13.5, gap=9)
notes(s, "This is the honest 'as-is' picture and the most important slide for a security conversation. Left side (green): the good news — the core data already lives on our own server in a single database file, passwords are encrypted, and uploaded photos are never even written to disk. Right side (red): the catch — the AI features currently send data to an outside AI provider. That is exactly what our policy does not allow. Be upfront: this is fixable, and the next slide shows how. Do not hide this; management will trust the plan more because we surfaced it.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The one big decision: external AI
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "THE KEY DECISION", "The AI question — and our options", 4)
text(s, 0.55, 1.55, 12.2, 0.5, [("Our rule: no external AI apps. Good news — the cost calculations don't need AI at all. AI is an optional helper.", 15, INK, True)])
opts = [
    ("Option A", "Switch AI OFF", ["Simplest & fully compliant","All cost engines still work","Lose: AI chat, photo-to-parts,\nAI summaries"], GREEN, "Recommended to start"),
    ("Option B", "Run AI in-house", ["Keep AI, nothing leaves our network","Needs our own AI server (hardware)","More setup effort"], BLUE, "Best long-term"),
    ("Option C", "Approved internal AI", ["Use a company-approved AI service","Keeps AI features","Depends on what IT allows"], PURPLE, "If available"),
]
for i,(tag,title,pts,c,foot) in enumerate(opts):
    x = 0.55 + i*4.15
    box(s, x, 2.25, 3.85, 4.05, fill=CARD, line=CARD_BD, line_w=1.2)
    box(s, x, 2.25, 3.85, 0.6, fill=c, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, 2.30, 3.85, 0.5, [(f"{tag}: {title}", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb = s.shapes.add_textbox(Inches(x+0.25), Inches(3.0), Inches(3.4), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for j,p in enumerate(pts):
        par = tf.paragraphs[0] if j==0 else tf.add_paragraph(); par.space_after=Pt(8)
        r=par.add_run(); r.text="•  "+p; r.font.size=Pt(13); r.font.color.rgb=INK; r.font.name="Segoe UI"
    chip(s, x+0.25, 5.75, 3.35, foot, c, h=0.42, size=11)
notes(s, "Frame the single most important decision simply. Our policy says no external AI apps. The reassuring headline: the actual cost calculations are done entirely by our own maths — they do not need AI. AI is only a convenience layer (a chat helper, reading a photo into a parts list, writing summaries). So we have three clean choices. Recommend starting with Option A (switch AI off) to get compliant and live quickly, then move to Option B (our own in-house AI) later if the team wants those helpers back. Option C if IT already offers an approved internal AI. Emphasise: none of these options block the core value of the tool.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Target architecture (flowchart)
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "THE TARGET SETUP", "How it will run — all inside our walls", 5)
def node(x,y,w,h,label,c,sub=None):
    box(s,x,y,w,h,fill=CARD,line=c,line_w=2.0)
    box(s,x,y,0.16,h,fill=c,line=None)
    text(s,x+0.28,y+ (0.14 if sub else h/2-0.2),w-0.4,0.5,[(label,13.5,INK,True)])
    if sub: text(s,x+0.28,y+0.55,w-0.4,0.5,[(sub,10.5,GREY,False)])
# firewall band
box(s, 3.35, 1.55, 9.45, 5.35, fill=RGBColor(0xEE,0xF4,0xFF), line=RGBColor(0xC7,0xD8,0xF5), line_w=1.5)
text(s, 3.5, 1.62, 9.0, 0.35, [("🔒  Inside our corporate network / firewall", 12.5, BLUE_DK, True)])
node(0.55,3.15,2.5,1.2,"Team browsers",GREY,"Engineers & buyers")
node(3.7,3.15,2.7,1.2,"Secure gateway",TEAL,"HTTPS + firewall + limits")
node(6.9,2.35,2.7,1.1,"CostVision app",BLUE,"The tool + calculations")
node(6.9,3.75,2.7,1.1,"Our database",GREEN,"Accounts & saved work")
node(10.0,2.35,2.6,1.1,"In-house AI",PURPLE,"Optional (Option B)")
node(10.0,3.75,2.6,1.1,"Encrypted backups",AMBER,"Nightly, retained")
# arrows
for (x1,y) in [(3.05,3.75)]: arrow(s, 3.08, 3.55, 0.55, 0.4, BLUE)
arrow(s, 6.42, 3.55, 0.42, 0.4, BLUE)
box(s, 8.05, 2.55, 0.42, 0.35, fill=PURPLE, line=None, shape=MSO_SHAPE.RIGHT_ARROW)
box(s, 8.05, 3.95, 0.42, 0.35, fill=GREEN, line=None, shape=MSO_SHAPE.RIGHT_ARROW)
text(s, 0.55, 5.9, 12.0, 0.8, [("Nothing needs the public internet. Optional external feeds (news, live pricing) can be turned off for a fully sealed install.", 13, GREY, False)])
notes(s, "Walk through the picture left to right in plain language. Our people open the tool in a normal web browser. Their request goes through a secure gateway (this handles encryption, blocks attacks, and limits abuse). Behind it sits the CostVision app, which does the calculations, and our own database, which holds accounts and saved work. Optionally, an in-house AI box and encrypted nightly backups. The big blue box is the key point: everything lives inside our firewall. Nothing has to touch the public internet — and any optional outside feeds can be switched off entirely for a fully sealed, air-gapped install.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Security controls (grid)
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "KEEPING IT SAFE", "How we protect data, accounts & the system", 6)
text(s, 0.55, 1.5, 12.2, 0.4, [("Green = already built in.   Blue = we add for production.", 13, GREY, True)])
items = [
 ("Encrypted passwords", "Stored scrambled, never in plain text", GREEN, "Built in"),
 ("Private per-user data", "You only see your own projects (tested)", GREEN, "Built in"),
 ("Photos not saved to disk", "Board images used in memory only", GREEN, "Built in"),
 ("Attack protection", "Security headers + request rate limits", GREEN, "Built in"),
 ("Company login (SSO)", "Use our existing corporate sign-in + MFA", BLUE, "To add"),
 ("Roles & permissions", "Admin / engineer / viewer access levels", BLUE, "To add"),
 ("Encryption everywhere", "HTTPS in transit, encrypted storage at rest", BLUE, "To add"),
 ("Audit trail", "Record who did what, sent to IT monitoring", BLUE, "To add"),
 ("Backups & recovery", "Nightly encrypted backups + restore drills", BLUE, "To add"),
]
cols=3; cw=4.0; ch=1.35; x0=0.55; y0=2.05
for i,(t,d,c,tag) in enumerate(items):
    r=i//cols; col=i%cols
    x=x0+col*(cw+0.15); y=y0+r*(ch+0.15)
    box(s,x,y,cw,ch,fill=CARD,line=CARD_BD,line_w=1.0)
    box(s,x,y,0.13,ch,fill=c,line=None)
    text(s,x+0.28,y+0.12,cw-0.4,0.4,[(t,13,INK,True)])
    text(s,x+0.28,y+0.55,cw-0.4,0.6,[(d,10.8,GREY,False)])
    text(s,x+cw-1.15,y+0.1,1.0,0.3,[(tag,9,c,True)],align=PP_ALIGN.RIGHT)
notes(s, "This is the reassurance slide for the security-minded. Do not read every box — summarise: a lot of protection is ALREADY built into the tool (the green items — encrypted passwords, private per-user data that we have actually tested, photos never written to disk, and standard attack protection). The blue items are the standard enterprise additions we will layer on for production: logging in with our existing company account and MFA, proper role-based access, encryption end to end, an audit trail feeding IT's monitoring, and reliable backups. Message: this follows normal corporate security practice — no surprises.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Roadmap (timeline)
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "THE PLAN", "Phased roll-out — low risk, step by step", 7)
phases = [
 ("Phase 0","Make it compliant","Switch AI to off or in-house; seal external feeds","~1 wk", BLUE),
 ("Phase 1","Secure pilot","Install on one server, HTTPS, backups; 3–5 pilot users","2–3 wks", TEAL),
 ("Phase 2","Company login & roles","Corporate sign-on, MFA, access levels, audit log","3–4 wks", GREEN),
 ("Phase 3","Scale & harden","Security test, wider roll-out, resilience","Ongoing", PURPLE),
]
# timeline line
box(s, 0.9, 3.55, 11.5, 0.06, fill=RGBColor(0xC7,0xD3,0xE6), line=None)
n=len(phases); span=11.2/n
for i,(tag,title,desc,dur,c) in enumerate(phases):
    cx = 1.0 + i*span + span/2
    # dot
    box(s, cx-0.16, 3.42, 0.32, 0.32, fill=c, line=WHITE, line_w=2.0, shape=MSO_SHAPE.OVAL)
    up = (i%2==0)
    cy = 1.85 if up else 4.05
    box(s, cx-1.35, cy, 2.7, 1.55, fill=CARD, line=CARD_BD, line_w=1.2)
    box(s, cx-1.35, cy, 2.7, 0.12, fill=c, line=None)
    text(s, cx-1.15, cy+0.16, 2.3, 0.35, [(tag+"  ·  "+dur, 11, c, True)])
    text(s, cx-1.15, cy+0.5, 2.3, 0.4, [(title, 13.5, INK, True)])
    text(s, cx-1.15, cy+0.92, 2.3, 0.6, [(desc, 10.5, GREY, False)])
    # connector
    box(s, cx-0.02, (3.4 if up else 3.6), 0.04, (0.05 if up else 0.45) if not up else 0.05, fill=c, line=None)
text(s, 0.55, 6.35, 12.2, 0.6, [("We are live and compliant after Phase 1. Everything after that adds convenience and scale — not core capability.", 13.5, INK, True)])
notes(s, "The plan, told as four simple steps on a timeline. Phase 0 (about a week): make it compliant — switch the AI off or in-house and seal any external feeds. Phase 1 (2–3 weeks): a proper secure install on one server with HTTPS and backups, tried by a small pilot group. Phase 2 (3–4 weeks): plug into our normal company login with MFA and set who can see what. Phase 3: ongoing hardening, a security test, and wider roll-out. The key reassurance at the bottom: we are already live and compliant after Phase 1 — the later phases add polish and scale, not core function, so value arrives early and risk stays low.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Benefits vs Considerations
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "THE TRADE-OFFS", "Benefits vs. things to consider", 8)
box(s, 0.55, 1.7, 6.0, 4.95, fill=RGBColor(0xEC,0xF7,0xF0), line=RGBColor(0xBF,0xE3,0xCE), line_w=1.2)
text(s, 0.8, 1.9, 5.5, 0.5, [("✅  Benefits", 18, GREEN, True)])
bullets(s, 0.8, 2.6, 5.5, 3.9, [
    "Our data never leaves our network",
    "Consistent, defensible costs across the team",
    "Faster quoting and stronger supplier negotiation",
    "No per-user external software fees",
    "We control updates, access and security",
    "Already tested & reliable (automated checks on every change)",
], size=13.5, gap=11)
box(s, 6.75, 1.7, 6.0, 4.95, fill=RGBColor(0xFC,0xF4,0xE9), line=RGBColor(0xF0,0xDD,0xBC), line_w=1.2)
text(s, 7.0, 1.9, 5.5, 0.5, [("⚠️  Things to consider", 18, AMBER, True)])
bullets(s, 7.0, 2.6, 5.5, 3.9, [
    "Needs a server and IT support to host it",
    "AI helpers need a decision (off, in-house, or approved)",
    "In-house AI (if chosen) needs extra hardware",
    "Someone owns updates, backups & access reviews",
    "Rate data should be reviewed for accuracy over time",
], size=13.5, gap=11)
notes(s, "Be balanced and honest — management trusts a plan that names its costs. Left: the benefits, led by the one that matters most for this audience — our data stays with us. Then the business wins: consistent numbers, faster quoting, no per-seat external fees, and full control. Right: the honest considerations — it needs a server and IT ownership, the AI question needs a decision, in-house AI would need hardware, and someone must own upkeep and keep the rate data current. Nothing here is a blocker; these are normal ownership costs for any internal system. Presenting both sides builds credibility.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Simple cost/benefit bar chart
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "WHY IT'S WORTH IT", "Effort now vs. value returned", 9)
chart_data = CategoryChartData()
chart_data.categories = ['Phase 0\nCompliant', 'Phase 1\nPilot live', 'Phase 2\nEnterprise', 'Phase 3\nScaled']
chart_data.add_series('Effort', (2, 4, 5, 3))
chart_data.add_series('Value delivered', (4, 8, 9, 10))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.7), Inches(8.0), Inches(4.9), chart_data)
chart = gframe.chart
chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.TOP; chart.legend.include_in_layout = False
chart.has_title = False
plot = chart.plots[0]; plot.gap_width = 90
chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = RGBColor(0xC9,0xD6,0xEC)
chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = BLUE
for ax in [chart.category_axis, chart.value_axis]:
    ax.tick_labels.font.size = Pt(11); ax.tick_labels.font.name = "Segoe UI"
# side takeaways
box(s, 9.0, 1.9, 3.75, 4.4, fill=CARD, line=CARD_BD, line_w=1.2)
text(s, 9.2, 2.1, 3.4, 0.4, [("Takeaways", 15, INK, True)])
bullets(s, 9.2, 2.65, 3.4, 3.5, [
    ("Value climbs faster than effort", GREEN, True),
    "Biggest jump is early (going live)",
    "Later phases keep adding value at modest effort",
    "One-time setup, long-term payoff",
], size=12.5, gap=10)
notes(s, "A simple visual to make the investment case without spreadsheets. The light bars are the effort in each phase; the blue bars are the value we get. The story the chart tells: value rises faster than effort, and the biggest single jump comes early — when the pilot goes live. After that, each phase keeps adding value for modest extra effort. Bottom line for management: this is a mostly one-time setup with a long-term payoff, and we start getting return almost immediately.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Decisions we need + next steps
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "WHAT WE NEED FROM YOU", "Decisions & next steps", 10)
box(s, 0.55, 1.7, 6.0, 4.9, fill=CARD, line=CARD_BD, line_w=1.2)
box(s, 0.55, 1.7, 6.0, 0.55, fill=BLUE, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.55, 1.75, 6.0, 0.45, [("3 decisions we need", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.85, 2.5, 5.5, 3.9, [
    ("1.  AI approach: OFF, in-house, or approved internal?", INK, True),
    ("2.  Who hosts it — which server / IT team owns it?", INK, True),
    ("3.  Pilot group: which 3–5 people go first?", INK, True),
], size=15, gap=16)
box(s, 6.75, 1.7, 6.0, 4.9, fill=CARD, line=CARD_BD, line_w=1.2)
box(s, 6.75, 1.7, 6.0, 0.55, fill=GREEN, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 6.75, 1.75, 6.0, 0.45, [("Immediate next steps", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 7.05, 2.5, 5.5, 3.9, [
    "Approve Phase 0 (make it compliant)",
    "IT & security review the plan",
    "Stand up one server for the pilot",
    "Confirm pilot users & success measures",
    "Target: pilot live within ~4 weeks",
], size=14, gap=13)
text(s, 0.55, 6.75, 12.2, 0.5, [("Ask: approval to begin Phase 0 and name a pilot group. Low cost, low risk, quick to value.", 14, BLUE_DK, True)], align=PP_ALIGN.CENTER)
notes(s, "Close with a clear ask so the meeting ends in a decision, not just discussion. Left: the three decisions we genuinely need from them — the AI approach, who hosts it, and who is in the pilot. Right: the concrete next steps and a realistic target of a live pilot in about four weeks. End on the one-line ask at the bottom: approve Phase 0 and name a pilot group. Reinforce the theme — low cost, low risk, quick to value, and our data stays ours. Then open for questions.")

out = "/home/user/leamington-marathi/CostVision-Implementation-Roadmap.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
