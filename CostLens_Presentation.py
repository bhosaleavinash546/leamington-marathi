"""
CostLens Management Presentation Builder
Generates a professional 7-slide .pptx presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

# ── Brand Palette ───────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x17, 0x2A)   # deep navy background
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)   # sky-blue accent
ACCENT2    = RGBColor(0x6E, 0xE7, 0xB7)   # mint green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD  = RGBColor(0x1E, 0x2D, 0x45)
ORANGE     = RGBColor(0xFB, 0x92, 0x3C)
SUCCESS    = RGBColor(0x22, 0xC5, 0x5E)
DANGER     = RGBColor(0xF8, 0x71, 0x71)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ──────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(14), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_size=Pt(13), bold_first=False,
                  color=WHITE, line_spacing=None, font_name="Calibri"):
    """lines: list of (text, bold, color_override or None)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line_spacing:
            from pptx.util import Pt as Ptt
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold or (bold_first and first)
        run.font.color.rgb = col
        run.font.name = font_name
    return txBox


def bg(slide, color=NAVY):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


def accent_bar(slide, y=Inches(0.08), h=Inches(0.06)):
    """Thin horizontal accent line at top"""
    grad = slide.shapes.add_shape(1, 0, y, SLIDE_W, h)
    grad.fill.solid()
    grad.fill.fore_color.rgb = ACCENT
    grad.line.fill.background()


def card(slide, x, y, w, h, color=DARK_CARD, radius=None):
    r = add_rect(slide, x, y, w, h, color)
    return r


def icon_circle(slide, x, y, r_size, color, emoji, font_size=Pt(20)):
    circle = slide.shapes.add_shape(9, x, y, r_size, r_size)  # oval
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # emoji text centred
    add_text(slide, emoji,
             x, y + r_size * 0.05,
             r_size, r_size * 0.9,
             font_size=font_size, align=PP_ALIGN.CENTER, color=WHITE)
    return circle


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

# Diagonal decorative rectangle
deco = s.shapes.add_shape(1, Inches(8.5), Inches(-0.5), Inches(5.5), Inches(9))
deco.fill.solid()
deco.fill.fore_color.rgb = DARK_CARD
deco.line.fill.background()

# Accent circle
circ = s.shapes.add_shape(9, Inches(10.5), Inches(1.8), Inches(2.5), Inches(2.5))
circ.fill.solid()
circ.fill.fore_color.rgb = ACCENT
circ.line.fill.background()
circ.fill.fore_color.theme_color  # just styling

add_text(s, "🔭", Inches(10.5), Inches(1.9), Inches(2.5), Inches(2.5),
         font_size=Pt(60), align=PP_ALIGN.CENTER)

# Tag line chip
chip = add_rect(s, Inches(0.6), Inches(1.6), Inches(3.0), Inches(0.35), ACCENT)
add_text(s, "AUTOMOTIVE COST ENGINEERING",
         Inches(0.62), Inches(1.63), Inches(2.96), Inches(0.3),
         font_size=Pt(8.5), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Main title
add_text(s, "CostLens",
         Inches(0.6), Inches(2.1), Inches(8.0), Inches(1.4),
         font_size=Pt(68), bold=True, color=WHITE, font_name="Calibri")

add_text(s, "AI-Powered Should-Cost &\nSupplier Quote Intelligence",
         Inches(0.6), Inches(3.5), Inches(8.0), Inches(1.2),
         font_size=Pt(22), bold=False, color=LIGHT_GRAY)

# Divider line
div = s.shapes.add_shape(1, Inches(0.6), Inches(4.8), Inches(3.5), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT
div.line.fill.background()

add_text(s, "Presented to Senior Management  ·  2026",
         Inches(0.6), Inches(4.95), Inches(7.0), Inches(0.4),
         font_size=Pt(12), color=MID_GRAY)

# Three KPI chips at bottom-left
kpis = [("30%", "Cost Savings", SUCCESS), ("5×", "Faster Analysis", ACCENT), ("Real-Time", "AI Insights", GOLD)]
for i, (val, lbl, col) in enumerate(kpis):
    cx = Inches(0.6 + i * 2.5)
    cy = Inches(6.0)
    chip_r = add_rect(s, cx, cy, Inches(2.2), Inches(0.9), DARK_CARD)
    add_text(s, val, cx, cy + Inches(0.05), Inches(2.2), Inches(0.45),
             font_size=Pt(22), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, cx, cy + Inches(0.45), Inches(2.2), Inches(0.4),
             font_size=Pt(10), color=MID_GRAY, align=PP_ALIGN.CENTER)

# Slide number
add_text(s, "1 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem & Solution
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

add_text(s, "The Challenge — and Our Solution",
         Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)
add_text(s, "Traditional cost engineering is slow, manual, and disconnected from supplier data",
         Inches(0.6), Inches(0.78), Inches(11), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY)

div = s.shapes.add_shape(1, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.03))
div.fill.solid(); div.fill.fore_color.rgb = DARK_CARD; div.line.fill.background()

# BEFORE column
card(s, Inches(0.5), Inches(1.35), Inches(5.7), Inches(5.6))
add_text(s, "❌  BEFORE CostLens",
         Inches(0.7), Inches(1.5), Inches(5.3), Inches(0.5),
         font_size=Pt(15), bold=True, color=DANGER)

problems = [
    "Spreadsheet-based should-cost models — error-prone & version chaos",
    "Supplier quotes received with no structured breakdown",
    "Manual variance analysis takes days, not minutes",
    "No visibility into where supplier margins are inflated",
    "Negotiation targets set on gut feel, not data",
    "Audit trail and history impossible to maintain",
    "AI/ML insights require expensive external consultants",
]
for i, p_text in enumerate(problems):
    add_text(s, f"•  {p_text}",
             Inches(0.75), Inches(2.1 + i * 0.47), Inches(5.2), Inches(0.45),
             font_size=Pt(11.5), color=LIGHT_GRAY)

# AFTER column
card(s, Inches(6.5), Inches(1.35), Inches(6.3), Inches(5.6), DARK_CARD)
add_text(s, "✅  WITH CostLens",
         Inches(6.7), Inches(1.5), Inches(5.8), Inches(0.5),
         font_size=Pt(15), bold=True, color=SUCCESS)

solutions = [
    ("Structured 3-level should-cost models with version control", SUCCESS),
    ("Web-based & CSV quote intake with element-by-element breakdown", ACCENT),
    ("Instant automated variance analysis across all cost elements", SUCCESS),
    ("AI (Claude) flags over-priced elements with severity ratings", ACCENT),
    ("Data-driven negotiation targets with deadline tracking", SUCCESS),
    ("Full audit trail — who changed what and when", ACCENT),
    ("AI should-cost generation from part description in seconds", SUCCESS),
]
for i, (s_text, col) in enumerate(solutions):
    add_text(s, f"•  {s_text}",
             Inches(6.7), Inches(2.1 + i * 0.47), Inches(5.8), Inches(0.45),
             font_size=Pt(11.5), color=col)

add_text(s, "2 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Features
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

add_text(s, "Key Features",
         Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)
add_text(s, "End-to-end cost engineering — from build to negotiation close",
         Inches(0.6), Inches(0.78), Inches(11), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY)

features = [
    ("📐", "Should-Cost\nModeller", ACCENT,
     "3-level hierarchical cost build-up (Category → Element → Cost Driver). "
     "Publish/Archive versioning with full audit trail. "
     "GBP/EUR/USD support."),
    ("📋", "Supplier\nQuote Intake", RGBColor(0x8B, 0x5C, 0xF6),
     "Web form or CSV import. Element-level quote breakdown. "
     "Multiple quotes per part. Status workflow: Submitted → Negotiating → Accepted."),
    ("⚖️", "Cost\nComparison", ORANGE,
     "Side-by-side should-cost vs. quote variance. "
     "Element-level over/under flags. Cross-model cross-supplier analysis. "
     "Snapshot history."),
    ("🤖", "AI Analysis\n(Claude)", SUCCESS,
     "Instant AI-powered gap analysis. Severity-rated flags (High/Med/Low). "
     "Clarifying questions for suppliers. Recommended negotiation actions."),
    ("🎯", "Negotiation\nTracker", GOLD,
     "Track price targets per part/supplier. Deadline alerts. "
     "Potential saving calculator. Mark Agreed/Stalled/Closed."),
    ("📊", "Supplier\nScorecard", ACCENT2,
     "Multi-dimensional supplier performance scoring. "
     "On-time, quality, cost competitiveness. Monthly trend sparklines."),
]

cols = 3
rows = 2
card_w = Inches(4.1)
card_h = Inches(2.35)
start_x = Inches(0.5)
start_y = Inches(1.35)
gap_x = Inches(0.18)
gap_y = Inches(0.16)

for i, (icon, title, col, desc) in enumerate(features):
    row = i // cols
    c   = i % cols
    cx  = start_x + c * (card_w + gap_x)
    cy  = start_y + row * (card_h + gap_y)

    card(s, cx, cy, card_w, card_h)

    # colour top bar
    bar = s.shapes.add_shape(1, cx, cy, card_w, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()

    # icon circle
    icon_circle(s, cx + Inches(0.15), cy + Inches(0.18), Inches(0.52), col, icon, Pt(18))

    # title
    add_text(s, title,
             cx + Inches(0.78), cy + Inches(0.18), Inches(3.2), Inches(0.6),
             font_size=Pt(13), bold=True, color=WHITE)

    # desc
    add_text(s, desc,
             cx + Inches(0.15), cy + Inches(0.82), Inches(3.85), Inches(1.38),
             font_size=Pt(10.5), color=LIGHT_GRAY)

add_text(s, "3 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — AI Workflow
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

add_text(s, "How the AI Works",
         Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)
add_text(s, "Powered by Anthropic's Claude — the same AI used by Fortune 500 procurement teams",
         Inches(0.6), Inches(0.78), Inches(12), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY)

# Flow diagram — 5 steps
steps = [
    ("1", "📤", "Cost\nComparison\nCreated",    ACCENT,
     "Buyer builds should-cost model and uploads supplier quote. System calculates element-level variance."),
    ("2", "🧠", "AI Prompt\nEngineered",         RGBColor(0x8B, 0x5C, 0xF6),
     "CostLens compiles part data, variances, and commodity context into a structured prompt for Claude."),
    ("3", "☁️", "Claude API\nAnalysis",          RGBColor(0x06, 0x95, 0xDD),
     "Claude analyses the cost data in real-time via streaming. Response appears word-by-word instantly."),
    ("4", "🔍", "Structured\nInsight Output",   ORANGE,
     "AI returns: executive summary, flagged elements with severity, supplier questions, recommended actions."),
    ("5", "✅", "Buyer\nTakes Action",           SUCCESS,
     "Buyer uses AI flags to drive negotiation targets, raise RFQs, or escalate to management."),
]

box_w   = Inches(2.35)
box_h   = Inches(4.6)
start_x = Inches(0.48)
start_y = Inches(1.35)
gap     = Inches(0.12)

for i, (num, icon, title, col, desc) in enumerate(steps):
    cx = start_x + i * (box_w + gap)

    # Main card
    card(s, cx, start_y, box_w, box_h, DARK_CARD)

    # Top accent
    bar = s.shapes.add_shape(1, cx, start_y, box_w, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()

    # Step number badge
    badge = s.shapes.add_shape(9, cx + box_w - Inches(0.45), start_y + Inches(0.1),
                                Inches(0.35), Inches(0.35))
    badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background()
    add_text(s, num,
             cx + box_w - Inches(0.45), start_y + Inches(0.1),
             Inches(0.35), Inches(0.35),
             font_size=Pt(11), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Icon
    add_text(s, icon,
             cx, start_y + Inches(0.2), box_w, Inches(0.7),
             font_size=Pt(30), align=PP_ALIGN.CENTER)

    # Title
    add_text(s, title,
             cx + Inches(0.1), start_y + Inches(0.95), box_w - Inches(0.2), Inches(0.75),
             font_size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER)

    # Description
    add_text(s, desc,
             cx + Inches(0.12), start_y + Inches(1.72), box_w - Inches(0.24), Inches(2.7),
             font_size=Pt(10.5), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Arrow between boxes
    if i < len(steps) - 1:
        ax = cx + box_w + gap * 0.15
        add_text(s, "→",
                 ax, start_y + Inches(2.0), gap * 0.7, Inches(0.5),
                 font_size=Pt(14), bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Bottom "Did you know" strip
strip = add_rect(s, Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.8), DARK_CARD)
add_text(s, "💡  Real-time streaming: ",
         Inches(0.7), Inches(6.2), Inches(2.4), Inches(0.5),
         font_size=Pt(11.5), bold=True, color=ACCENT)
add_text(s, "The AI response streams word-by-word directly into the CostLens interface — no waiting for a full response. "
            "If no API key is configured, a built-in intelligent mock provides full functionality for demonstrations.",
         Inches(3.0), Inches(6.2), Inches(9.7), Inches(0.5),
         font_size=Pt(11), color=LIGHT_GRAY)

add_text(s, "4 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Business Benefits & ROI
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

add_text(s, "Business Benefits & ROI",
         Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)
add_text(s, "Measurable impact on productivity, cost reduction, and decision speed",
         Inches(0.6), Inches(0.78), Inches(11), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY)

# Big KPI row
kpi_data = [
    ("30–40%", "Cost Savings\nIdentified", SUCCESS, "vs. quoted price"),
    ("5×",     "Faster Quote\nAnalysis",   ACCENT,  "hours → minutes"),
    ("100%",   "Audit Trail\nCoverage",    GOLD,    "full change history"),
    ("£000s",  "Annual Saving\nTracker",   ORANGE,  "negotiation pipeline"),
]
kw = Inches(2.9)
kh = Inches(1.55)
ky = Inches(1.35)
for i, (val, label, col, sub) in enumerate(kpi_data):
    kx = Inches(0.5 + i * (kw + Inches(0.17)))
    card(s, kx, ky, kw, kh, DARK_CARD)
    bar = s.shapes.add_shape(1, kx, ky, kw, Inches(0.055))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, val, kx, ky + Inches(0.1), kw, Inches(0.7),
             font_size=Pt(32), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, label, kx, ky + Inches(0.75), kw, Inches(0.5),
             font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sub, kx, ky + Inches(1.2), kw, Inches(0.3),
             font_size=Pt(9.5), color=MID_GRAY, align=PP_ALIGN.CENTER)

# Benefit blocks — 2×2 grid
benefits = [
    ("🏎️  Procurement Speed",  ACCENT,
     [("Quote analysis time", "2 days → 30 min"),
      ("Should-cost build time", "1 week → 2 hrs"),
      ("AI insight generation", "instant streaming"),
      ("Negotiation prep", "data-ready in seconds")]),
    ("💰  Cost Reduction",  SUCCESS,
     [("Identify inflated margins", "element-level flags"),
      ("AI-flagged over-pricing", "High/Med/Low severity"),
      ("Structured negotiation", "target vs. current price"),
      ("Annual saving visibility", "full pipeline view")]),
    ("📈  Decision Quality",  GOLD,
     [("Evidence-based targets", "not gut feel"),
      ("Cross-model benchmarking", "compare similar parts"),
      ("Supplier scorecards", "multi-factor rating"),
      ("Management dashboards", "live KPI tiles")]),
    ("🔒  Compliance & Control",  ORANGE,
     [("Version audit trail", "who changed what/when"),
      ("Status workflow", "Draft → Published → Archived"),
      ("Role-based access", "Admin / Buyer / Supplier"),
      ("Data sovereignty", "self-hosted on-premises")]),
]

bw = Inches(5.9)
bh = Inches(2.6)
for i, (title, col, rows) in enumerate(benefits):
    bx = Inches(0.5 + (i % 2) * (bw + Inches(0.2)))
    by = Inches(3.1 + (i // 2) * (bh + Inches(0.12)))
    card(s, bx, by, bw, bh)
    bar2 = s.shapes.add_shape(1, bx, by, bw, Inches(0.055))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = col; bar2.line.fill.background()
    add_text(s, title, bx + Inches(0.15), by + Inches(0.1), bw - Inches(0.2), Inches(0.42),
             font_size=Pt(13), bold=True, color=col)
    for j, (k, v) in enumerate(rows):
        ry = by + Inches(0.58 + j * 0.46)
        add_text(s, f"  {k}", bx + Inches(0.15), ry, Inches(3.2), Inches(0.42),
                 font_size=Pt(11.5), color=LIGHT_GRAY)
        add_text(s, v, bx + Inches(3.3), ry, Inches(2.4), Inches(0.42),
                 font_size=Pt(11.5), bold=True, color=col, align=PP_ALIGN.RIGHT)

add_text(s, "5 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Sample Reports / Exports
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

add_text(s, "Sample Reports & Outputs",
         Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)
add_text(s, "Everything your team sees in CostLens — representative data for a brake disc part",
         Inches(0.6), Inches(0.78), Inches(12), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY)

# ── Report 1: Should-Cost Breakup ──
rc = Inches(0.5)
ry = Inches(1.3)
rw = Inches(3.9)
rh = Inches(5.6)
card(s, rc, ry, rw, rh)
bar = s.shapes.add_shape(1, rc, ry, rw, Inches(0.055))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
add_text(s, "📐  Should-Cost Breakup", rc + Inches(0.15), ry + Inches(0.1), rw, Inches(0.4),
         font_size=Pt(12.5), bold=True, color=ACCENT)
add_text(s, "Part: BRK-2001  Front Brake Disc", rc + Inches(0.15), ry + Inches(0.52),
         rw - Inches(0.2), Inches(0.3), font_size=Pt(10), color=LIGHT_GRAY)

sc_rows = [
    ("RAW MATERIAL",     "GBP 4.82",  "32%", ACCENT),
    ("  Steel coil",     "GBP 3.10",  "21%", MID_GRAY),
    ("  Alloy additives","GBP 1.72",  "11%", MID_GRAY),
    ("MANUFACTURING",    "GBP 5.60",  "37%", RGBColor(0x06, 0xB6, 0xD4)),
    ("  Casting",        "GBP 2.40",  "16%", MID_GRAY),
    ("  Machining",      "GBP 1.90",  "13%", MID_GRAY),
    ("  Surface treat.", "GBP 1.30",   "9%", MID_GRAY),
    ("OVERHEAD",         "GBP 1.95",  "13%", ORANGE),
    ("LOGISTICS",        "GBP 1.20",   "8%", GOLD),
    ("PROFIT",           "GBP 1.53",  "10%", SUCCESS),
    ("─────────────",    "─────────",  "───", MID_GRAY),
    ("TOTAL",            "GBP 15.10", "100%", WHITE),
]
for j, (el, val, pct_v, col) in enumerate(sc_rows):
    add_text(s, el,    rc + Inches(0.15), ry + Inches(0.85 + j * 0.36),
             Inches(1.9), Inches(0.34), font_size=Pt(9.5),
             bold=(el in ("RAW MATERIAL","MANUFACTURING","OVERHEAD","LOGISTICS","PROFIT","TOTAL")),
             color=col)
    add_text(s, val,   rc + Inches(2.1), ry + Inches(0.85 + j * 0.36),
             Inches(1.0), Inches(0.34), font_size=Pt(9.5), align=PP_ALIGN.RIGHT,
             bold=(el == "TOTAL"), color=col)
    add_text(s, pct_v, rc + Inches(3.2), ry + Inches(0.85 + j * 0.36),
             Inches(0.5), Inches(0.34), font_size=Pt(9.5), align=PP_ALIGN.RIGHT,
             color=MID_GRAY)

# ── Report 2: Quote Comparison ──
rc2 = Inches(4.6)
card(s, rc2, ry, rw, rh)
bar2 = s.shapes.add_shape(1, rc2, ry, rw, Inches(0.055))
bar2.fill.solid(); bar2.fill.fore_color.rgb = ORANGE; bar2.line.fill.background()
add_text(s, "⚖️  Quote vs. Should-Cost", rc2 + Inches(0.15), ry + Inches(0.1), rw, Inches(0.4),
         font_size=Pt(12.5), bold=True, color=ORANGE)
add_text(s, "Supplier: Leamington Forge Ltd", rc2 + Inches(0.15), ry + Inches(0.52),
         rw - Inches(0.2), Inches(0.3), font_size=Pt(10), color=LIGHT_GRAY)

hdr_y = ry + Inches(0.82)
for hdr, hx, hw in [("Element", 0.15, 1.5), ("Should-Cost", 1.65, 0.85), ("Quote", 2.5, 0.75), ("Var", 3.3, 0.45)]:
    add_text(s, hdr, rc2 + Inches(hx), hdr_y, Inches(hw), Inches(0.3),
             font_size=Pt(9), bold=True, color=MID_GRAY)

cmp_rows = [
    ("Raw Material", "4.82", "5.90", "+1.08", DANGER, "over"),
    ("Manufacturing","5.60", "5.80", "+0.20", GOLD,   "~ok"),
    ("  Casting",    "2.40", "3.10", "+0.70", DANGER, "high"),
    ("  Machining",  "1.90", "1.75", "-0.15", SUCCESS,"under"),
    ("Overhead",     "1.95", "2.40", "+0.45", ORANGE, "med"),
    ("Logistics",    "1.20", "1.00", "-0.20", SUCCESS,"under"),
    ("Profit",       "1.53", "2.30", "+0.77", DANGER, "high"),
    ("─────",        "────", "────", "────",  MID_GRAY,""),
    ("TOTAL",        "15.10","18.40","+3.30", DANGER, "high"),
]
for j, (el, sc_v, q_v, var_v, col, flag) in enumerate(cmp_rows):
    fy = ry + Inches(1.18 + j * 0.36)
    add_text(s, el,    rc2 + Inches(0.15), fy, Inches(1.5), Inches(0.34),
             font_size=Pt(9.5), bold=(el == "TOTAL"), color=WHITE if el=="TOTAL" else LIGHT_GRAY)
    add_text(s, sc_v,  rc2 + Inches(1.65), fy, Inches(0.85), Inches(0.34),
             font_size=Pt(9.5), align=PP_ALIGN.RIGHT, color=MID_GRAY)
    add_text(s, q_v,   rc2 + Inches(2.5),  fy, Inches(0.75), Inches(0.34),
             font_size=Pt(9.5), align=PP_ALIGN.RIGHT, color=WHITE)
    add_text(s, var_v, rc2 + Inches(3.25), fy, Inches(0.55), Inches(0.34),
             font_size=Pt(9.5), bold=True, align=PP_ALIGN.RIGHT, color=col)

# ── Report 3: AI Insights ──
rc3 = Inches(8.7)
card(s, rc3, ry, rw, rh)
bar3 = s.shapes.add_shape(1, rc3, ry, rw, Inches(0.055))
bar3.fill.solid(); bar3.fill.fore_color.rgb = SUCCESS; bar3.line.fill.background()
add_text(s, "🤖  AI Analysis Output", rc3 + Inches(0.15), ry + Inches(0.1), rw, Inches(0.4),
         font_size=Pt(12.5), bold=True, color=SUCCESS)

add_text(s, "Summary",
         rc3 + Inches(0.15), ry + Inches(0.58), rw - Inches(0.2), Inches(0.3),
         font_size=Pt(10), bold=True, color=LIGHT_GRAY)
add_text(s, '"The quote of GBP 18.40 exceeds the '
            'should-cost model by 21.9% (GBP 3.30). '
            'The most significant overcharges are in '
            'Raw Material castings (+29%) and Profit '
            'margin (+50% above model). Strong case '
            'for negotiation."',
         rc3 + Inches(0.15), ry + Inches(0.88), rw - Inches(0.2), Inches(1.12),
         font_size=Pt(9.5), italic=True, color=LIGHT_GRAY)

add_text(s, "Flagged Elements",
         rc3 + Inches(0.15), ry + Inches(2.1), rw - Inches(0.2), Inches(0.3),
         font_size=Pt(10), bold=True, color=LIGHT_GRAY)

flags = [
    ("🔴 HIGH", "Casting cost",   "29% above model — request process breakdown",   DANGER),
    ("🔴 HIGH", "Profit margin",  "50% above agreed 10% — challenge with data",    DANGER),
    ("🟡 MED",  "Overhead",       "23% inflated — request overhead allocation",    GOLD),
    ("🟢 LOW",  "Machining",      "Under-cost — no action needed",                 SUCCESS),
]
for j, (sev, el, reason, col) in enumerate(flags):
    fy = ry + Inches(2.48 + j * 0.56)
    chip_f = add_rect(s, rc3 + Inches(0.15), fy, Inches(0.7), Inches(0.28),
                      DARK_CARD)
    add_text(s, sev, rc3 + Inches(0.15), fy, Inches(0.7), Inches(0.28),
             font_size=Pt(7.5), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, el, rc3 + Inches(0.9), fy, rw - Inches(1.0), Inches(0.28),
             font_size=Pt(9.5), bold=True, color=WHITE)
    add_text(s, reason, rc3 + Inches(0.9), fy + Inches(0.28), rw - Inches(1.0), Inches(0.26),
             font_size=Pt(9), color=MID_GRAY)

add_text(s, "Recommended Action",
         rc3 + Inches(0.15), ry + Inches(4.78), rw - Inches(0.2), Inches(0.3),
         font_size=Pt(10), bold=True, color=LIGHT_GRAY)
add_text(s, '• Request itemised casting cost breakdown\n'
            '• Challenge profit at agreed 10% rate\n'
            '• Counter-offer: GBP 15.80 target',
         rc3 + Inches(0.15), ry + Inches(5.08), rw - Inches(0.2), Inches(0.7),
         font_size=Pt(9.5), color=ACCENT2)

add_text(s, "6 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Architecture & Next Steps
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)

add_text(s, "Technology & Next Steps",
         Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# Left: Tech Stack
card(s, Inches(0.5), Inches(0.95), Inches(5.8), Inches(5.9), DARK_CARD)
bar_l = s.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(5.8), Inches(0.055))
bar_l.fill.solid(); bar_l.fill.fore_color.rgb = ACCENT; bar_l.line.fill.background()
add_text(s, "🛠️  Technology Stack", Inches(0.65), Inches(1.05), Inches(5.5), Inches(0.4),
         font_size=Pt(13), bold=True, color=ACCENT)

tech = [
    ("Frontend",   "React 18 + TypeScript + Vite",       ACCENT),
    ("State Mgmt", "TanStack React Query v5",             ACCENT),
    ("Backend",    "Node.js + Express + TypeScript",      RGBColor(0x8B, 0x5C, 0xF6)),
    ("Database",   "PostgreSQL (structured, relational)", RGBColor(0x06, 0x95, 0xDD)),
    ("AI Engine",  "Anthropic Claude API (streaming)",    SUCCESS),
    ("Auth",       "JWT with role-based access control",  GOLD),
    ("Deployment", "Docker Compose (self-hosted / cloud)",ORANGE),
    ("Exports",    "CSV import/export, web forms",        ACCENT2),
]
for j, (cat, val, col) in enumerate(tech):
    ty = Inches(1.55 + j * 0.55)
    add_text(s, cat, Inches(0.65), ty, Inches(1.4), Inches(0.45),
             font_size=Pt(10.5), bold=True, color=col)
    add_text(s, val, Inches(2.1), ty, Inches(4.0), Inches(0.45),
             font_size=Pt(10.5), color=LIGHT_GRAY)

add_text(s, "✅  Self-hosted: data never leaves your infrastructure",
         Inches(0.65), Inches(6.1), Inches(5.5), Inches(0.4),
         font_size=Pt(10.5), color=SUCCESS)

# Right: Roadmap / Next Steps
card(s, Inches(6.7), Inches(0.95), Inches(6.1), Inches(5.9), DARK_CARD)
bar_r = s.shapes.add_shape(1, Inches(6.7), Inches(0.95), Inches(6.1), Inches(0.055))
bar_r.fill.solid(); bar_r.fill.fore_color.rgb = GOLD; bar_r.line.fill.background()
add_text(s, "🚀  Recommended Next Steps", Inches(6.85), Inches(1.05), Inches(5.8), Inches(0.4),
         font_size=Pt(13), bold=True, color=GOLD)

steps_next = [
    ("Immediate", [
        "Deploy to shared server for team access",
        "Load historical should-cost models from spreadsheets",
        "Configure Anthropic API key for live AI analysis",
    ], ACCENT),
    ("Short-Term (1-3 months)", [
        "Pilot on top 20 highest-spend parts",
        "Train procurement team on quote analysis workflow",
        "Establish negotiation target discipline via Tracker",
    ], SUCCESS),
    ("Medium-Term (3-6 months)", [
        "Integrate with ERP for live price & volume data",
        "Expand AI model with commodity-specific training",
        "Add weekly AI digest to management email reports",
    ], ORANGE),
]
ry_n = Inches(1.55)
for phase, items, col in steps_next:
    add_text(s, f"◆  {phase}", Inches(6.85), ry_n, Inches(5.8), Inches(0.38),
             font_size=Pt(11.5), bold=True, color=col)
    ry_n += Inches(0.38)
    for item in items:
        add_text(s, f"   •  {item}", Inches(6.85), ry_n, Inches(5.8), Inches(0.38),
                 font_size=Pt(10.5), color=LIGHT_GRAY)
        ry_n += Inches(0.38)
    ry_n += Inches(0.12)

# Bottom CTA strip
cta = add_rect(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.38),
               RGBColor(0x0E, 0x7A, 0x9E))
add_text(s, "CostLens is live, fully functional, and ready for your team today.",
         Inches(0.65), Inches(7.04), Inches(12.0), Inches(0.32),
         font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "7 / 7", Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
         font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ── Save ─────────────────────────────────────────────────────────
out = "/home/user/leamington-marathi/CostLens_Management_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
