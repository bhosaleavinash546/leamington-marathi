#!/usr/bin/env python3
"""
CostVision and CAPEE — the business case in four slides, for senior management.

  1  Our should-cost today: capable, but hands-on. What we do well, where the
     manual effort sits, and the value model with JLR's numbers left blank.
  2  Option 1 — automate the data entry inside CAPEE. 3D model only, no AI, so
     it needs no AI approval and can start now.
  3  Option 2 — what CostVision does. Six capabilities, stated as what they add
     rather than how they are built.
  4  What we get from each. Benefits in kind, never in pounds.
  5  What changes for the business, what each option needs, and the ask.

RULES THIS FILE FOLLOWS.

1. NO AI IN OPTION 1. Option 1 reads the 3D model only. Drawing reading needs a
   language model and is out of scope for it; the tool's default setting is
   already "rules only, no AI call" and that is the setting slide 2 describes.
   This is the point that lets Option 1 start under today's policy, and slide 2
   says so on its face. Option 2 does use an AI assistant, and the approval to
   use AI is listed as the first thing it needs — stated, not buried.

2. THE AI SETS UP, IT NEVER PRICES. Slide 3 says this in plain words. The
   assistant picks the process route and the machine and fills what the
   geometry gives; the money is still worked out by `computeUniversalStack` on
   the same rate library. If that ever stops being true the slide is wrong.

3. NO INVENTED SAVING. Every benefit figure on slide 4 is an empty box for JLR
   to fill, with the arithmetic printed beside it. We have never timed an
   engineer costing a part in CAPEE and never compared a CostVision cost with a
   price JLR paid.

4. NO MACHINE TIMINGS. An earlier draft put throughput seconds on the Option 2
   slide. It answered a question nobody in that room was asking. Capability
   belongs there; the timings live in the 21-slide pack.

5. ONLY WHAT IS IN THE SOFTWARE, verified September 2026: 12 to 69 input boxes
   by commodity; the quote teardown runs live (a machined part quoted 28% above
   should-cost returned the gap by bucket, recovery levers and supplier
   questions); the comparison table shows ten countries side by side out of
   twenty in the rate library; calibration from actuals, scenario, sensitivity,
   landed cost and the DFM savings levers are all wired through to the UI.

Regenerate:  python3 build_capee_business_case_pptx.py
Output:      CostVision-CAPEE-Business-Case.pptx
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pptx_fixup import finalise

# ── Palette: identical to build_blueprint_pptx.py so the pack reads as one ────
INDIGO  = RGBColor(0x1D, 0x6F, 0xB8)
BLUE    = RGBColor(0x1D, 0x6F, 0xB8)
DARK    = RGBColor(0x16, 0x32, 0x5C)
BODY    = RGBColor(0x3A, 0x43, 0x56)
MUTED   = RGBColor(0x6B, 0x72, 0x80)
BG      = RGBColor(0xF4, 0xF7, 0xFB)
PANEL   = RGBColor(0xFF, 0xFF, 0xFF)
PANEL2  = RGBColor(0xE8, 0xF1, 0xFA)
GREENBG = RGBColor(0xEA, 0xF6, 0xEF)
AMBERBG = RGBColor(0xFC, 0xF3, 0xE3)
REDBG   = RGBColor(0xFB, 0xEC, 0xEA)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
AMBER   = RGBColor(0xB7, 0x79, 0x1F)
RED     = RGBColor(0xB0, 0x3A, 0x2E)
VIOLET  = RGBColor(0x6B, 0x3F, 0xA0)
LINE    = RGBColor(0xDC, 0xE3, 0xEE)
NAVY    = RGBColor(0x16, 0x32, 0x5C)
ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)
HERO_SUB = RGBColor(0xCA, 0xDC, 0xFC)
HERO_DIM = RGBColor(0x8F, 0xA3, 0xCC)

TITLE_FONT = 'Cambria'

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def box(slide, x, y, w, h, fill=None, line=None, round_=False, radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _h(v, minimum=Inches(0.12)):
    """Clamp a derived height to a positive value.

    Several helpers below size an inner text box as `h - <chrome>`. When a
    caller passes a short card the result goes zero or negative, which is
    invalid OOXML: PowerPoint refuses the file and offers to repair it, while
    LibreOffice silently tolerates it. A text box does not clip its contents, so
    clamping changes nothing visually.
    """
    return v if v > minimum else minimum


_PICTO = re.compile('([\U0001F300-\U0001FAFF☀-⛿✀-✒✙-➿️]+)')
EMOJI_FONT = 'Segoe UI Emoji'


def _emit_runs(p, t, size, color, bold, italic, base_font='Calibri'):
    for part in _PICTO.split(t):
        if not part:
            continue
        run = p.add_run(); run.text = part
        f = run.font
        f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
        if _PICTO.fullmatch(part): f.name = EMOJI_FONT
        elif base_font: f.name = base_font


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for r in para:
            t, size, color, bold = r[0], r[1], r[2], r[3]
            italic = r[4] if len(r) > 4 else False
            _emit_runs(p, t, size, color, bold, italic, base_font=font)
    return tb


def logo(slide, x=Inches(0.35), y=Inches(0.22), scale=1.0, on_dark=False):
    s = scale
    badge = box(slide, x, y, Inches(0.42 * s), Inches(0.42 * s),
                fill=ON_DARK if on_dark else INDIGO, round_=True, radius=0.28)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'cv'
    r.font.size = Pt(17 * s); r.font.bold = True; r.font.name = 'Calibri'
    r.font.color.rgb = NAVY if on_dark else ON_DARK
    text(slide, x + Inches(0.52 * s), y - Inches(0.03 * s), Inches(2.6), Inches(0.32),
         [[('CostVision', 18 * s, ON_DARK if on_dark else BLUE, True)]])
    # No 'AI COST INTELLIGENCE' strapline in this deck: slide 2's headline is that
    # Option 1 involves no AI, and a strapline arguing the opposite in the corner of
    # every slide is an own goal in front of an audience where AI is not yet approved.


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def header(title, kicker=None):
    slide = prs.slides.add_slide(BLANK)
    box(slide, 0, 0, W, H, fill=BG)
    logo(slide)
    box(slide, 0, Inches(0.78), W, Pt(2.2), fill=NAVY)
    if kicker:
        text(slide, Inches(0.45), Inches(0.95), Inches(11.5), Inches(0.3),
             [[(kicker.upper(), 11, BLUE, True)]])
        ty = Inches(1.22)
    else:
        ty = Inches(1.02)
    text(slide, Inches(0.45), ty, Inches(12.4), Inches(0.6), [[(title, 27, DARK, True)]],
         font=TITLE_FONT)
    return slide


def card(slide, x, y, w, h, accent, title, lines, title_size=13, body_size=10.5,
         fill=PANEL, title_color=None):
    """White card with a coloured left rule — the pack's workhorse block."""
    box(slide, x, y, w, h, fill=fill, line=LINE, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    text(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.3),
         [[(title, title_size, title_color or DARK, True)]])
    if lines:
        runs = [[(ln[0], body_size, ln[1] if len(ln) > 1 else BODY,
                  ln[2] if len(ln) > 2 else False)] for ln in lines]
        text(slide, x + Inches(0.22), y + Inches(0.48), w - Inches(0.4), _h(h - Inches(0.6)),
             runs, space_after=3, line_spacing=1.12)


def table(slide, x, y, w, cols, rows, col_w, head_fill=NAVY, row_h=Inches(0.34),
          size=9.5, head_size=9.5):
    """Simple native table drawn from shapes — keeps full colour control."""
    cx = x
    for i, c in enumerate(cols):
        box(slide, cx, y, col_w[i], Inches(0.36), fill=head_fill)
        text(slide, cx + Inches(0.09), y + Inches(0.09), col_w[i] - Inches(0.16), Inches(0.24),
             [[(c, head_size, ON_DARK, True)]])
        cx += col_w[i]
    ry = y + Inches(0.36)
    for r_i, row in enumerate(rows):
        cx = x
        band = PANEL if r_i % 2 == 0 else BG
        for i, cell in enumerate(row):
            val, colr, bold = (cell if isinstance(cell, tuple) else (cell, BODY, False))
            box(slide, cx, ry, col_w[i], row_h, fill=band, line=LINE)
            text(slide, cx + Inches(0.09), ry + Inches(0.07), col_w[i] - Inches(0.16),
                 row_h - Inches(0.1), [[(val, size, colr, bold)]])
            cx += col_w[i]
        ry += row_h
    return ry


def callout(slide, x, y, w, h, fill, accent, title, body, tsize=12, bsize=10.5):
    box(slide, x, y, w, h, fill=fill, line=None, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    text(slide, x + Inches(0.24), y + Inches(0.13), w - Inches(0.45), Inches(0.28),
         [[(title, tsize, accent, True)]])
    text(slide, x + Inches(0.24), y + Inches(0.44), w - Inches(0.45), _h(h - Inches(0.55)),
         [[(body, bsize, BODY, False)]], space_after=3, line_spacing=1.14)


def footer(slide, txt):
    text(slide, Inches(0.45), Inches(7.02), Inches(12.4), Inches(0.24),
         [[(txt, 8, MUTED, False)]])


# ── Icons ────────────────────────────────────────────────────────────────────
# 12 of the 16 icons in assets/workflow-deck/icons are WHITE artwork on a
# transparent background — measured, average RGB 255,255,255. Dropped straight
# onto a light card they are invisible. So a white icon always sits inside a
# filled circle. check / times / warn / arrow are coloured and work bare.
import os
ICON_DIR = 'assets/workflow-deck/icons'
_COLOURED_ICONS = {'check', 'times', 'warn', 'arrow'}


def icon_badge(slide, name, cx, cy, d=Inches(0.62), fill=INDIGO, pad=0.26):
    """Coloured circle with a white icon centred inside it."""
    path = os.path.join(ICON_DIR, f'{name}.png')
    if name not in _COLOURED_ICONS:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
        c.fill.solid(); c.fill.fore_color.rgb = fill
        c.line.fill.background(); c.shadow.inherit = False
    if os.path.exists(path):
        inset = int(d * pad)
        slide.shapes.add_picture(path, cx + inset, cy + inset, d - 2 * inset, d - 2 * inset)


def step_circle(slide, n, cx, cy, d=Inches(0.5), fill=INDIGO, size=17):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.margin_left = tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(size); r.font.bold = True; r.font.name = 'Calibri'
    r.font.color.rgb = ON_DARK


def chevron(slide, x, y, w, h, label, sub, fill, text_col=ON_DARK):
    """One block of a left-to-right process flow."""
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = text_col; r.font.name = 'Calibri'
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(8.5); r2.font.color.rgb = text_col; r2.font.name = 'Calibri'
    return shp


def flow_step(slide, x, y, w, h, icon, title, sub, accent):
    """Icon-topped step card used in the end-to-end summary slides."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, w, Inches(0.06), fill=accent)
    icon_badge(slide, icon, x + (w - Inches(0.62)) / 2, y + Inches(0.18), fill=accent)
    text(slide, x + Inches(0.08), y + Inches(0.92), w - Inches(0.16), Inches(0.34),
         [[(title, 10.5, DARK, True)]], align=PP_ALIGN.CENTER)
    text(slide, x + Inches(0.08), y + Inches(1.28), w - Inches(0.16), _h(h - Inches(1.34)),
         [[(sub, 8.8, BODY, False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)


def arrow_between(slide, x, y, w=Inches(0.3)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.22))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xB8, 0xC4, 0xD4)
    a.line.fill.background(); a.shadow.inherit = False


def lane(slide, x, y, w, h, label, colour, items, label_w=Inches(1.5)):
    """Swim-lane: who does what."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE)
    box(slide, x, y, label_w, h, fill=colour)
    tf_box = box(slide, x, y, label_w, h, fill=None)
    tf = tf_box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = ON_DARK; r.font.name = 'Calibri'
    cw = (w - label_w) / len(items)
    for i, it in enumerate(items):
        cx = x + label_w + cw * i
        if i:
            box(slide, cx, y + Inches(0.06), Pt(0.75), h - Inches(0.12), fill=LINE)
        text(slide, cx + Inches(0.12), y + Inches(0.13), cw - Inches(0.24), h - Inches(0.26),
             [[(it, 9.2, BODY, False)]], line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)


def chip(slide, x, y, w, h, label, formula, accent):
    """A named formula. The label says what it is, the line under it says how it
    is worked out — so the room can check the working rather than trust a total."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, w, Inches(0.06), fill=accent)
    text(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.26),
         [[(label, 10.8, DARK, True)]])
    text(slide, x + Inches(0.18), y + Inches(0.42), w - Inches(0.36), Inches(0.3),
         [[(formula, 13.5, accent, True)]])


def band(slide, x, y, w, h, icon, title, body, accent, title_size=12.5, body_size=10.2):
    """Full-width row: icon, heading, one line of plain explanation."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    icon_badge(slide, icon, x + Inches(0.27), y + (h - Inches(0.6)) / 2, d=Inches(0.6), fill=accent)
    text(slide, x + Inches(1.1), y + Inches(0.16), w - Inches(1.35), Inches(0.28),
         [[(title, title_size, DARK, True)]])
    text(slide, x + Inches(1.1), y + Inches(0.47), w - Inches(1.35), _h(h - Inches(0.55)),
         [[(body, body_size, BODY, False)]], line_spacing=1.13)


BLANKBOX = '________'

# ═══════════ 1 · OUR SHOULD-COST TODAY: CAPABLE, BUT HANDS-ON ═══════════════
s = header('Our should-cost today: capable, but hands-on', 'Where we are')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('We already cost parts properly. On the left is that job as it runs today. On the right is '
        'what it is worth, in our own numbers.', 11.5, BODY, False)]])

card(s, Inches(0.45), Inches(1.96), Inches(6.1), Inches(4.5), AMBER,
     'OUR SHOULD-COST TODAY · capable, but hands-on',
     [('Trusted and thorough', GREEN, True),
      ('A rigorous bottom-up method, across every commodity.',),
      ('',),
      ('Hands-on setup', DARK, True),
      ('Material, process and machine chosen by a person,',),
      ('part by part.',),
      ('',),
      ('Manual data entry', DARK, True),
      ('Geometry and tolerances read off the 3D model and the',),
      ('drawings, then typed in — 12 to 69 values per part.',),
      ('',),
      ('Updated by hand', DARK, True),
      ('Rates and inputs refreshed manually as economics move.',)])

text(s, Inches(6.78), Inches(1.98), Inches(6.1), Inches(0.3),
     [[('What it is worth: five numbers, four of them ours', 11.5, MUTED, True)]])
rows = [
    (('A', DARK, True), 'Minutes to cost one part today', (BLANKBOX, INDIGO, True)),
    (('B', DARK, True), 'Of those, minutes reading the model and typing', (BLANKBOX, INDIGO, True)),
    (('C', DARK, True), 'Parts we cost in a year', (BLANKBOX, INDIGO, True)),
    (('D', DARK, True), 'Parts in a basket we never cost individually', (BLANKBOX, INDIGO, True)),
    (('E', DARK, True), 'Share of those values the software fills itself', (BLANKBOX, GREEN, True)),
]
table(s, Inches(6.78), Inches(2.26), Inches(6.1),
      ['', 'What it is', 'Our number'], rows,
      [Inches(0.45), Inches(3.95), Inches(1.7)], row_h=Inches(0.34), size=9.8)

chips = [('Hours spent typing today', 'B × C ÷ 60', RED),
         ('Hours freed a year', 'B × E × C ÷ 60', GREEN),
         ('Extra parts that time would cost', 'B × E × C ÷ A', INDIGO)]
cy = Inches(4.44)
for lbl, fml, col in chips:
    box(s, Inches(6.78), cy, Inches(6.1), Inches(0.52), fill=PANEL, line=LINE, round_=True)
    box(s, Inches(6.78), cy, Inches(0.075), Inches(0.52), fill=col)
    text(s, Inches(7.02), cy + Inches(0.12), Inches(3.5), Inches(0.28),
         [[(lbl, 10.5, DARK, True)]])
    text(s, Inches(10.55), cy + Inches(0.10), Inches(2.2), Inches(0.3),
         [[(fml, 12.5, col, True)]], align=PP_ALIGN.RIGHT)
    cy += Inches(0.56)
# The assumption behind the middle row, said rather than left to be discovered.
text(s, Inches(6.78), Inches(6.12), Inches(6.1), Inches(0.3),
     [[('Assumes typing time falls in step with the values the software fills; the trial tests '
        'that.', 8.5, MUTED, False)]], line_spacing=1.1)

callout(s, Inches(0.45), Inches(6.56), Inches(12.43), Inches(0.82), GREENBG, GREEN,
        'The step forward: the same defensible should-cost, with the manual work automated',
        'Nothing about the method changes. The judgement, the build-up and the numbers we would '
        'defend to a supplier all stay as they are. What we are proposing to remove is the typing, '
        'and then the limit on how many parts we can get through.')
notes(s, "I want to start by saying what we do well, because none of this is a criticism of it. The "
         "left-hand card is our should-cost as it runs today. The method is trusted and thorough: a "
         "rigorous bottom-up build-up, applied across every commodity we buy, and it is the number "
         "we would defend in front of a supplier. Nothing I am about to propose changes that. What "
         "I want to talk about is the effort behind it. The setup is hands-on: a person chooses the "
         "material, the process and the machine, part by part. The data entry is manual: somebody "
         "opens the 3D model and the drawings, reads the geometry and the tolerances, and types "
         "them in. And when the economics move, somebody goes back and updates the rates by hand. "
         "The line at the bottom of that card is the size of it: between twelve and sixty-nine "
         "values typed per part depending on the commodity, and those counts come straight off our "
         "own input forms. None of that is wrong. It is just slow, and it is the reason we cost the "
         "parts we have time for rather than the parts we would like to. Now the right-hand side, "
         "and I want to be straight about why it is empty. I have not put a saving on this slide, "
         "because we have never timed this job and we have never compared this tool against a price "
         "we have actually paid. Any number I put there would be my estimate dressed up as a "
         "measurement. Four of these five are ours and we can get them quickly: how long a part "
         "takes, how much of that is reading and typing, how many parts we cost a year, and how "
         "many in a basket we never cost at all. Time a handful of parts in one commodity and you "
         "have the first two; the others are in our own records and the programme data. The fifth, "
         "E, is what the trial measures. Put those numbers in and the three boxes underneath fill "
         "themselves in: the hours we spend typing today, the hours we would get back, and how many "
         "more parts that time would cost. That is a business case we own rather than one I have "
         "handed you. And the green strip is the whole idea in one line: same method, same "
         "defensible answer, with the manual work taken out of it.")

# ═══════════ 2 · OPTION 1 · AUTOMATE THE DATA ENTRY, INSIDE CAPEE ═══════════
s = header('Option 1: automate the data entry, inside CAPEE', 'Option 1 · start now')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('CAPEE keeps doing the costing. The engineer stops typing the input and starts checking it. '
        'Reads the 3D model only, so no AI is involved.', 11.5, BODY, False)]])
steps = [
    ('upload', 'Attach the 3D model', 'STEP, IGES or STL, from the CAPEE screen', INDIGO),
    ('ruler', 'The software measures it', 'On a JLR server, with no internet connection', INDIGO),
    ('cog', 'The values fill in', 'Fixed rules. Same part, same numbers, always', INDIGO),
    ('person', 'The engineer confirms', 'Only what a 3D model cannot show', AMBER),
    ('calc', 'CAPEE costs the part', 'Exactly as it does now. It just gets the numbers', GREEN),
]
sw, gap = Inches(2.334), Inches(0.19)
x = Inches(0.45)
for i, (ic, t_, sub, c_) in enumerate(steps):
    flow_step(s, x, Inches(1.98), sw, Inches(1.9), ic, t_, sub, c_)
    if i < len(steps) - 1:
        arrow_between(s, x + sw + Inches(0.02), Inches(2.82), gap - Inches(0.04))
    x += sw + gap

card(s, Inches(0.45), Inches(4.02), Inches(6.1), Inches(2.05), GREEN,
     'WHAT THE 3D MODEL GIVES, WITH NOTHING TYPED',
     [('Weight, volume, overall size, surface area.',),
      ('Wall thickness, holes, pockets, bosses, bends,',),
      ('gear teeth, draft, machined face count.',),
      ('',),
      ('Fixed rules turn those measurements into the', DARK, True),
      ('values the CAPEE form asks for.', DARK, True)], fill=GREENBG)
card(s, Inches(6.78), Inches(4.02), Inches(6.1), Inches(2.05), AMBER,
     'WHAT A PERSON STILL SUPPLIES',
     [('Which material it is. A 3D model cannot tell steel',),
      ('from aluminium, and the same shape weighs about',),
      ('three times more in one than the other.',),
      ('',),
      ('Tolerance class, finish, heat treatment, coating.', DARK, True),
      ('The tool asks. It never guesses and never pre-ticks.', DARK, True)], fill=AMBERBG)
callout(s, Inches(0.45), Inches(6.17), Inches(12.43), Inches(1.06), PANEL2, INDIGO,
        'Why this one can start now',
        'The 3D model is measured by ordinary engineering software, the same kind of geometry '
        'engine that sits inside a CAD package, and the values are filled in by fixed rules. There '
        'is no AI in this option, so it needs no AI approval. The file is read on a JLR server and '
        'never leaves it, and CAPEE still does the costing on the same screens.')
notes(s, "This is the smaller of the two and it is deliberately narrow. CAPEE does not move. It "
         "still does the costing, on the same screens, with the same maths, and nobody has to learn "
         "a new tool. The only thing that changes is where the input comes from. The engineer "
         "attaches the 3D model from inside CAPEE. The software measures it on a JLR server with no "
         "internet connection, and fixed rules turn those measurements into the values the form "
         "wants. The engineer then confirms the handful of things a 3D model genuinely cannot show, "
         "and CAPEE costs the part. The green card is what we get for free once the model is "
         "attached: weight, size, surface area, wall thickness, the hole and pocket counts, and so "
         "on. The amber card is the honest limit. The model cannot tell us what the part is made "
         "of. Steel and aluminium look identical in a 3D file and the weight differs by about three "
         "times, so somebody answers that, or it comes off the part master. Same for tolerance "
         "class, finish and heat treatment. The tool asks rather than guessing. The point of the "
         "blue strip is the one I would make to this room: there is no AI in this option at all, so "
         "it does not wait on an AI decision. We could start it now.")

# ═══════════ 3 · WHAT COSTVISION DOES ══════════════════════════════════════
def tile(slide, x, y, w, h, icon, title, body, accent):
    """One capability: icon, heading, one plain sentence."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, w, Inches(0.06), fill=accent)
    icon_badge(slide, icon, x + Inches(0.26), y + Inches(0.26), d=Inches(0.6), fill=accent)
    text(slide, x + Inches(0.26), y + Inches(1.0), w - Inches(0.5), Inches(0.34),
         [[(title, 11.8, DARK, True)]])
    text(slide, x + Inches(0.26), y + Inches(1.36), w - Inches(0.5), _h(h - Inches(1.44)),
         [[(body, 9.8, BODY, False)]], line_spacing=1.14)


s = header('Option 2: what CostVision does', 'Option 2 · the capability')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('What the tool does. Five of these work today, one part at a time — running them over a '
        'whole basket is what Option 2 builds.', 11.5, BODY, False)]])
tiles = [
    ('cog', 'Cost a whole basket, unattended',
     'A list of parts and their 3D models goes in. Every part costed and reported '
     'with nobody sitting there. This is the one Option 2 builds.', AMBER),
    ('eye', 'A board photo becomes a costed BOM',
     'Take a picture of a circuit board. The tool works out what is on it, builds '
     'the bill of materials and the board spec, and costs it.', VIOLET),
    ('clip', 'An RFQ pack becomes costed lines',
     'The tool pulls the line items out of the pack, costs each one, ranks them by '
     'money and drafts the negotiation brief.', VIOLET),
    ('person', 'An assistant does the setting up',
     'It picks the process route and the machine and fills in what the geometry '
     'gives, then asks one question instead of a hundred.', INDIGO),
    ('coins', 'Supplier quotes taken apart',
     'Put their price next to our own build-up and see which part of it the gap is '
     'in, with the questions to put to them.', INDIGO),
    ('press', 'Design changes priced, not guessed',
     'It flags what makes a part expensive and puts a number on each change, '
     'ranked by the money it would save.', INDIGO),
]
tw, tg = Inches(4.01), Inches(0.2)
for i, (ic, t_, b_, c_) in enumerate(tiles):
    tx = Inches(0.45) + (tw + tg) * (i % 3)
    ty = Inches(2.0) + Inches(2.15) * (i // 3)
    tile(s, tx, ty, tw, Inches(2.05), ic, t_, b_, c_)
callout(s, Inches(0.45), Inches(6.32), Inches(12.43), Inches(0.95), PANEL2, INDIGO,
        'All six end in the same engine',
        'The AI reads, sorts and sets up — the photo, the pack, the drawing, the process route. It '
        'never sets a price. The money is worked out by the same fixed rules and the same rate book '
        'we use today, so every number can still be explained line by line.')
notes(s, "This is the slide I would spend the most time on, because it is the part of the case that "
         "is easy to under-sell. Six things the tool does, and one distinction I want to make "
         "before I go through them: five of these work today, one part at a time, in the tool as it "
         "stands. The first one does not, and that is why it is the amber tile. A list of parts and "
         "their models goes in and every part comes back costed, with a report per part and per "
         "basket, without anybody sitting there — that is what Option two builds, and today we cost "
         "the parts we have hours for and estimate the rest from a sample. Next to it is the one that usually gets a reaction in "
         "the room. You take a photograph of a circuit board and the tool works out what is on it, "
         "builds the bill of materials, works out the board specification and costs it. That is a "
         "job that takes an electronics buyer a long time and most of us cannot do at all. Third "
         "along, an RFQ pack goes in and comes back as costed lines, ranked by where the money is, "
         "with a first draft of the negotiation brief written for you. Bottom left is the "
         "assistant, and I want to be precise about it. It does the setting up: it picks the "
         "process route and the machine and fills in what the geometry gives it, and then it comes "
         "back to a person with one question rather than a hundred. Next to that is the one our "
         "buyers will care about most. Put a supplier's price next to our own build-up and the tool "
         "tells you which part of it the gap sits in and what to ask them about it. And the last "
         "one turns design-for-cost from a list into a decision: it flags what makes a part "
         "expensive and puts a number on each change, ordered by what it saves. The blue strip at "
         "the bottom is the sentence I would want people to leave with. The AI reads and sorts and "
         "sets up. It never sets a price. The money is still worked out by the same fixed rules and "
         "the same rate book we use now, so every number can be explained line by line, exactly as "
         "it can today.")

# ═══════════ 4 · WHAT WE GET FROM EACH ═════════════════════════════════════
def wide_card(slide, x, y, w, h, accent, title, col1, col2):
    """One card with two columns of text inside it.

    Option 2 has twice as much to say as Option 1, and giving it twice the
    width says so before anyone reads a word.
    """
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    text(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.3),
         [[(title, 13, DARK, True)]])
    cw = (w - Inches(0.62)) / 2
    for i, lines in enumerate((col1, col2)):
        runs = [[(ln[0], 10.5, ln[1] if len(ln) > 1 else BODY,
                  ln[2] if len(ln) > 2 else False)] for ln in lines]
        text(slide, x + Inches(0.22) + (cw + Inches(0.18)) * i, y + Inches(0.5),
             cw, _h(h - Inches(0.6)), runs, space_after=3, line_spacing=1.12)


s = header('What we get from each', 'The benefits')
text(s, Inches(0.45), Inches(1.66), Inches(12.4), Inches(0.3),
     [[('Stated in kind, not in pounds. Slide 1 turns these into hours once we fill in our own four '
        'numbers.', 11.5, BODY, False)]])

card(s, Inches(0.45), Inches(1.98), Inches(3.95), Inches(4.6), INDIGO,
     'OPTION 1 · input into CAPEE',
     [('Checking instead of typing', INDIGO, True),
      ('The values come off the 3D model.',),
      ('',),
      ('The same answer every time', INDIGO, True),
      ('One method, whoever costs the part.',),
      ('',),
      ('More parts for the same team', INDIGO, True),
      ('The time saved goes back into costing.',),
      ('',),
      ('Nothing new to learn', INDIGO, True),
      ('CAPEE does not change.',),
      ('',),
      ('It can start now', INDIGO, True),
      ('No AI, so no approval to wait for.',)])

wide_card(s, Inches(4.63), Inches(1.98), Inches(8.25), Inches(4.6), VIOLET,
          'OPTION 2 · CostVision end to end',
          [('It runs on its own', VIOLET, True),
           ('Parts go in, costed reports come out. A',),
           ('person answers only what the model cannot show.',),
           ('',),
           ('The same method for everything we make', VIOLET, True),
           ('Metal, plastic, rubber, composites,',),
           ('electronics and assemblies.',),
           ('',),
           ('We can cost electronics', VIOLET, True),
           ('A photo of a circuit board gives us the parts',),
           ('list, the board spec and a price.',),
           ('',),
           ('Software costed too', VIOLET, True),
           ('43 modules, from powertrain and driver',),
           ('assistance through to cloud.',)],
          [('Where a supplier price differs', VIOLET, True),
           ('Their price next to our build-up, with the gap',),
           ('shown in each part of the cost.',),
           ('',),
           ('And what to do about it', VIOLET, True),
           ('What to go after, what to ask them, and what',),
           ('it is worth over a year.',),
           ('',),
           ('A whole RFQ pack costed', VIOLET, True),
           ('The line items are pulled out, priced and put',),
           ('in order of money.',),
           ('',),
           ('Design changes with a price on them', VIOLET, True),
           ('What makes the part expensive, and what',),
           ('each change would save.',)])

callout(s, Inches(0.45), Inches(6.62), Inches(12.43), Inches(0.83), GREENBG, GREEN,
        'The two build on each other',
        "Option 1's benefits arrive first and are kept when Option 2 lands. Option 2's rest on the "
        'accuracy figure that Option 1 produces, which is why these are an order and not a choice.')
notes(s, "This is what we get for it, written in kind rather than in pounds, for the reason I gave "
         "on the first slide: we have never timed this job and never checked the tool against a "
         "price we paid, so a pound figure from me would be a guess. The left column is Option one, "
         "and it is short on purpose. The job turns from typing into checking, because the values "
         "come off the model. It gives the same answer every time, so parts costed by different "
         "engineers are comparable. The time freed goes back into costing more parts with the team "
         "we already have. Nobody learns a new tool. And it can start now, because there is no AI "
         "in it to approve. The right column is wider because Option two genuinely has more in it. "
         "It runs on its own — parts go in and costed reports come out, and a person answers only "
         "what the model cannot show. It uses the same method for everything we make, so a plastic "
         "clip and a machined bracket are built up the same way and can sit in the same report. It "
         "costs electronics from a photograph of the board, which is work we largely cannot do "
         "today. It costs software as well: forty-three modules, from powertrain and driver "
         "assistance through to the cloud back end, which matters more every programme. Then the "
         "second half of that column is the part I would put to purchasing. It shows where a "
         "supplier's price differs from our own build-up and which part of the cost the gap sits "
         "in. Then it tells us what to do about it: what to go after, what to ask them, and what it "
         "is worth over a year. A whole RFQ pack can go in and come back priced and put in order of "
         "money. And design changes come back with a price on them rather than as a list of "
         "suggestions. The green strip is the sequencing point. Option one's benefits arrive first "
         "and we keep them. Option two's depend on the accuracy figure Option one produces, which "
         "is why these are an order rather than a choice.")


# ═══════════ 5 · WHAT WE ARE ASKING FOR ════════════════════════════════════
def strip(slide, x, y, w, h, label, body, accent):
    """One short statement: bold label, then the sentence, on a single row."""
    box(slide, x, y, w, h, fill=PANEL, line=LINE, round_=True)
    box(slide, x, y, Inches(0.075), h, fill=accent)
    text(slide, x + Inches(0.24), y + Inches(0.15), Inches(2.9), Inches(0.28),
         [[(label, 10.8, accent, True)]])
    text(slide, x + Inches(3.3), y + Inches(0.15), w - Inches(3.55), Inches(0.28),
         [[(body, 10.2, BODY, False)]])


s = header('What we are asking for', 'The decision')
card(s, Inches(0.45), Inches(1.72), Inches(6.1), Inches(1.05), GREEN,
     'OPTION 1 · a decision today',
     [('Approve the trial. One commodity, 30 to 50 parts we have',),
      ('already bought, inside CAPEE.',)])
card(s, Inches(6.78), Inches(1.72), Inches(6.1), Inches(1.05), VIOLET,
     'OPTION 2 · no decision today',
     [('Start the AI approval in parallel, so it is not the thing',),
      ('holding us up when the trial finishes.',)])

card(s, Inches(0.45), Inches(2.89), Inches(6.1), Inches(1.32), INDIGO,
     'OPTION 1 NEEDS',
     [('A Linux server, and a way to pass numbers into CAPEE',),
      ('30 to 50 parts where we know the price we paid',),
      ('An engineer from the commodity team for the trial',)])
card(s, Inches(6.78), Inches(2.89), Inches(6.1), Inches(1.32), VIOLET,
     'OPTION 2 NEEDS',
     [('Everything above, plus approval to use AI',),
      ('Our own rate book loaded, with a change history',),
      ('A record of every costing, and the run itself built',)])

callout(s, Inches(0.45), Inches(4.33), Inches(12.43), Inches(0.9), GREENBG, GREEN,
        'What the trial gives us',
        'The two numbers nobody in this company can state today: how close the tool gets to a price '
        'we actually paid, and how much of the input it fills on its own. Until we have them there '
        'is no honest conversation to be had about Option 2.')

strip(s, Inches(0.45), Inches(5.36), Inches(12.43), Inches(0.58),
      'The rule that does not bend',
      'The AI reads, sorts and sets up. It never sets a price — the money comes from the same '
      'fixed rules we use today.', INDIGO)
strip(s, Inches(0.45), Inches(5.98), Inches(12.43), Inches(0.58),
      'If the trial disappoints',
      'Never checked against a price JLR has paid. If the answer is poor we stop, having spent '
      'only our own time.', AMBER)
strip(s, Inches(0.45), Inches(6.60), Inches(12.43), Inches(0.58),
      'Nobody has sized this yet',
      'How long it takes waits on one answer from IT about how CAPEE is built. The sizing comes '
      'back with it.', AMBER)
notes(s, "So this is the decision, and I have split it into what I want today and what I do not. On "
         "the left, Option one: approve the trial. One commodity, thirty to fifty parts we have "
         "already bought, run inside CAPEE. On the right, Option two: I am not asking you to decide "
         "it. What I am asking is that we start the AI approval alongside the trial, because that "
         "is the long pole and there is no sense discovering at the end of the trial that we now "
         "have to begin it. The two cards underneath are what each one needs. Option one needs a "
         "server, a way to pass numbers into CAPEE, some parts where we know what we paid, and an "
         "engineer from the commodity team for the trial. Option two needs all of that plus the AI "
         "approval, our own rate book with a change history, and a few things built around the "
         "engine. The green strip is what we get out of the trial and it is the reason to do it: "
         "two numbers that nobody in this company can state today. How close the tool gets to a "
         "price we actually paid, and how much of the input it fills on its own. Until we have "
         "those, any conversation about Option two is people trading opinions. Then three short "
         "things I want to say plainly rather than have asked. The first is the governance point "
         "and it is the one I would repeat if you take nothing else away: the AI reads, sorts and "
         "sets up, and it never sets a price. The money comes out of the same fixed rules and the "
         "same rate book as today, so every number can be defended line by line exactly as it can "
         "now. The second is the exit. We have never compared this against a price we have paid. If "
         "the trial says the tool is not close enough, we stop, and what we will have spent is our "
         "own engineering time — that is precisely why I am asking for a trial and not a rollout. "
         "And the third is the honest answer to how long: nobody has sized it. It waits on one "
         "answer from IT about how CAPEE is built, and the sizing comes back with that answer "
         "rather than being invented now.")


# ───────────────────────────────────────────────────────────────────────────
OUT = 'CostVision-CAPEE-Business-Case.pptx'
prs.save(OUT)
finalise(OUT)


def assert_powerpoint_can_open(path):
    """Fail the build on the OOXML faults that make PowerPoint offer to repair.

    LibreOffice opens files PowerPoint rejects, so converting to PDF proves
    nothing about whether the deck will open on a colleague's laptop. A shape
    with a zero or negative extent is the fault this build actually hit: a short
    callout made an inner text box `h - chrome` wide, which went negative, and
    PowerPoint refused the file while LibreOffice rendered it happily.
    """
    import zipfile, re
    import xml.etree.ElementTree as ET
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    z = zipfile.ZipFile(path)
    faults = []
    if z.testzip() is not None:
        faults.append('zip archive is corrupt')
    for n in sorted(x for x in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', x)):
        root = ET.fromstring(z.read(n))
        for ext in root.iter('{%s}ext' % A):
            cx, cy = int(ext.get('cx', '1')), int(ext.get('cy', '1'))
            if cx <= 0 or cy <= 0:
                faults.append(f'{n}: shape extent cx={cx} cy={cy}')
        ids = [int(e.get('id')) for e in root.iter('{%s}cNvPr' % NS) if e.get('id')]
        for d in sorted({i for i in ids if ids.count(i) > 1}):
            faults.append(f'{n}: duplicate shape id {d}')
    if faults:
        raise SystemExit('DECK IS INVALID - PowerPoint would ask to repair it:\n  '
                         + '\n  '.join(faults))
    return len([x for x in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', x)])


n_slides = assert_powerpoint_can_open(OUT)
print(f'{OUT}  -  {n_slides} slides, validated')
