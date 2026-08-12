#!/usr/bin/env python3
"""BrainSpark — the whole-platform director presentation.

    python3 scripts/make-platform-deck.py

Lived outside the repo until now, which is why it went five months without a
refresh and nobody noticed: a generator you cannot find is a deck that silently
goes stale. Screenshots live in scripts/deck-assets/ and are committed with it.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
import copy

# ── Palette ─────────────────────────────────────────────────────────────────
# Light theme palette (legacy names kept, values remapped)
BG     = RGBColor(0xF3, 0xF5, 0xF9)   # page background — soft light gray
NAVY   = RGBColor(0xF3, 0xF5, 0xF9)   # (legacy) page fill == BG
NAVY2  = RGBColor(0xFF, 0xFF, 0xFF)   # panels / header — white
NAVY3  = RGBColor(0xE4, 0xEA, 0xF1)   # chips / gridlines — light
BORDER = RGBColor(0xD8, 0xE0, 0xEA)   # hairline card/panel border
TEAL   = RGBColor(0x0D, 0x94, 0x88)   # accents deepened to ~600 for contrast on white
TEALD  = RGBColor(0x0B, 0x76, 0x6E)
GOLD   = RGBColor(0xB4, 0x76, 0x09)
EMER   = RGBColor(0x05, 0x96, 0x69)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ROSE   = RGBColor(0xE1, 0x1D, 0x48)
INK    = RGBColor(0x16, 0x27, 0x3C)   # primary text — dark
WHITE  = INK                          # (legacy) primary text now dark
WT     = RGBColor(0xFF, 0xFF, 0xFF)   # text/glyphs on colored fills
SLATE  = RGBColor(0x55, 0x64, 0x77)   # muted text
SLATED = RGBColor(0x83, 0x90, 0xA1)   # dim text
LIGHT  = RGBColor(0x37, 0x45, 0x58)   # body text
CARD   = RGBColor(0xFF, 0xFF, 0xFF)   # cards — white

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

FONT = "Segoe UI"
FONTL = "Segoe UI Light"
FONTSB = "Segoe UI Semibold"

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=NAVY):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    # flat design (no injected effect XML) — maximum PowerPoint compatibility
    _ = shadow
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        if isinstance(para, tuple): para = [para]
        for seg in para:
            t, size, color, bold = seg[0], seg[1], seg[2], seg[3]
            font = seg[4] if len(seg) > 4 else FONT
            ital = seg[5] if len(seg) > 5 else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = font; r.font.italic = ital
    return tb

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

_slide_no = [0]

def header(s, kicker, title, idx=None):
    """Slide chrome. `idx` is COUNTED, not typed: a deck that renumbers itself
    is a deck somebody will actually insert a slide into."""
    if idx is None:
        _slide_no[0] += 1
        idx = _slide_no[0]
    rect(s, 0, 0, EMU_W, Inches(1.35), fill=NAVY2)
    rect(s, 0, Inches(1.35), EMU_W, Pt(3), fill=TEAL)
    rect(s, Inches(0.55), Inches(0.42), Inches(0.11), Inches(0.55), fill=GOLD)
    txt(s, Inches(0.85), Inches(0.28), Inches(10.5), Inches(0.35),
        [(kicker.upper(), 12, TEAL, True, FONTSB)])
    txt(s, Inches(0.85), Inches(0.60), Inches(11.3), Inches(0.6),
        [(title, 27, WHITE, True, FONTSB)])
    # slide number chip
    rect(s, Inches(12.35), Inches(0.5), Inches(0.55), Inches(0.4), fill=NAVY3, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(12.35), Inches(0.53), Inches(0.55), Inches(0.35),
        [(f"{idx:02d}", 13, SLATE, True)], align=PP_ALIGN.CENTER)

def footer(s):
    txt(s, Inches(0.85), Inches(7.06), Inches(8), Inches(0.3),
        [[("BrainSpark", 9, TEAL, True), ("  ·  AI-Powered Cost Engineering Platform", 9, SLATED, False)]])
    txt(s, Inches(9.5), Inches(7.06), Inches(3), Inches(0.3),
        [("Confidential — Internal", 9, SLATED, False)], align=PP_ALIGN.RIGHT)

def chip(s, x, y, w, text, color, tcolor=WHITE, size=11):
    rect(s, x, y, w, Inches(0.42), fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, y+Inches(0.02), w, Inches(0.38), [(text, size, tcolor, True, FONTSB)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def icon_badge(s, x, y, glyph, color, d=Inches(0.62)):
    c = rect(s, x, y, d, d, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, y, d, d, [(glyph, 18, WT, True, "Arial")], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

def card(s, x, y, w, h, glyph, title, body, accent):
    rect(s, x, y, w, h, fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, Inches(0.09), h, fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    icon_badge(s, x+Inches(0.28), y+Inches(0.26), glyph, accent, d=Inches(0.55))
    txt(s, x+Inches(1.0), y+Inches(0.26), w-Inches(1.2), Inches(0.5),
        [(title, 15, WHITE, True, FONTSB)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(0.32), y+Inches(0.95), w-Inches(0.6), h-Inches(1.05),
        [(body, 11.5, SLATE, False)], space=1.06)

def arrow(s, x, y, w, h=Inches(0.5), color=TEAL):
    rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.CHEVRON)

# Photorealistic imagery (Canva-generated, user-selected). Source ratio 1587×2245
# (w/h ≈ 0.707). pic() centre-crops vertically to fill the target box exactly,
# biased by `focus` (0 = keep top … 1 = keep bottom) toward the subject, and
# frames it with a hairline border.
import os
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deck-assets")
IMGS = {k: os.path.join(IMG_DIR, f"deck-{k}.png") for k in ("hero", "demo", "close")}
SRC_RATIO = 1587 / 2245

def pic(s, key, x, y, w, h, focus=0.5):
    path = IMGS[key]
    if not os.path.exists(path):
        return None   # image not supplied — slide falls back to plain background
    p = s.shapes.add_picture(path, x, y, width=w, height=h)
    box_ratio = w / h
    if box_ratio > SRC_RATIO:            # box is wider than source → crop top/bottom
        keep = SRC_RATIO / box_ratio     # fraction of source height kept
        crop = 1 - keep
        p.crop_top = crop * focus
        p.crop_bottom = crop * (1 - focus)
    else:                                # box narrower → crop left/right evenly
        keep = box_ratio / SRC_RATIO
        crop = 1 - keep
        p.crop_left = crop / 2
        p.crop_right = crop / 2
    ln = rect(s, x, y, w, h, fill=None, line=BORDER, line_w=1.0)
    return p

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, NAVY)
# decorative side band
rect(s, 0, 0, Inches(0.35), EMU_H, fill=TEAL)
rect(s, Inches(0.35), 0, Inches(0.12), EMU_H, fill=GOLD)
# photorealistic hero: EV rolling chassis (Canva-generated, user-selected)
pic(s, "hero", Inches(8.95), Inches(1.35), Inches(3.95), Inches(4.65), focus=0.68)
txt(s, Inches(1.1), Inches(1.5), Inches(7.6), Inches(0.5),
    [("AI-POWERED COST ENGINEERING PLATFORM", 15, TEAL, True, FONTSB)])
txt(s, Inches(1.05), Inches(2.15), Inches(7.7), Inches(1.6),
    [[("Brain", 66, WHITE, True, FONTSB), ("Spark", 66, TEAL, True, FONTSB)]])
txt(s, Inches(1.1), Inches(3.5), Inches(7.6), Inches(1.0),
    [("Turning whole-vehicle cost-reduction ideas into defensible,\nquote-ready numbers — in minutes, not weeks.", 19, LIGHT, False, FONTL)], space=1.1)
# stat strip (compressed left of the hero image)
stats = [("13+", "vehicle systems"), ("11", "innovation methods"), ("1,600+", "curated ideas"), ("100%", "engine-verified")]
sx = Inches(1.1)
for v, l in stats:
    rect(s, sx, Inches(5.15), Inches(1.83), Inches(1.05), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, sx, Inches(5.3), Inches(1.83), Inches(0.45), [(v, 20, GOLD, True, FONTSB)], align=PP_ALIGN.CENTER)
    txt(s, sx, Inches(5.8), Inches(1.83), Inches(0.35), [(l, 9.5, SLATE, False)], align=PP_ALIGN.CENTER)
    sx += Inches(1.95)
txt(s, Inches(1.1), Inches(6.6), Inches(11), Inches(0.4),
    [[("Prepared for: ", 12, SLATED, False), ("Director, Cost Engineering", 12, WHITE, True),
      ("      |      Presenter: ", 12, SLATED, False), ("Avinash Bhosale", 12, WHITE, True)]])
notes(s, "Open warmly: 'Thanks for making time. I want to walk you through BrainSpark — a cost-engineering platform I've been building on weekends. In about 20 minutes I'll show what it does, why the numbers can be trusted, and then a quick live demo.' Emphasise the one-line promise on screen: defensible, quote-ready numbers in minutes, not weeks. Keep this slide up while people settle.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Why this matters", "The Cost-Engineering Challenge"); footer(s)
pains = [
    ("●", "Slow to insight", "Manual should-costing and VAVE workshops take days to weeks per part — ideas arrive after the design is frozen.", ROSE),
    ("◆", "Hard to defend", "Spreadsheet estimates and gut-feel savings don't survive scrutiny in a sourcing or programme review.", GOLD),
    ("■", "Fragmented knowledge", "OEM best practice, teardown data and commodity prices live in silos and in people's heads.", PURPLE),
    ("▲", "Narrow coverage", "Tools tend to cover one commodity; nobody sees the whole vehicle at once.", BLUE),
]
x = Inches(0.6)
for g, t, b, c in pains:
    card(s, x, Inches(1.75), Inches(2.95), Inches(3.0), g, t, b, c)
    x += Inches(3.08)
rect(s, Inches(0.6), Inches(5.05), Inches(12.13), Inches(1.5), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(5.05), Inches(0.12), Inches(1.5), fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.0), Inches(5.25), Inches(11.4), Inches(1.1),
    [[("The opportunity:  ", 16, TEAL, True, FONTSB),
      ("bring AI ideation and a deterministic cost engine together, so we identify savings across the ", 15, LIGHT, False),
      ("whole vehicle", 15, GOLD, True),
      (" and back every number with real engineering — fast enough to act while the design is still open.", 15, LIGHT, False)]],
    anchor=MSO_ANCHOR.MIDDLE, space=1.08)
notes(s, "Frame the pain the Director already feels. Two moves matter here: first, cost ideas usually land too late to change anything; second, when they do land, the numbers are hard to defend. Add that knowledge is fragmented and tools are narrow. Then pivot to the opportunity line at the bottom — that's exactly the gap BrainSpark fills. Don't rush; let the problem land before showing the solution.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS BRAINSPARK
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Overview", "What is BrainSpark?"); footer(s)
txt(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.9),
    [[("A single platform that generates ", 18, LIGHT, False), ("cost-reduction & value-improvement ideas", 18, TEAL, True),
      (" across the whole vehicle — and prices every one of them on a deterministic should-cost engine.", 18, LIGHT, False)]], space=1.1)
feat = [
    ("★", "AI Idea Generation", "System-aware ideas grounded in OEM & teardown benchmarks", TEAL),
    ("●", "Should-Cost Engine", "Bottom-up costs from rates, mass, geometry & live commodities", GOLD),
    ("3D", "CAD-to-Cost", "Upload a model or drawing → instant feature-based cost + 3D view", BLUE),
    ("◆", "Innovation Methods", "TRIZ, Value Engineering, DFA/DFMA, Design-to-Cost built in", PURPLE),
    ("■", "Idea Marketplace", "Thousands of verified, OEM-benchmarked ideas to reuse", EMER),
    ("▲", "Decision Analytics", "Monte-Carlo, ROI, savings pipeline & dashboards", ROSE),
]
x, y = Inches(0.6), Inches(2.75)
for i, (g, t, b, c) in enumerate(feat):
    cx = x + (i % 3) * Inches(4.08)
    cy = y + (i // 3) * Inches(1.95)
    card(s, cx, cy, Inches(3.9), Inches(1.78), g, t, b, c)
notes(s, "This is the elevator pitch in one slide. Say it plainly: 'BrainSpark generates cost ideas across the whole car, and immediately prices each one on a deterministic engine.' Then gesture at the six tiles — these are the building blocks I'll go one level deeper on in the next few slides. The key phrase to plant early is 'deterministic should-cost engine', because trust is our differentiator.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE USP / TRUST ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "The core differentiator", "Our USP: AI Proposes, the Engine Verifies"); footer(s)
txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.55),
    [("Most AI tools estimate savings. BrainSpark checks every figure against real engineering — so nothing is hallucinated.", 15, SLATE, False)])
# three-stage trust flow
boxes = [
    ("AI", "AI PROPOSES", "Claude generates ideas & alternatives grounded in benchmarks and live data", TEAL),
    ("£", "ENGINE VERIFIES", "Deterministic should-cost engine re-computes every number from first principles", GOLD),
    ("✓", "STAMPED RESULT", "Each idea is marked confirmed / contradicted — defensible, quote-ready", EMER),
]
bx = Inches(0.6)
for i, (g, t, b, c) in enumerate(boxes):
    rect(s, bx, Inches(2.5), Inches(3.7), Inches(2.5), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    icon_badge(s, bx+Inches(1.5), Inches(2.75), g, c, d=Inches(0.75))
    txt(s, bx, Inches(3.6), Inches(3.7), Inches(0.4), [(t, 15, c, True, FONTSB)], align=PP_ALIGN.CENTER)
    txt(s, bx+Inches(0.3), Inches(4.05), Inches(3.1), Inches(0.9), [(b, 11.5, SLATE, False)], align=PP_ALIGN.CENTER, space=1.05)
    if i < 2:
        arrow(s, bx+Inches(3.75), Inches(3.5), Inches(0.55), color=GOLD)
    bx += Inches(4.08)
rect(s, Inches(0.6), Inches(5.35), Inches(12.13), Inches(1.15), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(5.35), Inches(0.12), Inches(1.15), fill=EMER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.0), Inches(5.5), Inches(11.5), Inches(0.85),
    [[("Why it matters:  ", 15, EMER, True, FONTSB),
      ("the same engine runs behind every screen — ideas, CAD, marketplace — so a saving shown to a supplier is the same number our engineers can reproduce and stand behind.", 14, LIGHT, False)]],
    anchor=MSO_ANCHOR.MIDDLE, space=1.06)
notes(s, "This is the most important slide — slow down. The trap with AI is confident, wrong numbers. We avoid that with a simple rule: the AI is allowed to have ideas, but it is never allowed to invent the cost. A separate deterministic engine recomputes every figure and stamps it confirmed or contradicted. Tell the Director: this is what lets us put a number in front of a supplier and defend it. If they remember one thing from today, it should be 'AI proposes, the engine verifies.'")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHOLE-VEHICLE COVERAGE
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Breadth", "Whole-Vehicle System Coverage"); footer(s)
txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
    [("One platform spans the full vehicle — powertrain to trim, ICE to BEV — with curated levers and benchmarks for each system.", 14, SLATE, False)])
systems = [
    ("BIW / Body Structure", BLUE), ("Powertrain — ICE", GOLD), ("Powertrain — BEV / EDU", TEAL),
    ("Battery Pack & BMS", EMER), ("Chassis & Suspension", PURPLE), ("Braking & Wheel-end", ROSE),
    ("Interior Systems", BLUE), ("Exterior & Trim", TEAL), ("ADAS & Sensors", GOLD),
    ("Thermal & HVAC", EMER), ("E/E & Wiring", PURPLE), ("Transmission & Driveline", ROSE),
]
x, y = Inches(0.6), Inches(2.35)
for i, (t, c) in enumerate(systems):
    cx = x + (i % 4) * Inches(3.06)
    cy = y + (i // 4) * Inches(1.12)
    rect(s, cx, cy, Inches(2.9), Inches(0.95), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, cx+Inches(0.18), cy+Inches(0.26), Inches(0.16), Inches(0.42), fill=c)
    txt(s, cx+Inches(0.5), cy, Inches(2.35), Inches(0.95), [(t, 12.5, WHITE, True, FONTSB)], anchor=MSO_ANCHOR.MIDDLE, space=1.0)
txt(s, Inches(0.6), Inches(6.05), Inches(12), Inches(0.5),
    [[("Plus commodity models for ", 12.5, SLATE, False),
      ("machining, stamping, casting, moulding, forging, wiring harness, electrical-steel laminations, glass & PCB assembly.", 12.5, LIGHT, True)]], space=1.05)
notes(s, "The point of this slide is breadth. Cost tools are usually single-commodity; BrainSpark sees the whole car. Walk across a couple of systems the Director cares about most — probably BIW and battery for automotive — and mention that under each system there's a curated set of proven cost levers, not generic advice. The bottom line is the commodity coverage: we can cost most of what a vehicle is made of.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURE: AI IDEA GENERATION
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 1", "AI Idea Generation"); footer(s)
txt(s, Inches(0.6), Inches(1.6), Inches(6.4), Inches(0.5), [("Pick a system, set the context, generate.", 15, TEAL, True, FONTSB)])
pts = [
    ("System-aware & grounded", "Ideas tuned to system, vehicle type, volume and region — with live web search for current prices and innovations."),
    ("Learns from your decisions", "Ideas you approve or confirm in production feed back as positive examples; proven marketplace precedents ground every run."),
    ("Nine analysis lenses", "Layer TRIZ, Value Engineering, DFA, Design-to-Cost, biomimicry and more onto the same part."),
    ("Every idea costed & ranked", "Engine-checked saving, difficulty and payback — ranked by verified value, with every ranking factor visible."),
]
y = Inches(2.25)
for t, b in pts:
    icon_badge(s, Inches(0.6), y, "▪", TEAL, d=Inches(0.42))
    txt(s, Inches(1.2), y-Inches(0.05), Inches(5.9), Inches(0.4), [(t, 14, WHITE, True, FONTSB)])
    txt(s, Inches(1.2), y+Inches(0.32), Inches(5.9), Inches(0.7), [(b, 11.5, SLATE, False)], space=1.03)
    y += Inches(1.12)
# right visual: idea card mock
rect(s, Inches(7.3), Inches(1.9), Inches(5.4), Inches(4.7), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(7.6), Inches(2.1), Inches(4.8), Inches(0.4), [("EXAMPLE IDEA CARD", 11, SLATED, True, FONTSB)])
rect(s, Inches(7.6), Inches(2.55), Inches(4.8), Inches(1.75), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.8), Inches(2.7), Inches(4.4), Inches(0.7), [("One-piece hot-stamped B-pillar tailored blank", 13.5, WHITE, True, FONTSB)], space=1.0)
txt(s, Inches(7.8), Inches(3.42), Inches(4.4), Inches(0.4), [("Deletes 14 welds + 2 dies · −0.9 kg/veh", 11, SLATE, False)])
chip(s, Inches(7.8), Inches(3.85), Inches(1.5), "£2.04M/yr", EMER, WT, 12)
chip(s, Inches(9.45), Inches(3.85), Inches(1.35), "High", GOLD, WT, 11)
chip(s, Inches(10.95), Inches(3.85), Inches(1.25), "✓ Verified", TEAL, WT, 11)
txt(s, Inches(7.6), Inches(4.55), Inches(4.8), Inches(0.4), [("ENGINE CHECK", 11, SLATED, True, FONTSB)])
rect(s, Inches(7.6), Inches(4.95), Inches(4.8), Inches(1.35), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.8), Inches(5.1), Inches(4.4), Inches(1.1),
    [[("Should-cost recomputed: ", 11.5, SLATE, False), ("£17.10", 11.5, EMER, True), (" vs ", 11.5, SLATE, False), ("£23.40", 11.5, ROSE, True), (" baseline.\n", 11.5, SLATE, False),
      ("Saving confirmed by the deterministic engine.", 11.5, LIGHT, True)]], space=1.05)
notes(s, "Now we go one level down into the first capability. The story: an engineer picks a system — say Body Structure — sets volume, region and vehicle type, and hits generate. Before it answers, it searches for current prices and recent innovations, so ideas are grounded in today's reality, not old training data. The new bit since you last saw this: it learns. Ideas the team approves or confirms in production feed back into the next generation as positive examples, and proven ideas from our own library ground every run — so the tool gets more 'us' the more we use it. Point at the example card: an annual saving, a difficulty flag, and the engine's own recomputed cost underneath. That green 'verified' tag is the trust architecture in action.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURE: SHOULD-COST ENGINE
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 2", "Deterministic Should-Cost Engine"); footer(s)
txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
    [("Bottom-up piece-price from real drivers — no black box, fully reproducible.", 15, TEAL, True, FONTSB)])
drivers = [
    ("■", "Material", "Density, live commodity €/kg, buy-to-fly & scrap recovery", BLUE),
    ("●", "Process", "Machine rate × cycle time, labour, setup, secondary ops", PURPLE),
    ("▲", "Tooling", "Amortised die / mould / fixture over programme volume", GOLD),
    ("◆", "Overheads", "Factory burden, SG&A, packaging & freight by region", TEAL),
]
x = Inches(0.6)
for g, t, b, c in drivers:
    card(s, x, Inches(2.3), Inches(2.95), Inches(2.35), g, t, b, c)
    x += Inches(3.08)
# capability strip
strip = [("9-region", "labour & overhead library"), ("Live prices", "ECB FX + commodity feeds"),
         ("Monte-Carlo", "P10 / P50 / P90 bands"), ("Self-calibrating", "learns from real quotes")]
x = Inches(0.6)
for v, l in strip:
    rect(s, x, Inches(4.95), Inches(2.95), Inches(1.35), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, Inches(5.15), Inches(2.95), Inches(0.5), [(v, 18, GOLD, True, FONTSB)], align=PP_ALIGN.CENTER)
    txt(s, x, Inches(5.72), Inches(2.95), Inches(0.5), [(l, 11, SLATE, False)], align=PP_ALIGN.CENTER, space=1.0)
    x += Inches(3.08)
notes(s, "This is the engine behind every number. Stress that it's bottom-up and transparent: material from mass and live commodity prices, process from machine rate times cycle time, tooling amortised over the programme, then overheads and margin by region. Four things to name-drop: nine costing regions, live commodity and FX feeds, Monte-Carlo confidence bands so we show a range not a false-precision point, and self-calibration — feed it a real supplier quote and it tunes itself. This is what makes the output audit-proof.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURE: CAD-TO-COST
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 3", "Feature-Based CAD-to-Cost"); footer(s)
txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5),
    [("Upload a STEP / STL model or a 2D drawing — cost is driven by the geometry itself, the way aPriori-class tools work.", 14.5, TEAL, True, FONTSB)])
# pipeline
steps = [("↑","Upload","STEP / STL /\ndrawing / PDF"),("◆","Extract","Volume, area,\nholes, thickness"),
         ("●","Cost","Feature-based\ncycle & material"),("3D","Visualise","Interactive\n3D viewer"),("▲","Report","DFMA score +\nshould-cost")]
bx = Inches(0.6)
for i,(g,t,b) in enumerate(steps):
    rect(s, bx, Inches(2.45), Inches(2.15), Inches(1.95), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    icon_badge(s, bx+Inches(0.78), Inches(2.62), g, TEAL, d=Inches(0.6))
    txt(s, bx, Inches(3.3), Inches(2.15), Inches(0.35), [(t,13,WHITE,True,FONTSB)], align=PP_ALIGN.CENTER)
    txt(s, bx, Inches(3.68), Inches(2.15), Inches(0.65), [(b,10.5,SLATE,False)], align=PP_ALIGN.CENTER, space=1.0)
    if i<4: arrow(s, bx+Inches(2.18), Inches(3.15), Inches(0.32), color=GOLD)
    bx += Inches(2.48)
rect(s, Inches(0.6), Inches(4.75), Inches(12.13), Inches(1.65), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(4.75), Inches(0.12), Inches(1.65), fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.0), Inches(4.92), Inches(11.5), Inches(0.45), [("Why feature-based beats a mass estimate:", 14, BLUE, True, FONTSB)])
txt(s, Inches(1.0), Inches(5.38), Inches(11.5), Inches(0.95),
    [[("It reads the ", 12.5, LIGHT, False), ("actual geometry", 12.5, WHITE, True),
      (" — removal volume, surface area, holes, wall thickness, press tonnage — so it tells a slow, dear titanium billet apart from a cheap aluminium one, and a deep-drawn panel apart from a flat bracket. Live for machining and stamping today.", 12.5, LIGHT, False)]], space=1.06)
notes(s, "This is the slide engineers get excited about. Instead of guessing cost from weight, it reads the geometry. Upload a CAD model, it pulls out the features that actually drive cost — how much metal is removed, surface area to finish, number and depth of holes, press tonnage for a stamping — and builds the cycle time from that. Walk the five-step pipeline left to right and mention the 3D viewer, which makes it tangible in the demo. Honest scope note: machining and stamping are live today; more commodities are on the roadmap.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PCB & ELECTRONICS COST INTELLIGENCE (NEW)
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 3b — new", "Electronics: PCB Cost Intelligence"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5),
    [("Photograph a circuit board — get a multi-country should-cost with live component prices and a full cost breakdown.", 14.5, TEAL, True, FONTSB)])
# left: the pipeline
pcb_steps = [
    ("↑", "Capture", "Up to 5 HD photos, or import the BOM from Excel / PDF"),
    ("◆", "Read", "AI vision identifies components, markings and package types"),
    ("£", "Price", "Live DigiKey / Octopart part prices where available — tagged LIVE vs AI-estimated"),
    ("●", "Cost", "Assembly, test, fabrication and overheads modelled across 12 manufacturing countries"),
]
y = Inches(2.2)
for g, t, b in pcb_steps:
    icon_badge(s, Inches(0.6), y, g, BLUE, d=Inches(0.46))
    txt(s, Inches(1.25), y-Inches(0.04), Inches(5.6), Inches(0.4), [(t, 13.5, WHITE, True, FONTSB)])
    txt(s, Inches(1.25), y+Inches(0.33), Inches(5.65), Inches(0.75), [(b, 11, SLATE, False)], space=1.03)
    y += Inches(1.08)
# right: outputs card
rect(s, Inches(7.1), Inches(1.95), Inches(5.6), Inches(4.6), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(7.4), Inches(2.15), Inches(5.0), Inches(0.4), [("WHAT YOU GET", 11, SLATED, True, FONTSB)])
outs = [
    ("China vs India vs Vietnam + 9 more", "side-by-side landed cost per country", TEAL),
    ("P10 / P50 / P90 + tornado chart", "which drivers move the price most", GOLD),
    ("Industry-standard cost breakdown", "open-book waterfall: material → conversion → NRE → overheads → landed", PURPLE),
    ("Engine-verified insights", "DFM, sourcing and optimisation ideas — each re-costed before it's shown", EMER),
]
yy = Inches(2.6)
for t, b, c in outs:
    icon_badge(s, Inches(7.4), yy, "▪", c, d=Inches(0.36))
    txt(s, Inches(7.95), yy-Inches(0.05), Inches(4.6), Inches(0.4), [(t, 12.5, WHITE, True, FONTSB)], space=1.0)
    txt(s, Inches(7.95), yy+Inches(0.3), Inches(4.6), Inches(0.6), [(b, 10.5, SLATE, False)], space=1.02)
    yy += Inches(0.98)
notes(s, "This one is brand new since the last update, and it fills a real gap — electronics. The workflow is deliberately simple: photograph the board, or drop in the BOM as Excel or PDF, and the AI reads the components off the image. Where a part number resolves, we pull a live distributor price and tag it LIVE; where it doesn't, the estimate is clearly tagged as AI — we never blur the two. Then the engine costs assembly, test and fabrication across twelve manufacturing countries side by side, with a proper sensitivity view and an open-book cost breakdown a supplier would recognise. If procurement has ever asked 'what should this ECU cost from Vietnam versus China?' — this answers it in minutes.")

# ════════════════════════════════════════════════════════════════════════════
# DFM / DFA STUDIO
#
# This deck predated the tool entirely — the largest capability gap in it.
# Figures here are the ones the DFM deck counts from the code: 248 rules across
# 35 process families, 33 thresholds from a first-hand primary source, and a
# geometry gate of 200 analytic checks.
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 4 — new", "DFM / DFA Studio: Manufacturability, Measured"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5),
    [("Upload a STEP file. Every rule your chosen process breaks — marked on the faces that broke it, priced by the same engines that quote the part.", 14.5, TEAL, True, FONTSB)])
dfm_steps = [
    ("↑", "Read", "OpenCascade measures draft, undercuts, wall thickness and tool reach from the B-rep — no AI in the measurement"),
    ("◆", "Judge", "248 rules across 35 process families, thresholds resolved for YOUR alloy, each one citing its source"),
    ("●", "Locate", "Every finding paints the actual faces that broke it; click a face to jump back to its finding"),
    ("▲", "Decide", "Findings priced per part and per year, written as actions with an owner-role and a target date"),
]
y = Inches(2.2)
for g, t, b in dfm_steps:
    icon_badge(s, Inches(0.6), y, g, GOLD, d=Inches(0.46))
    txt(s, Inches(1.25), y-Inches(0.04), Inches(5.6), Inches(0.4), [(t, 13.5, WHITE, True, FONTSB)])
    txt(s, Inches(1.25), y+Inches(0.33), Inches(5.65), Inches(0.75), [(b, 11, SLATE, False)], space=1.03)
    y += Inches(1.08)
rect(s, Inches(7.1), Inches(1.95), Inches(5.6), Inches(4.6), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(7.4), Inches(2.15), Inches(5.0), Inches(0.4), [("WHY IT IS BELIEVABLE", 11, SLATED, True, FONTSB)])
dfm_outs = [
    ("Three outcomes, not two", "pass, fail, or NOT EVALUATED with the reason — a rule that could not run is never a pass", TEAL),
    ("200 / 200 geometry gate", "every fixture is a shape whose truth is arithmetic; CI fails on any regression", EMER),
    ("Nine primary standards, read", "NADCA, SFSA, DuPont, Covestro, ISO 8062-4, DIN 16742 — 33 thresholds now first-hand", GOLD),
    ("It reads the drawing too", "2D tolerances judged by the same rules, then cross-checked against the 3D model", PURPLE),
]
yy = Inches(2.6)
for t, b, c in dfm_outs:
    icon_badge(s, Inches(7.4), yy, "▪", c, d=Inches(0.36))
    txt(s, Inches(7.95), yy-Inches(0.05), Inches(4.6), Inches(0.4), [(t, 12.5, WHITE, True, FONTSB)], space=1.0)
    txt(s, Inches(7.95), yy+Inches(0.3), Inches(4.6), Inches(0.6), [(b, 10.5, SLATE, False)], space=1.02)
    yy += Inches(0.98)
notes(s, "This is the newest major tool and it closes the loop between design and cost. You upload a STEP file, tell it the material and the process, and it measures the part — draft angles, undercuts, wall thickness, whether a cutter can physically reach a surface. That measurement is pure geometry; there is no AI anywhere in it. Then 248 rules across 35 process families judge it, with the threshold resolved for your specific alloy, and every rule citing where its number came from. The thing suppliers react to is on the right of the third row: a finding does not just say 'undercut' — it paints the actual faces that broke the rule, and you can click a painted face to jump back to the finding. Two points of discipline worth stressing to a sceptic. First, three outcomes: pass, fail, or NOT EVALUATED with the reason — a rule that could not run is never quietly counted as a pass. Second, the accuracy gate: two hundred analytic fixtures whose correct answer is arithmetic, and CI fails on any regression. And since most of our supply base still puts requirements on a 2D drawing, it reads the drawing as well and tells you where the drawing and the model disagree.")

# ════════════════════════════════════════════════════════════════════════════
# HORIZON — technology foresight
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 5 — new", "Horizon: Which Technology, and When"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5),
    [("Sourcing decisions outlive the technology they were made on. Horizon dates the change — with evidence, and with a band around every year.", 14.5, TEAL, True, FONTSB)])
hz_steps = [
    ("◆", "Curated register", "180 technologies across every commodity, ICE through BEV — each with adoption, cost trajectory and a confidence tier"),
    ("↑", "Evidence, computed", "Patent-filing velocity from PatentsView and cited retrieval — the numbers come from engines, never from the model"),
    ("●", "Segment lenses", "One click focuses the whole register on our battleground — off-road, luxury SUV, software-defined vehicle"),
    ("▲", "Prediction ledger", "Old forecasts are snapshotted and scored against what happened. A foresight tool that never grades itself is astrology"),
]
y = Inches(2.2)
for g, t, b in hz_steps:
    icon_badge(s, Inches(0.6), y, g, PURPLE, d=Inches(0.46))
    txt(s, Inches(1.25), y-Inches(0.04), Inches(5.6), Inches(0.4), [(t, 13.5, WHITE, True, FONTSB)])
    txt(s, Inches(1.25), y+Inches(0.33), Inches(5.65), Inches(0.75), [(b, 11, SLATE, False)], space=1.03)
    y += Inches(1.08)
rect(s, Inches(7.1), Inches(1.95), Inches(5.6), Inches(4.6), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(7.4), Inches(2.15), Inches(5.0), Inches(0.4), [("WHAT ONE QUERY RETURNS", 11, SLATED, True, FONTSB)])
hz_outs = [
    ("Three horizon lanes", "act now, plan for it, or watch — so a roadmap can be built from it", TEAL),
    ("Crossing years, with a band", "when a technology reaches a quarter and half its market — never a single false-precision year", GOLD),
    ("Regulatory radar", "dated obligations with their LEGAL status, not just the headline", ROSE),
    ("A board-ready report", "expert-panel critique and a branded multi-page PDF, one click each", EMER),
]
yy = Inches(2.6)
for t, b, c in hz_outs:
    icon_badge(s, Inches(7.4), yy, "▪", c, d=Inches(0.36))
    txt(s, Inches(7.95), yy-Inches(0.05), Inches(4.6), Inches(0.4), [(t, 12.5, WHITE, True, FONTSB)], space=1.0)
    txt(s, Inches(7.95), yy+Inches(0.3), Inches(4.6), Inches(0.6), [(b, 10.5, SLATE, False)], space=1.02)
    yy += Inches(0.98)
notes(s, "The other tool this deck did not cover. The problem it solves is one we all live with: a sourcing decision commits us for seven years, and we make it on today's technology assumptions. Horizon dates the change. Behind it is a curated register of a hundred and eighty technologies across every commodity, from ICE through to battery-electric, each carrying an adoption curve, a cost trajectory and a confidence tier. The evidence layer is the part that separates it from asking a chatbot: patent-filing velocity is computed from the real PatentsView database, and retrieved sources are cited or the finding is dropped server-side. Numbers come from engines; the AI writes the prose. Two things to point at on the right. Crossing years always come with a band — anyone who gives you a single year for the future is selling false precision. And the prediction ledger: it snapshots its own old forecasts and scores them against what actually happened. A foresight tool that never checks its own predictions is astrology, and I would rather show you our error than hide it.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ACCURACY CHART (feature vs mass)
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Proof point", "Feature-Based Accuracy vs a Mass Estimate"); footer(s)
txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
    [("Held-out benchmark — average error vs reference prices (lower is better). Same parts, two methods.", 13.5, SLATE, False)])
# Hand-drawn grouped bar chart — no embedded workbook, fully native shapes.
base_y, top_y, cx0 = Inches(6.05), Inches(2.65), Inches(1.55)
plot_h = int(base_y) - int(top_y)
scale = plot_h / 80.0
for gv in (0, 20, 40, 60, 80):
    gy = Emu(int(int(base_y) - scale * gv))
    rect(s, cx0, gy, Inches(6.35), Pt(0.75), fill=NAVY3)
    txt(s, cx0 - Inches(0.62), gy - Inches(0.12), Inches(0.5), Inches(0.25), [(str(gv), 10, SLATE, False)], align=PP_ALIGN.RIGHT)
def _bar(xc, val, color):
    h = Emu(int(scale * val)); w = Inches(0.9)
    rect(s, xc, Emu(int(base_y) - h), w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, xc - Inches(0.15), Emu(int(base_y) - h) - Inches(0.36), w + Inches(0.3), Inches(0.3),
        [(f"{val:.1f}%", 12, WHITE, True, FONTSB)], align=PP_ALIGN.CENTER)
_bar(Inches(2.15), 43.7, SLATED); _bar(Inches(3.2), 34.5, TEAL)
_bar(Inches(5.0), 74.8, SLATED); _bar(Inches(6.05), 38.9, TEAL)
txt(s, Inches(1.9), base_y + Inches(0.12), Inches(2.35), Inches(0.35), [("Machining parts", 12.5, LIGHT, True, FONTSB)], align=PP_ALIGN.CENTER)
txt(s, Inches(4.75), base_y + Inches(0.12), Inches(2.35), Inches(0.35), [("Stamping parts", 12.5, LIGHT, True, FONTSB)], align=PP_ALIGN.CENTER)
txt(s, cx0 - Inches(0.62), top_y - Inches(0.42), Inches(6.5), Inches(0.3), [("Mean Absolute % Error", 12, SLATE, False)])
rect(s, Inches(1.95), Inches(6.62), Inches(0.28), Inches(0.2), fill=SLATED)
txt(s, Inches(2.32), Inches(6.57), Inches(2.1), Inches(0.3), [("Mass estimate (old)", 11, SLATE, False)])
rect(s, Inches(4.55), Inches(6.62), Inches(0.28), Inches(0.2), fill=TEAL)
txt(s, Inches(4.92), Inches(6.57), Inches(3.0), Inches(0.3), [("Feature-based (BrainSpark)", 11, LIGHT, False)])
# takeaways
rect(s, Inches(8.35), Inches(2.25), Inches(4.4), Inches(4.4), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.65), Inches(2.45), Inches(3.9), Inches(0.4), [("WHAT THIS SHOWS", 12, GOLD, True, FONTSB)])
for i,(t) in enumerate(["Feature-based cuts error on both — stamping error halves (75% → 39%).",
                        "Core should-cost engine: 100% hit-rate, 8.3% MAPE on 16 held-out parts (CI-gated).",
                        "Validated on held-out parts the model never calibrated against.",
                        "Method chosen for honesty, not to overfit the benchmark."]):
    yy = Inches(3.0)+i*Inches(0.9)
    icon_badge(s, Inches(8.65), yy, "✓", EMER, d=Inches(0.38))
    txt(s, Inches(9.2), yy-Inches(0.02), Inches(3.35), Inches(0.85), [(t, 11.5, LIGHT, False)], space=1.03)
notes(s, "Proof, not claims. This is a held-out benchmark — parts the model wasn't tuned on. The grey bars are the old weight-based estimate; the teal bars are the new feature-based method. The biggest jump is on stamping, where the old approach was worst — error halves from about 75% down to 39%; machining improves from 44% to 34%. And note the second bullet: the core should-cost engine itself — the one behind every number in the tool — scores a 100% hit-rate with 8.3% average error on sixteen held-out reference parts, and CI fails if that regresses. Be candid: the CAD-feature figures are engineering-grade estimates, not supplier quotes, and I deliberately stopped tuning to avoid overfitting the benchmark. That honesty is the point — we show ranges and we don't pretend to precision we don't have.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — INNOVATION METHODS (now 11 methods)
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capability 6", "Eleven Structured Innovation Methods — Built In"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
    [("Real methodology as deterministic logic — not a chatbot prompt. The AI turns each method's output into concrete, engine-checked embodiments.", 13, SLATE, False)])
methods = [
    ("TRIZ", "40 inventive principles + contradiction matrix", TEAL),
    ("Value Engineering", "Value index = worth ÷ cost, per function", GOLD),
    ("FAST Function-Cost Matrix", "Component cost mapped onto functions — attack poor value", PURPLE),
    ("DFA / DFMA", "Boothroyd-Dewhurst min-parts & efficiency", BLUE),
    ("Design-to-Cost", "Allocates the cost gap across buckets", EMER),
    ("Spec & Tolerance Challenge", "Engine-costed relaxations — CTQs stay locked", ROSE),
    ("Teardown Delta", "Attribute-by-attribute vs a benchmark part", TEAL),
    ("SCAMPER + Morphological", "Lateral prompts + Zwicky concept space", GOLD),
    ("Effects & Trends · Circularity", "Cheaper physical effects, next-gen jumps, EU ELV design", BLUE),
]
x, y = Inches(0.6), Inches(2.2)
for i,(t,b,c) in enumerate(methods):
    cx = x + (i%3)*Inches(4.08); cy = y + (i//3)*Inches(1.5)
    rect(s, cx, cy, Inches(3.9), Inches(1.36), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, cx, cy, Inches(3.9), Inches(0.56), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy+Inches(0.28), Inches(3.9), Inches(0.28), fill=c)
    txt(s, cx+Inches(0.25), cy, Inches(3.45), Inches(0.56), [(t, 12.5, WT, True, FONTSB)], anchor=MSO_ANCHOR.MIDDLE, space=0.95)
    txt(s, cx+Inches(0.25), cy+Inches(0.66), Inches(3.45), Inches(0.65), [(b, 10.5, LIGHT, False)], space=1.02)
notes(s, "A differentiator versus generic AI tools: these are genuine engineering methods implemented as deterministic code, then the AI dresses the output into concrete ideas for your specific part. TRIZ actually walks the contradiction matrix; DFA computes the real Boothroyd-Dewhurst minimum part count. Three are new since the last update and worth calling out: the FAST function-cost matrix — the classical value-engineering core, mapping every component's cost onto the functions it serves; the Spec and Tolerance Challenge, where relaxation savings are computed by the cost engine itself and safety-critical characteristics are locked in code so the AI can never propose touching them; and Teardown Delta — compare your part attribute-by-attribute against a benchmark and turn every gap into an idea. Method from theory, numbers from the engine, wording from the AI.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SELF-IMPROVING AI: DEEP MODE + LEARNING LOOP (NEW)
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "New — how it gets better", "Deep Mode & the Learning Loop"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5),
    [("The same architecture the research labs use for scientific discovery: propose, critique, rank, repair — verified at every step.", 14, TEAL, True, FONTSB)])
# left: deep mode pipeline
rect(s, Inches(0.6), Inches(2.15), Inches(6.1), Inches(4.35), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(0.9), Inches(2.32), Inches(5.5), Inches(0.4), [("DEEP MODE — opt-in, per analysis", 12, PURPLE, True, FONTSB)])
deep = [
    ("1", "Expert panel critiques", "Manufacturing, supplier-commercial and quality personas — each reviews every idea from its own discipline"),
    ("2", "Tournament ranking", "Ideas judged head-to-head, chess-style — the strongest rise on merit, not on wording"),
    ("3", "Verified repair", "Challenged or engine-contradicted ideas get one fix attempt — accepted only if the engine re-confirms it"),
]
yy = Inches(2.85)
for n, t, b in deep:
    icon_badge(s, Inches(0.9), yy, n, PURPLE, d=Inches(0.44))
    txt(s, Inches(1.52), yy-Inches(0.04), Inches(5.0), Inches(0.4), [(t, 13, WHITE, True, FONTSB)])
    txt(s, Inches(1.52), yy+Inches(0.33), Inches(5.0), Inches(0.75), [(b, 10.5, SLATE, False)], space=1.03)
    yy += Inches(1.18)
# right: learning loop
rect(s, Inches(6.95), Inches(2.15), Inches(5.78), Inches(4.35), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(7.25), Inches(2.32), Inches(5.2), Inches(0.4), [("THE LEARNING LOOP — always on", 12, EMER, True, FONTSB)])
loop = [
    ("Learns your taste", "approvals and production-confirmed savings become positive examples for the next run"),
    ("Targets the gaps", "a coverage map of 1,600+ verified ideas shows where thinking is thin — the AI aims there"),
    ("Kills duplicates", "near-identical ideas are auto-merged; restatements of known ideas get flagged"),
    ("Measured, not asserted", "a benchmark harness A/B-tests every pipeline change before it ships"),
]
yy = Inches(2.85)
for t, b in loop:
    icon_badge(s, Inches(7.25), yy, "▪", EMER, d=Inches(0.36))
    txt(s, Inches(7.8), yy-Inches(0.05), Inches(4.75), Inches(0.4), [(t, 12.5, WHITE, True, FONTSB)])
    txt(s, Inches(7.8), yy+Inches(0.3), Inches(4.75), Inches(0.6), [(b, 10.5, SLATE, False)], space=1.02)
    yy += Inches(0.92)
notes(s, "This slide answers the question every director should ask about AI tools: 'does it get better, or is this as good as it gets?' Two answers. Deep Mode is an opt-in switch for the analyses that matter: a panel of three AI experts — a manufacturing engineer, a commercial manager and a quality lead — critiques every idea from its own angle, the ideas then compete head-to-head like a chess tournament, and anything the panel or the engine knocks down gets exactly one repair attempt, which only survives if the engine re-confirms it. It costs a few times more tokens, which is why it's a switch and not the default. The learning loop runs all the time: the tool learns from what the team approves, aims at the white space in our idea library, and merges duplicates. And the last point matters most — we don't claim any of this made the ideas better, we measure it with a benchmark harness before each change ships. That's engineering discipline applied to the AI itself.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — MARKETPLACE + ANALYTICS
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Capabilities 7 & 8", "Idea Marketplace + Decision Analytics"); footer(s)
# left: marketplace
rect(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(4.85), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
icon_badge(s, Inches(0.9), Inches(1.95), "★", EMER, d=Inches(0.6))
txt(s, Inches(1.65), Inches(1.98), Inches(4.6), Inches(0.6), [("Idea Marketplace", 18, WHITE, True, FONTSB)], anchor=MSO_ANCHOR.MIDDLE)
for i,t in enumerate(["1,600+ ideas — each labelled Curated (benchmark-sourced) or Community (user-submitted)",
                      "Coverage map shows where the library is thin — the AI targets those gaps",
                      "Theme clusters + duplicate warnings keep it clean and searchable",
                      "Every idea carries an engine-checked should-cost"]):
    yy = Inches(2.85)+i*Inches(0.82)
    icon_badge(s, Inches(0.95), yy, "▪", EMER, d=Inches(0.36))
    txt(s, Inches(1.5), yy-Inches(0.02), Inches(4.75), Inches(0.75), [(t, 12.5, LIGHT, False)], space=1.03)
# right: analytics
rect(s, Inches(6.8), Inches(1.7), Inches(5.93), Inches(4.85), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
icon_badge(s, Inches(7.1), Inches(1.95), "▲", ROSE, d=Inches(0.6))
txt(s, Inches(7.85), Inches(1.98), Inches(4.6), Inches(0.6), [("Decision Analytics", 18, WHITE, True, FONTSB)], anchor=MSO_ANCHOR.MIDDLE)
for i,t in enumerate(["Monte-Carlo confidence bands on every cost",
                      "ROI & business-case calculator",
                      "Savings pipeline with G0–G3 stage-gate scorecards",
                      "Patent check with real, citable patent records"]):
    yy = Inches(2.85)+i*Inches(0.82)
    icon_badge(s, Inches(7.15), yy, "▪", ROSE, d=Inches(0.36))
    txt(s, Inches(7.7), yy-Inches(0.02), Inches(4.75), Inches(0.75), [(t, 12.5, LIGHT, False)], space=1.03)
notes(s, "Two capabilities that make it a platform rather than a gadget. On the left, the marketplace is an institutional memory — over sixteen hundred ideas an engineer can pull from instead of starting blank, each carrying an engine-checked cost. And note the honesty detail: every idea is labelled by origin — Curated means it came from benchmark and teardown curation, Community means one of us submitted it and it passed review. We never dress seeded content up as community wisdom; that's the same no-faking rule the cost figures follow. The library polices itself too: a coverage map shows where our thinking is thin, duplicate submissions get caught at the door, and ideas cluster into browsable themes. On the right, analytics turn ideas into decisions: confidence bands, an ROI calculator, and a savings pipeline with stage-gate scorecards — promoting an idea past an unfinished checklist warns you exactly what was skipped. Plus a patent check that cites real patent records, with the standing caveat that it's awareness, not legal advice.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — END-TO-END WORKFLOW
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "How it works", "End-to-End Workflow"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.45),
    [("From a part or a CAD file to a tracked, defensible saving — one continuous flow.", 14, SLATE, False)])
flow = [("1","Select / Upload","Choose a system or drop a CAD model","↑",TEAL),
        ("2","Configure","Volume, region, currency, lenses","◆",GOLD),
        ("3","Generate + Cost","AI ideates · engine prices each","●",BLUE),
        ("4","Verify","Engine stamps confirmed / contradicted","✓",EMER),
        ("5","Decide & Track","ROI, pipeline, export to sourcing","▲",ROSE)]
bx = Inches(0.55)
for i,(n,t,b,g,c) in enumerate(flow):
    rect(s, bx, Inches(2.35), Inches(2.28), Inches(2.5), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, bx, Inches(2.35), Inches(2.28), Inches(0.55), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx, Inches(2.6), Inches(2.28), Inches(0.3), fill=c)
    txt(s, bx+Inches(0.25), Inches(2.35), Inches(0.6), Inches(0.55), [(n, 20, WT, True, FONTSB)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, bx+Inches(0.8), Inches(2.35), Inches(1.4), Inches(0.55), [(g, 16, WT, True, "Arial")], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    txt(s, bx+Inches(0.18), Inches(3.05), Inches(2.0), Inches(0.5), [(t, 13.5, WHITE, True, FONTSB)], space=1.0)
    txt(s, bx+Inches(0.18), Inches(3.62), Inches(1.95), Inches(1.1), [(b, 11, SLATE, False)], space=1.05)
    if i<4: arrow(s, bx+Inches(2.3), Inches(3.35), Inches(0.26), color=GOLD)
    bx += Inches(2.53)
rect(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.15), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.9), Inches(5.5), Inches(11.6), Inches(0.85),
    [[("Feedback loop:  ", 14, TEAL, True, FONTSB),
      ("real supplier quotes feed back into the engine's calibration, so every quote makes the next estimate sharper.", 13.5, LIGHT, False)]],
    anchor=MSO_ANCHOR.MIDDLE, space=1.05)
notes(s, "Tie it all together as one flow. Select a system or upload a CAD file; set the context; the AI generates ideas while the engine prices each; the engine verifies; and you decide and track. Then the loop at the bottom is the clever part — when we later get a real supplier quote, it feeds back and calibrates the engine, so the tool gets more accurate the more we use it. This is the slide to keep on screen right before the live demo, because the demo will follow exactly these five steps.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — AI CAPABILITIES
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Under the hood", "AI Capabilities"); footer(s)
ai = [
    ("★","Agentic reasoning","The AI calls the cost engine as a tool, tests alternatives, and learns from rejections — it explores, it doesn't just answer.",TEAL),
    ("◆","Grounded web search","Live search for current commodity prices and recent OEM innovations before it generates.",GOLD),
    ("■","Multimodal input","Reads 3D models, 2D drawings and multi-sheet PDF drawing packs via vision.",BLUE),
    ("▲","Structured output","Forced tool-schemas guarantee clean, typed data — never free-text guesses.",PURPLE),
    ("●","Guardrails","Prompt-injection defences; untrusted part text treated as data, not instructions.",ROSE),
    ("=","Model tiering","Fast model for cheap steps, flagship model for the hard reasoning — cost-efficient.",EMER),
]
x, y = Inches(0.6), Inches(1.7)
for i,(g,t,b,c) in enumerate(ai):
    cx = x + (i%3)*Inches(4.08); cy = y + (i//3)*Inches(2.5)
    card(s, cx, cy, Inches(3.9), Inches(2.32), g, t, b, c)
notes(s, "For a technically-curious Director, here's what the AI actually does — and what keeps it safe. Agentic reasoning means it uses the cost engine as a tool, tries alternatives and learns when the engine rejects an impossible combination. It searches the web for live prices before answering. It reads CAD and PDF drawings, not just text. Outputs are forced into strict schemas so we get clean data. And there are guardrails: anything a user types is treated as data, never as instructions, so the tool can't be hijacked. Model tiering keeps running costs sensible.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TECHNOLOGY STACK
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Engineering", "Technology & Tools Used"); footer(s)
groups = [
    ("AI & Reasoning", TEAL, ["Anthropic Claude (Opus / Sonnet)", "Agentic tool-use loop", "Expert-panel & tournament ranking", "Structured JSON output"]),
    ("Frontend", BLUE, ["React + TypeScript + Vite", "Tailwind CSS UI", "Three.js 3D CAD viewer", "PWA / mobile (Capacitor)"]),
    ("Backend & Data", GOLD, ["Node.js + Express", "SQLite (better-sqlite3)", "OpenCascade CAD geometry", "FX, commodity, DigiKey & patent APIs"]),
    ("Cost Engine", EMER, ["Deterministic should-cost core", "Machining, stamping & PCB models", "Monte-Carlo simulation", "280+ automated tests & benchmarks"]),
]
x = Inches(0.6)
for t, c, items in groups:
    rect(s, x, Inches(1.75), Inches(2.95), Inches(4.75), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, Inches(1.75), Inches(2.95), Inches(0.75), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, Inches(2.15), Inches(2.95), Inches(0.35), fill=c)
    txt(s, x, Inches(1.75), Inches(2.95), Inches(0.75), [(t, 14.5, WT, True, FONTSB)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, it in enumerate(items):
        yy = Inches(2.75)+i*Inches(0.92)
        icon_badge(s, x+Inches(0.22), yy, "▪", c, d=Inches(0.34))
        txt(s, x+Inches(0.66), yy-Inches(0.03), Inches(2.2), Inches(0.85), [(it, 11.5, LIGHT, False)], space=1.02)
    x += Inches(3.08)
notes(s, "Keep this light unless the Director is technical. Headline: it's built on a modern, production-grade stack, not a prototype held together with string. Claude does the reasoning; React and Three.js give a polished UI with a real 3D viewer; Node and SQLite run the backend with live price and FX feeds; and the cost engine is proper software — unit-tested and benchmarked. If they ask 'is this maintainable and safe', the answer is yes: standard tools, tested code, and it runs self-contained.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — THE ENGINEERING EFFORT (NEW)
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "Behind the scenes", "The Engineering Effort Behind the Tool"); footer(s)
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5),
    [("Engineer-directed, AI-accelerated development — the same human + AI leverage the tool itself delivers.", 14.5, TEAL, True, FONTSB)])
# top stat strip
estats = [("95,000+", "lines of production code"), ("279", "source files"), ("758", "automated tests"), ("CI-gated", "accuracy benchmarks")]
sx = Inches(0.6)
for v, l in estats:
    rect(s, sx, Inches(2.15), Inches(2.95), Inches(1.05), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, sx, Inches(2.27), Inches(2.95), Inches(0.5), [(v, 24, GOLD, True, FONTSB)], align=PP_ALIGN.CENTER)
    txt(s, sx, Inches(2.82), Inches(2.95), Inches(0.35), [(l, 10.5, SLATE, False)], align=PP_ALIGN.CENTER)
    sx += Inches(3.08)
# left: drawn horizontal bar breakdown (no embedded chart — max compatibility)
rect(s, Inches(0.6), Inches(3.42), Inches(6.8), Inches(3.25), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(0.9), Inches(3.56), Inches(6.2), Inches(0.35), [("WHERE THE CODE LIVES", 11, SLATED, True, FONTSB)])
code = [
    ("Frontend app (React / TypeScript)", 20.4, BLUE),
    ("Cost engines + AI pipeline (backend)", 11.3, TEAL),
    ("Curated knowledge bases (13 domains)", 5.5, GOLD),
    ("Automated tests", 3.0, EMER),
    ("Benchmarks & eval harnesses", 1.8, PURPLE),
]
by = Inches(4.0); bar_x = Inches(0.9); max_w = 4.05; scale_k = max_w / 20.4
for t, v, c in code:
    txt(s, bar_x, by-Inches(0.02), Inches(6.1), Inches(0.3), [(t, 10.5, LIGHT, False)])
    rect(s, bar_x, by+Inches(0.22), Inches(v*scale_k), Inches(0.16), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bar_x+Inches(v*scale_k)+Inches(0.1), by+Inches(0.13), Inches(1.2), Inches(0.3), [(f"{v:.1f}k", 10.5, WHITE, True, FONTSB)])
    by += Inches(0.47)
# right: what the discipline buys
rect(s, Inches(7.65), Inches(3.42), Inches(5.08), Inches(3.25), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
txt(s, Inches(7.95), Inches(3.56), Inches(4.5), Inches(0.35), [("WHAT THAT DISCIPLINE BUYS", 11, SLATED, True, FONTSB)])
buys = [
    "Every cost engine ships with unit tests — 758 run on every change",
    "Accuracy proven on held-out parts, not tuned to look good",
    "AI changes are A/B-measured; a debt report flags unmeasured claims",
    "Daily automatic backups, a decision log & an honest known-gaps register",
    "Directed, reviewed & verified by an engineer — AI did the typing",
]
yy = Inches(3.96)
for t in buys:
    icon_badge(s, Inches(7.95), yy, "✓", EMER, d=Inches(0.32))
    txt(s, Inches(8.42), yy-Inches(0.02), Inches(4.15), Inches(0.55), [(t, 10, LIGHT, False)], space=1.02)
    yy += Inches(0.545)
notes(s, "This slide is about the effort and the standard — and be upfront about how it was built, because it's the strongest part of the story: I built this with AI-assisted development. I set the direction, made the engineering decisions, reviewed and tested everything; AI did the heavy lifting on the code. In other words, I used exactly the human-plus-AI way of working this tool offers our engineers — and this is what one engineer can produce with it on personal time. Walk the numbers: over forty-two thousand lines of production code, including five and a half thousand lines of curated engineering knowledge — the levers and benchmarks that make ideas automotive-grade rather than generic — and roughly one line of test or benchmark code for every eight lines of product. Then land on the discipline column: two hundred and eighty-five tests on every change, accuracy proven on held-out parts, AI changes A/B-measured with a debt report that flags anything still unmeasured, and it's run like a product — daily automatic backups, a written decision log, and an honest register of what's NOT production-ready yet. That last point is worth saying out loud: I keep a list of the gaps — monitoring, deployment, spend tracking — because knowing your limits is part of the engineering. AI accelerated the build; discipline governed it.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — BENEFITS
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "The value", "Benefits & Business Impact"); footer(s)
ben = [
    ("●","Speed","Costed ideas in minutes","From part or CAD to a defensible number in minutes, not the days a manual should-cost takes.",TEAL),
    ("✓","Confidence","Numbers that survive review","Every figure is engine-verified and reproducible — defensible in sourcing and programme gates.",EMER),
    ("■","Breadth","See the whole vehicle","One platform across all systems and commodities — no more single-tool blind spots.",BLUE),
    ("★","Leverage","Capture know-how","Turns scattered OEM best-practice into a reusable, searchable asset for the whole team.",GOLD),
]
x = Inches(0.6)
for g,t,st,b,c in ben:
    rect(s, x, Inches(1.75), Inches(2.95), Inches(4.35), fill=CARD, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    icon_badge(s, x+Inches(1.17), Inches(2.05), g, c, d=Inches(0.7))
    txt(s, x, Inches(2.9), Inches(2.95), Inches(0.4), [(t, 17, WHITE, True, FONTSB)], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.2), Inches(3.35), Inches(2.55), Inches(0.5), [(st, 12, c, True, FONTSB)], align=PP_ALIGN.CENTER, space=1.0)
    txt(s, x+Inches(0.28), Inches(4.0), Inches(2.4), Inches(2.0), [(b, 11.5, SLATE, False)], align=PP_ALIGN.CENTER, space=1.06)
    x += Inches(3.08)
txt(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
    [[("Net effect:  ", 14, GOLD, True, FONTSB), ("more cost-down ideas, found earlier, defended better — and captured as reusable team knowledge.", 14, LIGHT, True)]],
    align=PP_ALIGN.CENTER)
notes(s, "Translate features into value the Director cares about. Speed: ideas arrive while the design is still open, so they can actually be actioned. Confidence: the numbers survive a sourcing review because they're reproducible. Breadth: no blind spots — we see the whole car. Leverage: the know-how that usually walks out the door at 5pm becomes a searchable team asset. The one-line net effect at the bottom is your summary sentence — say it out loud.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DEMO
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, NAVY)
rect(s, 0, 0, Inches(0.35), EMU_H, fill=TEAL)
rect(s, Inches(0.35), 0, Inches(0.12), EMU_H, fill=GOLD)
# photorealistic engineering-workspace image (Canva-generated, user-selected)
pic(s, "demo", Inches(8.95), Inches(1.3), Inches(3.85), Inches(5.2), focus=0.38)
txt(s, Inches(1.1), Inches(1.5), Inches(7.6), Inches(0.5), [("SEE IT LIVE", 15, TEAL, True, FONTSB)])
txt(s, Inches(1.05), Inches(2.1), Inches(7.7), Inches(1.0), [("Live Demo", 52, WHITE, True, FONTSB)])
txt(s, Inches(1.1), Inches(3.3), Inches(7.6), Inches(0.5), [("A 5-minute walk through the real workflow:", 17, LIGHT, False, FONTL)])
demo = ["Generate cost-reduction ideas for a live vehicle system",
        "Upload a CAD part → feature-based should-cost + 3D view",
        "Photograph a PCB → AI reads the BOM → 12-country cost",
        "Run a should-cost, close a quote gap with targeted ideas",
        "Show the savings pipeline, scorecards & ROI business case"]
y = Inches(3.95)
for i,t in enumerate(demo):
    icon_badge(s, Inches(1.1), y, str(i+1), GOLD, d=Inches(0.5))
    txt(s, Inches(1.8), y+Inches(0.02), Inches(6.9), Inches(0.5), [(t, 14.5, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.62)
notes(s, "Transition to the live tool. Keep the demo tight — five beats, mirroring the workflow slide: generate ideas for a system, upload a CAD part to show feature-based cost and the 3D viewer, then the crowd-pleaser — photograph a circuit board and watch it read the components and cost it across twelve countries. Fourth beat: run a should-cost with a supplier quote and show the gap-closure ideas, so they see the engine verify in real time. Finish on the pipeline with its scorecards and ROI so it lands as a business tool, not a toy. If anything is slow or offline, narrate what it would show and move on — don't fight the tech in the room. Have a backup screen-recording ready just in case.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — ROADMAP
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s); header(s, "What's next", "Roadmap"); footer(s)
phases = [("NOW","Live today", EMER, ["DFM / DFA Studio — 248 rules, 35 families","Horizon foresight — 180 technologies dated","Self-improving AI ideation + Deep Mode","PCB suite: photo → BOM → 12-country cost","Marketplace, scorecards & patent check"]),
          ("NEXT","3–6 months", GOLD, ["Weld / assembly cost models","Empirical confidence bands","Commonality & variant reduction","Wider 2D-drawing coverage (GD&T datums)"]),
          ("LATER","Vision", TEAL, ["Composite & paint models","Team collaboration & approvals","ERP / PLM integration","Portfolio-level cost intelligence"])]
x = Inches(0.6)
for tag, when, c, items in phases:
    rect(s, x, Inches(1.8), Inches(3.95), Inches(4.7), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, Inches(1.8), Inches(3.95), Inches(0.95), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, Inches(2.4), Inches(3.95), Inches(0.35), fill=c)
    txt(s, x+Inches(0.35), Inches(1.9), Inches(3.3), Inches(0.5), [(tag, 22, WT, True, FONTSB)])
    txt(s, x+Inches(0.35), Inches(2.38), Inches(3.3), Inches(0.35), [(when, 12, WT, True, FONTSB)])
    for i, it in enumerate(items):
        yy = Inches(3.0)+i*Inches(0.68)
        icon_badge(s, x+Inches(0.28), yy, "▪", c, d=Inches(0.32))
        txt(s, x+Inches(0.7), yy-Inches(0.02), Inches(3.1), Inches(0.6), [(it, 11.5, LIGHT, False)], space=1.02)
    x += Inches(4.08)
notes(s, "Show momentum and a credible path — and note how much moved from 'next' into 'now' since the last version of this deck: the whole DFM/DFA Studio and the Horizon foresight tool are both new, and the feature-based casting and moulding models that sat under 'next' last time shipped inside the Studio. 'Next' is the near-term engineering: more commodity models and tighter empirical confidence bands. 'Later' is the vision — collaboration, PLM/ERP integration and portfolio-level cost intelligence. The subtext for the Director: this isn't finished-and-frozen, it's a living platform with visible velocity — and with a little support it accelerates.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — ASK / CLOSE
# ════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, NAVY)
rect(s, 0, 0, Inches(0.35), EMU_H, fill=TEAL)
rect(s, Inches(0.35), 0, Inches(0.12), EMU_H, fill=GOLD)
# photorealistic closing still life: machined gear + PCB (Canva-generated, user-selected)
pic(s, "close", Inches(9.35), Inches(2.9), Inches(3.4), Inches(3.4), focus=0.42)
txt(s, Inches(1.1), Inches(1.15), Inches(8), Inches(0.5), [("WHERE WE GO FROM HERE", 15, TEAL, True, FONTSB)])
txt(s, Inches(1.05), Inches(1.75), Inches(10.5), Inches(0.9), [("Let's put it to work", 46, WHITE, True, FONTSB)])
asks = [("▲","Run a live pilot","Pick one active programme or part family and trial BrainSpark on it."),
        ("■","Benchmark it","Validate its numbers against a few known supplier quotes together."),
        ("◆","Sponsor & shape","A little time and direction to prioritise the features that matter most to us.")]
y = Inches(3.05)
for g,t,b in asks:
    icon_badge(s, Inches(1.1), y, g, GOLD, d=Inches(0.62))
    txt(s, Inches(1.95), y-Inches(0.05), Inches(7.2), Inches(0.45), [(t, 18, WHITE, True, FONTSB)])
    txt(s, Inches(1.95), y+Inches(0.42), Inches(7.2), Inches(0.45), [(b, 13.5, SLATE, False)])
    y += Inches(1.15)
rect(s, Inches(1.1), Inches(6.5), Inches(11.2), Inches(0.03), fill=NAVY3)
txt(s, Inches(1.1), Inches(6.65), Inches(11), Inches(0.5),
    [[("Thank you.  ", 15, TEAL, True, FONTSB), ("Questions & discussion welcome   ·   ", 14, LIGHT, False), ("Avinash Bhosale · Cost Engineering", 14, SLATE, False)]])
notes(s, "Close with a clear, low-risk ask — don't leave it open-ended. Three concrete next steps: run a small pilot on one real programme, benchmark it together against a handful of known quotes so trust is earned not asserted, and give it light sponsorship to steer the roadmap. Then hand over to questions. End on the collaborative note: this is built to support the team's work, and the fastest way to prove its worth is to try it on something real.")

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..",
                   "BrainSpark_Director_Presentation.pptx")
prs.save(OUT)
print("wrote", os.path.relpath(OUT))
print("Saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
