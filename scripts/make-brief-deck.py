#!/usr/bin/env python3
"""BrainSpark — the 2-slide director brief.

    python3 scripts/make-brief-deck.py

Same design system as the platform deck. Its Lucide PNGs are committed under
scripts/deck-assets/icons/ so this runs on a clean checkout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF3, 0xF5, 0xF9)
NAVY2  = RGBColor(0xFF, 0xFF, 0xFF)
NAVY3  = RGBColor(0xE4, 0xEA, 0xF1)
BORDER = RGBColor(0xD8, 0xE0, 0xEA)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
GOLD   = RGBColor(0xB4, 0x76, 0x09)
EMER   = RGBColor(0x05, 0x96, 0x69)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ROSE   = RGBColor(0xE1, 0x1D, 0x48)
INK    = RGBColor(0x16, 0x27, 0x3C)
WT     = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x55, 0x64, 0x77)
SLATED = RGBColor(0x83, 0x90, 0xA1)
LIGHT  = RGBColor(0x37, 0x45, 0x58)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]
FONT, FONTL, FONTSB = "Segoe UI", "Segoe UI Light", "Segoe UI Semibold"

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        if isinstance(para, tuple): para = [para]
        for seg in para:
            t, size, color, bold = seg[0], seg[1], seg[2], seg[3]
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = seg[4] if len(seg) > 4 else FONT
    return tb

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

import os
ICON_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deck-assets", "icons")

def badge(s, x, y, glyph, color, d=Inches(0.5)):
    rect(s, x, y, d, d, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, y, d, d, [(glyph, 15, WT, True, "Arial")], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def ibadge(s, x, y, name, color, d=Inches(0.5)):
    """Colored rounded badge with a rendered Lucide icon (white line-art PNG)."""
    rect(s, x, y, d, d, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p = os.path.join(ICON_DIR, f"{name}.png")
    if os.path.exists(p):
        pad = int(d * 0.21)
        s.shapes.add_picture(p, x + pad, y + pad, width=d - 2 * pad, height=d - 2 * pad)

def mini_icon(s, x, y, name, color, d=Inches(0.3)):
    """Small tinted icon chip (no heavy badge) for dense grids."""
    rect(s, x, y, d, d, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p = os.path.join(ICON_DIR, f"{name}.png")
    if os.path.exists(p):
        pad = int(d * 0.18)
        s.shapes.add_picture(p, x + pad, y + pad, width=d - 2 * pad, height=d - 2 * pad)

def header(s, kicker, title):
    rect(s, 0, 0, EMU_W, Inches(1.2), fill=NAVY2)
    rect(s, 0, Inches(1.2), EMU_W, Pt(3), fill=TEAL)
    rect(s, Inches(0.55), Inches(0.34), Inches(0.11), Inches(0.52), fill=GOLD)
    txt(s, Inches(0.85), Inches(0.22), Inches(11), Inches(0.32), [(kicker.upper(), 11.5, TEAL, True, FONTSB)])
    txt(s, Inches(0.85), Inches(0.52), Inches(11.5), Inches(0.6), [(title, 25, INK, True, FONTSB)])
    txt(s, Inches(0.85), Inches(7.08), Inches(8), Inches(0.3),
        [[("BrainSpark", 9, TEAL, True), ("  ·  AI-Powered Cost Engineering Platform  ·  Avinash Bhosale", 9, SLATED, False)]])
    txt(s, Inches(9.9), Inches(7.08), Inches(2.6), Inches(0.3), [("Confidential — Internal", 9, SLATED, False)], align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — AI IDEA GENERATION (the headline capability)
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "The headline capability", "AI Idea Generation — Grounded, Learning, Verified")
txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.6),
    [[("Pick a vehicle system, set volume, region and vehicle type — get a full set of ", 14, LIGHT, False),
      ("costed, ranked, engine-checked cost-reduction ideas", 14, TEAL, True, FONTSB),
      (" in minutes. ", 14, LIGHT, False), ("AI proposes; the engine verifies; nothing is invented.", 14, GOLD, True, FONTSB)]], space=1.06)

# left: how the generation engine works
gen = [
    ("globe", "Grounded, not guessed", "Live web search for current prices and innovations + a curated knowledge base of proven levers across 13 vehicle domains, with OEM & teardown benchmarks.", TEAL),
    ("brain", "Learns your organisation", "Ideas you approve or confirm in production feed back as positive examples; the AI targets gaps in the 1,600+ idea library and auto-merges near-duplicates.", GOLD),
    ("users", "Deep Mode (opt-in)", "A 3-persona expert panel — manufacturing, supplier-commercial, quality — critiques every idea; a tournament ranks them; weak ideas get one engine-verified repair.", PURPLE),
    ("badge", "Ranked & defensible", "Explainable ranking: annual value × payback × validation quality × engine cross-check — every factor visible on the idea card, every boost labelled.", EMER),
]
y = Inches(2.15)
for g, t, b, c in gen:
    ibadge(s, Inches(0.6), y, g, c, d=Inches(0.5))
    txt(s, Inches(1.28), y-Inches(0.04), Inches(6.0), Inches(0.34), [(t, 13.5, INK, True, FONTSB)])
    txt(s, Inches(1.28), y+Inches(0.31), Inches(6.02), Inches(0.85), [(b, 10.5, SLATE, False)], space=1.05)
    y += Inches(1.17)

# right: example idea card (mirrors the real UI)
rect(s, Inches(7.6), Inches(2.1), Inches(5.15), Inches(4.55), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
mini_icon(s, Inches(7.9), Inches(2.24), "bulb", GOLD, d=Inches(0.32))
txt(s, Inches(8.32), Inches(2.28), Inches(4.2), Inches(0.32), [("WHAT AN IDEA LOOKS LIKE", 10.5, SLATED, True, FONTSB)])
rect(s, Inches(7.9), Inches(2.66), Inches(4.55), Inches(1.62), fill=BG, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.1), Inches(2.8), Inches(4.15), Inches(0.62), [("One-piece hot-stamped B-pillar tailored blank", 12.5, INK, True, FONTSB)], space=1.0)
txt(s, Inches(8.1), Inches(3.42), Inches(4.15), Inches(0.3), [("Deletes 14 welds + 2 dies · −0.9 kg/vehicle", 10, SLATE, False)])
for cx, cw, ct, cc in [(8.1, 1.35, "£2.04M/yr", EMER), (9.55, 1.0, "High", GOLD), (10.65, 1.35, "✓ Verified", TEAL)]:
    rect(s, Inches(cx), Inches(3.78), Inches(cw), Inches(0.36), fill=cc, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(cx), Inches(3.8), Inches(cw), Inches(0.32), [(ct, 10.5, WT, True, FONTSB)], align=PP_ALIGN.CENTER)
txt(s, Inches(7.9), Inches(4.45), Inches(4.6), Inches(0.3), [("ENGINE CROSS-CHECK", 10.5, SLATED, True, FONTSB)])
rect(s, Inches(7.9), Inches(4.78), Inches(4.55), Inches(0.92), fill=BG, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.1), Inches(4.9), Inches(4.2), Inches(0.72),
    [[("Should-cost recomputed:  ", 10.5, SLATE, False), ("£17.10", 10.5, EMER, True, FONTSB), (" vs ", 10.5, SLATE, False), ("£23.40", 10.5, ROSE, True, FONTSB), (" baseline — saving confirmed by the deterministic engine.", 10.5, LIGHT, False)]], space=1.04)
txt(s, Inches(7.9), Inches(5.85), Inches(4.6), Inches(0.7),
    [[("Every idea carries provenance: ", 10, SLATE, False), ("engine-verified · prior-art · similar-to-approved · panel critiques", 10, TEAL, True, FONTSB), (" — boosts are never silent.", 10, SLATE, False)]], space=1.05)
txt(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.32),
    [[("Also in the platform:  Should-Cost Engine · CAD → Cost · PCB → BOM → Cost · Marketplace & Savings Pipeline", 10.5, SLATED, False)]], align=PP_ALIGN.CENTER)
notes(s, "This is the star of the show, so spend the time here. The flow in one sentence: pick a system — say Body Structure — set volume, region and vehicle type, and in minutes you get a full set of costed, ranked ideas. Then the four things that make it different from asking a chatbot. First, it's grounded: it searches live prices and carries a curated knowledge base of proven levers for thirteen vehicle domains — so ideas cite real OEM practice, not generic advice. Second, it learns US: what we approve and confirm in production feeds back as positive examples, and it deliberately aims at the white space in our own idea library. Third, Deep Mode: a panel of three AI experts — manufacturing, commercial, quality — critiques every idea, they compete in a tournament, and anything weak gets one repair attempt that only survives if the engine re-confirms it. Fourth, the output is defensible: point at the example card — the saving chip, and below it the engine's own recomputed cost. That green verified tag is the whole philosophy: AI proposes, the engine verifies. Mention the bottom strip only in passing — the demo will show the rest of the platform.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INNOVATION STUDIO: 11 STRUCTURED METHODS
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Innovation Studio", "Eleven Engineering Methods — Not Chat Prompts")
txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.55),
    [[("Classical idea-generation methodology implemented as ", 13.5, LIGHT, False), ("deterministic logic", 13.5, TEAL, True, FONTSB),
      (" — the method computes the analysis, the AI turns it into concrete embodiments for your part, the engine checks every figure.", 13.5, LIGHT, False)]], space=1.06)

methods = [
    ("target", "TRIZ", "40 inventive principles + contradiction matrix", TEAL),
    ("scale", "Value Engineering", "Value index = worth ÷ cost per function", GOLD),
    ("network", "FAST Function-Cost Matrix", "Component cost mapped onto functions", PURPLE),
    ("boxes", "DFA / DFMA", "Boothroyd-Dewhurst minimum part count", BLUE),
    ("crosshair", "Design-to-Cost", "Cost gap allocated across buckets", EMER),
    ("sliders", "Spec & Tolerance Challenge", "Engine-costed relaxations · CTQs locked", ROSE),
    ("compare", "Teardown Delta", "Attribute-by-attribute vs a benchmark", TEAL),
    ("shuffle", "SCAMPER", "Seven-verb structured creativity sweep", GOLD),
    ("grid", "Morphological Analysis", "Zwicky concept-space recombination", PURPLE),
    ("zap", "Effects & Evolution Trends", "Cheaper physical effects · next-gen jumps", BLUE),
    ("recycle", "Design for Circularity", "Cut cost and meet EU end-of-life rules", EMER),
    ("layers", "+ 9 one-click lenses", "Layer onto any analysis — incl. Biomimicry", SLATED),
]
x0, y0 = Inches(0.6), Inches(2.05)
for i, (ic, t, b, c) in enumerate(methods):
    cx = x0 + (i % 4) * Inches(3.19); cy = y0 + (i // 4) * Inches(1.08)
    rect(s, cx, cy, Inches(3.05), Inches(0.96), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    mini_icon(s, cx+Inches(0.16), cy+Inches(0.16), ic, c, d=Inches(0.3))
    txt(s, cx+Inches(0.58), cy+Inches(0.12), Inches(2.4), Inches(0.42), [(t, 11, INK, True, FONTSB)], space=0.95)
    txt(s, cx+Inches(0.16), cy+Inches(0.56), Inches(2.78), Inches(0.36), [(b, 9, SLATE, False)], space=0.98)

# bottom row: real accuracy mini-chart + proof/next-step card
rect(s, Inches(0.6), Inches(5.45), Inches(6.2), Inches(1.3), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
mini_icon(s, Inches(0.82), Inches(5.58), "chart", TEAL, d=Inches(0.3))
txt(s, Inches(1.24), Inches(5.6), Inches(5.4), Inches(0.3), [("HELD-OUT ACCURACY — FEATURE-BASED VS MASS ESTIMATE (MAPE, lower is better)", 8.5, SLATED, True, FONTSB)])
acc = [("Machining", 43.7, 34.5), ("Stamping", 74.8, 38.9)]
ay = Inches(5.95); bar_x = Inches(1.75); bar_max = 3.6; scale_k = bar_max / 80.0
for label, mass, feat in acc:
    txt(s, Inches(0.82), ay+Inches(0.01), Inches(0.9), Inches(0.25), [(label, 9, LIGHT, True, FONTSB)])
    rect(s, bar_x, ay, Inches(mass*scale_k), Inches(0.115), fill=SLATED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bar_x+Inches(mass*scale_k)+Inches(0.07), ay-Inches(0.045), Inches(1.1), Inches(0.22), [(f"{mass:.1f}%  mass", 8, SLATED, False)])
    rect(s, bar_x, ay+Inches(0.16), Inches(feat*scale_k), Inches(0.115), fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bar_x+Inches(feat*scale_k)+Inches(0.07), ay+Inches(0.115), Inches(1.4), Inches(0.22), [(f"{feat:.1f}%  feature", 8, TEAL, True, FONTSB)])
    ay += Inches(0.42)
rect(s, Inches(7.0), Inches(5.45), Inches(5.73), Inches(1.3), fill=NAVY2, line=BORDER, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
mini_icon(s, Inches(7.22), Inches(5.58), "shield", EMER, d=Inches(0.3))
txt(s, Inches(7.64), Inches(5.6), Inches(4.9), Inches(0.3), [("PROOF, NOT PROMISES", 8.5, SLATED, True, FONTSB)])
txt(s, Inches(7.22), Inches(5.95), Inches(5.3), Inches(0.35),
    [[("Core engine: ", 10.5, LIGHT, False), ("100% hit-rate · 8.3% avg. error", 10.5, INK, True, FONTSB), (" (CI-gated) · 758 tests", 10.5, LIGHT, False)]], space=1.0)
mini_icon(s, Inches(7.22), Inches(6.32), "play", GOLD, d=Inches(0.3))
txt(s, Inches(7.64), Inches(6.36), Inches(5.0), Inches(0.32),
    [[("Next: ", 10.5, GOLD, True, FONTSB), ("5-minute live demo, then a pilot on one programme.", 10.5, LIGHT, False)]])
notes(s, "This slide answers 'how is this different from asking a chatbot for ideas?' The answer: methodology as software. Each of these eleven is a genuine engineering method implemented as deterministic code — TRIZ actually walks the contradiction matrix, DFA actually computes the Boothroyd-Dewhurst minimum part count, the FAST matrix enforces that component costs map onto functions and sum correctly. The AI's job is the last mile: turning the method's output into concrete embodiments for your specific part, which the cost engine then checks. Call out the three newest, because they're the most 'cost-engineering-native': the FAST function-cost matrix — the classical VE core; the Spec and Tolerance Challenge, where relaxation savings are computed by the engine itself and safety-critical characteristics are locked in code so the AI can never propose touching them; and Teardown Delta — benchmark your part attribute-by-attribute and turn every gap into an idea. The twelfth tile reminds them all of this is also available as one-click lenses on any analysis. Then the bottom band: the proof numbers are measured and CI-gated, not marketing — and go straight into the demo from here.")

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..",
                   "BrainSpark_Director_Brief.pptx")
prs.save(OUT)
print("wrote", os.path.relpath(OUT))
print("Saved:", len(prs.slides._sldIdLst), "slides")
