"""
CostVision — 16-Slide Executive Presentation Generator
Produces a professional .pptx file with dark theme, data tables,
two-column layouts and branded colour scheme.
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

from pptx_fixup import finalise

# ─── Brand colours — light theme, matched to CostVision-Workflow-Explained ────
# The palette is deliberately the same set of hex values as build_workflow_deck.mjs
# so the two decks read as one pack when they are presented back to back: a very
# light blue-grey page, white cards, navy headings, and accents dark enough to
# stay legible as TEXT on white (the dark-theme accents were tuned for the
# opposite job and turn to pastel mush on paper and on a bright projector).
#
# The token NAMES are historical — TEXT_W is the primary text colour, not white.
# They are kept because ~300 call sites use them, so the theme is one edit here.
BG          = RGBColor(0xF4, 0xF7, 0xFB)   # page
SURFACE     = RGBColor(0xF7, 0xF9, 0xFC)   # tinted panel / alternating table row
SURFACE2    = RGBColor(0xFF, 0xFF, 0xFF)   # card
BORDER      = RGBColor(0xDC, 0xE3, 0xEE)   # hairline rule
ACCENT_B    = RGBColor(0x1D, 0x6F, 0xB8)   # blue
ACCENT_G    = RGBColor(0x2E, 0x8B, 0x57)   # green
ACCENT_P    = RGBColor(0x6B, 0x3F, 0xA0)   # purple
ORANGE      = RGBColor(0xB7, 0x79, 0x1F)   # amber
RED         = RGBColor(0xB0, 0x3A, 0x2E)   # red
TEXT_W      = RGBColor(0x16, 0x32, 0x5C)   # primary — navy
TEXT_G      = RGBColor(0x3A, 0x43, 0x56)   # body — slate
TEXT_D      = RGBColor(0x6B, 0x72, 0x80)   # secondary — muted
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# The title slide stays a dark navy panel, exactly as the Workflow deck's does —
# a light deck still wants one dark plate to open on. These are the only colours
# used against it, and they are the ONLY place white type is correct.
HERO_BG     = RGBColor(0x16, 0x32, 0x5C)   # navy plate
HERO_PANEL  = RGBColor(0x1E, 0x40, 0x70)   # slightly lifted centre
HERO_TEXT   = RGBColor(0xFF, 0xFF, 0xFF)
HERO_SUB    = RGBColor(0xCA, 0xDC, 0xFC)
HERO_DIM    = RGBColor(0x8F, 0xA3, 0xCC)

# Workflow sets titles in Cambria and everything else in Calibri. Matching that
# is most of what makes the two decks look like siblings.
TITLE_FONT  = 'Cambria'

# Slide dimensions: 16:9 widescreen
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank layout

# ─── Helper functions ─────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Fill background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

# ── Emoji font handling ──────────────────────────────────────────────────────
# Calibri carries no colour-emoji coverage, so an emoji inside a Calibri run is
# left to font fallback — which resolves differently on every machine, and to an
# empty box where the fallback is missing. Splitting the run and naming an emoji
# font for the pictographic part makes the deck render the same on somebody
# else's laptop as it does here.
#
# Dingbats (U+2700-27BF, the tick and cross) are deliberately NOT routed here:
# Calibri draws them cleanly as typographic marks, and an emoji font would turn
# a neat tick into a colour sticker.
# U+2700-27BF is split deliberately: the tick/cross family (U+2713-2718)
# renders as clean typographic marks in Calibri, but its neighbours (the
# pencil, the question mark) render as colour emoji on most systems and so
# belong with the pictographs.
_PICTO = re.compile('([\U0001F300-\U0001FAFF\u2600-\u26FF\u2700-\u2712\u2719-\u27BF\uFE0F]+)')
EMOJI_FONT = 'Segoe UI Emoji'


def _emit_runs(p, t, size, color, bold, italic, base_font='Calibri'):
    """Add `t` to paragraph `p`, giving pictographs an emoji font."""
    for part in _PICTO.split(t):
        if not part:
            continue
        run = p.add_run()
        run.text = part
        f = run.font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
        if _PICTO.fullmatch(part):
            f.name = EMOJI_FONT
        elif base_font:
            f.name = base_font


def txb(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT_W, align=PP_ALIGN.LEFT,
        wrap=True, italic=False, font=None):
    """Add a text box."""
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    _emit_runs(p, text, size, color, bold, italic, base_font=font)
    return tf_box

def rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, x, y, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def accent_bar(slide, x, y, w=Inches(0.5), h=Inches(0.04), color=ACCENT_B):
    """Thin horizontal accent bar."""
    r = rect(slide, x, y, w, h, color)
    return r

def slide_header(slide, slide_num, section_label, title_text, subtitle_text=""):
    """Standard slide header with top bar, slide number, section label, title."""
    # Top navy strip — the Workflow deck's masthead
    rect(slide, 0, 0, W, Inches(0.10), HERO_BG)

    # Slide number (top-left)
    txb(slide, f"SLIDE {slide_num:02d} / 19", Inches(0.35), Inches(0.14), Inches(2), Inches(0.35),
        size=7.5, color=TEXT_D, bold=True)

    # CostVision logo (top-right)
    txb(slide, "CostVision", W - Inches(1.9), Inches(0.12), Inches(1.6), Inches(0.35),
        size=11, bold=True, color=ACCENT_B, align=PP_ALIGN.RIGHT)

    # Horizontal rule under top bar
    rect(slide, 0, Inches(0.52), W, Inches(0.008), BORDER)

    # Section label
    txb(slide, section_label.upper(), Inches(0.5), Inches(0.68), Inches(8), Inches(0.28),
        size=8, bold=True, color=ACCENT_B)

    # Main title
    txb(slide, title_text, Inches(0.5), Inches(0.9), Inches(12), Inches(0.6),
        size=26, bold=True, color=TEXT_W, font=TITLE_FONT)

    # Subtitle
    if subtitle_text:
        txb(slide, subtitle_text, Inches(0.5), Inches(1.45), Inches(12), Inches(0.38),
            size=12, color=TEXT_G)

    # Footer rule
    rect(slide, 0, H - Inches(0.45), W, Inches(0.008), BORDER)
    txb(slide, "CostVision — AI-Powered Should-Cost Intelligence  |  Avinash Bhosale  |  Confidential",
        Inches(0.4), H - Inches(0.4), W - Inches(0.8), Inches(0.35),
        size=7.5, color=TEXT_D, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body, accent=ACCENT_B, icon="",
         title_pt=9.5, body_pt=8.5):
    """Card with coloured left border, title, body."""
    # Card background
    r = rect(slide, x, y, w, h, SURFACE2, BORDER, Pt(0.5))
    # Accent left border
    rect(slide, x, y, Inches(0.06), h, accent)
    # Icon + Title
    title_text = f"{icon}  {title}" if icon else title
    txb(slide, title_text, x + Inches(0.14), y + Inches(0.1), w - Inches(0.2), Inches(0.34),
        size=title_pt, bold=True, color=TEXT_W)
    # Body
    txb(slide, body, x + Inches(0.14), y + Inches(0.42), w - Inches(0.2), h - Inches(0.54),
        size=body_pt, color=TEXT_G, wrap=True)

def stat_card(slide, x, y, w, h, number, label, color=ACCENT_B):
    """Stat card with big number."""
    rect(slide, x, y, w, h, SURFACE2, BORDER, Pt(0.5))
    txb(slide, number, x, y + Inches(0.1), w, Inches(0.55),
        size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, label, x, y + Inches(0.62), w, Inches(0.38),
        size=8, color=TEXT_G, align=PP_ALIGN.CENTER, wrap=True)

def pill(slide, x, y, text, color=ACCENT_B):
    """Small pill badge."""
    w = Inches(2.2)
    h = Inches(0.28)
    rect(slide, x, y, w, h, SURFACE2, color, Pt(0.8))
    txb(slide, text, x + Inches(0.1), y + Inches(0.02), w - Inches(0.15), h,
        size=7.5, color=color, bold=True)

def bullet_block(slide, x, y, w, h, items, title=None, title_color=ACCENT_B):
    """Block of bullet points."""
    yy = y
    if title:
        txb(slide, title, x, yy, w, Inches(0.3),
            size=9.5, bold=True, color=title_color)
        yy += Inches(0.3)
    for item in items:
        # Bullet dot
        rect(slide, x + Inches(0.05), yy + Inches(0.13), Inches(0.06), Inches(0.06), ACCENT_B)
        txb(slide, item, x + Inches(0.22), yy, w - Inches(0.25), Inches(0.36),
            size=8.5, color=TEXT_G, wrap=True)
        yy += Inches(0.33)

def flow_box(slide, x, y, w, h, num, icon, title, body, color=ACCENT_B):
    """Step box for workflow."""
    rect(slide, x, y, w, h, SURFACE2, BORDER, Pt(0.5))
    # Step number circle area
    rect(slide, x, y, w, Inches(0.04), color)
    # Number
    txb(slide, str(num), x, y + Inches(0.06), w, Inches(0.3),
        size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Icon
    txb(slide, icon, x, y + Inches(0.33), w, Inches(0.3),
        size=14, color=TEXT_W, align=PP_ALIGN.CENTER)
    # Title
    txb(slide, title, x + Inches(0.06), y + Inches(0.62), w - Inches(0.12), Inches(0.32),
        size=7.5, bold=True, color=TEXT_W, align=PP_ALIGN.CENTER)
    # Body
    txb(slide, body, x + Inches(0.06), y + Inches(0.92), w - Inches(0.12), h - Inches(1.0),
        size=7.5, color=TEXT_G, align=PP_ALIGN.CENTER, wrap=True)

def comm_card(slide, x, y, w, h, icon, name, sub, color=BORDER):
    """Small commodity card."""
    rect(slide, x, y, w, h, SURFACE2, color, Pt(0.8))
    txb(slide, icon, x, y + Inches(0.18), w, Inches(0.40),
        size=17, color=TEXT_W, align=PP_ALIGN.CENTER)
    txb(slide, name, x, y + Inches(0.64), w, Inches(0.26),
        size=9.0, bold=True, color=TEXT_W, align=PP_ALIGN.CENTER)
    txb(slide, sub, x, y + Inches(0.94), w, Inches(0.58),
        size=8.0, color=TEXT_D, align=PP_ALIGN.CENTER)

# ─── TABLE HELPER ─────────────────────────────────────────────────────────────

def add_table(slide, rows, cols, x, y, w, h, header_row, data_rows,
              col_widths=None, header_colors=None):
    """Add a styled table."""
    table = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    # Header row
    for c, text in enumerate(header_row):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HERO_BG
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(7.5)
        run.font.bold = True
        run.font.color.rgb = HERO_TEXT

    # Data rows
    for r, row_data in enumerate(data_rows):
        bg = SURFACE2 if r % 2 == 0 else SURFACE
        for c, text in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(7.5)
            # Colour cost column green
            if c == 3:
                run.font.color.rgb = ACCENT_G
                run.font.bold = True
            elif c == 0:
                run.font.color.rgb = ACCENT_B
                run.font.bold = True
            else:
                run.font.color.rgb = TEXT_G


def notes(slide, text):
    """Attach plain-text speaker notes to a slide (humanised, conversational)."""
    slide.notes_slide.notes_text_frame.text = text.strip()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

# Navy plate. This is the one slide in the deck that is NOT light, and so the
# one place HERO_TEXT/HERO_SUB (white and pale blue) are the correct colours —
# everywhere else those would be invisible.
rect(slide, 0, 0, W, H, HERO_BG)
# The lifted panel is sized to the content it holds (wordmark down to the second
# pill row) and centred on the slide — it used to run to 6.5" with the bottom
# third empty, and sat 0.09" off centre.
# It closes above the capability pills, which are wider than it is — the pills
# then sit cleanly on the navy plate instead of straddling the panel edge.
_HP_W = Inches(9.53)
rect(slide, (W - _HP_W) / 2, Inches(0.85), _HP_W, Inches(3.30), HERO_PANEL)
rect(slide, 0, 0, W, Inches(0.08), ACCENT_B)

# Main logo / wordmark
txb(slide, "CostVision", Inches(0.6), Inches(1.0), Inches(12), Inches(1.4),
    size=72, bold=True, color=HERO_TEXT, align=PP_ALIGN.CENTER, font=TITLE_FONT)

# Accent tagline bar
rect(slide, Inches(4.5), Inches(2.45), Inches(4.4), Inches(0.06), HERO_SUB)

txb(slide, "AI-Powered Should-Cost Intelligence for Modern Engineering",
    Inches(0.6), Inches(2.55), Inches(12), Inches(0.6),
    size=17, color=HERO_SUB, align=PP_ALIGN.CENTER)

txb(slide, "Designed & Developed by  Avinash Bhosale",
    Inches(0.6), Inches(3.2), Inches(12), Inches(0.45),
    size=13, bold=True, color=HERO_TEXT, align=PP_ALIGN.CENTER)

txb(slide, "Cost Engineering & Digital Innovation",
    Inches(0.6), Inches(3.65), Inches(12), Inches(0.35),
    size=10, color=HERO_DIM, align=PP_ALIGN.CENTER)

# Capability pills row
pill_data = [
    ("21 Commodity Models",  ACCENT_G, Inches(1.0)),
    ("AI Agent (NLP)",       ACCENT_B, Inches(3.4)),
    ("CAD-to-Cost AI",       ACCENT_P, Inches(5.8)),
    ("DFM / DFA Layer",      ORANGE,   Inches(8.2)),
    ("20 Regions · 10 FX",  ACCENT_B, Inches(10.6)),
]
for label, col, px in pill_data:
    rect(slide, px, Inches(4.55), Inches(2.1), Inches(0.36), SURFACE2, col, Pt(0.8))
    txb(slide, label, px + Inches(0.08), Inches(4.57), Inches(1.95), Inches(0.32),
        size=8, bold=True, color=col, align=PP_ALIGN.CENTER)

# Second row pills
pill2_data = [
    ("Assembly BOM Rollup",    ACCENT_G, Inches(1.8)),
    ("Wright's Law Learning Curve", ACCENT_P, Inches(4.2)),
    ("Supplier Quote Comparison",   ORANGE,   Inches(6.9)),
    ("Cloud Team Sync",             ACCENT_B, Inches(9.5)),
]
for label, col, px in pill2_data:
    rect(slide, px, Inches(5.05), Inches(2.5), Inches(0.33), SURFACE2, col, Pt(0.8))
    txb(slide, label, px + Inches(0.08), Inches(5.07), Inches(2.35), Inches(0.29),
        size=7.5, bold=True, color=col, align=PP_ALIGN.CENTER)

# Footer
rect(slide, 0, H - Inches(0.45), W, Inches(0.008), HERO_PANEL)
txb(slide, "CONFIDENTIAL — Management Review  |  July 2026",
    Inches(0.4), H - Inches(0.4), W - Inches(0.8), Inches(0.35),
    size=7.5, color=HERO_DIM, align=PP_ALIGN.CENTER)

notes(slide,
    "Welcome, and thanks for making the time. What you're looking at is CostVision — a "
    "should-cost platform I've built to answer one deceptively simple question: what should this "
    "part actually cost to make? Not what a supplier quotes, but what the engineering and the "
    "economics say it should be. "
    "The one principle I want you to hold onto for the whole session is this: the AI never sets a "
    "price. It reads the input — a CAD model, a photo of a board, a plain-English description — and "
    "it classifies and interprets. Every pound of cost after that is deterministic arithmetic we can "
    "trace line by line. That's what makes the number defensible in a negotiation. "
    "Across the next few slides I'll show you the breadth — twenty-one commodities, twenty regions — "
    "and then the part I'm proudest of: as of this year the tool checks its own homework and learns "
    "from real quotes, and I've proven that on five real automotive CAD parts.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 2, "The Challenge", "Why CostVision? Traditional Costing is Broken",
             "Manual effort, weeks of delay, and supplier dependency cost organisations millions annually.")

# 6 problem cards in 3×2 grid
problems = [
    (RED,    "⏱", "Weeks, Not Minutes",
     "Manual should-cost takes 2–4 weeks per part. Sourcing decisions stall. Programme timelines slip."),
    (ORANGE, "🔒", "Supplier Dependency",
     "Without independent cost intelligence, engineers accept quotes at face value — no leverage, no transparency."),
    (RED,    "📊", "Inconsistent Methodology",
     "Every engineer calculates cost differently. No standard model. No auditable assumptions. No cross-team comparability."),
    (ORANGE, "🌍", "No Regional Visibility",
     "Labour, machine, material and energy rates vary dramatically by country. Manual adjustment is error-prone."),
    (RED,    "🎯", "Late-Stage Cost Surprises",
     "Cost issues discovered at SOP — not at concept. DFM problems missed. Expensive redesigns. Margin erosion at launch."),
    (ORANGE, "💸", "No Supplier Quote Benchmark",
     "Buyers lack a defensible floor price for negotiation. Supplier margins are opaque and unchallenged."),
]

# The card rows and the stat bar below them are sized from a common bottom
# target, so the block finishes just above the footer instead of leaving an inch
# of empty page under it.
cols = 3
sx, sy = Inches(0.45), Inches(2.0)
gap = Inches(0.12)
cw = (W - sx * 2 - gap * (cols - 1)) / cols
STAT_BAR_H = Inches(0.86)
_rows2 = -(-len(problems) // cols)
ch = (Inches(6.95) - STAT_BAR_H - Inches(0.14) - sy - gap * (_rows2 - 1)) / _rows2

for i, (col, ico, title, body) in enumerate(problems):
    r, c = divmod(i, cols)
    cx = sx + c * (cw + gap)
    cy = sy + r * (ch + gap)
    card(slide, cx, cy, cw, ch, f"{ico}  {title}", body, accent=col,
         title_pt=10.5, body_pt=9.5)

# Industry stat bar
_bar_y = sy + _rows2 * (ch + gap) + Inches(0.02)
rect(slide, sx, _bar_y, W - sx * 2, STAT_BAR_H, SURFACE2, BORDER, Pt(0.5))
# "80%" at 36 pt is wider than the 1.1" box, so wrapping broke it across two
# lines and dropped the "%" out of the panel. Centred and unwrapped it stays on
# one line inside the panel and clear of the sentence beside it.
txb(slide, "80%", sx + Inches(0.15), _bar_y + Inches(0.06), Inches(1.1), Inches(0.72),
    size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, wrap=False)
txb(slide, "of part cost is locked in at the design stage — yet most cost analysis happens after design freeze."
    "  CostVision shifts cost intelligence to where it matters: concept phase.",
    sx + Inches(1.35), _bar_y + Inches(0.14), Inches(10.8), Inches(0.62),
    size=10, color=TEXT_G)
assert _bar_y + STAT_BAR_H <= Inches(6.98), 'slide 2 stat bar runs into the footer'

notes(slide,
    "Before I show you the tool, let me be honest about the problem it's solving, because most of "
    "us in this room have lived it. A should-cost done by hand takes two to four weeks per part. By "
    "the time it lands, the sourcing decision has often already been made. And because everyone "
    "builds their model slightly differently, you can't compare one engineer's number to another's — "
    "there's no single method, no auditable trail. "
    "The stat at the bottom is the one that really matters: roughly eighty percent of a part's cost "
    "is locked in at the design stage, but almost all of our cost analysis happens after the design "
    "is frozen — when changing anything is expensive. So we discover the cost problem at exactly the "
    "point we can no longer cheaply fix it. "
    "Everything CostVision does is aimed at pulling that intelligence forward, to the concept phase, "
    "where a change is still just a conversation.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 3, "The Solution", "What is CostVision?",
             "A world-class should-cost engine with AI at every layer — from natural language input to DFM/DFA output.")

cards_data = [
    (ACCENT_B, "🤖", "AI Agent (Natural Language)",
     "Describe a part in plain English. AI auto-builds the complete cost model: material, routing, machines, operations."),
    (ACCENT_P, "📐", "CAD-to-Cost Automation",
     "Upload STEP / IGES / photo. AI extracts geometry, infers material, selects process model — full cost in minutes."),
    (ACCENT_G, "🧠", "21 Engineering-Grade Cost Models",
     "Physics-based models for every major commodity. Routing logic, yield curves, tooling amortisation — all encoded."),
    (ORANGE,   "🔗", "Assembly BOM Rollup",
     "Cost entire assemblies — not just single parts. Multi-level BOM with per-component commodity models."),
    (ACCENT_B, "📉", "Learning Curve / Wright's Law",
     "Project cost at any volume milestone. Configure learning curve %. Essential for LTA negotiations."),
    (ACCENT_P, "💰", "Supplier Quote Comparison",
     "Log supplier quotes. Compare vs should-cost instantly. Identify margin gaps. Generate negotiation talking points."),
    (ACCENT_G, "🔬", "DFM / DFA Intelligence",
     "Bolt-on AI scores manufacturability 1–10. Critical issues, saving potential, actionable recommendations."),
    (ORANGE,   "🌍", "20 Regions · 10 Currencies",
     "Real-time labour, machine, material and energy rates. Currency auto-switches when region is selected."),
    (ACCENT_B, "📈", "Export, Scenarios & Team Sync",
     "6-sheet Excel, PDF report, A/B/C scenario comparison, tornado sensitivity chart, cloud team sync."),
]

cw = Inches(4.1)
ch = Inches(1.45)
sx, sy = Inches(0.45), Inches(1.98)
gap = Inches(0.1)

for i, (col, ico, title, body) in enumerate(cards_data):
    r, c = divmod(i, 3)
    cx = sx + c * (cw + gap)
    cy = sy + r * (ch + gap)
    card(slide, cx, cy, cw, ch, f"{ico}  {title}", body, accent=col)

notes(slide,
    "So here's the whole platform on one slide. I'll come back to several of these in detail, so "
    "don't feel you need to read every card. "
    "The three I'd point you to first are the entry points, because they're what make it fast. You "
    "can describe a part in plain English and the AI builds the model for you. You can upload a CAD "
    "file or even a photo, and it measures the geometry and infers the process. Or, if you prefer, "
    "you fill in the form yourself. Whichever door you come in, you land on the same engine — "
    "twenty-one engineering-grade commodity models underneath. "
    "The rest of the cards are the things that make it a real tool rather than a calculator: assembly "
    "roll-ups, learning curves for volume pricing, supplier-quote comparison, DFM scoring, twenty "
    "regions with automatic currency, and full export. The theme to notice is that the AI is doing "
    "the reading and interpreting — the cost maths stays deterministic underneath.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI Agent: Natural Language Costing
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 4, "Game-Changing Capability", "AI Agent — Natural Language Should-Costing",
             "The fastest path from part concept to full cost breakdown: just describe the part in plain English.")

# Left column — how it works
lx, ly, lw = Inches(0.45), Inches(2.0), Inches(5.8)
txb(slide, "How It Works", lx, ly, lw, Inches(0.3), size=10, bold=True, color=ACCENT_B)

steps = [
    ("1. Describe", "Type a part description: 'Aluminium bracket, 6082-T6, approximately 200×100×50mm, 3 machined holes, anodised finish, 5,000/year UK production'"),
    ("2. AI Analyses", "Claude AI interprets the description, infers geometry, selects the most cost-effective manufacturing route, and chooses appropriate machines and labour grades"),
    ("3. Auto-Populate", "The cost model is automatically populated: material weight, 3 machining operations, cycle times, OEE, tooling NRE — all filled without manual input"),
    ("4. Full Result", "Click Calculate. Full 8-bucket cost breakdown, AI insights, DFM/DFA score, sensitivity analysis, and export — all ready in under 60 seconds"),
]

yy = Inches(2.35)
for title, body in steps:
    rect(slide, lx, yy, lw, Inches(1.02), SURFACE2, BORDER, Pt(0.5))
    rect(slide, lx, yy, Inches(0.06), Inches(1.02), ACCENT_B)
    txb(slide, title, lx + Inches(0.14), yy + Inches(0.08), lw - Inches(0.2), Inches(0.28),
        size=9, bold=True, color=TEXT_W)
    txb(slide, body, lx + Inches(0.14), yy + Inches(0.34), lw - Inches(0.2), Inches(0.62),
        size=8, color=TEXT_G, wrap=True)
    yy += Inches(1.08)

# Right column — chat mockup
rx = Inches(6.6)
ry = Inches(2.0)
rw = Inches(6.3)

# Chat bubbles
chats = [
    ("USER",  "Describe a 4-layer FR4 PCB, 100×80mm, ENIG finish, flying probe test, 10,000/yr, UK manufacture"),
    ("AI",    "Got it. Analysing your PCB specification...\n• 4-layer FR4 standard, 100×80 mm board\n• ENIG surface finish (+22% vs HASL)\n• Flying probe electrical test\n• UK pricing region applied\n• NRE amortised over 10,000 boards"),
    ("USER",  "What is the should-cost?"),
    ("AI",    "Should-Cost Result: £ 4.18 per board\n• PCB Fab (material + process): £ 3.24\n• Test (flying probe): £ 0.71\n• NRE amortised: £ 0.23\nDFM Score: 8.5/10 — Good manufacturability"),
]

# The panel height is DERIVED from the transcript rather than typed. It used to
# be a fixed 4.92", the transcript overran it, and the "Key Benefits" heading
# clamped back on top of the last AI bubble while the second row of benefits
# dropped through the panel floor onto the footer.
BUB_PAD, BUB_LINE, BUB_GAP = 0.26, 0.21, 0.08
_transcript = sum(BUB_PAD + (m.count('\n') + 1) * BUB_LINE for _, m in chats) \
    + BUB_GAP * (len(chats) - 1)
rh = Inches(0.44 + _transcript + 0.14)
rect(slide, rx, ry, rw, rh, SURFACE, BORDER, Pt(0.5))

# Mockup title bar
rect(slide, rx, ry, rw, Inches(0.36), SURFACE2)
txb(slide, "●  CostVision AI Agent", rx + Inches(0.14), ry + Inches(0.07),
    rw - Inches(0.2), Inches(0.25), size=8.5, bold=True, color=ACCENT_B)

cy2 = ry + Inches(0.44)
for role, msg in chats:
    is_user = role == "USER"
    # The AI reply used to be a dark blue bubble against a dark page. On a light
    # page the same distinction is made the other way round: the user's turn is
    # the plain white bubble, the AI's is a pale blue tint.
    bg_col = SURFACE2 if is_user else RGBColor(0xE8, 0xF1, 0xFA)
    border_col = BORDER if is_user else ACCENT_B
    lines = msg.count('\n') + 1
    bh = Inches(BUB_PAD + lines * BUB_LINE)
    if is_user:
        rect(slide, rx + Inches(0.14), cy2, rw - Inches(0.28), bh, bg_col, border_col, Pt(0.5))
        txb(slide, f"You: {msg}", rx + Inches(0.24), cy2 + Inches(0.06),
            rw - Inches(0.5), bh - Inches(0.1), size=7.5, color=TEXT_G)
    else:
        rect(slide, rx + Inches(0.14), cy2, rw - Inches(0.28), bh, bg_col, border_col, Pt(0.5))
        txb(slide, f"AI: {msg}", rx + Inches(0.24), cy2 + Inches(0.06),
            rw - Inches(0.5), bh - Inches(0.1), size=7.5, color=TEXT_W)
    cy2 += bh + Inches(0.08)

# Key benefits — one strip under the panel, label and all four on a single line,
# so it cannot collide with the transcript however long that grows.
benefits = ["⚡ Zero form-filling", "🎯 Any commodity",
            "✓ Assumptions visible", "📊 Full audit trail"]
by = ry + rh + Inches(0.12)
txb(slide, "Key Benefits", rx, by, Inches(0.95), Inches(0.22),
    size=8.5, bold=True, color=ACCENT_G, wrap=False)
_bw = (rw - Inches(1.05)) / len(benefits)
for i, b in enumerate(benefits):
    txb(slide, b, rx + Inches(1.05) + i * _bw, by, _bw, Inches(0.22),
        size=7.5, color=TEXT_G, wrap=False)
assert by + Inches(0.22) <= H - Inches(0.45), 'benefits strip runs into the footer'

notes(slide,
    "This is the feature that gets the biggest reaction in a live demo, so let me walk the flow on "
    "the right. Someone types a description — here, a four-layer PCB with an ENIG finish, ten "
    "thousand a year, made in the UK. The AI reads that, fills in the whole model — layer count, "
    "finish premium, test method, the NRE spread over the batch — and comes back with four pounds "
    "eighteen a board, plus a manufacturability score. "
    "The thing I always stress at this point: the AI didn't invent that price. It interpreted the "
    "words into engineering inputs, and the deterministic engine did the arithmetic. So every one of "
    "those assumptions is visible and editable — nothing is a black box. "
    "What it buys you is zero manual form-filling, it works for any commodity, and there's a full "
    "audit trail behind every figure. That last point is what makes the number safe to take into a "
    "supplier meeting.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CAD-to-Cost Automation
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 5, "Automation", "CAD-to-Cost: Geometry to Should-Cost Automatically",
             "Upload a STEP file or part photo — AI extracts features, infers material, and generates a full cost model.")

# 3 stat boxes, then the supported-formats panel beside them. These used to be
# laid on the same 1.98"–3.00" band at the same x, so the formats bar was drawn
# straight over all three stat cards and hid them completely.
STAT_Y, STAT_H = Inches(1.98), Inches(1.0)
stats = [("10×", "Faster than manual costing", ACCENT_B),
         ("±8%", "Typical model accuracy", ACCENT_G),
         ("0", "Manual routing steps required", ORANGE)]
_sw5, _sg5 = Inches(1.75), Inches(0.1)
for i, (num, lbl, col) in enumerate(stats):
    stat_card(slide, Inches(0.45) + i * (_sw5 + _sg5), STAT_Y, _sw5, STAT_H, num, lbl, col)
_fmt_x = Inches(0.45) + len(stats) * (_sw5 + _sg5) + Inches(0.1)
assert _fmt_x >= Inches(0.45) + 3 * (_sw5 + _sg5), 'formats panel overlaps the stat cards'

# What AI extracts
ex_items = [
    "Part geometry: envelope dimensions, volume, wall thickness, surface area",
    "Feature recognition: holes, pockets, ribs, undercuts, threads, datum features",
    "Material inference: visual texture & density → grade selection (e.g. Al 6082-T6)",
    "Process selection: geometry → most cost-effective manufacturing route",
    "Cycle time: material removal rate, tool path length, fixture & setup strategy",
    "Complexity: Simple / Moderate / Complex → machine type & OEE selection",
    "DFM pre-check: thin walls, tight radii, undercuts flagged before costing begins",
]
rect(slide, Inches(0.45), Inches(3.12), Inches(6.0), Inches(3.8), SURFACE2, BORDER, Pt(0.5))
rect(slide, Inches(0.45), Inches(3.12), Inches(6.0), Inches(0.04), ACCENT_B)
txb(slide, "What AI Extracts from CAD / Photo", Inches(0.6), Inches(3.16),
    Inches(5.7), Inches(0.3), size=9.5, bold=True, color=TEXT_W)
yy = Inches(3.5)
for item in ex_items:
    rect(slide, Inches(0.65), yy + Inches(0.13), Inches(0.06), Inches(0.06), ACCENT_B)
    txb(slide, item, Inches(0.84), yy, Inches(5.4), Inches(0.36),
        size=8, color=TEXT_G, wrap=True)
    yy += Inches(0.34)

# Result breakdown mockup (right side)
rx = Inches(6.85)
ry = Inches(3.12)
rw = Inches(6.0)
rh = Inches(3.8)
rect(slide, rx, ry, rw, rh, SURFACE, BORDER, Pt(0.5))
rect(slide, rx, ry, rw, Inches(0.04), ACCENT_P)
txb(slide, "AI Analysis Result — Aluminium Bracket", rx + Inches(0.15), ry + Inches(0.08),
    rw - Inches(0.25), Inches(0.3), size=9.5, bold=True, color=TEXT_W)

# Detection results grid
det = [("AI Detected Material", "Al 6082-T6", ACCENT_B),
       ("Suggested Process", "3-Axis CNC Milling + Drilling", ACCENT_G),
       ("Envelope (L×W×H)", "220 × 140 × 85 mm", TEXT_W),
       ("Est. Net Weight", "0.82 kg", TEXT_W),
       ("Complexity", "Moderate — 3 setups", ORANGE),
       ("DFM Pre-check", "1 flag: 3 setups → consolidate", ORANGE)]
for i, (lbl, val, col) in enumerate(det):
    dx = rx + Inches(0.15) + (i % 2) * Inches(2.9)
    dy = ry + Inches(0.5) + (i // 2) * Inches(0.52)
    rect(slide, dx, dy, Inches(2.7), Inches(0.46), SURFACE2, BORDER, Pt(0.3))
    txb(slide, lbl, dx + Inches(0.08), dy + Inches(0.03), Inches(2.54), Inches(0.18),
        size=7.5, color=TEXT_D)
    txb(slide, val, dx + Inches(0.08), dy + Inches(0.22), Inches(2.54), Inches(0.2),
        size=8.5, bold=True, color=col)

# Cost bar
rect(slide, rx + Inches(0.15), ry + Inches(2.08), rw - Inches(0.3), Inches(0.52), SURFACE2, BORDER, Pt(0.3))
txb(slide, "Should-Cost Preview", rx + Inches(0.3), ry + Inches(2.12), Inches(3), Inches(0.22),
    size=8, color=TEXT_G)
txb(slide, "£ 28.40", rx + Inches(3.8), ry + Inches(2.1), Inches(1.9), Inches(0.3),
    size=14, bold=True, color=ACCENT_G, align=PP_ALIGN.RIGHT)

# mini bar chart (5 segments)
bar_data = [(0.38, ACCENT_B, "Mat 38%"), (0.32, ACCENT_P, "Proc 32%"),
            (0.12, ACCENT_G, "Lab 12%"), (0.08, ORANGE, "Tool 8%"),
            (0.10, TEXT_D, "OH 10%")]
bx_start = rx + Inches(0.15)
for pct, col, lbl in bar_data:
    bw_seg = (rw - Inches(0.3)) * pct
    rect(slide, bx_start, ry + Inches(2.66), bw_seg, Inches(0.16), col)
    bx_start += bw_seg

txb(slide, "  ".join(f"■ {l}" for _, _, l in bar_data),
    rx + Inches(0.15), ry + Inches(2.86), rw - Inches(0.3), Inches(0.22),
    size=7.5, color=TEXT_D)

# File types supported — right of the stat cards, aligned to the panel below it
_fmt_w = Inches(12.85) - _fmt_x
rect(slide, _fmt_x, STAT_Y, _fmt_w, STAT_H, SURFACE2, BORDER, Pt(0.5))
rect(slide, _fmt_x, STAT_Y, Inches(0.06), STAT_H, ACCENT_P)
txb(slide, "Supported formats", _fmt_x + Inches(0.2), STAT_Y + Inches(0.12),
    _fmt_w - Inches(0.3), Inches(0.25), size=9, color=TEXT_W, bold=True)
txb(slide, ".STEP   .IGES   .STP   .IGS   .STL   .JPG   .PNG   .HEIC",
    _fmt_x + Inches(0.2), STAT_Y + Inches(0.42), _fmt_w - Inches(0.3), Inches(0.25),
    size=9, color=ACCENT_B, bold=True)
txb(slide, "Drag-drop or file picker — CAD geometry measured by a full kernel, photos read by vision.",
    _fmt_x + Inches(0.2), STAT_Y + Inches(0.68), _fmt_w - Inches(0.3), Inches(0.25),
    size=8, color=TEXT_D)

notes(slide,
    "This is the automation path, and there's one design decision here I really want to land: the "
    "geometry is the ground truth. When you upload a STEP file, a proper CAD kernel measures it — the "
    "real volume, the bounding box, the B-rep faces, the holes and pockets. The AI then interprets "
    "that measured geometry; it doesn't get to overrule it. If the AI ever suggests a number that "
    "contradicts the measured volume or weight, we treat that as a bug and the sanity layer catches "
    "it. "
    "A quick note for the technically minded, because it matters for deployment: measuring STEP and "
    "IGES files needs a heavyweight geometry kernel, and we now ship that inside the deployed "
    "container — the build that carries it is verified in our continuous-integration pipeline on "
    "every change, so it isn't a 'works on my laptop' story. STL files take a lighter, pure-code fast "
    "path. "
    "Net effect: geometry to a full should-cost in minutes, ten times faster than doing the routing "
    "by hand, with the machine doing the tedious measurement and a human still owning the judgement.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 3D CAD Viewer (latest capabilities)
# ══════════════════════════════════════════════════════════════════════════════
TEAL = RGBColor(0x0E, 0x80, 0x74)   # Workflow deck's teal
slide = add_slide()
slide_header(slide, 6, "New in 2026 · Engineering-Grade Viewer",
             "3D CAD Viewer — Inspect, Measure & Analyse",
             "A SaaS-grade CAD viewer built into CostVision — STEP / IGES / STL, in the browser, no separate CAD seat.")

_vg = [
    (ACCENT_B, "Views & Navigation",
     "Iso / Front / Top / Right, Fit, orientation view-cube, smooth logarithmic-depth zoom (no z-fighting on metre-scale parts), maximize / full-window."),
    (ACCENT_G, "Display & DFM Analysis",
     "Shaded / wireframe, bounding-box dims, grid toggle, colour-by-face-type, per-component colours + legend, draft & undercut heatmap, wall-thickness heatmap."),
    (ACCENT_P, "Structure & Assembly",
     "Collapsible model tree (bodies · features · faces), section planes X/Y/Z, exploded view (Radial / X / Y / Z), per-part Move & Rotate — isolated, the rest of the assembly stays put."),
    (ORANGE, "Measure & Inspect",
     "Distance, radius / diameter, angle, point X/Y/Z, face-to-face perpendicular — exact B-rep readouts, CSV export, snapshot straight into the cost report."),
    (TEAL, "SaaS-grade UX",
     "Docked collapsible tree, drag-to-reorder tool groups, slide-away toolbar, professional line-icon set — fits cleanly beside the app sidebar at any width."),
]
_gx, _gy, _gw, _gh, _ggap = Inches(0.45), Inches(2.02), Inches(3.98), Inches(2.18), Inches(0.24)
for i, (col, t, b) in enumerate(_vg):
    r, c = divmod(i, 3)
    card(slide, _gx + c * (_gw + _ggap), _gy + r * (_gh + Inches(0.2)), _gw, _gh, t, b, accent=col)
# 6th cell — summary stat
_sx = _gx + 2 * (_gw + _ggap); _sy = _gy + 1 * (_gh + Inches(0.2))
rect(slide, _sx, _sy, _gw, _gh, SURFACE2, BORDER, Pt(0.5))
rect(slide, _sx, _sy, Inches(0.06), _gh, ACCENT_G)
txb(slide, "One viewer, many jobs", _sx + Inches(0.16), _sy + Inches(0.12), _gw - Inches(0.24), Inches(0.3), size=10, bold=True, color=ACCENT_G)
txb(slide, "24 tools · 6 groups\nSTEP / IGES / STL\nRuns in the browser\nB-rep-accurate geometry\n\nEverything measured here\nfeeds the cost engine →",
    _sx + Inches(0.16), _sy + Inches(0.5), _gw - Inches(0.24), _gh - Inches(0.6), size=9, color=TEXT_G, wrap=True)

notes(slide,
    "This is the 3D CAD viewer, and it's now genuinely engineering-grade — it lives inside CostVision, in the "
    "browser, so nobody needs a separate CAD seat just to look at a part. Left to right: full navigation with a "
    "view-cube and smooth zoom that no longer shatters on big automotive parts; a display and DFM column that "
    "colours faces by type, colours each component of an assembly, and heat-maps draft and wall thickness — those "
    "last two are pure manufacturability signals. Then structure and assembly: a proper collapsible tree, section "
    "planes, exploded views along any axis, and the ability to move or rotate a single component in isolation. "
    "Measurement is exact B-rep, not mesh approximation, and it exports. And the whole thing has SaaS-grade polish "
    "— docked tree, draggable tools, a clean icon set. But the punchline is the box on the right: everything the "
    "viewer measures is the same geometry that feeds the cost engine, which is the next slide.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Automatic data capture → Cost engine
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 7, "CAD-to-Cost · Zero Manual Entry",
             "What's Auto-Captured — Fed to the Cost Engine",
             "Geometry is measured, not typed. AI only classifies material & process; every £ is deterministic and traceable.")

_hdr = ["Auto-captured from geometry (no typing)", "→ Feeds this cost driver"]
_rows = [
    ["Volume — exact, from the CAD kernel (OCCT)", "Material mass · stock size"],
    ["Weight per material = volume × density", "Raw-material £  (Al / steel / iron / Cu / Ti)"],
    ["Bounding box / envelope", "Machine sizing · stock · setups"],
    ["Body / component count", "Assembly cost · BOM line count"],
    ["Hole & boss table — Ø, depth, count", "Secondary machining: drill / bore / tap"],
    ["Wall thickness — min / mean / heatmap", "Moulding & casting feasibility · cooling"],
    ["Draft & undercut analysis", "Tooling complexity — slides for undercuts"],
    ["Face-type areas (planar / cylindrical)", "Machining area · paint / coat area"],
    ["CNC cycle-time & setup estimate", "Process-time baseline"],
]
add_table(slide, len(_rows) + 1, 2, Inches(0.45), Inches(2.05), Inches(7.35), Inches(4.55),
          _hdr, _rows, col_widths=[Inches(4.15), Inches(3.2)])

# Right column — worked example + golden rule
_rx, _rw = Inches(8.05), Inches(4.85)
rect(slide, _rx, Inches(2.05), _rw, Inches(2.72), SURFACE2, BORDER, Pt(0.5))
rect(slide, _rx, Inches(2.05), Inches(0.06), Inches(2.72), ACCENT_B)
txb(slide, "⚙  Worked example — Spur Gear (m3 · z38)", _rx + Inches(0.16), Inches(2.14), _rw - Inches(0.26), Inches(0.3),
    size=10, bold=True, color=TEXT_W)
txb(slide,
    "20MnCr5 case-hardening steel — captured automatically:\n"
    "   Volume 264.9 cm³   ·   Weight 2.08 kg\n"
    "   Envelope 120 × 120 × 30 mm   ·   1 body\n\n"
    "→ Should-cost, India ex-works (GBP):\n"
    "   Material £3.36 · Process £11.40 · Labour £1.36\n"
    "   Overhead £1.95 · Margin £1.83  + heat-treat £1.77\n"
    "   =  £21.90 / part   (≈ ₹2,398)   ·   £10.53 / kg",
    _rx + Inches(0.16), Inches(2.5), _rw - Inches(0.26), Inches(2.2), size=9, color=TEXT_G, wrap=True)

rect(slide, _rx, Inches(4.95), _rw, Inches(1.65), SURFACE2, ACCENT_G, Pt(0.9))
rect(slide, _rx, Inches(4.95), Inches(0.06), Inches(1.65), ACCENT_G)
txb(slide, "🔒  The golden rule", _rx + Inches(0.16), Inches(5.04), _rw - Inches(0.26), Inches(0.3),
    size=10, bold=True, color=ACCENT_G)
txb(slide,
    "AI never sets a price. It only reads and classifies — material family, process route, feature intent. "
    "Every number above is deterministic arithmetic in the cost engine, bounded by the measured geometry and "
    "traceable to the rate library. That is what makes the output defensible.",
    _rx + Inches(0.16), Inches(5.4), _rw - Inches(0.26), Inches(1.15), size=9, color=TEXT_G, wrap=True)

notes(slide,
    "This is the slide to dwell on, because it's the heart of the automation. On the left is exactly what the "
    "tool pulls off the CAD model with no human typing anything in — and next to each item, the cost driver it "
    "feeds. Volume and weight come straight from the CAD kernel, so the material bucket is exact, not a guess — "
    "and weight is just volume times the density of whichever material you pick. The bounding box sizes the "
    "machine and the stock. Body count gives you assembly and BOM scope. The hole-and-boss table becomes drilling "
    "and boring operations. Wall thickness tells you if a part can even be moulded or cast. Draft analysis flags "
    "undercuts that need slides in the tool. Face areas drive machining and paint area. And there's a first-cut "
    "cycle-time estimate. On the right is a real worked example — the gear we just costed: two kilos of gear "
    "steel, captured automatically, costing about twenty-two pounds ex-works from India, all in sterling. "
    "And the box underneath is the guardrail that makes all of this trustworthy in front of a supplier: the AI "
    "never sets the price. It classifies; the engine does deterministic arithmetic on the measured geometry, and "
    "every rate traces back to the library. That's the difference between a defensible should-cost and a guess.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — End-to-End Workflow
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 8, "Workflow", "AI Agent — 8-Step End-to-End Flow",
             "Three entry points. One intelligent engine. Full cost intelligence in under 5 minutes.")

# Entry point banner
rect(slide, Inches(0.45), Inches(1.98), Inches(12.45), Inches(0.5), SURFACE2, BORDER, Pt(0.5))
txb(slide, "Entry Point A: 🤖 AI Agent (plain English description)   |   "
    "Entry Point B: 📐 CAD / Photo Upload   |   Entry Point C: ✏️ Manual Form Entry",
    Inches(0.6), Inches(2.04), Inches(12.1), Inches(0.36),
    size=9, color=TEXT_G, align=PP_ALIGN.CENTER)

# 8 flow boxes
flow_items = [
    ("1", "🏭", "Select\nCommodity",      "Choose from 21 manufacturing commodities or let AI Agent auto-detect"),
    ("2", "🤖", "AI Agent /\nCAD / Manual", "Describe part, upload file, or enter parameters directly"),
    ("3", "🔍", "AI Feature\nExtraction",  "AI reads geometry, identifies features, infers material & complexity"),
    ("4", "🛤️", "AI Routing &\nModel Select", "AI sequences operations, selects machines, labour grades & cycle times"),
    ("5", "⚙️", "Should-Cost\nCalculation",  "8-bucket model: Material+Process+Labour+Tooling+OH+Margin"),
    ("6", "💡", "AI Insights\nGenerated",  "Benchmark comparisons, regional cost index, supplier margin analysis"),
    ("7", "🔬", "DFM / DFA\nAnalysis",    "Manufacturability + assembly scores, critical issues, saving potential"),
    ("8", "📤", "Export &\nShare",         "6-sheet Excel, PDF report, scenario save, cloud team sync"),
]

n = len(flow_items)
fw = (W - Inches(0.9)) / n
fh = Inches(3.8)
fx = Inches(0.45)
fy = Inches(2.6)
colors = [ACCENT_B, ACCENT_B, ACCENT_P, ACCENT_P, ORANGE, ACCENT_G, ACCENT_G, ACCENT_B]

for i, (num, ico, title, body) in enumerate(flow_items):
    flow_box(slide, fx + i*fw, fy, fw - Inches(0.04), fh,
             num, ico, title, body, colors[i])
    # Arrow between boxes
    if i < n - 1:
        txb(slide, "→", fx + (i+1)*fw - Inches(0.2), fy + Inches(1.7),
            Inches(0.25), Inches(0.3), size=9, bold=True, color=ACCENT_B)

# Tags row
tags = ["⚡ Minutes, not weeks", "✓ Auditable assumptions", "🤖 Claude AI powered",
        "🌍 20 regions", "📊 A/B/C scenarios", "📐 Assembly BOM", "📉 Learning curve"]
tx = Inches(0.45)
for tag in tags:
    rect(slide, tx, Inches(6.55), Inches(1.72), Inches(0.3), SURFACE2, ACCENT_B, Pt(0.5))
    txb(slide, tag, tx + Inches(0.08), Inches(6.58), Inches(1.58), Inches(0.24),
        size=7.5, color=ACCENT_B)
    tx += Inches(1.78)

notes(slide,
    "This slide is really just to show that all three entry points — plain English, CAD upload, or "
    "manual entry — funnel into the same eight-step flow. I won't read the boxes out; the shape is "
    "what matters. "
    "The important beat is in the middle: steps three and four, where the AI extracts the features "
    "and sequences the operations. That's the part that used to take an experienced cost engineer "
    "days, and it's where most of the time saving comes from. Everything after it — the eight-bucket "
    "calculation, the insights, the DFM check, the export — is deterministic and fast. "
    "So the headline is the first tag at the bottom: minutes, not weeks, but with auditable "
    "assumptions the whole way through. Speed without a black box.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 21 Commodities
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 9, "Full Coverage", "21 Manufacturing Commodities — One Platform",
             "Every major manufacturing process — all modelled to engineering depth with full routing, yield, and tooling logic.")

# 21 commodity cards
commodities = [
    ("🛠️", "CNC Machining",        "3/4/5-axis, turning, grinding"),
    ("🪨",  "Casting",              "HPDC, gravity, investment, sand"),
    ("🏗️", "Cast + Machine",       "Casting + post-ops + HT"),
    ("🔨",  "Forging",              "Hot, warm, cold die forging"),
    ("🪗",  "Sheet Metal",          "Progressive/transfer die"),
    ("✂️",  "Sheet Metal Fab",      "Laser · punch · press brake · MIG"),
    ("⚙️",  "Gear Cutting",         "Hob · shape · skive · broach · grind"),
    ("🏺",  "Injection Moulding",   "Thermoplastics, multi-cavity"),
    ("🫧",  "Blow Moulding",        "HDPE, PET, extrusion blow"),
    ("📏",  "Extrusion",            "Profile, pipe, sheet, rod"),
    ("♨️",  "Thermoforming",        "Vacuum, pressure, twin-sheet"),
    ("🔄",  "Rotational Moulding",  "Large hollow shapes, LLDPE"),
    ("🧱",  "Rubber Moulding",      "Compression, injection, LSR"),
    ("🪢",  "Composites",           "CFRP/GFRP, autoclave, RTM, AFP"),
    ("🎨",  "Painting / Coating",   "E-coat, powder, wet, anodise"),
    ("🚗",  "BIW Assembly",         "Spot weld, MIG, hem, bond"),
    ("🖥️", "PCB Fabrication",      "2–16L, HDI, rigid-flex, RF/μwave"),
    ("🔌",  "PCBA / SMD",           "SMT, TH, reflow, conformal coat"),
    ("🔗",  "Wiring Harness",       "Cut-strip-crimp, test, sub-asm"),
    ("🧩",  "Assembly BOM Rollup",  "Multi-part BOM, full cost rollup"),
    ("🤖",  "AI Agent",             "Natural language → any commodity"),
]

# Two cards sharing a glyph reads as a mistake from the back of the room — ⚙️,
# 🔩 and 🔗 were each doing duty on two commodities until the render showed it.
assert len({icon for icon, _, _ in commodities}) == len(commodities), \
    'two commodity cards share an icon'

# 6 wide x 4 rows. Eight columns fitted the page but squeezed the subtitle to
# 6 pt, which no projector makes legible; the card must be wide enough to carry
# 7.5 pt type, so the column count follows the type, not the other way round.
cols_c = 7                      # 21 cards = exactly 3 full rows, no orphans
cw_c = (W - Inches(0.9)) / cols_c - Inches(0.055)
_rows_c = -(-len(commodities) // cols_c)
ch_c = (Inches(6.95) - Inches(1.98) - Inches(0.065) * (_rows_c - 1)) / _rows_c
sx_c = Inches(0.45)
sy_c = Inches(1.98)

for i, (ico, name, sub) in enumerate(commodities):
    r, c = divmod(i, cols_c)
    cx2 = sx_c + c * (cw_c + Inches(0.055))
    cy2 = sy_c + r * (ch_c + Inches(0.065))
    col = ACCENT_B if i < 19 else ACCENT_G   # last three are cross-cutting, not a process
    comm_card(slide, cx2, cy2, cw_c, ch_c, ico, name, sub, color=col)

notes(slide,
    "Breadth is the point of this slide. Twenty-one commodities, from CNC machining and casting and "
    "forging, through the whole plastics family — injection, blow, extrusion, thermoforming, "
    "rotomoulding — into rubber and composites, then electronics: PCB fabrication, assembly, wiring "
    "harness, and body-in-white. "
    "Why does that matter? Because a real product isn't one process. A door module is stampings plus "
    "a harness plus a PCBA plus paint. If your costing tool only does machining, you're stuck the "
    "moment the part is plastic. CostVision costs all of them with the same rigour and the same "
    "eight-bucket structure, which is exactly what lets the assembly roll-up add them together into "
    "one number. "
    "And the last tile is the AI agent — the natural-language front door into any of the other "
    "twenty.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Should-Cost Model Depth
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 10, "Engineering Depth", "What Sits Behind Every Number",
             "Twelve layers, all of them arithmetic you can follow. Nothing is estimated by feel.")

depth_cards = [
    (ACCENT_B, "📦 Materials",
     "30+ grades — steels, aluminium, plastics, rubber, composites, board laminates. Priced in the local currency."),
    (ACCENT_G, "🏭 Machines & Rates",
     "50+ machine types with an hourly rate you can edit per project and per region."),
    (ORANGE,   "🛤️ Process Routing",
     "The order of operations, and for each one: which machine, how long, who runs it."),
    (ACCENT_P, "🔩 Tooling & NRE",
     "Dies, moulds and fixtures spread over the year's volume, with tool life accounted for."),
    (ACCENT_G, "📉 Yield & Scrap",
     "What actually comes out good first time, driven by the features that make a part hard to build."),
    (ACCENT_B, "🌍 Regional Rates",
     "20 regions, each with its own labour, energy and machine rates."),
    (ORANGE,   "📊 Benchmarks & Sensitivity",
     "Checked against published cost ranges, and ranked by which inputs move the number most."),
    (ACCENT_P, "📉 Learning Curve",
     "What the part costs as volume builds — the figure you need for a long-term agreement."),
    (ACCENT_G, "💰 Quote Comparison",
     "Put the supplier's quote next to ours and see the gap, in their currency or yours."),
    (ACCENT_B, "🤖 Stated Assumptions",
     "Where something isn't known, the AI fills it in and says so — every assumption is on the page."),
    (ORANGE,   "⚙️ Gear Cutting",
     "Hobbing, shaping, skiving and grinding times worked out from the gear itself, not a rule of thumb."),
    (ACCENT_P, "🚢 Landed Cost",
     "Duty, carbon levy and shipping terms — the cost at our door, not at the supplier's gate."),
]

# Card width is DERIVED from the slide, never typed. A hardcoded 3.88" put
# column 4 at 12.39"-16.27" on a 13.33" slide, so cards 3 and 7 — "Tooling &
# NRE Amortisation" and "Learning Curve" — rendered entirely off the screen.
# The speaker note tells the presenter to point at one of them.
COLS_D = 4
MARGIN_D = Inches(0.45)
gap_d = Inches(0.1)
cw_d = (W - MARGIN_D * 2 - gap_d * (COLS_D - 1)) / COLS_D
# 12 cards over 4 columns is exactly 3 full rows — no orphan row — and the card
# height is derived so the block finishes just above the footer instead of
# leaving 1.5" of dead slide beneath it.
sx_d = MARGIN_D
sy_d = Inches(1.98)
BOTTOM_D = Inches(6.95)
rows_d = -(-len(depth_cards) // COLS_D)
ch_d = (BOTTOM_D - sy_d - gap_d * (rows_d - 1)) / rows_d
assert sx_d + (COLS_D - 1) * (cw_d + gap_d) + cw_d <= W, 'depth grid overflows the slide'
assert sy_d + (rows_d - 1) * (ch_d + gap_d) + ch_d <= BOTTOM_D + Inches(0.01)
# The card body sits at a fixed offset below the title, so a title that wraps to
# a second line lands ON the body text. "Learning Curve (Wright's Law)" (31 ch)
# did exactly that; "Tooling & NRE Amortisation" (28 ch) is the longest measured
# to stay on one line at 10.5 pt in a 3.03" card. Budget set from the render.
TITLE_CH_MAX = 29
for _, _t, _ in depth_cards:
    assert len(_t) <= TITLE_CH_MAX, f'card title wraps onto the body: {_t!r}'

for i, (col, title, body) in enumerate(depth_cards):
    r, c = divmod(i, COLS_D)
    cx_d = sx_d + c * (cw_d + gap_d)
    cy_d = sy_d + r * (ch_d + gap_d)
    card(slide, cx_d, cy_d, cw_d, ch_d, title, body, accent=col,
         title_pt=10.5, body_pt=9.5)

notes(slide,
    "If anyone in the room is going to be sceptical, it's usually about depth — is this a real "
    "engineering model, or a spreadsheet with a nice front end? This slide is the answer. "
    "Underneath every commodity there's a material database with the real grades and FX-adjusted "
    "prices, a machine and rate library you can edit per project and region, and proper routing "
    "logic — each operation has a machine, a cycle time, an OEE, a setup, a labour grade. Yield and "
    "scrap are modelled, not guessed. Tooling and NRE are amortised over real volume with tool-life "
    "curves. "
    "The two I'd draw a manager's eye to are benchmarking and the learning curve. Every cost element "
    "is checked against published benchmark ranges, so an outlier gets flagged rather than quietly "
    "shipped. And Wright's Law lets you project cost as volume grows — which is exactly the "
    "conversation you need for a long-term agreement. This is the layer that makes the output "
    "defensible rather than just plausible.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Advanced Features
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 11, "Advanced Capabilities", "Beyond a Cost Number",
             "Five tools that turn the estimate into something you can negotiate with.")

# Left col
lx9 = Inches(0.45)
lw9 = Inches(6.0)

adv_left = [
    (ACCENT_P, "📉 Learning Curve",
     "Cost falls as volume builds. Set the curve, read the cost at any year.\n\n"
     "• Quote year-three cost, not prototype cost\n"
     "• Built for long-term agreements and make-vs-buy"),
    (ACCENT_G, "🧩 Assembly Roll-Up",
     "Cost a whole module, not just one part.\n\n"
     "• Mix any of the 21 commodities in one build\n"
     "• Every component keeps its own cost model\n"
     "• Body + paint + harness + PCBA = one total"),
]

yy9 = Inches(2.0)
for col, title, body in adv_left:
    rect(slide, lx9, yy9, lw9, Inches(2.3), SURFACE2, BORDER, Pt(0.5))
    rect(slide, lx9, yy9, Inches(0.06), Inches(2.3), col)
    txb(slide, title, lx9 + Inches(0.14), yy9 + Inches(0.1), lw9 - Inches(0.2), Inches(0.3),
        size=10, bold=True, color=TEXT_W)
    txb(slide, body, lx9 + Inches(0.14), yy9 + Inches(0.42), lw9 - Inches(0.2), Inches(1.8),
        size=8.5, color=TEXT_G, wrap=True)
    yy9 += Inches(2.42)

# Right col
rx9 = Inches(6.85)
rw9 = Inches(6.0)

adv_right = [
    (ORANGE, "💰 Quote Comparison",
     "Log the supplier's quote. Any currency, converted for you.\n"
     "• Tells you the gap in words: \"18% above should-cost\"\n"
     "• Gives you the talking points to open with"),
    (ACCENT_B, "📊 What Moves the Number",
     "Every driver flexed ±10%, ranked by impact.\n"
     "• Shows the few inputs worth checking before you negotiate\n"
     "• Chart exports straight into your supplier pack"),
    (ACCENT_G, "🧪 Scenarios Side by Side",
     "Baseline, target and stretch on one screen.\n"
     "• See exactly where the money moved\n"
     "• Save and share with the team"),
]

yy9r = Inches(2.0)
for col, title, body in adv_right:
    h = Inches(1.52)
    rect(slide, rx9, yy9r, rw9, h, SURFACE2, BORDER, Pt(0.5))
    rect(slide, rx9, yy9r, Inches(0.06), h, col)
    txb(slide, title, rx9 + Inches(0.14), yy9r + Inches(0.08), rw9 - Inches(0.2), Inches(0.28),
        size=9.5, bold=True, color=TEXT_W)
    txb(slide, body, rx9 + Inches(0.14), yy9r + Inches(0.34), rw9 - Inches(0.2), h - Inches(0.45),
        size=8, color=TEXT_G, wrap=True)
    yy9r += Inches(1.62)

notes(slide,
    "These are the features that turn a cost number into commercial leverage. "
    "Top left, the learning curve — as volume doubles, cost comes down by a set percentage, so you "
    "can quote the cost at year-three volume, not just the prototype cost. Below it, the assembly "
    "roll-up I keep coming back to: build a multi-level BOM from any mix of commodities and get a "
    "total with the per-component breakdown intact. "
    "On the right is the negotiation toolkit. Log a supplier's quote and the tool tells you, in "
    "plain terms, 'this supplier is eighteen percent above should-cost,' and hands you the talking "
    "points to say so. The tornado chart shows which inputs actually move the number, so you know "
    "what to validate before you sit down. And scenarios let you hold baseline, target, and stretch "
    "side by side. This is where the tool stops being an estimator and starts being a negotiator's "
    "brief.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DFM / DFA Intelligence
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 12, "AI Intelligence Layer", "DFM / DFA — How to Make the Part Cheaper",
             "Runs after the cost is settled. It reads the result and advises — it never changes the price.")

# Score banner
rect(slide, Inches(0.45), Inches(2.0), Inches(12.45), Inches(0.78), SURFACE2, BORDER, Pt(0.5))
txb(slide, "Manufacturability Score (DFM)", Inches(0.65), Inches(2.06),
    Inches(3.5), Inches(0.28), size=8.5, bold=True, color=TEXT_W)
txb(slide, "8.5 / 10", Inches(0.65), Inches(2.34), Inches(2), Inches(0.36),
    size=20, bold=True, color=ACCENT_G)
txb(slide, "Assembly Efficiency (DFA)", Inches(4.5), Inches(2.06),
    Inches(3.5), Inches(0.28), size=8.5, bold=True, color=TEXT_W)
txb(slide, "7.0 / 10", Inches(4.5), Inches(2.34), Inches(2), Inches(0.36),
    size=20, bold=True, color=ORANGE)
txb(slide, "Total Potential Saving", Inches(8.4), Inches(2.06),
    Inches(2.5), Inches(0.28), size=8.5, bold=True, color=TEXT_W)
txb(slide, "~19.3%", Inches(8.4), Inches(2.34), Inches(2.5), Inches(0.36),
    size=20, bold=True, color=ACCENT_G)
txb(slide, "Scored out of 10  ·  saving is the top three fixes combined",
    Inches(10.8), Inches(2.2), Inches(2.0), Inches(0.54),
    size=7.5, color=TEXT_D, wrap=True)

# Four quadrant cards
dfm_dfa = [
    (ACCENT_B, "🔧 What It Finds — Making the Part",
     "Too much material going to scrap.\n"
     "Machines running well below what they should.\n"
     "Tooling, overhead or supplier margin out of line.\n"
     "Each issue graded critical, major or minor."),
    (ACCENT_G, "🧩 What It Finds — Building It",
     "Assemblies done by hand that a robot could do.\n"
     "Too many setups — every one adds cost and variation.\n"
     "Too many different fasteners, so tools keep changing."),
    (ORANGE,   "💡 What to Do About It",
     "Six to eight actions per part, with the saving on each:\n\n"
     "Automate the manual operations   5–20%\n"
     "Move to a near-net shape   5–15%\n"
     "Raise machine uptime   5–12%\n"
     "Plus quick wins: more volume, open-book overhead, re-bid."),
    (ACCENT_P, "🤖 Ask the AI to Go Deeper",
     "One click, after the rule-based findings are on screen.\n\n"
     "Why the cost sits where it does.\n"
     "The three things to fix first.\n"
     "What to say to the supplier, and a different process to consider."),
]

# Height derived from a bottom target: a typed 2.0" put the second row of cards
# at 7.10", straight through the footer rule.
cw_d2 = Inches(6.0)
_top_d2, _gap_d2 = Inches(3.0), Inches(0.12)
ch_d2 = (Inches(6.95) - _top_d2 - _gap_d2) / 2
for i, (col, title, body) in enumerate(dfm_dfa):
    r, c = divmod(i, 2)
    dx = Inches(0.45) + c * (cw_d2 + _gap_d2)
    dy = _top_d2 + r * (ch_d2 + _gap_d2)
    card(slide, dx, dy, cw_d2, ch_d2, title, body, accent=col,
         title_pt=10.5, body_pt=9.5)
assert _top_d2 + 2 * ch_d2 + _gap_d2 <= Inches(6.96), 'DFM grid runs into the footer'

notes(slide,
    "One thing I was careful about here: the DFM and DFA layer runs after the should-cost, and it "
    "never touches the calculation. It reads the finished cost structure and asks manufacturing "
    "questions of it. That separation is deliberate — I didn't want a 'design advice' feature "
    "quietly nudging the price. "
    "What it gives you is a manufacturability score and an assembly-efficiency score out of ten, and "
    "then something more useful than a grade: a ranked list of levers with rough savings against "
    "each. Automate a high-labour operation, five to twenty percent. Move to a near-net shape, five "
    "to fifteen percent on material. Increase volume to dilute the tooling NRE — quick win, low "
    "risk. "
    "And if you want more than the rule-based read, there's a deep-analysis button that has the AI "
    "write the root-cause commentary and a negotiation strategy. So the cost number comes with a "
    "'here's how to make it cheaper' attached — at the concept stage, where it's still free to act "
    "on. "
    "One important update since the last time I showed you this. Your engineering team told me, "
    "quite bluntly, that these findings were too generic — that reading the cost and guessing "
    "backwards isn't good enough. They were right. So there's now a second layer underneath this "
    "one that opens the 3D model itself and measures every face on the part. Nineteen rules, and "
    "every one of them cites a published standard rather than my opinion. It doesn't say 'tooling "
    "looks high' — it says which specific faces on the part cannot come out of the die, and you can "
    "click one and see it highlighted on the model. That is the difference between a hint and an "
    "argument you can take to a supplier.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Regional & Currency Intelligence
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 13, "Global Coverage", "20 Manufacturing Regions — 10 Currencies",
             "Select a region and every rate adjusts automatically — labour, machines, energy, material and currency.")

# Regions table
regions = [
    ("🇬🇧 UK",         "GBP £", "Skilled £26/hr",   "Industry / Aerospace baseline"),
    ("🇩🇪 Germany",    "EUR €", "Skilled €32/hr",   "Premium, high automation"),
    ("🇵🇱 Poland",     "EUR €", "Skilled €11/hr",   "Low-cost EU option"),
    ("🇷🇴 Romania",    "EUR €", "Skilled €7/hr",    "Ultra-low EU labour"),
    ("🇨🇳 China",      "CNY ¥", "Skilled ¥28/hr",   "Global volume manufacturing"),
    ("🇮🇳 India",      "INR ₹", "Skilled ₹320/hr",  "Engineering & electronics hub"),
    ("🇲🇽 Mexico",     "MXN $", "Skilled MXN82/hr", "Near-shore for North America"),
    ("🇺🇸 USA",        "USD $", "Skilled $32/hr",   "Premium North American rates"),
    ("🇹🇭 Thailand",   "THB ฿", "Skilled ฿180/hr",  "Electronics & automotive"),
    ("🇻🇳 Vietnam",    "VND ₫", "Skilled ₫52k/hr",  "Lowest-cost electronics hub"),
    ("🇧🇷 Brazil",     "BRL R$","Skilled R$22/hr",  "South America manufacturing"),
    ("🇰🇷 S. Korea",   "KRW ₩", "Skilled ₩22k/hr",  "Electronics & automotive"),
]

rect(slide, Inches(0.45), Inches(1.98), Inches(7.5), Inches(4.98), SURFACE, BORDER, Pt(0.5))
txb(slide, "Key Manufacturing Regions (sample)", Inches(0.6), Inches(2.04),
    Inches(5), Inches(0.28), size=9, bold=True, color=TEXT_W)

headers = ["Region", "Currency", "Sample Labour Rate", "Position"]
col_ws = [Inches(1.5), Inches(0.9), Inches(1.9), Inches(2.7)]
add_table(slide, len(regions)+1, 4,
          Inches(0.48), Inches(2.36), Inches(7.44), Inches(4.5),
          headers, regions, col_widths=col_ws)

# Right side — key points
rx11 = Inches(8.2)
ry11 = Inches(1.98)
rw11 = Inches(4.7)

key_points = [
    (ACCENT_B, "Auto-Currency Switch",
     "When you select a manufacturing region (e.g. China), the display currency automatically switches to CNY. No manual FX conversion needed."),
    (ACCENT_G, "Full Rate Adjustment",
     "Every rate adjusts: skilled/semi-skilled/engineer/inspector labour, electricity, gas, machine rates, material prices — all recalculated for the selected region."),
    (ACCENT_P, "Editable Rate Library",
     "All rates are editable. Override any rate with project-specific data, supplier quotes, or latest market prices. Resets to calibrated baseline at any time."),
    (ORANGE,   "Regional Comparison",
     "AI Insights tab shows estimated cost in all 20 regions simultaneously — instant make-vs-buy regional arbitrage analysis for any part."),
]

yy11 = ry11
for col, title, body in key_points:
    rect(slide, rx11, yy11, rw11, Inches(1.1), SURFACE2, BORDER, Pt(0.5))
    rect(slide, rx11, yy11, Inches(0.06), Inches(1.1), col)
    txb(slide, title, rx11 + Inches(0.14), yy11 + Inches(0.08), rw11 - Inches(0.2), Inches(0.28),
        size=9, bold=True, color=TEXT_W)
    txb(slide, body, rx11 + Inches(0.14), yy11 + Inches(0.34), rw11 - Inches(0.2), Inches(0.72),
        size=8, color=TEXT_G, wrap=True)
    yy11 += Inches(1.18)

notes(slide,
    "A part made in Poland doesn't just have cheaper labour — it has a different machine rate, "
    "different energy cost, a different currency. Getting that right by hand is where manual "
    "should-costs quietly go wrong. "
    "Here you pick the region and everything reprices at once: the four labour tiers, electricity "
    "and gas, machine rates, material prices, and the display currency flips automatically — pick "
    "China and you're reading yuan. And because the AI Insights tab estimates the same part across "
    "all twenty regions at once, you get instant regional arbitrage: 'this is X in the UK, Y in "
    "Mexico, Z in Vietnam' — a make-versus-buy answer backed by rates, not a gut feel. "
    "Every rate is editable too, so if you have a real supplier quote or a fresher market price, you "
    "drop it in and the model uses it.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Agentic vs Autonomous Agentic AI (concept, with examples)
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 14, "The Concept · Explained",
             "Agentic vs Autonomous Agentic AI — With Examples",
             "Two levels of \"agentic\": on the left the AI acts because you asked; on the right it acts on its own, unattended.")

def _agentic_panel(x, accent, head, tag, examples):
    w, y, h = Inches(6.05), Inches(2.0), Inches(4.4)
    rect(slide, x, y, w, h, SURFACE2, BORDER, Pt(0.5))
    rect(slide, x, y, Inches(0.06), h, accent)
    txb(slide, head, x + Inches(0.2), y + Inches(0.12), w - Inches(0.35), Inches(0.34),
        size=14, bold=True, color=accent)
    txb(slide, tag, x + Inches(0.2), y + Inches(0.5), w - Inches(0.35), Inches(0.3),
        size=8.5, color=TEXT_G, italic=True)
    ry, rh = y + Inches(0.9), Inches(0.74)
    for lab, sub in examples:
        rect(slide, x + Inches(0.18), ry, w - Inches(0.36), rh, SURFACE, BORDER, Pt(0.3))
        rect(slide, x + Inches(0.18), ry, Inches(0.05), rh, accent)
        txb(slide, lab, x + Inches(0.34), ry + Inches(0.07), w - Inches(0.52), Inches(0.24),
            size=10, bold=True, color=TEXT_W)
        txb(slide, sub, x + Inches(0.34), ry + Inches(0.31), w - Inches(0.52), rh - Inches(0.36),
            size=8, color=TEXT_G, wrap=True)
        ry += rh + Inches(0.1)

_agentic_panel(Inches(0.45), ACCENT_B, "Agentic AI",
    "Assisted — you trigger it and stay in the loop; a human approves",
    [("AI Agent", "Describe a part in plain English — it builds the full cost model for you."),
     ("CAD / photo → cost", "Upload a STEP file — it measures geometry and infers material & process."),
     ("Negotiation coach", "Open a part — it drafts the buyer's counter-argument and target price."),
     ("Rate-data assistant", "Ask a costing question — it answers from your own rates, with citations.")])
_agentic_panel(Inches(6.83), ACCENT_P, "Autonomous Agentic AI",
    "Self-directed — runs unattended on a schedule, within limits you set",
    [("Savings monitor", "Runs on the server, compares paid vs should-cost, opens findings itself — £0.5M/yr unattended."),
     ("Self-audit", "Re-checks every estimate for known errors and corrects within bounds — unasked."),
     ("Calibration & drift", "Learns from logged quotes, re-derives factors, watches for drift continuously."),
     ("Outcome-weighted ranking", "Learns which findings actually convert and re-prioritises its own queue.")])

rect(slide, Inches(0.45), Inches(6.5), Inches(12.45), Inches(0.5), SURFACE2, ACCENT_G, Pt(0.6))
txb(slide, "The common thread:  you set the boundaries and every action stays glass-box and auditable — "
    "autonomy never means the AI sets a price in secret.",
    Inches(0.72), Inches(6.6), Inches(12.0), Inches(0.32), size=9.5, color=TEXT_G)

notes(slide,
    "This is the slide to slow down on if the room is new to the word 'agentic', because there are really two levels "
    "and people run them together. On the left is agentic AI in the everyday sense: the AI takes actions and uses "
    "tools, but you started the task and you stay in the loop. You describe a part and it builds the cost model; you "
    "upload a CAD file and it reads the geometry and picks the process; you open a part and it drafts the negotiation "
    "argument; you ask a question and it answers from our own rate data, with citations. Genuinely useful — but it's "
    "waiting for you to ask. On the right is the step that surprises people: autonomous agentic AI, where nobody is at "
    "the keyboard. A monitor runs on our server on a schedule, compares what we pay against what things should cost, "
    "and opens findings on its own — that's the half a million pounds a year it surfaced unattended in the demo. The "
    "self-audit re-checks every estimate for known mistakes and corrects them without being asked; the calibration "
    "keeps learning from logged quotes and watching for drift; the agent even re-prioritises its own queue toward the "
    "findings that actually convert to cash. The line at the bottom is the one I'd underline: the difference between "
    "the two columns is only who starts the action — in both, the human sets the boundaries and every action stays "
    "glass-box and auditable. Autonomy here never means the AI quietly sets a price; it means it does the watching for "
    "us, and shows its working the moment it finds something worth acting on.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — 2026 Agentic Intelligence: self-audit, learn-from-actuals, validation
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 15, "New in 2026 · Agentic Intelligence",
             "The Tool Now Checks Its Own Work — and Learns",
             "It re-checks every estimate before you see it, and it learns from the quotes you have already had.")

# Left column — three new-capability cards
lx12, lw12 = Inches(0.45), Inches(5.55)
agentic_cards = [
    (ACCENT_G, "🔍 It checks its own homework",
     "Every estimate is re-checked against the mistakes this tool has actually made before — wrong "
     "process, impossible machining time, a weight that disagrees with the CAD. It corrects them and "
     "shows its working. The measured geometry always wins over the AI."),
    (ACCENT_B, "📈 It learns from your quotes",
     "Load the supplier quotes you already have. The engine corrects itself against them, and the "
     "estimate gets tighter as your history grows. Every estimate now says how well calibrated it is."),
    (ORANGE, "🛠️ It picks the right machine",
     "One rule now sizes the machine for the part on every commodity — presses, die-casting, "
     "moulding, extrusion — instead of a fixed guess per commodity."),
]
yy12 = Inches(2.0)
for col, title, body in agentic_cards:
    h = Inches(1.55)
    rect(slide, lx12, yy12, lw12, h, SURFACE2, BORDER, Pt(0.5))
    rect(slide, lx12, yy12, Inches(0.06), h, col)
    txb(slide, title, lx12 + Inches(0.14), yy12 + Inches(0.08), lw12 - Inches(0.2), Inches(0.28),
        size=9.5, bold=True, color=TEXT_W)
    txb(slide, body, lx12 + Inches(0.14), yy12 + Inches(0.36), lw12 - Inches(0.2), h - Inches(0.46),
        size=8, color=TEXT_G, wrap=True)
    yy12 += Inches(1.62)

# Right column — validation on 5 real automotive CAD parts
rx12 = Inches(6.35)
rw12 = Inches(6.55)
rect(slide, rx12, Inches(2.0), rw12, Inches(0.04), ACCENT_G)
txb(slide, "Proven on five real automotive CAD parts", rx12, Inches(2.1), rw12, Inches(0.3),
    size=10, bold=True, color=ACCENT_G)
txb(slide, "Each was mis-costed on the first pass; the self-audit caught it and corrected it.",
    rx12, Inches(2.42), rw12, Inches(0.28), size=8, color=TEXT_G)

headers12a = ["Real Part", "AI First Pass", "What Was Wrong", "Corrected"]
data12a = [
    ("Fuel tank",         "£216.97", "Treated as cast",       "£24.62 blow-mould"),
    ("Front bumper",      "£29.18",  "Machined, not moulded", "£7.79 injection"),
    ("Servo horn",        "£333.00", "Bulk, not net-shape",   "£6.40 net-shape"),
    ("Stub axle",         "no result", "Classed as the wrong process", "Forged steel 8.1 kg"),
    ("Seat cross-member", "no result", "Machine sizing gap",    "£1.11 runs clean"),
]
col_ws12a = [Inches(1.55), Inches(1.15), Inches(1.95), Inches(1.9)]
add_table(slide, len(data12a) + 1, 4,
          rx12, Inches(2.78), rw12, Inches(2.55),
          headers12a, data12a, col_widths=col_ws12a)

# Credibility strip under the table
rect(slide, rx12, Inches(5.5), rw12, Inches(1.42), SURFACE2, BORDER, Pt(0.5))
rect(slide, rx12, Inches(5.5), Inches(0.06), Inches(1.42), ACCENT_B)
txb(slide, "Built to be trusted, not just believed", rx12 + Inches(0.16), Inches(5.58),
    rw12 - Inches(0.25), Inches(0.28), size=9.5, bold=True, color=TEXT_W)
cred_items = [
    "1,777 automated tests — if the logic breaks, the build fails, not the demo",
    "Every estimate is re-checked before you ever see the number",
    "The CAD engine is tested in the same container we deploy",
]
yyc = Inches(5.9)
for it in cred_items:
    rect(slide, rx12 + Inches(0.2), yyc + Inches(0.11), Inches(0.06), Inches(0.06), ACCENT_B)
    txb(slide, it, rx12 + Inches(0.38), yyc, rw12 - Inches(0.55), Inches(0.3),
        size=8, color=TEXT_G, wrap=True)
    yyc += Inches(0.32)

notes(slide,
    "This is the slide I most want to talk you through, because it's what's new this year and it's "
    "the honest answer to 'can I trust the number?' "
    "Earlier versions of this tool were confident and sometimes wrong. So I built a self-audit — a "
    "deterministic layer that re-checks every estimate against a list of mistakes I've actually seen "
    "it make. Did it treat a hollow part as a solid casting? Is the machining time longer than it "
    "would physically take to finish the part? Does the weight disagree with the CAD? When it finds "
    "one of those, it applies a bounded correction and shows its working. Crucially, the AI is never "
    "allowed to overrule the measured geometry — the geometry wins. "
    "Alongside that, the tool now learns. You can bulk-import your real supplier quotes, and it "
    "builds per-segment correction factors and a confidence band, so the estimate gets tighter as "
    "your own history grows. And the machine-sizing that used to be hard-coded for a couple of "
    "processes is now universal — it picks the right press or tonnage for the part across "
    "commodities. "
    "The right-hand side is the proof, and I'm deliberately showing you the failures. Five real "
    "automotive CAD parts, every one mis-costed on the first pass. A fuel tank that came out at two "
    "hundred and seventeen pounds because it was treated as a casting — the audit spotted it's a "
    "thin hollow part, re-routed it to blow moulding, and it dropped to twenty-five. A servo horn "
    "that priced from a solid billet at three hundred and thirty-three pounds, corrected to six "
    "forty on a net shape. A stub axle that was mis-classified until the audit forced it to forged "
    "steel. And a seat cross-member that used to crash the engine outright and now runs clean at a "
    "pound eleven. "
    "Underneath all of it: just over a thousand automated tests, the self-audit on every estimate, "
    "and the CAD engine shipping in a container we build-test in CI. That's what lets me stand here "
    "and say the number is defensible — it's checked, it's proven, and it's honest about what it got "
    "wrong.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Demo: Real-World Cost Benchmarks
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 16, "Live Demo Data", "Real-World Cost Benchmarks — Luxury SUV Programme",
             "All figures generated by CostVision's should-cost engine. UK manufacturing region, GBP.")

headers12 = ["Commodity", "Demo Part", "Key Inputs", "Should-Cost", "AI Top Insight", "DFM Flag"]
data12 = [
    ("⚙️ Machining",     "Al Gearbox Bracket",    "Al 6082-T6 · 0.85 kg · 3 setups",       "£ 28.40", "Process 44% → near-net-shape saves £8",    "Major: 3 setups → 5-axis"),
    ("🪨 HPDC Casting",  "Transmission Housing",  "Al A380 · 3.2 kg · 50k/yr",             "£ 18.70", "Die 9% OK. Margin 21% → negotiate",        "Good — no critical issues"),
    ("🪗 Sheet Metal",   "Rear Door Inner",       "DC04 1.2mm · 4.1 kg · 60k/yr",          "£ 22.10", "Blank nesting 64% → improve saves £2.80",  "Critical: nesting <65%"),
    ("🏺 Injection",     "PP Bumper Bracket",     "PP-TD20 · 0.42 kg · 100k/yr",           "£  3.85", "Mould 8% — good at this volume",           "Good — runner waste 15%"),
    ("🧱 Rubber",        "EPDM Door Seal",        "Extrusion · 2.8m · 85 Shore A",         "£  4.10", "Process 42% — cure cycle dominant",        "Major: cure time optimise"),
    ("🖥️ PCB Fab",      "4-Layer FR4 ECU Board", "100×80mm · ENIG · Flying probe",        "£  4.18", "NRE 6% per board — increase batch",        "Good — standard technology"),
    ("🔌 PCBA",          "ECU PCBA 85 comp.",     "SMT + TH · 4 TH joints",               "£ 32.60", "Labour 38% — convert TH to SMT saves £4",  "Major: high TH labour"),
    ("🔗 Wiring Harness","Door Module Harness",   "24 circuits · 2.1m avg",               "£ 11.80", "Labour 58% — automation saves 20%",        "Critical: labour >50%"),
    ("🪢 Composites",    "CFRP Floor Pan",        "Prepreg CFRP · 1.8 kg · autoclave",    "£ 185.00","Labour 44% — evaluate ATL/AFP",            "Major: manual layup cost"),
    ("🔨 Forging",       "Steel Control Arm",     "42CrMo4 · 2.1 kg · 80k/yr",           "£  9.20", "Die 11% — acceptable at volume",           "Good — closed-die efficient"),
]

col_ws12 = [Inches(1.4), Inches(1.55), Inches(1.95), Inches(0.9), Inches(2.95), Inches(1.9)]
add_table(slide, 11, 6,
          Inches(0.45), Inches(1.98), Inches(12.45), Inches(4.78),
          headers12, data12, col_widths=col_ws12)

# Notes row
rect(slide, Inches(0.45), Inches(6.92), Inches(12.45), Inches(0.36), SURFACE2, BORDER, Pt(0.3))
txb(slide, "All costs in GBP  |  AI insights generated per part  |  DFM/DFA score per part  |  Export to 6-sheet Excel or PDF  |  Supplier quote comparison available for each",
    Inches(0.6), Inches(6.96), Inches(12.1), Inches(0.28),
    size=7.5, color=TEXT_D, align=PP_ALIGN.CENTER)

notes(slide,
    "To make the last slide concrete, here's a spread of real outputs across the commodities on a "
    "luxury-SUV programme — every figure came out of the engine, not off a rate card. "
    "I won't read the table; I'd rather you notice the pattern in the two right-hand columns. Each "
    "part comes with an AI insight that points at the biggest lever and a DFM flag that tells you "
    "whether it's clean or needs attention. The machined bracket says process is forty-four percent, "
    "go near-net-shape and save eight pounds. The wiring harness says labour is over half, so it's an "
    "automation candidate. The sheet-metal door flags that blank nesting is under sixty-five percent. "
    "That's the whole pitch in one table: not just a price, but a costed price with a reason and a "
    "next action attached to it — consistently, across ten very different processes.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Export, Reporting & Team Collaboration
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 17, "Reporting & Collaboration", "Getting the Number Out of the Room",
             "Reports people can act on, and one shared baseline for the whole team.")

export_cards = [
    (ACCENT_G, "📊 Excel — Six Sheets",
     "The whole build-up, not just the total.\n\n"
     "Summary · material · every operation ·\n"
     "machine rates · labour rates ·\n"
     "and a sheet recording every assumption."),
    (ACCENT_B, "📄 PDF Report",
     "The version you send out.\n\n"
     "Cover with the part photo, the cost summary,\n"
     "the operation-by-operation build,\n"
     "and the savings and DFM findings."),
    (ACCENT_P, "☁️ One Shared Baseline",
     "Everyone works from the same numbers.\n\n"
     "Saved centrally, not on a laptop.\n"
     "No arguing over which spreadsheet is current."),
    (ORANGE,   "📸 Part Photos",
     "Attach a photo to any calculation.\n\n"
     "It appears on the report cover — and the\n"
     "AI can cost from the photo itself."),
    (ACCENT_G, "🔑 Secure Access",
     "Enterprise sign-in, as you would expect.\n\n"
     "Email verification, password reset,\n"
     "and sign-in attempts capped."),
    (ACCENT_B, "❓ Built-In Help",
     "Nobody needs training to start.\n\n"
     "Getting started, the AI agent, CAD and photo,\n"
     "a glossary and troubleshooting — all in the app."),
]

cw_e = Inches(4.1)
ch_e = Inches(2.15)
sx_e = Inches(0.45)
sy_e = Inches(1.98)
gap_e = Inches(0.12)

for i, (col, title, body) in enumerate(export_cards):
    r, c = divmod(i, 3)
    cx_e = sx_e + c * (cw_e + gap_e)
    cy_e = sy_e + r * (ch_e + gap_e)
    card(slide, cx_e, cy_e, cw_e, ch_e, title, body, accent=col)

notes(slide,
    "A cost number that lives in one engineer's head isn't worth much, so this slide is about "
    "getting it out and shared. "
    "The Excel export is six sheets — not just the total, but the full build-up: material, every "
    "operation, the rate library snapshot, and a traceability sheet that records the assumptions and "
    "the AI notes with a version and date. That last sheet is what makes an estimate auditable "
    "months later. The PDF is the boardroom version, cover photo to DFM recommendations. "
    "And it's genuinely multi-user — scenarios live in a database with proper authentication, so the "
    "whole team works off one baseline instead of emailing spreadsheets around and losing track of "
    "which version is current. It scales from one engineer to a global team without changing how it "
    "works.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Business Benefits & ROI
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 18, "Value Delivered", "What This Is Worth",
             "The same claims, in numbers a finance team recognises.")

# Stat row
stats14 = [
    ("70–90%", "Reduction in costing time",                     ACCENT_B),
    ("±8%",    "Typical cost model accuracy",                   ACCENT_G),
    ("15–25%", "Typical supplier price reduction",              ORANGE),
    ("3×",     "Faster sourcing decisions",                     ACCENT_P),
    ("£15M",   "On £500M of spend, a 3% improvement",           ACCENT_G),
]
sw = Inches(2.4)
sh = Inches(1.1)
for i, (num, lbl, col) in enumerate(stats14):
    stat_card(slide, Inches(0.45) + i * (sw + Inches(0.1)), Inches(1.98), sw, sh, num, lbl, col)

# Benefits detail cards
benefits14 = [
    (ACCENT_B, "⚡ Speed",
     "Two to four weeks becomes ten minutes. Engineers spend the day deciding, not building spreadsheets."),
    (ACCENT_G, "🎯 Accuracy",
     "Built from the physics of the process and checked against real shop-floor data — not a rule of thumb."),
    (ORANGE,   "💰 Negotiation",
     "A floor price you can defend, part by part. Supplier margin becomes visible instead of assumed."),
    (ACCENT_P, "📐 Early Warning",
     "Most of the cost is fixed at design. Problems surface while changing them is still free."),
    (ACCENT_G, "🌍 Where to Build It",
     "Compare 20 regions in one click. Make-vs-buy answered with real rates."),
    (ACCENT_B, "📊 One Method",
     "The same model everywhere, with every assumption on record. Leadership sees one set of numbers."),
    (ORANGE,   "📉 The Right Volume",
     "See what the part costs as volume grows, and when a new process starts to pay."),
    (ACCENT_P, "☁️ Shared Knowledge",
     "One baseline the whole team works from. Cost know-how stays with the business, not one engineer."),
]

# Card height derived from a bottom target — a typed 1.3" left 1.4" of empty
# page under the grid once the copy was cut back.
sx_b, sy_b, gap_b = Inches(0.45), Inches(3.25), Inches(0.1)
COLS_B = 4
cw_b = (W - sx_b * 2 - gap_b * (COLS_B - 1)) / COLS_B
_rows_b = -(-len(benefits14) // COLS_B)
ch_b = (Inches(6.95) - sy_b - gap_b * (_rows_b - 1)) / _rows_b

for i, (col, title, body) in enumerate(benefits14):
    r, c = divmod(i, COLS_B)
    card(slide, sx_b + c * (cw_b + gap_b), sy_b + r * (ch_b + gap_b),
         cw_b, ch_b, title, body, accent=col, title_pt=10.5, body_pt=9.5)

notes(slide,
    "Let me translate all of that into the numbers a manager actually cares about. "
    "Speed first: what took two to four weeks now takes minutes. That's the seventy-to-ninety "
    "percent time reduction, and it means engineers spend their day on decisions instead of "
    "spreadsheets. On accuracy, the models are calibrated to real shop-floor data and now corrected "
    "by the self-audit — which is exactly what tightens that plus-or-minus-eight-percent band over "
    "time. "
    "But the line that pays for the whole thing is negotiation power. When you walk into a supplier "
    "meeting with a defensible floor price, a fifteen-to-twenty-five percent reduction is a "
    "realistic ask — and the far-right stat makes it concrete: on five hundred million of spend, a "
    "three percent improvement is fifteen million pounds. This tool doesn't need to be perfect to "
    "pay for itself many times over; it needs to be defensible and fast, and it is both.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Roadmap & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, 19, "Vision & Next Steps", "Intelligent, Instant, Integrated — What Comes Next",
             "CostVision is live today. Here is what comes next — and how to get started.")

# Current capabilities
rect(slide, Inches(0.45), Inches(2.0), Inches(5.9), Inches(4.92), SURFACE2, BORDER, Pt(0.5))
rect(slide, Inches(0.45), Inches(2.0), Inches(5.9), Inches(0.04), ACCENT_G)
txb(slide, "📍 Current Capabilities — Live Today", Inches(0.6), Inches(2.08),
    Inches(5.6), Inches(0.3), size=9.5, bold=True, color=ACCENT_G)

live_items = [
    "21 commodities costed to engineering depth",
    "Describe a part in plain English — the model builds itself",
    "Upload a CAD file or a photo, get a cost in minutes",
    "Every estimate re-checked, and corrected, before you see it",
    "Learns from the supplier quotes you already have",
    "Whole assemblies, not just single parts",
    "Learning curve, quote comparison and scenarios",
    "20 regions and 10 currencies, switched automatically",
    "Excel and PDF reports, shared across the team",
]

# The pitch is derived so the list always fills its panel — the old fixed 0.34"
# was tuned to thirteen items and would leave a gap at any other count.
yy15 = Inches(2.44)
_pitch15 = (Inches(6.80) - yy15) / len(live_items)
for item in live_items:
    rect(slide, Inches(0.65), yy15 + Inches(0.15), Inches(0.06), Inches(0.06), ACCENT_G)
    txb(slide, item, Inches(0.84), yy15, Inches(5.3), Inches(0.3),
        size=9, color=TEXT_G, wrap=True)
    yy15 += _pitch15

# Roadmap
rx15 = Inches(6.6)
ry15 = Inches(2.0)
rw15 = Inches(6.3)

phases = [
    (ORANGE,   "Phase 2 — H2 2026",
     "Connect to Teamcenter and SAP\n"
     "Cost-target dashboard for programme teams\n"
     "Compare several suppliers side by side"),
    (ACCENT_P, "Phase 3 — 2027",
     "Calibrate itself from live ERP actuals\n"
     "Draft the RFQ straight from the should-cost\n"
     "Portfolio view for the leadership team"),
    (ACCENT_B, "Phase 4 — 2027+",
     "Cost that updates as the design changes\n"
     "AI suggestions for redesigning the part\n"
     "Supplier risk scoring · mobile app for plant visits"),
]

yy15r = ry15
for col, title, body in phases:
    rect(slide, rx15, yy15r, rw15, Inches(1.06), SURFACE2, BORDER, Pt(0.5))
    rect(slide, rx15, yy15r, Inches(0.06), Inches(1.06), col)
    txb(slide, title, rx15 + Inches(0.14), yy15r + Inches(0.06), rw15 - Inches(0.2), Inches(0.28),
        size=9, bold=True, color=col)
    txb(slide, body, rx15 + Inches(0.14), yy15r + Inches(0.32), rw15 - Inches(0.2), Inches(0.72),
        size=8, color=TEXT_G, wrap=True)
    yy15r += Inches(1.14)

# Pilot invite + contact
rect(slide, rx15, yy15r, rw15, Inches(1.42), SURFACE2, BORDER, Pt(0.5))
rect(slide, rx15, yy15r, Inches(0.06), Inches(1.42), ACCENT_G)
txb(slide, "🚀  Pilot Rollout — Next Steps", rx15 + Inches(0.14), yy15r + Inches(0.08),
    rw15 - Inches(0.2), Inches(0.28), size=9, bold=True, color=ACCENT_G)
pilot_text = (
    "• Two or three commodities, one programme team\n"
    "• Four to six weeks, on real parts against real quotes\n"
    "• We call it a success at 15% price reduction found per part\n"
    "• Avinash Bhosale — Cost Engineering & Digital Innovation"
)
txb(slide, pilot_text, rx15 + Inches(0.14), yy15r + Inches(0.38), rw15 - Inches(0.2), Inches(0.92),
    size=8.5, color=TEXT_G, wrap=True)

# Vision quote at bottom
rect(slide, Inches(0.45), Inches(7.08), W - Inches(0.9), Inches(0.36), SURFACE2, BORDER, Pt(0.3))
txb(slide, '"Every engineer should have instant access to accurate, defensible should-cost intelligence — at concept stage, across every commodity, in every region."  — Avinash Bhosale',
    Inches(0.65), Inches(7.11), W - Inches(1.3), Inches(0.28),
    size=8, color=TEXT_D, italic=True, align=PP_ALIGN.CENTER)

notes(slide,
    "So where does that leave us. The left column is not a wish list — it's live today. Twenty-one "
    "commodities, the AI front doors, and this year's additions: the self-audit on every estimate, "
    "learning from your actuals, universal machine-sizing, and the whole thing deployable to the "
    "cloud with the CAD engine and a thousand-plus tests behind it. That's a working platform, not a "
    "prototype. "
    "The right-hand column is honest about what's next. Near-term it's about fitting into your "
    "existing systems — PLM and ERP connectors, a cost-target waterfall for programme teams. "
    "Further out, the calibration learns automatically from a live actuals feed, it starts "
    "generating the RFQ itself, and eventually a digital-twin cost model that updates as the design "
    "moves. "
    "But I don't want the roadmap to distract from the ask, which is small and concrete. It's at the "
    "bottom of the box: a four-to-six week pilot on two or three commodity families, one programme "
    "team, measured against real supplier quotes, with the success bar set at finding at least "
    "fifteen percent per part. If it clears that bar on your parts, the business case makes itself. "
    "That's what I'd like to agree today. Thank you — I'm happy to take questions.")


# ─── Save ─────────────────────────────────────────────────────────────────────
output_path = "/home/user/leamington-marathi/CostVision-Executive-Presentation.pptx"
prs.save(output_path)
# python-pptx cannot declare the notes master, and PowerPoint will not open a
# deck that has notes slides without one. See pptx_fixup.py.
_fixed = finalise(output_path)
print(f"Saved: {output_path}" + (f"  [fixed: {', '.join(_fixed)}]" if _fixed else ""))
print(f"Slides: {len(prs.slides)}")
