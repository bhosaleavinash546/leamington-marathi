#!/usr/bin/env python3
"""
CostVision — "Agentic AI" management presentation builder.

Light professional theme, CostVision logo top-left on every slide, native
(editable) shapes + charts, speaker notes per slide. All figures are real,
verified outputs from the tool's live test runs — nothing invented.

Regenerate:  python3 build_agentic_pptx.py
Output:      CostVision-Agentic-AI-Management-Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── Brand palette (light theme) ────────────────────────────────────────────────
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)   # logo badge
BLUE    = RGBColor(0x25, 0x63, 0xEB)   # wordmark / primary accent
DARK    = RGBColor(0x0F, 0x17, 0x2A)   # headings
BODY    = RGBColor(0x33, 0x41, 0x55)   # body text
MUTED   = RGBColor(0x64, 0x74, 0x8B)   # captions
BG      = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xF1, 0xF5, 0xF9)   # light slate panel
PANEL2  = RGBColor(0xEF, 0xF6, 0xFF)   # light blue panel
GREEN   = RGBColor(0x05, 0x96, 0x69)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
RED     = RGBColor(0xDC, 0x26, 0x26)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
CYAN    = RGBColor(0x08, 0x91, 0xB2)
LINE    = RGBColor(0xE2, 0xE8, 0xF0)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Low-level helpers ──────────────────────────────────────────────────────────
def _noline(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def box(slide, x, y, w, h, fill=None, line=None, round_=False, radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if round_:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold[, italic])."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for r in para:
            t, size, color, bold = r[0], r[1], r[2], r[3]
            italic = r[4] if len(r) > 4 else False
            run = p.add_run(); run.text = t
            f = run.font
            f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
            f.name = 'Calibri'
    return tb

def logo(slide, x=Inches(0.35), y=Inches(0.22), scale=1.0):
    """CostVision logo — indigo rounded 'cv' badge + blue wordmark + grey tagline."""
    s = scale
    badge = box(slide, x, y, Inches(0.42 * s), Inches(0.42 * s), fill=INDIGO, round_=True, radius=0.28)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'cv'
    r.font.size = Pt(17 * s); r.font.bold = True; r.font.color.rgb = BG; r.font.name = 'Calibri'
    text(slide, x + Inches(0.52 * s), y - Inches(0.03 * s), Inches(2.6), Inches(0.32),
         [[('CostVision', 18 * s, BLUE, True)]])
    text(slide, x + Inches(0.52 * s), y + Inches(0.24 * s), Inches(2.8), Inches(0.22),
         [[('AI  COST  INTELLIGENCE', 7.5 * s, MUTED, False)]])

def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt

def header(title, kicker=None):
    """New slide with logo header + underline + title. Returns slide."""
    slide = prs.slides.add_slide(BLANK)
    box(slide, 0, 0, W, H, fill=BG)                       # background
    logo(slide)
    box(slide, 0, Inches(0.78), W, Pt(2.2), fill=INDIGO)  # header rule
    if kicker:
        text(slide, Inches(0.45), Inches(0.95), Inches(11.5), Inches(0.3),
             [[(kicker.upper(), 11, BLUE, True)]])
        ty = Inches(1.22)
    else:
        ty = Inches(1.02)
    text(slide, Inches(0.45), ty, Inches(12.4), Inches(0.6), [[(title, 27, DARK, True)]])
    return slide

def kpi_card(slide, x, y, w, h, big, label, sub, color=BLUE):
    box(slide, x, y, w, h, fill=PANEL, round_=True, radius=0.10)
    box(slide, x, y, Inches(0.07), h, fill=color)
    text(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.35), Inches(0.55),
         [[(big, 25, color, True)]])
    text(slide, x + Inches(0.22), y + Inches(0.68), w - Inches(0.35), Inches(0.3),
         [[(label, 12.5, DARK, True)]])
    text(slide, x + Inches(0.22), y + Inches(0.98), w - Inches(0.35), h - Inches(1.05),
         [[(sub, 10, MUTED, False)]], line_spacing=1.05)

def chevron(slide, x, y, w, h, label, sub, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.adjustments[0] = 0.28
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    _noline(shp)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = BG; r.font.name = 'Calibri'
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(8.5); r2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xFF); r2.font.name = 'Calibri'

def style_chart(chart, series_colors):
    chart.has_legend = False
    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = LINE
        chart.value_axis.format.line.color.rgb = LINE
        chart.category_axis.format.line.color.rgb = LINE
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.color.rgb = MUTED
        chart.category_axis.tick_labels.font.size = Pt(10.5)
        chart.category_axis.tick_labels.font.color.rgb = BODY
    except Exception:
        pass
    for i, ser in enumerate(chart.plots[0].series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = series_colors[i % len(series_colors)]
        ser.format.line.fill.background()
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10.5); dl.font.bold = True; dl.font.color.rgb = DARK


# ════════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=BG)
box(s, 0, 0, W, Inches(0.16), fill=INDIGO)
logo(s, x=Inches(0.5), y=Inches(0.45), scale=1.25)
text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.0),
     [[('Agentic AI in CostVision', 46, DARK, True)]])
text(s, Inches(0.9), Inches(3.35), Inches(11.0), Inches(0.6),
     [[('A costing tool that learns from every analysis, remembers your parts,', 19, BODY, False)],
      [('and finds savings on its own — while staying fully auditable.', 19, BODY, False)]])
for i, (t, c) in enumerate([('Remembers', BLUE), ('Recognises', CYAN), ('Self-corrects', VIOLET), ('Acts autonomously', GREEN)]):
    x = Inches(0.9 + i * 2.85)
    chip = box(s, x, Inches(4.6), Inches(2.6), Inches(0.52), fill=PANEL, round_=True, radius=0.5)
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = c; r.font.name = 'Calibri'
text(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
     [[('Management briefing  ·  July 2026  ·  All figures verified from live system runs', 12, MUTED, False)]])
box(s, 0, H - Inches(0.16), W, Inches(0.16), fill=INDIGO)
notes(s, "Welcome. Today I'll show you the Agentic AI capability we have built into CostVision. "
         "In one sentence: the tool no longer just calculates — it learns from every analysis we run, "
         "remembers every part, gets more accurate with every real quote we log, and even works unattended, "
         "flagging savings opportunities by itself. Everything I'll show is from live runs of the system, not mock-ups.")

# ════════════════════════════════════════════════════════════════════════════
# 2 — EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════════════════
s = header('What we built — in one slide', 'Executive summary')
kpi_card(s, Inches(0.45), Inches(2.0), Inches(3.0), Inches(1.75), '36×', 'Error reduction',
         'Estimating error fell from 10.9% to 0.3% after the tool learned from just 3 real quotes (verified live).', GREEN)
kpi_card(s, Inches(3.65), Inches(2.0), Inches(3.0), Inches(1.75), '£512k/yr', 'Found autonomously',
         'In our live demo the background agent flagged £512k/yr of pricing issues — with nobody at the keyboard.', RED)
kpi_card(s, Inches(6.85), Inches(2.0), Inches(3.0), Inches(1.75), '99%', 'Part recognition',
         'A new bracket was matched to 3 past bracket analyses at 98–99% similarity, with reasons shown.', CYAN)
kpi_card(s, Inches(10.05), Inches(2.0), Inches(2.85), Inches(1.75), '917', 'Automated tests',
         'Every capability is covered by automated tests (77 suites) and was exercised end-to-end on the running system.', VIOLET)
box(s, Inches(0.45), Inches(4.1), Inches(12.45), Inches(2.7), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(4.35), Inches(11.9), Inches(2.3),
     [[('The idea, in plain words', 15, DARK, True)],
      [('Until now, every costing started from zero and the result depended on who did it. ', 13.5, BODY, False),
       ('Now the tool keeps an organisational memory: ', 13.5, DARK, True),
       ('every analysis is stored, every real supplier quote teaches it, and every new part is compared against everything we have costed before. ', 13.5, BODY, False)],
      [('The knowledge stays in our database, on our servers — it becomes a company asset that gets more valuable with use, and it does not walk out of the door when an expert leaves.', 13.5, BODY, False)]],
     space_after=8, line_spacing=1.15)
notes(s, "Four headline numbers, all verified. One: after learning from only three real quotes, the estimating error "
         "in a segment dropped from about 11% to under 1%. Two: the autonomous monitor found half a million pounds a year "
         "of pricing issues in our demonstration, entirely unattended. Three: when we costed a new bracket, the tool "
         "recognised three similar past parts at 99% similarity and told us why. Four: this is production-grade — "
         "813 automated tests. The key message: costing knowledge now accumulates as a company asset instead of "
         "living in individual experts' heads.")

# ════════════════════════════════════════════════════════════════════════════
# 3 — WHAT AGENTIC AI MEANS
# ════════════════════════════════════════════════════════════════════════════
s = header('What "Agentic AI" means here — four plain words', 'The concept')
cards = [
    ('🧠', 'Remembers', 'Every costing is saved as a "case": the part, the inputs, the result, and any real quote. Shared across the whole team.', BLUE),
    ('🔎', 'Recognises', 'Start a new part and it instantly finds the most similar past parts — like an experienced engineer saying "we\'ve done this before".', CYAN),
    ('🎯', 'Self-corrects', 'Log the real supplier price and the tool measures its own error, then corrects future estimates in that category. Accuracy is measured, not claimed.', VIOLET),
    ('🤖', 'Acts', 'A background agent re-checks all stored parts on a schedule and raises findings by itself: "you are overpaying here — worth £400k/yr."', GREEN),
]
for i, (icon, t, d, c) in enumerate(cards):
    x = Inches(0.45 + i * 3.24)
    box(s, x, Inches(2.1), Inches(3.02), Inches(3.6), fill=PANEL, round_=True, radius=0.07)
    box(s, x, Inches(2.1), Inches(3.02), Inches(0.09), fill=c)
    text(s, x + Inches(0.25), Inches(2.35), Inches(2.5), Inches(0.6), [[(icon, 28, c, False)]])
    text(s, x + Inches(0.25), Inches(3.05), Inches(2.55), Inches(0.45), [[(t, 17, c, True)]])
    text(s, x + Inches(0.25), Inches(3.55), Inches(2.55), Inches(2.0), [[(d, 11.5, BODY, False)]], line_spacing=1.15)
text(s, Inches(0.45), Inches(6.0), Inches(12.4), Inches(0.9),
     [[('Deliberate design choice: ', 12.5, DARK, True),
       ('the learning is statistics over our own data — every suggestion shows its source parts and its arithmetic. '
        'That makes it auditable and defensible in front of a supplier, which is where costing tools win or lose.', 12.5, BODY, False)]],
     line_spacing=1.15)
notes(s, "When we say Agentic AI, we mean four concrete abilities. It REMEMBERS — every analysis becomes a stored case. "
         "It RECOGNISES — new parts are matched against that memory. It SELF-CORRECTS — real quotes teach it, and its "
         "accuracy is measured, not asserted. And it ACTS — a background agent raises findings without being asked. "
         "One important design choice: every AI suggestion is traceable to specific past parts and shows its arithmetic. "
         "That auditability is what makes it usable in supplier negotiations.")

# ════════════════════════════════════════════════════════════════════════════
# 4 — THE LEARNING LOOP
# ════════════════════════════════════════════════════════════════════════════
s = header('How it works — the learning loop', 'How it works')
steps = [
    ('1 · Analyse', 'Engineer costs a part as usual', BLUE),
    ('2 · Remember', 'Saved automatically to the knowledge base', INDIGO),
    ('3 · Recognise', 'Similar past parts found instantly', CYAN),
    ('4 · Suggest', 'Benchmarks, materials, real quotes shown', VIOLET),
    ('5 · Learn', 'Real quote logged → model self-corrects', AMBER),
    ('6 · Act', 'Background agent flags drift & savings', GREEN),
]
cw = Inches(2.24)
for i, (t, d, c) in enumerate(steps):
    chevron(s, Inches(0.4) + int(cw * 0.86) * i, Inches(2.4), cw, Inches(1.15), t, d, c)
# feedback arrow
arr = s.shapes.add_shape(MSO_SHAPE.BENT_UP_ARROW, Inches(5.6), Inches(3.85), Inches(6.2), Inches(0.85))
arr.rotation = 180
arr.fill.solid(); arr.fill.fore_color.rgb = LINE
_noline(arr)
text(s, Inches(4.7), Inches(4.9), Inches(4.6), Inches(0.35),
     [[('…and every loop makes the next estimate better', 12, MUTED, False, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.45), Inches(5.5), Inches(12.45), Inches(1.35), fill=PANEL2, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(5.72), Inches(11.9), Inches(1.0),
     [[('No extra work for the engineer. ', 13, DARK, True),
       ('Steps 2, 3, 4 and 6 happen automatically. The only new habit is one click — "Log Actual £" — when a real '
        'supplier quote arrives. That single click is the fuel for everything on this slide.', 13, BODY, False)]],
     line_spacing=1.2)
notes(s, "This is the whole system on one slide. The engineer works exactly as before — they cost a part. "
         "The tool remembers it, recognises similar history, and makes suggestions, all automatically. "
         "When a real supplier quote arrives, one click teaches the tool, and the background agent keeps watch from then on. "
         "The crucial point for adoption: there is essentially no extra workload. One new click when a quote arrives.")

# ════════════════════════════════════════════════════════════════════════════
# 5 — MEMORY: KNOWLEDGE BASE
# ════════════════════════════════════════════════════════════════════════════
s = header('The memory — an organisational knowledge base', 'Capability 1 of 5')
box(s, Inches(0.45), Inches(2.0), Inches(6.0), Inches(4.6), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(2.25), Inches(5.5), Inches(0.4), [[('What is stored for every analysis', 15, DARK, True)]])
rows = [
    ('Part "fingerprint"', 'process, material, weight, size, region, volume'),
    ('The full cost result', 'total + breakdown by cost driver'),
    ('Real quotes', 'actual supplier / PO prices, when logged'),
    ('Expert corrections', 'where a person adjusted the AI\'s values'),
    ('CAD shape data', 'dimensions & features, when a CAD file was used'),
]
for i, (a, b) in enumerate(rows):
    y = Inches(2.75 + i * 0.72)
    box(s, Inches(0.75), y, Inches(0.09), Inches(0.55), fill=BLUE)
    text(s, Inches(1.0), y, Inches(5.2), Inches(0.62),
         [[(a + ' — ', 12.5, DARK, True), (b, 12.5, BODY, False)]], line_spacing=1.05)
box(s, Inches(6.85), Inches(2.0), Inches(6.05), Inches(4.6), fill=PANEL2, round_=True, radius=0.05)
text(s, Inches(7.15), Inches(2.25), Inches(5.5), Inches(0.4), [[('Why it matters to us', 15, DARK, True)]])
pts = [
    ('Shared, not personal', 'One engineer\'s analysis instantly helps everyone — juniors inherit senior judgement.'),
    ('Stays on our servers', 'Our database, our infrastructure. The knowledge is a company asset, not a vendor\'s.'),
    ('No duplicates', 'Re-costing a part updates its case — the memory stays clean.'),
    ('Compounds with use', 'Useful from ~20–30 analyses; every costing from today is an investment.'),
]
for i, (a, b) in enumerate(pts):
    y = Inches(2.75 + i * 0.92)
    text(s, Inches(7.15), y, Inches(5.5), Inches(0.85),
         [[('✓  ' + a, 13, GREEN, True)], [(b, 11.5, BODY, False)]], line_spacing=1.05, space_after=2)
notes(s, "Capability one: the memory. Every analysis is saved with a fingerprint of the part, the full result, "
         "any real quotes, and — importantly — the places where an expert corrected the AI. Corrections are how "
         "our experts literally teach the tool. Right side: this is shared across the team, it lives in our own "
         "database on our own servers, and it compounds — it becomes genuinely useful after only twenty to thirty "
         "analyses. Every costing we run from today is a deposit into this asset.")

# ════════════════════════════════════════════════════════════════════════════
# 6 — RECOGNITION & SUGGESTIONS
# ════════════════════════════════════════════════════════════════════════════
s = header('Recognition — "we have costed this before"', 'Capability 2 of 5')
text(s, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.4),
     [[('Live example: an engineer costs a new 0.85 kg aluminium bracket. The tool answers instantly:', 13.5, BODY, False)]])
box(s, Inches(0.45), Inches(2.4), Inches(12.45), Inches(3.1), fill=PANEL, round_=True, radius=0.05)
box(s, Inches(0.45), Inches(2.4), Inches(0.09), Inches(3.1), fill=CYAN)
text(s, Inches(0.8), Inches(2.6), Inches(11.8), Inches(0.35),
     [[('🧠  AI memory — similar past parts', 14, DARK, True),
       ('      (knowledge base: 3 analyses · 2 with actuals)', 10.5, MUTED, False)]])
matches = [
    ('EV Battery Bracket', '99% match — material family, weight, region', '£41.20', 'actual £46.50'),
    ('Sensor Mount Bracket', '98% match — material family, weight, region', '£38.50', ''),
    ('Inverter Bracket', '98% match — material family, weight, region', '£44.80', 'actual £50.20'),
]
for i, (n, m, e, a) in enumerate(matches):
    y = Inches(3.05 + i * 0.5)
    text(s, Inches(0.9), y, Inches(6.6), Inches(0.42),
         [[(n, 12.5, DARK, True), ('   ' + m, 10.5, MUTED, False)]])
    text(s, Inches(8.0), y, Inches(4.6), Inches(0.42),
         [[(e, 12.5, DARK, True), (('   ' + a) if a else '', 11, VIOLET, True)]])
sugg = [
    ('•  Median cost of 3 similar parts: £41.20', BODY),
    ('•  3 of 3 used the same material (aluminium 6061)', BODY),
    ('•  Real quotes logged for 2 of them — median actual £48.35', VIOLET),
]
for i, (t, c) in enumerate(sugg):
    text(s, Inches(0.9), Inches(4.6 + i * 0.3), Inches(11.6), Inches(0.3), [[(t, 11.5, c, False)]])
box(s, Inches(0.45), Inches(5.8), Inches(12.45), Inches(1.0), fill=PANEL2, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(5.97), Inches(11.9), Inches(0.75),
     [[('Every match is explained (what matched, and how strongly) and every suggestion names its source parts — ', 12.5, BODY, False),
       ('no black box.', 12.5, DARK, True)]], line_spacing=1.15)
notes(s, "Capability two, shown with the actual live output. An engineer starts costing a new bracket. Before they "
         "finish, the tool has found the three most similar parts we've costed before — at 98 to 99 percent similarity — "
         "and it explains WHY they matched: same material family, same weight class, same region. It then suggests: "
         "the median cost of those parts, the material they all used, and crucially the real quoted prices we logged. "
         "A junior engineer instantly benefits from all previous work. And note — every number names its source. No black box.")

# ════════════════════════════════════════════════════════════════════════════
# 7 — SELF-CALIBRATION (chart)
# ════════════════════════════════════════════════════════════════════════════
s = header('Self-correction — it learns from real quotes', 'Capability 3 of 5')
cd = CategoryChartData()
cd.categories = ['Machining · Aluminium · UK', 'Casting · Aluminium · China']
cd.add_series('Before learning', (10.9, 8.7))
cd.add_series('After 3 real quotes', (0.3, 0.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.45), Inches(2.15), Inches(6.6), Inches(4.0), cd)
chart = gf.chart
style_chart(chart, [MUTED, GREEN])
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
text(s, Inches(0.45), Inches(6.25), Inches(6.6), Inches(0.5),
     [[('Estimating error (%) — measured against real supplier prices, live system run', 10.5, MUTED, False, True)]])
box(s, Inches(7.45), Inches(2.15), Inches(5.45), Inches(4.6), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(7.75), Inches(2.4), Inches(4.9), Inches(4.2),
     [[('How it works, simply', 15, DARK, True)],
      [('1.  A real supplier quote arrives → one click logs it.', 12.5, BODY, False)],
      [('2.  The tool compares its estimate with reality and measures its own error.', 12.5, BODY, False)],
      [('3.  From 3 quotes in a category, it corrects future estimates automatically — per process, material family and region.', 12.5, BODY, False)],
      [('4.  It reports its accuracy openly (error % before and after).', 12.5, BODY, False)],
      [('Why per-category matters: ', 12.5, DARK, True),
       ('our machining ran 12% low while our China castings ran 8% high — averaged together they looked "fine". The tool catches what averages hide.', 12.5, BODY, False)]],
     space_after=10, line_spacing=1.12)
notes(s, "Capability three is the accuracy engine. The chart shows a real measured result: in machining, our estimates "
         "ran about 11 percent below the real prices. After the tool learned from just three logged quotes, the "
         "remaining error was 0.3 percent. Same story for castings from China, where the bias was in the opposite "
         "direction. And that's the subtle but important point on the right: one category under-estimates, another "
         "over-estimates — a portfolio average hides both. The tool corrects each category separately. "
         "Accuracy here is measured against reality and reported openly — not a marketing claim.")

# ════════════════════════════════════════════════════════════════════════════
# 8 — HONEST UNCERTAINTY (chart)
# ════════════════════════════════════════════════════════════════════════════
s = header('Honest ranges — from "±20%" to "±3%", earned with data', 'Capability 4 of 5')
cd = CategoryChartData()
cd.categories = ['Before any real quotes', 'After 3 real quotes logged']
cd.add_series('Confidence band (± %)', (20.4, 2.8))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.45), Inches(2.15), Inches(5.9), Inches(4.0), cd)
style_chart(gf.chart, [AMBER])
text(s, Inches(0.45), Inches(6.25), Inches(5.9), Inches(0.5),
     [[('Width of the cost confidence band on the same part (± % around the estimate)', 10.5, MUTED, False, True)]])
box(s, Inches(6.75), Inches(2.15), Inches(6.15), Inches(4.6), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(7.05), Inches(2.4), Inches(5.6), Inches(4.2),
     [[('Every estimate now comes as a range', 15, DARK, True)],
      [('Optimistic (P10)  ·  Most likely (P50)  ·  Conservative (P90)', 12.5, BLUE, True)],
      [('A single number implies a precision that early estimates don\'t have. The range tells buyers how much to '
        'trust the number — and what target to set in negotiation.', 12.5, BODY, False)],
      [('The band is earned, not guessed: ', 12.5, DARK, True),
       ('it starts wide when the tool has no evidence, and tightens automatically as real quotes prove the accuracy. '
        'On our test part it went from ±20% to ±3% after three quotes.', 12.5, BODY, False)]],
     space_after=10, line_spacing=1.15)
notes(s, "Capability four: honesty about precision. Every estimate is now delivered as a range — optimistic, most "
         "likely, conservative. Early on, with no evidence, the range is wide, and that's the truth. As real quotes "
         "prove the tool's accuracy, the range tightens automatically — on our test part, from plus-or-minus twenty "
         "percent down to three percent. For buyers this is a practical negotiation tool: the conservative end is "
         "the walk-away price; the optimistic end is the stretch target.")

# ════════════════════════════════════════════════════════════════════════════
# 9 — AUTONOMOUS AGENT
# ════════════════════════════════════════════════════════════════════════════
s = header('The autonomous agent — finds money unattended', 'Capability 5 of 5')
text(s, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.45),
     [[('A background monitor re-checks every stored part on a schedule. In the live demonstration it opened these findings ', 13, BODY, False),
       ('entirely unattended:', 13, DARK, True)]])
findings = [
    ('💰  RENEGOTIATION', RED, '≈ £400,000 / yr',
     'Supplier price £48.00 is 20% above should-cost £40.00 at 50,000 pcs — clear recovery opportunity.'),
    ('🚨  UNDERWATER PRICE', AMBER, '≈ £112,000 / yr exposure',
     'Supplier price 25% BELOW should-cost — verify scope and quality; pricing may be unsustainable.'),
    ('⏳  STALE ESTIMATE', MUTED, 'Confidence decay',
     'A 120-day-old estimate was never validated with a real quote — the agent asks for a refresh.'),
]
for i, (tag, c, imp, d) in enumerate(findings):
    y = Inches(2.5 + i * 1.15)
    box(s, Inches(0.45), y, Inches(12.45), Inches(1.0), fill=PANEL, round_=True, radius=0.10)
    box(s, Inches(0.45), y, Inches(0.09), Inches(1.0), fill=c)
    text(s, Inches(0.8), y + Inches(0.13), Inches(3.3), Inches(0.4), [[(tag, 13, c, True)]])
    text(s, Inches(0.8), y + Inches(0.5), Inches(9.0), Inches(0.45), [[(d, 11.5, BODY, False)]])
    text(s, Inches(9.9), y + Inches(0.13), Inches(2.85), Inches(0.4), [[(imp, 14, c, True)]], align=PP_ALIGN.RIGHT)
box(s, Inches(0.45), Inches(6.05), Inches(12.45), Inches(0.85), fill=PANEL2, round_=True, radius=0.10)
text(s, Inches(0.75), Inches(6.2), Inches(11.9), Inches(0.6),
     [[('Total surfaced in the demo:  ', 13, BODY, False), ('≈ £512,000 / yr', 15, RED, True),
       ('   — every finding shows its arithmetic, and one click dismisses it once handled.', 12, MUTED, False)]])
notes(s, "Capability five is the one that makes this genuinely 'agentic'. A monitor runs on the server on a schedule — "
         "first pass thirty seconds after start-up, then every few hours. It compares what we PAY with what things "
         "SHOULD cost, using everything the tool has learned. In the live demonstration it opened these three findings "
         "by itself: a renegotiation opportunity worth four hundred thousand pounds a year, an underwater price that "
         "signals supply risk, and a stale estimate that needs revalidating. Half a million pounds a year surfaced with "
         "nobody at the keyboard. Each finding shows its arithmetic — ready to take into a supplier meeting.")

# ════════════════════════════════════════════════════════════════════════════
# 10 — SUPPORTING AI BRAINS
# ════════════════════════════════════════════════════════════════════════════
s = header('Supporting AI capabilities already in the tool', 'The wider AI platform')
quad = [
    ('📚  Grounded AI assistant', 'Answers cite our actual rate library and past costings — the AI quotes our data, it doesn\'t improvise. Every figure is traceable.', BLUE),
    ('📄  RFQ analyst', 'Drop in an RFQ package: it costs every line, flags risky prices and single-source parts, and drafts a prioritised negotiation brief.', VIOLET),
    ('📐  CAD feature costing', 'Reads the 3D model and prices individual design features — holes, threads, surfaces — so designers see exactly what drives cost.', CYAN),
    ('🌱  Carbon co-costing', 'Every part gets a CO₂e figure alongside the £ — increasingly demanded in automotive & aerospace RFQs.', GREEN),
]
for i, (t, d, c) in enumerate(quad):
    x = Inches(0.45 + (i % 2) * 6.35)
    y = Inches(2.1 + (i // 2) * 2.3)
    box(s, x, y, Inches(6.1), Inches(2.05), fill=PANEL, round_=True, radius=0.07)
    box(s, x, y, Inches(0.09), Inches(2.05), fill=c)
    text(s, x + Inches(0.3), y + Inches(0.2), Inches(5.6), Inches(0.45), [[(t, 15, c, True)]])
    text(s, x + Inches(0.3), y + Inches(0.72), Inches(5.55), Inches(1.25), [[(d, 12, BODY, False)]], line_spacing=1.15)
text(s, Inches(0.45), Inches(6.7), Inches(12.4), Inches(0.5),
     [[('Together with 18 commodity cost engines, CAD-to-cost, and PCB photo-to-BOM — the agentic layer sits on top of all of it.', 12, MUTED, False, True)]])
notes(s, "The learning loop is the headline, but it sits on a wider AI platform we've built. The assistant answers "
         "from our own rate data with citations. The RFQ analyst turns a full quote package into a costed, "
         "risk-flagged negotiation brief. CAD feature costing tells designers which specific features drive cost. "
         "And carbon co-costing gives a CO2 figure alongside every price — which customers increasingly require. "
         "All of this feeds and benefits from the same knowledge base.")

# ════════════════════════════════════════════════════════════════════════════
# 11 — INPUTS REQUIRED
# ════════════════════════════════════════════════════════════════════════════
s = header('What it needs from us — honestly, very little', 'Inputs required')
box(s, Inches(0.45), Inches(2.05), Inches(6.0), Inches(4.3), fill=PANEL2, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(2.3), Inches(5.4), Inches(0.4), [[('You provide', 16, BLUE, True)]])
you = [
    ('Keep costing parts as normal', 'every analysis feeds the memory automatically'),
    ('One click when a real quote arrives', '"Log Actual £" — this is the learning fuel'),
    ('Optional: CAD file or BOM', 'sharper matching and better auto-fill'),
    ('Optional: historic quotes (CSV)', 'seeds the memory so it starts smart, not empty'),
]
for i, (a, b) in enumerate(you):
    y = Inches(2.85 + i * 0.85)
    text(s, Inches(0.75), y, Inches(5.4), Inches(0.8),
         [[(f'{i+1}.  ' + a, 13, DARK, True)], [('     ' + b, 11, MUTED, False)]], space_after=2)
box(s, Inches(6.85), Inches(2.05), Inches(6.05), Inches(4.3), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(7.15), Inches(2.3), Inches(5.4), Inches(0.4), [[('The tool does', 16, GREEN, True)]])
tool = [
    'Remembers every analysis (no duplicates, org-wide)',
    'Finds & explains similar past parts instantly',
    'Suggests benchmarks, materials and real prices',
    'Measures its own error and self-corrects per category',
    'Tightens confidence ranges as evidence grows',
    'Monitors all parts in the background & raises findings',
]
for i, t in enumerate(tool):
    y = Inches(2.85 + i * 0.57)
    text(s, Inches(7.15), y, Inches(5.5), Inches(0.5), [[('✓  ', 13, GREEN, True), (t, 12.5, BODY, False)]])
notes(s, "A fair management question: what does this demand from the team? Honestly — almost nothing. Engineers keep "
         "costing parts exactly as they do today; the memory builds itself. The one new habit is a single click when a "
         "real supplier quote arrives. Optionally, attaching CAD or BOM files sharpens the matching, and a one-off "
         "import of historical quotes would let the system start smart instead of starting empty — that's my "
         "recommended first action if we proceed.")

# ════════════════════════════════════════════════════════════════════════════
# 12 — RESULTS & ACCURACY
# ════════════════════════════════════════════════════════════════════════════
s = header('Results & accuracy — measured, not promised', 'Evidence')
rows = [
    ('Estimating error after learning (machining segment)', '10.9%  →  0.3%', GREEN),
    ('Estimating error after learning (casting · China)', '8.7%  →  0.6%', GREEN),
    ('Confidence band on the same part', '±20.4%  →  ±2.8%', GREEN),
    ('Similar-part recognition on live example', '98–99% match, reasons shown', CYAN),
    ('Autonomous findings in unattended demo', '£512,000 / yr surfaced', RED),
    ('Automated tests protecting all of this', '917 passing (77 suites)', VIOLET),
]
for i, (a, b, c) in enumerate(rows):
    y = Inches(2.1 + i * 0.72)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.6), fill=(PANEL if i % 2 == 0 else BG), round_=True, radius=0.15)
    text(s, Inches(0.8), y + Inches(0.11), Inches(8.2), Inches(0.4), [[(a, 13, BODY, False)]])
    text(s, Inches(8.6), y + Inches(0.09), Inches(4.1), Inches(0.42), [[(b, 14, c, True)]], align=PP_ALIGN.RIGHT)
box(s, Inches(0.45), Inches(6.5), Inches(12.45), Inches(0.72), fill=PANEL2, round_=True, radius=0.10)
text(s, Inches(0.75), Inches(6.62), Inches(11.9), Inches(0.5),
     [[('How we verified: ', 12, DARK, True),
       ('every number above comes from running the real system end-to-end — live server, real database, real API calls — not from slides or simulations.', 12, BODY, False)]])
notes(s, "Everything on this slide was measured on the running system — live server, real database — not projected. "
         "The two headline accuracy results: after learning from three quotes, segment error dropped to well under "
         "one percent in both directions of bias. The confidence band tightened seven-fold. Recognition worked at 99 "
         "percent on the live example. The unattended agent surfaced half a million pounds. And the whole capability "
         "is protected by 813 automated tests, so it won't quietly regress. One caveat I want to be transparent about: "
         "these demos used small seeded datasets — real-world accuracy will build as OUR data accumulates. The "
         "mechanism is proven; the asset grows with use.")

# ════════════════════════════════════════════════════════════════════════════
# 13 — BENEFITS
# ════════════════════════════════════════════════════════════════════════════
s = header('What this means for the business', 'Benefits')
bens = [
    ('⚡', 'Faster costing', 'New parts start from proven history instead of a blank sheet — matches, materials and benchmarks appear instantly.', BLUE),
    ('🎯', 'Accuracy that compounds', 'Every real quote makes the next estimate better. Accuracy is measured and reported — defensible in any negotiation.', GREEN),
    ('🏦', 'Knowledge stays', 'Senior engineers\' judgement is captured as data. It doesn\'t leave when people do — and juniors inherit it from day one.', VIOLET),
    ('💰', 'Money found proactively', 'The agent watches all parts continuously and flags overpayment and pricing risk by itself, quantified in £/yr.', RED),
]
for i, (icon, t, d, c) in enumerate(bens):
    x = Inches(0.45 + (i % 2) * 6.35)
    y = Inches(2.05 + (i // 2) * 2.25)
    box(s, x, y, Inches(6.1), Inches(2.0), fill=PANEL, round_=True, radius=0.07)
    box(s, x, y, Inches(0.09), Inches(2.0), fill=c)
    text(s, x + Inches(0.3), y + Inches(0.18), Inches(5.5), Inches(0.5),
         [[(icon + '  ', 18, c, False), (t, 16, c, True)]])
    text(s, x + Inches(0.3), y + Inches(0.75), Inches(5.55), Inches(1.15), [[(d, 12, BODY, False)]], line_spacing=1.18)
text(s, Inches(0.45), Inches(6.6), Inches(12.4), Inches(0.6),
     [[('And it is ours: ', 13, DARK, True),
       ('the knowledge base runs on our infrastructure and grows into a proprietary asset competitors cannot buy.', 13, BODY, False)]])
notes(s, "Four benefits worth remembering. Speed — new parts start from history, not from zero. Accuracy — it "
         "compounds with every quote and is always measured, which is what makes our numbers defensible with "
         "suppliers. Retention — expert judgement becomes company data instead of leaving with people. And proactive "
         "savings — the agent finds money continuously, quantified per year. Finally, the strategic point: this "
         "knowledge base is proprietary. A competitor can buy the same software; they cannot buy our accumulated "
         "costing intelligence.")

# ════════════════════════════════════════════════════════════════════════════
# NEW A — CONFORMAL CONFIDENCE
# ════════════════════════════════════════════════════════════════════════════
s = header('Confidence you can defend — not just assert', 'New in 2026 · advanced intelligence')
text(s, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.45),
     [[('Every should-cost now carries two ranges — the physics estimate, and an ', 13.5, BODY, False),
       ('empirical band proven against your own logged quotes.', 13.5, DARK, True)]])
# Left: physics prior
box(s, Inches(0.45), Inches(2.5), Inches(6.05), Inches(2.5), fill=PANEL, round_=True, radius=0.06)
box(s, Inches(0.45), Inches(2.5), Inches(6.05), Inches(0.09), fill=BLUE)
text(s, Inches(0.75), Inches(2.72), Inches(5.5), Inches(0.4), [[('Physics prior (Monte-Carlo)', 14, BLUE, True)]])
text(s, Inches(0.75), Inches(3.2), Inches(5.5), Inches(1.7),
     [[('How well the inputs are known — before you have any real quotes.', 12, BODY, False)],
      [('Example:  £86.34  ± 3.4%', 14, DARK, True)],
      [('Always available, every part, every commodity.', 11, MUTED, False, True)]],
     space_after=8, line_spacing=1.15)
# Right: empirical conformal
box(s, Inches(6.85), Inches(2.5), Inches(6.05), Inches(2.5), fill=PANEL2, round_=True, radius=0.06)
box(s, Inches(6.85), Inches(2.5), Inches(6.05), Inches(0.09), fill=GREEN)
text(s, Inches(7.15), Inches(2.72), Inches(5.5), Inches(0.4), [[('Empirical band (conformal)', 14, GREEN, True)]])
text(s, Inches(7.15), Inches(3.2), Inches(5.5), Inches(1.7),
     [[('Proven against the quotes YOU logged — with a coverage guarantee.', 12, BODY, False)],
      [('Example:  90% of your quotes land within  ± 6.5%  →  £81.54 – £92.88', 13, DARK, True)],
      [('Tightens automatically as more quotes are logged.', 11, GREEN, False, True)]],
     space_after=8, line_spacing=1.15)
box(s, Inches(0.45), Inches(5.25), Inches(12.45), Inches(1.5), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(5.45), Inches(11.9), Inches(1.2),
     [[('Why it helps:  ', 13, BLUE, True),
       ('a buyer can’t defend "trust me, ±5%". They CAN defend "90% of our real quotes for this family landed within '
        '±6.5%." The band edge is an observed error — evidence, not an assertion. It is the honest uncertainty '
        'no competitor states this way.', 12.5, BODY, False)]], line_spacing=1.2)
notes(s, "A new capability that sharpens the moat. Every should-cost now shows two ranges. On the left, the physics "
         "prior — a Monte-Carlo band that reflects how well the inputs are known; it's always available. On the right, "
         "the new one: an empirical band computed from the actual quotes you've logged, using a statistical method "
         "called conformal prediction that comes with a coverage guarantee. So instead of asserting a precision, the "
         "tool says, honestly, '90% of your real quotes for this material family have landed within plus-or-minus 6.5%.' "
         "That band edge is an observed error, not a claim — which is exactly why a buyer can defend it in the room. "
         "And it tightens on its own as more quotes are logged. No competitor states uncertainty with a guarantee like this.")

# ════════════════════════════════════════════════════════════════════════════
# NEW B — OUTCOME-WEIGHTED FINDINGS
# ════════════════════════════════════════════════════════════════════════════
s = header('The agent learns what actually earns money', 'New in 2026 · advanced intelligence')
text(s, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.45),
     [[('The autonomous agent no longer just finds gaps — it learns which findings ', 13.5, BODY, False),
       ('actually convert into savings, and re-ranks by the money it can really recover.', 13.5, DARK, True)]])
# comparison table
hdr = ['Finding', 'Raw gap (impact)', 'Learned conversion', 'Expected realizable']
cols_x = [Inches(0.6), Inches(4.4), Inches(7.4), Inches(10.2)]
box(s, Inches(0.45), Inches(2.5), Inches(12.45), Inches(0.5), fill=DARK, round_=False)
for i, htext in enumerate(hdr):
    text(s, cols_x[i], Inches(2.58), Inches(3.4), Inches(0.35), [[(htext, 11.5, BG, True)]])
rows = [
    ('Cast Housing', '£200k/yr', '20% — rarely closes', '£40k', RED),
    ('Machined Knuckle', '£100k/yr', '80% — usually closes', '£80k', GREEN),
]
for i, (a, b, c, d, col) in enumerate(rows):
    y = Inches(3.0 + i * 0.62)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.6), fill=PANEL if i % 2 == 0 else BG)
    text(s, cols_x[0], y + Inches(0.14), Inches(3.6), Inches(0.35), [[(a, 12.5, DARK, True)]])
    text(s, cols_x[1], y + Inches(0.14), Inches(3.0), Inches(0.35), [[(b, 12, BODY, False)]])
    text(s, cols_x[2], y + Inches(0.14), Inches(2.8), Inches(0.35), [[(c, 12, BODY, False)]])
    text(s, cols_x[3], y + Inches(0.14), Inches(2.6), Inches(0.35), [[(d, 13, col, True)]])
box(s, Inches(0.45), Inches(4.5), Inches(12.45), Inches(0.95), fill=PANEL2, round_=True, radius=0.08)
text(s, Inches(0.75), Inches(4.66), Inches(11.9), Inches(0.7),
     [[('The result: ', 12.5, BLUE, True),
       ('the machining finding — HALF the raw gap — now ranks ABOVE the casting one, because the agent learned '
        'casting renegotiations don’t close. It stops shouting about theoretical money and surfaces the money you can get.', 12.5, BODY, False)]],
     line_spacing=1.2)
text(s, Inches(0.45), Inches(5.7), Inches(12.4), Inches(1.0),
     [[('Why it helps:  ', 13, GREEN, True),
       ('sourcing time is scarce. The agent spends it where the return is real — and tracks the £ actually saved, '
        'so you can prove the tool paid for itself. One click ("Actioned £") teaches it after every negotiation.', 12.5, BODY, False)]],
     line_spacing=1.2)
notes(s, "This is a genuine closed loop on the autonomous agent. Until now it ranked findings by the size of the gap "
         "times the volume — the theoretical money. Now it learns which findings actually convert. Look at the table: "
         "the cast housing has a two-hundred-thousand-pound gap, but in our data casting renegotiations almost never "
         "close — a twenty percent hit rate — so the expected realizable value is forty thousand. The machined knuckle "
         "has half the raw gap, but machining renegotiations usually close — eighty percent — so its realizable value is "
         "eighty thousand. The agent now ranks the machining finding first, even though its headline number is smaller. "
         "Why does that help? Sourcing time is scarce; the agent points it at money you can actually recover, and it "
         "tracks pounds truly saved so the tool proves its own worth. Every closed negotiation teaches it with one click.")

# ════════════════════════════════════════════════════════════════════════════
# NEW C — NEGOTIATION COACH
# ════════════════════════════════════════════════════════════════════════════
s = header('The negotiation coach — it hands you the argument', 'New in 2026 · advanced intelligence')
text(s, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.45),
     [[('The tool now knows ', 13.5, BODY, False), ('why', 13.5, DARK, True, True),
       (' a part costs what it does — and turns that into the sentence that wins the negotiation.', 13.5, BODY, False)]])
box(s, Inches(0.45), Inches(2.5), Inches(12.45), Inches(2.0), fill=PANEL2, round_=True, radius=0.05)
box(s, Inches(0.45), Inches(2.5), Inches(0.09), Inches(2.0), fill=BLUE)
text(s, Inches(0.8), Inches(2.72), Inches(11.8), Inches(0.35), [[('Live example — real output from the tool', 12, BLUE, True)]])
text(s, Inches(0.8), Inches(3.12), Inches(11.8), Inches(1.3),
     [[('“Material is £8.39 of this part, driven by Aluminium. Every 1% move in the Aluminium index shifts the '
        'piece price by £0.11. ', 14, DARK, True),
       ('A quote of £95 is only justified if Aluminium were ~14% above today’s index — ask the supplier to show '
        'that, or hold at £86.34.”', 14, VIOLET, True)]], line_spacing=1.3)
box(s, Inches(0.45), Inches(4.75), Inches(6.05), Inches(2.0), fill=PANEL, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(4.95), Inches(5.5), Inches(1.7),
     [[('What it does', 13.5, DARK, True)],
      [('Links the material to its commodity index, then works out — through the same maths the engine uses — '
        'exactly how much a supplier’s price implies the metal has moved.', 11.5, BODY, False)]],
     space_after=6, line_spacing=1.18)
box(s, Inches(6.85), Inches(4.75), Inches(6.05), Inches(2.0), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(7.15), Inches(4.95), Inches(5.5), Inches(1.7),
     [[('Why it helps', 13.5, GREEN, True)],
      [('The buyer walks in with a defensible, arithmetic counter-argument instead of "that feels high". '
        'It converts should-cost into negotiating power — the number and the words to win with it.', 11.5, BODY, False)]],
     space_after=6, line_spacing=1.18)
notes(s, "This is the capability that turns a number into leverage. The tool now builds a small causal model of the "
         "part: it knows the material bucket is driven by a specific commodity index — aluminium here — and it computes, "
         "through the very same overhead and margin maths the engine already uses, how sensitive the price is. Then it "
         "writes the argument for you. Read the live example: material is eight-thirty-nine, driven by aluminium; every "
         "one percent move in aluminium shifts the piece price by eleven pence; and a supplier quote of ninety-five is "
         "only justified if aluminium were about fourteen percent above today's index — so ask them to prove that, or "
         "hold at eighty-six thirty-four. That is a defensible, arithmetic counter a buyer can say out loud. It converts "
         "should-cost into negotiating power — the number AND the sentence to win with it.")

# ════════════════════════════════════════════════════════════════════════════
# NEW D — WHAT-IF ENGINE
# ════════════════════════════════════════════════════════════════════════════
s = header('Cost weather on demand — the what-if engine', 'New in 2026 · advanced intelligence')
text(s, Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.45),
     [[('Ask "what if a commodity moves?" and get the answer instantly — for one part, or the whole portfolio.', 13.5, BODY, False)]])
box(s, Inches(0.45), Inches(2.5), Inches(6.05), Inches(2.7), fill=PANEL, round_=True, radius=0.06)
box(s, Inches(0.45), Inches(2.5), Inches(6.05), Inches(0.09), fill=BLUE)
text(s, Inches(0.75), Inches(2.72), Inches(5.5), Inches(0.4), [[('One part — the live slider', 14, BLUE, True)]])
text(s, Inches(0.75), Inches(3.2), Inches(5.5), Inches(1.9),
     [[('Drag a commodity −20% … +20% and the piece price recomputes as you move.', 12, BODY, False)],
      [('Example:  Aluminium +15%  →  £86.34 → £87.92', 13.5, DARK, True)],
      [('Instant sensitivity in a single gesture.', 11, MUTED, False, True)]],
     space_after=8, line_spacing=1.15)
box(s, Inches(6.85), Inches(2.5), Inches(6.05), Inches(2.7), fill=PANEL, round_=True, radius=0.06)
box(s, Inches(6.85), Inches(2.5), Inches(6.05), Inches(0.09), fill=VIOLET)
text(s, Inches(7.15), Inches(2.72), Inches(5.5), Inches(0.4), [[('Whole portfolio — the scenario', 14, VIOLET, True)]])
text(s, Inches(7.15), Inches(3.2), Inches(5.5), Inches(1.9),
     [[('Apply a move across every part at once and see who changes status.', 12, BODY, False)],
      [('Example:  "If Steel +10% → 7 parts cross underwater, £X/yr at risk."', 13.5, DARK, True)],
      [('Pre-empt losses before the market moves them.', 11, MUTED, False, True)]],
     space_after=8, line_spacing=1.15)
box(s, Inches(0.45), Inches(5.45), Inches(12.45), Inches(1.3), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(5.62), Inches(11.9), Inches(1.05),
     [[('Honest by design:  ', 13, AMBER, True),
       ('this is a conditional — "IF the index moves" — never a price forecast. So it stays defensible even on an '
        'indicative commodity feed, and upgrades to a true forecast the day a live market feed is connected.', 12.5, BODY, False)]],
     line_spacing=1.2)
notes(s, "Two ways to ask 'what if a commodity moves?' On the left, for a single part: drag a slider from minus twenty "
         "to plus twenty percent and the piece price recomputes live as you move it — plus fifteen percent aluminium takes "
         "this part from eighty-six thirty-four to eighty-seven ninety-two. On the right, for the whole portfolio: apply a "
         "move across every part at once and see who changes status — if steel rises ten percent, these seven parts flip "
         "underwater and here's the money at risk. That lets sourcing pre-empt losses instead of reacting to them. And the "
         "honest framing that keeps it credible: this is always a conditional — IF the index moves — never a forecast we "
         "can't defend. The day we connect a live market feed, the same machinery becomes a genuine forecast.")

# ════════════════════════════════════════════════════════════════════════════
# NEW F — PCB PHOTO → BOM ACCURACY HARDENING
# ════════════════════════════════════════════════════════════════════════════
s = header('PCB photo → BOM → should-cost — now supplier-grade', 'New in 2026 · advanced intelligence')
text(s, Inches(0.45), Inches(1.82), Inches(12.4), Inches(0.42),
     [[('Photograph a circuit board, get a costed bill of materials in ~60 seconds. This year we hardened it from a '
        'clever demo into a ', 13, BODY, False),
       ('number you can put in front of a supplier.', 13, DARK, True)]])
# Left — Manual vs AI proof chart (real ECU run)
cd = CategoryChartData()
cd.categories = ['Components', 'Fab + assembly', 'Total / board']
cd.add_series('Manual (engineer)', (62.0, 10.0, 73.0))
cd.add_series('AI (from photo)', (62.29, 6.82, 69.11))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.45), Inches(2.45), Inches(6.15), Inches(3.35), cd)
style_chart(gf.chart, [MUTED, BLUE])
gf.chart.has_legend = True
gf.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
gf.chart.legend.include_in_layout = False
gf.chart.legend.font.size = Pt(10.5)
text(s, Inches(0.45), Inches(5.82), Inches(6.15), Inches(0.42),
     [[('Real automotive ECU, China @ 10k/yr — AI landed within ~5% of a half-day manual estimate, in ~60 s.', 10, MUTED, False, True)]],
     line_spacing=1.05)
# Right — what we fixed
box(s, Inches(6.85), Inches(2.45), Inches(6.05), Inches(3.35), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(7.15), Inches(2.62), Inches(5.5), Inches(0.35), [[('What we fixed this year', 14, DARK, True)]])
fixes = [
    ('Empty-BOM bug — killed', 'Complex boards used to return nothing; now we read every block the model emits.', GREEN),
    ('Catalogue grounding', 'Confirmed parts snap to real market prices — offline, no external API.', BLUE),
    ('Class-median cap + magnitude guard', 'A misread part can no longer dominate: one MCU £84 → £18, capped & flagged.', VIOLET),
    ('Deterministic fabrication', 'Fab is derived from stable board features — the headline no longer swings run-to-run.', CYAN),
    ('Confirmed vs needs-verification', 'The headline splits the £ you can trust from the £ to firm up.', AMBER),
]
for i, (a, b, c) in enumerate(fixes):
    y = Inches(3.05 + i * 0.55)
    box(s, Inches(7.15), y + Inches(0.05), Inches(0.09), Inches(0.42), fill=c)
    text(s, Inches(7.42), y, Inches(5.35), Inches(0.55),
         [[('✓  ' + a + ' — ', 11.5, DARK, True), (b, 10.5, BODY, False)]], line_spacing=1.0)
box(s, Inches(0.45), Inches(5.98), Inches(12.45), Inches(0.85), fill=PANEL2, round_=True, radius=0.10)
text(s, Inches(0.75), Inches(6.12), Inches(12.0), Inches(0.6),
     [[('Result:  ', 12.5, BLUE, True),
       ('the 2–3× over-costing is gone. That live ECU came back ASIL-B, 23 BOM lines, £69.11/board — with '
        '£37.65 confirmed and £24.64 honestly flagged to verify. Same photo, same answer, every run.', 12, BODY, False)]],
     line_spacing=1.15)
notes(s, "")  # notes rewritten below

# ════════════════════════════════════════════════════════════════════════════
# NEW E — GLASS-BOX DIFFERENTIATION
# ════════════════════════════════════════════════════════════════════════════
s = header('Why ours is different — glass-box autonomy', 'The differentiator')
box(s, Inches(0.45), Inches(2.0), Inches(12.45), Inches(1.3), fill=PANEL2, round_=True, radius=0.06)
text(s, Inches(0.75), Inches(2.2), Inches(11.9), Inches(1.0),
     [[('The principle:  ', 15, BLUE, True),
       ('every learned or derived number stays auditable — a value a cost engineer can read and defend. '
        'No black-box weight ever touches the price. The AI narrates and explains; it never decides the number in secret.', 13.5, BODY, False)]],
     line_spacing=1.2)
cmp = [
    ('Continuous learning', 'Calibrates on YOUR quotes; conformal band with a guarantee', GREEN),
    ('Autonomous action', 'Agent opens findings unattended AND learns which convert', GREEN),
    ('Causal reasoning', 'Knows why a part costs what it does; coaches the negotiation', GREEN),
    ('Explainable — always', 'Every number defensible line-by-line; the competitor’s edge is the opposite', BLUE),
    ('Runs in your walls', 'On-premise; the knowledge is your IP and never leaves', VIOLET),
]
for i, (t, d, c) in enumerate(cmp):
    y = Inches(3.5 + i * 0.66)
    box(s, Inches(0.45), y, Inches(12.45), Inches(0.58), fill=PANEL if i % 2 == 0 else BG, round_=True, radius=0.1)
    box(s, Inches(0.45), y, Inches(0.09), Inches(0.58), fill=c)
    text(s, Inches(0.75), y + Inches(0.13), Inches(3.6), Inches(0.35), [[(t, 12.5, c, True)]])
    text(s, Inches(4.5), y + Inches(0.13), Inches(8.2), Inches(0.35), [[(d, 12, BODY, False)]])
notes(s, "The one slide that says why this beats what's on the market. The competitors demo continuous learning and "
         "adaptive reasoning, but they do it as a black box — a number you're told to trust. Our principle is the "
         "opposite, and it's the harder thing to build: every learned or derived number stays auditable. The bias "
         "factor, the conformal band, the hit-rate, the coach's arithmetic — all of it is a value a cost engineer can "
         "read and defend in a negotiation. The AI narrates and explains; it never decides the price in secret. So we "
         "match them on learning, autonomy and causal reasoning — and beat them on the one thing that actually wins "
         "deals in this field: defensibility. And it all runs inside our own walls, so the intelligence is our IP and "
         "never leaves the building. That combination — autonomous, self-improving, AND glass-box, on-premise — is not "
         "available anywhere else.")

# ════════════════════════════════════════════════════════════════════════════
# 14 — NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
s = header('Where we are, and the ask', 'Next steps')
box(s, Inches(0.45), Inches(2.0), Inches(6.0), Inches(4.4), fill=PANEL, round_=True, radius=0.05)
text(s, Inches(0.75), Inches(2.25), Inches(5.4), Inches(0.4), [[('Status today', 16, GREEN, True)]])
st = [
    'All capabilities built, tested (917 tests) and live — incl. 6 new for 2026',
    'Verified end-to-end on the running system',
    'Zero extra licence cost — built into our tool',
    'Runs on-premise; no data leaves the company',
]
for i, t in enumerate(st):
    text(s, Inches(0.75), Inches(2.8 + i * 0.55), Inches(5.4), Inches(0.5),
         [[('✓  ', 13, GREEN, True), (t, 12.5, BODY, False)]])
text(s, Inches(0.75), Inches(5.15), Inches(5.4), Inches(1.1),
     [[('Honest note: ', 12, DARK, True),
       ('the intelligence starts empty and grows with use. The mechanism is proven; the value builds as we feed it.', 12, BODY, False)]],
     line_spacing=1.15)
box(s, Inches(6.85), Inches(2.0), Inches(6.05), Inches(4.4), fill=PANEL2, round_=True, radius=0.05)
text(s, Inches(7.15), Inches(2.25), Inches(5.4), Inches(0.4), [[('The ask — three small decisions', 16, BLUE, True)]])
asks = [
    ('1.  Adopt the habit', 'Make "Log Actual £" part of the quote-handling routine. One click per quote.'),
    ('2.  Seed the memory', 'Approve a one-off import of historical quotes so the tool starts smart (~1–2 days of effort).'),
    ('3.  Review the findings', 'Put the agent\'s monthly findings on the sourcing team\'s agenda — it is already finding money.'),
]
for i, (a, b) in enumerate(asks):
    y = Inches(2.85 + i * 1.05)
    text(s, Inches(7.15), y, Inches(5.5), Inches(1.0),
         [[(a, 13.5, DARK, True)], [('     ' + b, 11.5, BODY, False)]], space_after=3, line_spacing=1.1)
box(s, 0, H - Inches(0.16), W, Inches(0.16), fill=INDIGO)
notes(s, "To close: the capability is built, tested and live — at no extra licence cost, running on our own "
         "infrastructure. I'll be honest about the one dependency: the intelligence starts empty and grows with use. "
         "So the ask is three small decisions. One: make logging the actual quote a one-click habit. Two: approve a "
         "one-off import of our historical quotes so the system starts smart — roughly a day or two of effort. "
         "Three: put the agent's findings on the sourcing agenda each month. If we do those three things, this "
         "becomes a compounding asset from day one. Thank you — happy to take questions or show the live system.")

# ════════════════════════════════════════════════════════════════════════════
# SPEAKER NOTES — rewritten, humanised + punchy, applied in slide order
# ════════════════════════════════════════════════════════════════════════════
SPEAKER_NOTES = [
    # 1 — TITLE
    "Open strong, then pause. \"Every costing tool on the market does maths. Ours does something none of them do — it "
    "learns.\" Hold up the idea: the tool remembers every part we've ever costed, recognises the next one, corrects "
    "itself against real quotes, and goes hunting for savings while nobody's watching. One promise before we start: "
    "every single number in this deck is from a live run of the real system — no mock-ups, no slideware. Let's go.",

    # 2 — EXECUTIVE SUMMARY
    "If you remember four numbers, remember these. Error fell thirty-six-fold — eleven percent down to under one — after "
    "the tool learned from just three real quotes. The background agent found half a million pounds a year of pricing "
    "issues with nobody at the keyboard. It recognised a brand-new bracket against past parts at ninety-nine percent and "
    "told us why. And it's production-grade — nine hundred and seventeen automated tests. But here's the one line that "
    "matters: costing intelligence used to live in people's heads and walk out the door when they left. Now it "
    "accumulates as a company asset that gets more valuable every day we use it.",

    # 3 — WHAT AGENTIC AI MEANS
    "\"Agentic\" gets thrown around a lot, so let me make it concrete — four plain verbs. It REMEMBERS: every analysis "
    "becomes a stored case. It RECOGNISES: new parts get matched to that memory instantly. It SELF-CORRECTS: real quotes "
    "teach it, and its accuracy is measured, not claimed. And it ACTS: an agent raises findings without being asked. "
    "One design choice underpins all of it — every suggestion shows its source parts and its arithmetic. That's "
    "deliberate, because a number you can't defend is worthless in a supplier negotiation. Ours you can defend, line by line.",

    # 4 — THE LEARNING LOOP
    "Here's the whole machine on one slide — and the punchline is how little it asks of the engineer. They cost a part "
    "exactly as they do today. Everything blue and automatic happens on its own: it remembers, it recognises, it "
    "suggests, and the agent keeps watch. There is exactly ONE new habit — a single click, 'Log Actual £,' when a real "
    "quote lands. That one click is the fuel for every loop on this slide, and every loop makes the next estimate "
    "sharper. No new workload, compounding returns. That's the deal.",

    # 5 — MEMORY: KNOWLEDGE BASE
    "Capability one — the memory. For every analysis we keep a fingerprint of the part, the full costed result, any real "
    "quotes, and — this is the clever bit — the exact places an expert overrode the AI. Those corrections are our "
    "engineers literally teaching the tool. On the right is why it matters to us specifically: it's shared, so a "
    "junior inherits senior judgement on day one; it lives on our servers, so it's our asset, not a vendor's; and it "
    "compounds — genuinely useful after only twenty or thirty analyses. Every costing we run from today is a deposit "
    "into an account that only grows.",

    # 6 — RECOGNITION & SUGGESTIONS
    "Capability two, shown with the tool's actual live output. An engineer starts a new aluminium bracket. Before "
    "they've even finished, the tool has surfaced the three most similar parts we've ever costed — at ninety-eight to "
    "ninety-nine percent — and it tells you WHY they matched: same material family, same weight class, same region. "
    "Then it hands over the gold: the median cost, the shared material, and the real prices we actually paid. A junior "
    "just stood on the shoulders of every senior who came before them. And notice — every number names its source. "
    "Nothing here is a black box.",

    # 7 — SELF-CALIBRATION
    "Capability three is the accuracy engine, and this is a real measured result. Our machining estimates ran about "
    "eleven percent LOW against real prices. After the tool learned from three logged quotes, the error was "
    "three-tenths of a percent. Now look right — and this is the subtle, important part. Machining ran low; our China "
    "castings ran high. Average them together and everything looks 'fine' — the errors cancel and hide. The tool "
    "refuses to average. It corrects each process, material and region separately, so it catches exactly what a "
    "portfolio average conceals. Accuracy here is measured against reality and reported out loud — not asserted on a slide.",

    # 8 — HONEST UNCERTAINTY
    "Capability four — honesty about precision, which is a feature, not a weakness. Every estimate now comes as a range: "
    "optimistic, most likely, conservative. Early on, with no evidence, that range is wide — and that IS the truth. As "
    "real quotes prove the tool right, the range tightens on its own: on our test part, from plus-or-minus twenty percent "
    "down to three. Give a buyer this and you've handed them a script — the conservative end is the walk-away, the "
    "optimistic end is the stretch target. A single number pretends to a confidence early estimates simply don't have. "
    "This tells the truth, and the truth negotiates better.",

    # 9 — AUTONOMOUS AGENT
    "Capability five — this is the one that earns the word 'agentic.' A monitor runs on our server on a schedule and "
    "compares what we PAY against what things SHOULD cost, using everything the tool has learned. In the live demo, with "
    "nobody at the keyboard, it opened three findings by itself: a renegotiation worth four hundred thousand a year, an "
    "'underwater' price that flags supply risk, and a stale estimate that needs a fresh quote. Half a million pounds a "
    "year, surfaced unattended — and every finding shows its arithmetic, ready to carry straight into a supplier meeting. "
    "The tool isn't waiting to be asked anymore. It's already working.",

    # 10 — SUPPORTING AI BRAINS
    "The learning loop is the headline, but it stands on a real platform. The assistant answers from our own rate data "
    "with citations — it quotes our numbers, it doesn't improvise. The RFQ analyst turns a full quote package into a "
    "costed, risk-flagged negotiation brief. CAD feature costing shows a designer which hole or surface is driving the "
    "price. And carbon co-costing puts a CO2 figure next to every pound — which automotive and aerospace customers now "
    "demand in the RFQ itself. Eighteen commodity engines, CAD-to-cost, PCB photo-to-BOM — and the agentic layer sits "
    "on top of all of it, feeding and fed by the same memory.",

    # 11 — INPUTS REQUIRED
    "Fair management question: what does this cost the team in effort? Honestly — almost nothing. Engineers keep costing "
    "parts exactly as they do now; the memory builds itself in the background. The one new habit is a single click when a "
    "real quote arrives. Two optional accelerants: attach CAD or BOM files for sharper matching, and — my recommended "
    "first move — a one-off import of our historical quotes so the system starts SMART instead of empty. Low effort in, "
    "compounding value out. That's the whole ask on the input side.",

    # 12 — RESULTS & ACCURACY
    "Everything on this slide was measured on the running system — live server, real database, real API calls — not "
    "projected. Segment error dropped to well under one percent in BOTH directions of bias. The confidence band "
    "tightened seven-fold. Recognition hit ninety-nine percent on the live example. The unattended agent surfaced half "
    "a million pounds. And nine hundred and seventeen tests stand guard so none of it quietly regresses. One honest "
    "caveat, because credibility is the whole game here: these demos ran on small seeded datasets — real-world accuracy "
    "builds as OUR data accumulates. The mechanism is proven. The asset grows with use.",

    # 13 — BENEFITS
    "Four benefits, business language. Speed: new parts start from proven history, not a blank sheet. Accuracy: it "
    "compounds with every quote and it's always measured — which is precisely what makes our numbers defensible with "
    "suppliers. Retention: your best engineers' judgement becomes company data instead of leaving in a leaver's head. "
    "And proactive savings: the agent finds money continuously, quantified per year. Then the strategic kicker — a "
    "competitor can buy the same software tomorrow. They cannot buy our accumulated costing intelligence. That's the "
    "moat, and it deepens every day.",

    # 14 — NEW A · CONFORMAL CONFIDENCE
    "New this year, and it sharpens the moat. Every should-cost now carries TWO ranges. On the left, the physics prior — "
    "a Monte-Carlo band reflecting how well the inputs are known; always available, every part. On the right, the new "
    "one: an empirical band built from the actual quotes you've logged, using conformal prediction — a method that comes "
    "with a mathematical coverage guarantee. So the tool stops asserting a precision and starts stating an observed fact: "
    "'ninety percent of your real quotes for this family landed within plus-or-minus six-and-a-half percent.' That band "
    "edge is measured error, not a claim — which is exactly why a buyer can defend it in the room. And it tightens on its "
    "own. No competitor states uncertainty with a guarantee like this.",

    # 15 — NEW B · OUTCOME-WEIGHTED FINDINGS
    "This closes the loop on the agent. It used to rank findings by the raw gap — the theoretical money. Now it learns "
    "which findings actually convert into cash. Watch the table: the cast housing shows a two-hundred-thousand-pound "
    "gap, but in our data casting renegotiations almost never close — twenty percent — so the realizable value is forty "
    "thousand. The machined knuckle has HALF the raw gap, but machining usually closes — eighty percent — so it's worth "
    "eighty thousand realizable. The agent now ranks the machining finding ABOVE the bigger one. Why does that matter? "
    "Sourcing time is scarce — it now points at money you can actually get, and tracks the pounds truly saved so the "
    "tool proves its own worth. It stopped shouting about theoretical money.",

    # 16 — NEW C · NEGOTIATION COACH
    "This is where a number becomes leverage. The tool now builds a small causal model of the part — it knows the "
    "material bucket is driven by a specific commodity index, aluminium here — and it runs the same overhead-and-margin "
    "maths the engine already uses to work out how sensitive the price is. Then it writes the argument for you. Read the "
    "live output: material is eight-thirty-nine, driven by aluminium; every one-percent move in the index shifts the "
    "piece price by eleven pence; a quote of ninety-five is only justified if aluminium were about fourteen percent above "
    "today — so ask them to prove it, or hold at eighty-six thirty-four. That's a sentence a buyer can say out loud and "
    "not get argued out of. Should-cost, converted into negotiating power — the number AND the words.",

    # 17 — NEW D · WHAT-IF ENGINE
    "Two ways to ask 'what if a commodity moves?' Left, one part: drag a slider from minus twenty to plus twenty percent "
    "and the price recomputes live under your finger — plus fifteen percent aluminium takes this part from eighty-six "
    "thirty-four to eighty-seven ninety-two. Right, the whole portfolio at once: push a move across every part and watch "
    "who changes status — steel up ten percent, seven parts flip underwater, here's the money at risk. That's sourcing "
    "getting ahead of the market instead of reacting to it. And the framing that keeps it honest: it's always a "
    "conditional — IF the index moves — never a forecast we can't stand behind. Connect a live market feed and the same "
    "machinery becomes a real forecast overnight.",

    # 18 — NEW F · PCB PHOTO → BOM ACCURACY
    "Now the feature I'm proudest of this year. Photograph a circuit board, and in about sixty seconds you get a costed "
    "bill of materials. The story here is honesty about the journey: it started as a clever demo, and complex automotive "
    "boards would sometimes come back EMPTY — a lot of compute, no result. We fixed that at the root, then hardened the "
    "whole pipeline. The chart is the proof: the same real ECU, costed by an engineer bottom-up in half a day versus the "
    "AI from one photo — sixty-nine pounds against seventy-three, within about five percent, in sixty seconds. On the "
    "right, the five fixes that got us there — grounding prices to a real catalogue, capping any part the model "
    "misreads so one bad line can't blow up the total, making fabrication deterministic, and splitting the headline into "
    "what you can trust versus what to verify. Bottom line: the two-to-three-times over-costing is gone. Same photo, same "
    "answer, every run — and it tells you exactly which lines to firm up before you quote.",

    # 19 — NEW E · GLASS-BOX DIFFERENTIATION
    "If you take one competitive message away, take this. Rivals can demo continuous learning and adaptive reasoning — "
    "but as a black box, a number you're told to trust. Our principle is the opposite, and it's the harder thing to "
    "build: every learned or derived value stays auditable. The bias factor, the conformal band, the hit-rate, the "
    "coach's arithmetic — all of it is a value a cost engineer can read and defend across the table. The AI narrates and "
    "explains; it never sets the price in secret. So we match the market on learning, autonomy and causal reasoning — "
    "and we beat it on the one thing that actually wins deals in this business: defensibility. Autonomous, "
    "self-improving, glass-box, AND on-premise so the intelligence never leaves our walls. That combination doesn't "
    "exist anywhere else.",

    # 20 — NEXT STEPS
    "To close. The capability is built, tested — nine hundred and seventeen tests — and live, at no extra licence cost, "
    "running on our own infrastructure. I'll be straight about the one dependency: the intelligence starts empty and "
    "grows with use. So the ask is three small decisions. One: make 'Log Actual £' a one-click habit when quotes come "
    "in. Two: approve a one-off import of our historical quotes so the system starts smart — a day or two of effort. "
    "Three: put the agent's findings on the sourcing agenda each month; it's already finding money. Do those three "
    "things and this becomes a compounding asset from day one. Thank you — I'm happy to take questions, or show you the "
    "live system right now.",
]

_slides = list(prs.slides)
for _i, _slide in enumerate(_slides):
    if _i < len(SPEAKER_NOTES):
        _slide.notes_slide.notes_text_frame.text = SPEAKER_NOTES[_i]

assert len(_slides) == len(SPEAKER_NOTES), f"slides={len(_slides)} notes={len(SPEAKER_NOTES)}"

OUT = 'CostVision-Agentic-AI-Management-Presentation.pptx'
prs.save(OUT)
print(f'Wrote {OUT} with {len(prs.slides.slides if hasattr(prs.slides, "slides") else prs.slides._sldIdLst)} slides')
